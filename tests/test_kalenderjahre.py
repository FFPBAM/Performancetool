"""Regressionstest: der Saeulen-Chart zeigt nur VOLLE Kalenderjahre (12.08.2026).

Gemeldet hat es Philip an der Pro-Broschuere: Im Chart "PERFORMANCE P.A.
(NACH KOSTEN) IM BENCHMARKVERGLEICH" stand ein Balken "2023" — die Strategie
"Muster FFPB Pro" laeuft aber erst seit dem 01.09.2023. Der Balken zeigte
122 Tage (+3,23 % gegen +5,11 % Benchmark) und stand als Jahreswert neben
den echten Jahren 2024 und 2025.

Ursache: compute_performance_data uebersprang ein Jahr nur dann, wenn es GAR
KEINE Daten hatte (`if sub.empty`). Ob die Daten das Jahr ABDECKEN, hat
niemand geprueft. Betroffen waren 7 der 19 Strategien — jede, deren Auflage
in das rollierende Fuenf-Jahres-Fenster faellt.

Das ist keine Kosmetik: Ein Teiljahr unter der Ueberschrift "p.a." ist eine
falsche Sachaussage in einem Kundendokument.

Geprueft wird in drei Schritten:

  Schritt 1 (numpy + pandas)  — analytics.compute_performance_data gegen von
                                Hand nachgerechnete Grenzfaelle: Auflagejahr,
                                beide Toleranzraender, Loch in der Historie,
                                Strategie ohne ein einziges volles Jahr.
  Schritt 2 (+ pandas)        — an den ECHTEN CSVs: jeder gelieferte Balken
                                deckt sein Kalenderjahr wirklich ab, und die
                                sieben bekannten Stummel sind weg.
  Schritt 3 (+ python-pptx,
             streamlit)       — am ECHTEN ARTEFAKT: die Kategorien des
                                Saeulen-Charts in gebauten Broschueren.

Ohne pptx/streamlit wird Schritt 3 sauber uebersprungen.

    python tests/test_kalenderjahre.py

Rueckgabewert 0 = bestanden, 1 = fehlgeschlagen.

BEWUSST NICHT geprueft wird das Streamlit-Tool: dort bleiben angebrochene
Jahre stehen (Entscheidung Philip, 12.08.2026 — der Berater waehlt den
Zeitraum selbst und sieht ihn neben dem Chart). Siehe den Kommentar in
streamlit_app.compute_bar_data.
"""

import glob
import os
import sys

import numpy as np
import pandas as pd

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)
sys.path.insert(0, os.path.join(WURZEL, "tests"))
os.chdir(WURZEL)   # Vorlage/ und Daten/ werden relativ geladen

from modules.analytics import (  # noqa: E402
    JAHR_RAND_TOLERANZ_TAGE, compute_performance_data,
)
from test_benchmark_erkennung import (  # noqa: E402
    FEE_TEST, lade_zeitreihe, neuester_tag,
)

DATEN = os.path.join(WURZEL, "Daten")
NS_C = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}


def _ist(titel, ist, soll):
    if ist == soll:
        print(f"    OK   — {titel}: {ist}")
        return 0
    print(f"    FEHLER — {titel}: {ist}, erwartet {soll}")
    return 1


# ───────────────────────────── Schritt 1 ──────────────────────────────────

def _reihe(start, ende, luecke=None):
    """Tagesreihe mit konstanter Rendite; `luecke` = (von, bis) wird
    herausgeschnitten. Die Werte sind egal — geprueft werden die JAHRE."""
    idx = pd.date_range(start, ende, freq="D")
    if luecke:
        idx = idx[(idx < pd.Timestamp(luecke[0])) | (idx > pd.Timestamp(luecke[1]))]
    df = pd.DataFrame({"ret_port": np.full(len(idx), 0.0002),
                       "ret_bm": np.full(len(idx), 0.0001)}, index=idx)
    return df


def _jahre(start, ende, luecke=None):
    erg = compute_performance_data(_reihe(start, ende, luecke), FEE_TEST)
    return erg["performance_pa"]["jahre"], erg["performance_pa"]


def schritt1_grenzfaelle():
    print("Schritt 1 — die Regel an von Hand nachgerechneten Faellen")
    f = 0
    f += _ist("Toleranz steht auf 3 Tagen", JAHR_RAND_TOLERANZ_TAGE, 3)

    # Der gemeldete Fall: Auflage 01.09.2023, Daten bis Juli 2026.
    # Fenster = 2021..2025; 2021/2022 fehlen ganz, 2023 ist angebrochen.
    j, _ = _jahre("2023-09-01", "2026-07-21")
    f += _ist("Auflage 01.09.2023 (der Pro-Fall)", j, [2024, 2025])

    # Kontrolle: volle Historie liefert weiterhin genau fuenf Jahre.
    j, pa = _jahre("2015-01-01", "2026-07-21")
    f += _ist("lange Historie -> volles Fenster", j, [2021, 2022, 2023, 2024, 2025])
    f += _ist("Benchmark-Liste so lang wie die Jahresliste",
              len(pa["benchmark"]), len(j))
    f += _ist("Referenz-Liste so lang wie die Jahresliste",
              len(pa["referenz"]), len(j))

    # ── Der Rand am JAHRESANFANG (neu) ────────────────────────────────────
    # Feiertagsstart: bis zum 4. Januar zaehlt das Jahr als vollstaendig.
    j, _ = _jahre("2021-01-04", "2026-07-21")
    f += _ist("Start 04.01.2021 — gerade noch drin", j[0], 2021)
    j, _ = _jahre("2021-01-05", "2026-07-21")
    f += _ist("Start 05.01.2021 — einen Tag zu spaet", j, [2022, 2023, 2024, 2025])

    # ── Der Rand am JAHRESENDE (bestand schon, darf nicht kippen) ─────────
    j, _ = _jahre("2015-01-01", "2025-12-31")
    f += _ist("Daten bis 31.12.2025 -> 2025 zaehlt",
              j, [2021, 2022, 2023, 2024, 2025])
    j, _ = _jahre("2015-01-01", "2025-12-28")
    f += _ist("Daten bis 28.12.2025 -> 2025 zaehlt noch",
              j, [2021, 2022, 2023, 2024, 2025])
    j, _ = _jahre("2015-01-01", "2025-12-27")
    f += _ist("Daten bis 27.12.2025 -> 2025 faellt raus",
              j, [2020, 2021, 2022, 2023, 2024])

    # ── Loch MITTEN in der Historie ───────────────────────────────────────
    # Deshalb prueft _ist_volles_jahr aus `sub` und nicht aus df.index.min().
    j, _ = _jahre("2015-01-01", "2026-07-21", luecke=("2023-01-01", "2023-12-31"))
    f += _ist("2023 fehlt komplett", j, [2021, 2022, 2024, 2025])
    j, _ = _jahre("2015-01-01", "2026-07-21", luecke=("2023-03-01", "2023-03-31"))
    f += _ist("2023 hat ein Loch im Maerz — Raender stehen, Jahr zaehlt",
              j, [2021, 2022, 2023, 2024, 2025])
    j, _ = _jahre("2015-01-01", "2026-07-21", luecke=("2023-12-20", "2023-12-31"))
    f += _ist("2023 bricht am 19.12. ab", j, [2021, 2022, 2024, 2025])

    # ── Der Grenzfall, den es bisher nicht gab: gar kein volles Jahr ──────
    j, pa = _jahre("2026-01-02", "2026-07-21")
    f += _ist("Auflage im laufenden Jahr -> LEER, kein Absturz", j, [])
    f += _ist("und die Wertelisten sind auch leer",
              (pa["referenz"], pa["benchmark"]), ([], []))
    return f


# ───────────────────────────── Schritt 2 ──────────────────────────────────

# Gemessen am 12.08.2026 (Datenstand 260721). Die sieben Strategien, deren
# Auflage in das Fenster faellt — genau hier stand vorher ein Stummelbalken.
# Schluessel: Portfolio-Name aus der CSV. Wert: (Auflagedatum, Jahre NACHHER).
# Aendert sich das Auflagedatum (neue Datenlieferung), wird der Fall
# uebersprungen statt falsch rot — die Regel selbst wird trotzdem geprueft.
BEKANNTE_FAELLE = {
    "Muster FFPB Pro":           ("2023-09-01", [2024, 2025]),
    "Muster FFPB Pro Dividende": ("2024-10-22", [2025]),
    "Comdirect 30":              ("2024-03-12", [2025]),
    "Comdirect 70":              ("2024-03-12", [2025]),
    "Comdirect 100":             ("2024-03-12", [2025]),
    "Muster SCHWEIZ Substanz":   ("2022-09-22", [2023, 2024, 2025]),
    "Muster SCHWEIZ Aktien":     ("2022-09-12", [2023, 2024, 2025]),
}


def schritt2_echte_daten():
    """Der Test, der den Fehler verhindert haette: nicht die Jahresliste
    gegen eine Erwartung, sondern JEDEN gelieferten Balken gegen die Daten,
    die ihn tragen."""
    print("\nSchritt 2 — jeder Balken an den echten CSVs nachgemessen")
    tag = neuester_tag()
    if tag is None:
        print("    FEHLER — keine CSVs in Daten/")
        return 1
    dateien = sorted(glob.glob(os.path.join(DATEN, f"*_{tag}_*.CSV")))
    if not dateien:
        print(f"    FEHLER — keine Dateien fuer Tag {tag}")
        return 1

    tol = pd.Timedelta(days=JAHR_RAND_TOLERANZ_TAGE)
    kopf = f"    {'Strategie':28s} {'Auflage':>10s}  Balken"
    print(f"    Datenstand {tag} — {len(dateien)} Strategien")
    print(kopf)
    print("    " + "-" * (len(kopf) - 4))

    fehler = 0
    geprueft = 0
    for pfad in dateien:
        name, df = lade_zeitreihe(pfad)
        pa = compute_performance_data(df, FEE_TEST)["performance_pa"]
        jahre = pa["jahre"]
        maengel = []

        # (a) Die eigentliche Zusage: jedes gelieferte Jahr ist abgedeckt.
        for jahr in jahre:
            sub = df[df.index.year == jahr]
            if sub.index.min() > pd.Timestamp(jahr, 1, 1) + tol:
                maengel.append(f"{jahr} beginnt erst {sub.index.min():%d.%m.}")
            if sub.index.max() < pd.Timestamp(jahr, 12, 31) - tol:
                maengel.append(f"{jahr} endet schon {sub.index.max():%d.%m.}")
            geprueft += 1

        # (b) Die Listen bleiben zueinander passend (sonst verrutschen die
        #     Balken gegen ihre Beschriftung).
        if len(pa["referenz"]) != len(jahre):
            maengel.append("Referenz-Liste hat eine andere Laenge")
        if pa["benchmark"] and len(pa["benchmark"]) != len(jahre):
            maengel.append("Benchmark-Liste hat eine andere Laenge")

        # (c) Die sieben bekannten Faelle namentlich festgenagelt.
        bekannt = BEKANNTE_FAELLE.get(name)
        start = df.index.min()
        if bekannt and str(start.date()) == bekannt[0] and jahre != bekannt[1]:
            maengel.append(f"erwartet {bekannt[1]}")

        fehler += 1 if maengel else 0
        marke = " <-" if name in BEKANNTE_FAELLE else ""
        print(f"    {name[:28]:28s} {start:%d.%m.%Y}  {jahre}{marke}")
        for m in maengel:
            print(f"        FEHLER — {m}")

    print(f"    {geprueft} Balken nachgemessen, {len(BEKANNTE_FAELLE)} "
          f"Faelle namentlich geprueft")
    return fehler


# ───────────────────────────── Schritt 3 ──────────────────────────────────

# (Anzeigename, Familie, erwartete Kategorien). Offensiv ist der Kontrollfall
# mit voller Historie — an ihm darf sich nichts geaendert haben.
ARTEFAKT_FAELLE = [
    ("Pro",                        "Thema", ["2024", "2025"]),
    ("Schweiz_substanzorientiert", "Thema", ["2023", "2024", "2025"]),
    ("Offensiv",                   "Thema",
     ["2021", "2022", "2023", "2024", "2025"]),
]
COMDIRECT_ERWARTET = ["2025"]


def _kategorien(prs, folien_idx):
    """Kategorien des Saeulen-Charts direkt aus der Chart-XML."""
    from modules.pptx_slides import SHAPE_WE_CHART_BAR
    for shape in prs.slides[folien_idx].shapes:
        if shape.name != SHAPE_WE_CHART_BAR:
            continue
        if not getattr(shape, "has_chart", False):
            continue
        cat = shape.chart._chartSpace.find(".//c:ser/c:cat", NS_C)
        if cat is None:
            return []
        return [v.text for v in cat.findall(".//c:pt/c:v", NS_C)]
    return None


def schritt3_artefakt():
    print("\nSchritt 3 — Kategorien im Chart gebauter Broschueren")
    try:
        from pptx import Presentation
    except ImportError as ex:
        print(f"    UEBERSPRUNGEN — {ex}")
        return 0
    try:
        import tempfile
        from modules.vorlagen_config import (
            FAMILIE_ALLE_STRATEGIEN, VORLAGEN_FAMILIEN,
        )
        from modules.portfolioanalyse import duration_info_aus_bestand
        from test_export_smoke import _bauen, _daten, _portfolio
    except ImportError as ex:
        print(f"    UEBERSPRUNGEN — {ex}")
        return 0

    d = _daten()
    ausgabe = tempfile.mkdtemp(prefix="ffpb_kalenderjahre_")
    fehler = 0

    # ── Themen-Familie: eine Broschuere je Strategie ──────────────────────
    _, cfg = VORLAGEN_FAMILIEN["Thema"]
    idx = cfg["block_positionen"]["wertentwicklung"] - 1
    for name, familie, soll in ARTEFAKT_FAELLE:
        if name not in d["d2c"]:
            print(f"    {name[:28]:28s} UEBERSPRUNGEN (nicht in den Daten)")
            continue
        ziel, _g, meldungen = _bauen([_portfolio(name, d)], familie, d,
                                     ausgabe, f"{name}.pptx")
        for m in meldungen:
            print(f"    ! BUILD-MELDUNG {name}: {m[:88]}")
            fehler += 1
        ist = _kategorien(Presentation(ziel), idx)
        fehler += _ist(f"{name} Kategorien", ist, soll)

    # ── comdirect: alle drei Strategien in EINER Broschuere, feste Bloecke ──
    from modules.portfolioanalyse import _familien_portfolios
    namen = FAMILIE_ALLE_STRATEGIEN["comdirect"]
    portfolios, fehlend = _familien_portfolios(
        namen, d["namen"], d["d2c"], d["pf_data"], duration_info_aus_bestand)
    if fehlend:
        print(f"    comdirect UEBERSPRUNGEN (fehlende Daten: {fehlend})")
        return fehler
    ziel, _g, meldungen = _bauen(portfolios, "comdirect", d, ausgabe,
                                 "comdirect.pptx")
    for m in meldungen:
        print(f"    ! BUILD-MELDUNG comdirect: {m[:88]}")
        fehler += 1
    prs = Presentation(ziel)
    _, cfg_cd = VORLAGEN_FAMILIEN["comdirect"]
    for k, block in enumerate(cfg_cd["feste_bloecke"]):
        if "wertentwicklung" not in block:
            continue
        ist = _kategorien(prs, block["wertentwicklung"] - 1)
        fehler += _ist(f"{portfolios[k][0]} Kategorien", ist, COMDIRECT_ERWARTET)
    return fehler


def main():
    fehler = schritt1_grenzfaelle() + schritt2_echte_daten() + schritt3_artefakt()
    print()
    if fehler:
        print(f"FEHLGESCHLAGEN — {fehler} Abweichung(en)")
        return 1
    print("BESTANDEN — der Saeulen-Chart zeigt nur Kalenderjahre, die die")
    print("            Zeitreihe vollstaendig abdeckt.")
    return 0


if __name__ == "__main__":
    sys.exit(main())

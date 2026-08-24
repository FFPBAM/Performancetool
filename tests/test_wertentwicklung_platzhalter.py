"""Pruefstein: keine Folie behaelt nach dem Befuellen die Zahlen ihrer Vorlage.

DER FEHLER, GEGEN DEN DIESER TEST GEBAUT IST (24.08.2026)

Fehlt einer Strategie die Zeitreihe, waehrend ihr Bestand vorhanden ist, baut
`pptx_export` die Broschuere trotzdem — und die Wertentwicklungs-Folie behaelt
die BEISPIELDATEN DER VORLAGE. Der Titel ist zu dem Zeitpunkt schon gesetzt,
die Folie sieht deshalb bearbeitet aus, waehrend Kennzahlen-Tabelle, Saeulen-
und Liniendiagramm auf Vorlagenstand bleiben. Bis zum 24.08.2026 geschah das
ausserdem LAUTLOS: kein Eintrag in `LAST_BUILD_ERRORS`, keine Meldung.

Am 24.08.2026 nachgemessen (Datenstand 260821): comdirect Folie 7 mit
2024: 5,36 % / 2025: 6,24 %, ESG Folie 17 mit -12,91 / 5,56 / 6,91 / 7,03 % —
bytegleich mit der unveraenderten .pptx, die Nachbarfolien korrekt ersetzt.
Es faellt also nicht auf. In einem KUNDENDOKUMENT stuenden damit
Fantasiezahlen als Wertentwicklung.

DREI ZWILLINGE, NICHT EINER

Dieselbe stille Zeile stand in `_build_we_data`, `_build_perf_data` und
`_build_rollierend_data`. Die Rolle "rollierend" steckt in der Familie
*Thema*, also in einer echten Kundenbroschuere.

DIE ANZAHL DER MELDUNGEN IST DER EIGENTLICHE PRUEFSTEIN

Nicht jede Vorlage kennt jede Folie: nur `Vorlage_FFPB` fuehrt die Rolle
"performance", nur `Vorlage_Thema` die Rolle "rollierend"; comdirect, CVV,
ESG und ETF fuehren allein "anlagevorschlag" und "wertentwicklung". Eine
Warnung ueber eine Folie, die es in dieser Broschuere gar nicht gibt, waere
schlimmer als keine — wer zwei Wochen lang eine unzutreffende Zeile
ueberliest, ueberliest bald auch die zutreffende. Schritt 3 haelt deshalb die
GENAUE Zahl je Familie fest und nicht bloss "es wurde irgendetwas gemeldet".

WARUM DER TEST DIE ZAHLEN NICHT IM TEXT SUCHT

Sie stehen nicht im Text. Die Werte liegen als Cache im nativen PowerPoint-
Chart ("Diagramm links" / "Diagramm rechts"); nur die Kennzahlen-Tabelle ist
Text. Der Fingerabdruck einer Folie besteht deshalb aus drei Teilen, und erst
wenn ALLE DREI mit der Vorlage uebereinstimmen, gilt die Folie als unbefuellt.

Das ist kein Detail, sondern der Grund, warum dieser Test sich nicht mit
`tests/test_kalenderjahre.py` beisst: Hat eine Strategie kein abgeschlossenes
Kalenderjahr, bleibt ABSICHTLICH nur der Saeulen-Chart auf Vorlagenstand,
waehrend die Kennzahlen-Tabelle befuellt wird. Eine Teil-Uebereinstimmung ist
damit harmlos und wird nur als HINWEIS ausgegeben. Schritt 4 nagelt die
Trennung der beiden Meldungsarten zusaetzlich fest.

DIE FOLIEN WERDEN UEBER SHAPE-NAMEN GEFUNDEN, NICHT UEBER NUMMERN

Muster von `tests/test_quelle_position.py`: In einer gebauten Broschuere sind
die Themen-Folien vervielfaeltigt und die Positionen verschoben. Ueber die
Shapes findet der Test sie trotzdem — und nimmt die Performance-Folie gleich
mit, die denselben Fehler hatte.

    python tests/test_wertentwicklung_platzhalter.py

Braucht python-pptx, streamlit und die echten Daten. Fehlt etwas davon, wird
sauber uebersprungen statt zu scheitern. Laufzeit rund 20 Sekunden.
"""

import io
import os
import sys
import glob
import traceback

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)
sys.path.insert(0, os.path.join(WURZEL, "tests"))

# Vorlage/ und Daten/ werden im Code relativ zum Arbeitsverzeichnis geladen
os.chdir(WURZEL)

try:
    from pptx import Presentation
except ImportError:
    print("UEBERSPRUNGEN — python-pptx nicht installiert")
    sys.exit(0)

try:
    # Die Bau-Helfer stehen im Smoketest und werden bereits von drei weiteren
    # Suiten importiert — sie bilden nach, was die Oberflaeche beim Klick auf
    # "PowerPoint erstellen" tut. Sie hier zu kopieren hiesse, zwei Fassungen
    # desselben Aufbaus zu pflegen.
    import test_export_smoke as smoke
    from modules import pptx_export
    from modules.pptx_export import generate_portfolioanalyse_pptx
    from modules.portfolioanalyse import (
        VORLAGEN_FAMILIEN, FAMILIE_ALLE_STRATEGIEN, _familien_portfolios,
        _vorlage_fuer_familie, duration_info_aus_bestand,
    )
    from modules.pptx_slides import (
        SHAPE_WE_TABLE, SHAPE_WE_CHART_BAR, SHAPE_WE_CHART_LINE,
        SHAPE_TITLE, SHAPE_TITLE_ALT, WE_COL_VALUE,
    )
except ImportError as ex:
    print(f"UEBERSPRUNGEN — Abhaengigkeit fehlt: {ex}")
    sys.exit(0)

# EIGENE SYMBOLE WERDEN NICHT MIT IMPORTIERT, SONDERN NACHGESCHLAGEN (#65).
#
# Stuenden sie oben im `from ... import`, faenge derselbe `except ImportError`
# sie ab — und der Pruefstein meldete UEBERSPRUNGEN statt FEHLER, sobald
# jemand die Meldung wieder entfernt. Ein Test, der sich beim Verschwinden
# seines Gegenstands selbst stilllegt, ist keiner. Genau dieser Fall ist hier
# der wahrscheinliche: Die drei Meldungen sind erst am 24.08.2026 entstanden.
_FEHLENDE = [n for n in ("MELDUNG_KEINE_ZEITREIHE", "_melde_ohne_zeitreihe",
                         "_build_we_data", "_build_perf_data",
                         "_build_rollierend_data")
             if not hasattr(pptx_export, n)]
if _FEHLENDE:
    print("Pruefstein: Wertentwicklungs-Folien tragen keine Vorlagenzahlen\n")
    print(f"    FEHLER — modules/pptx_export.py fehlt: {_FEHLENDE}")
    print("    Ohne diese Meldungen faellt eine fehlende Zeitreihe wieder "
          "lautlos aus, und die Folie zeigt die Beispieldaten der Vorlage.")
    print("\nFEHLGESCHLAGEN — 1 Punkt(e)")
    sys.exit(1)

MELDUNG_KEINE_ZEITREIHE = pptx_export.MELDUNG_KEINE_ZEITREIHE
_build_we_data = pptx_export._build_we_data
_build_perf_data = pptx_export._build_perf_data
_build_rollierend_data = pptx_export._build_rollierend_data


# ─────────────────────────────────────────────────────────────────────────────
# Fingerabdruck einer Folie
# ─────────────────────────────────────────────────────────────────────────────

def _serien(chart_shape):
    """Kategorien und alle Serienwerte eines Charts als vergleichbarer Wert."""
    try:
        chart = chart_shape.chart
    except Exception:
        return None
    try:
        kategorien = tuple(str(k) for k in chart.plots[0].categories)
    except Exception:
        kategorien = ()
    reihen = []
    for serie in chart.series:
        # Gerundet, weil dieselbe Zahl nach dem Schreiben durch python-pptx in
        # der letzten Stelle abweichen kann. Sechs Stellen sind bei Renditen
        # weit feiner als jeder echte Unterschied.
        reihen.append(tuple(None if v is None else round(float(v), 6)
                            for v in serie.values))
    return (kategorien, tuple(reihen))


def _tabellenspalte(tabellen_shape):
    """Die Wertspalte der Kennzahlen-Tabelle als Tupel von Zeichenketten."""
    try:
        tab = tabellen_shape.table
    except Exception:
        return None
    werte = []
    for zeile in tab.rows:
        zellen = list(zeile.cells)
        werte.append(zellen[WE_COL_VALUE].text.strip()
                     if len(zellen) > WE_COL_VALUE else "")
    return tuple(werte)


def _titel(slide):
    for name in (SHAPE_TITLE, SHAPE_TITLE_ALT):
        for shape in slide.shapes:
            if shape.name == name and shape.has_text_frame:
                return shape.text_frame.text.strip()
    return ""


def _fingerabdruecke(prs):
    """Alle Folien mit Kennzahlen-Tabelle UND beiden Diagrammen.

    Genau diese Konstellation tragen die Wertentwicklungs- und die
    Performance-Folien — beide Fill-Funktionen hatten denselben stillen
    Ausfall. Ueber die Shapes statt ueber Positionsnummern, damit der Test
    auch in einer Broschuere greift, in der Folien vervielfaeltigt wurden.
    """
    raus = []
    for nr, slide in enumerate(prs.slides, start=1):
        namen = {sh.name: sh for sh in slide.shapes}
        if not {SHAPE_WE_TABLE, SHAPE_WE_CHART_BAR,
                SHAPE_WE_CHART_LINE} <= set(namen):
            continue
        raus.append({
            "nr": nr,
            "titel": _titel(slide),
            "balken": _serien(namen[SHAPE_WE_CHART_BAR]),
            "linie": _serien(namen[SHAPE_WE_CHART_LINE]),
            "tabelle": _tabellenspalte(namen[SHAPE_WE_TABLE]),
        })
    return raus


def _teile(fa):
    return (fa["balken"], fa["linie"], fa["tabelle"])


def _uebereinstimmung(gebaut, vorlagen_teile):
    """Wie viele der drei Teile stimmen mit IRGENDEINER Vorlagenfolie ueberein?

    Mengenbasiert und nicht positionsbasiert, weil `_normalisiere_vorlage` die
    Folien der Standard-Vorlage umsortiert: dort ist Vorlagenfolie 11 die
    Wertentwicklung, im Bau steht sie an Position 8.

    Returns: 3 heisst "die Folie ist unberuehrt".
    """
    beste = 0
    for teile in vorlagen_teile:
        treffer = sum(1 for a, b in zip(_teile(gebaut), teile)
                      if a is not None and a == b)
        beste = max(beste, treffer)
    return beste


# ─────────────────────────────────────────────────────────────────────────────
# Bauen
# ─────────────────────────────────────────────────────────────────────────────

def _portfolios_der_familie(familie, d):
    """Alle Strategien einer Familie — damit JEDER Folienplatz belegt ist.

    WARUM VOLLSTAENDIG: Eine Vorlage haelt feste Bloecke (comdirect drei,
    cVV fuenf). Wer nur eine Strategie uebergibt, laesst die uebrigen Bloecke
    unberuehrt — sie tragen dann voellig zu Recht die Vorlagenzahlen und
    wuerden diesen Test grundlos rot faerben.
    """
    alle = FAMILIE_ALLE_STRATEGIEN.get(familie)
    if alle:
        portfolios, fehlend = _familien_portfolios(
            alle, d["namen"], d["d2c"], d["pf_data"], duration_info_aus_bestand)
        return (None if fehlend else portfolios), fehlend
    name = next((n for n in d["namen"] if n in d["d2c"]), None)
    if name is None:
        return None, ["keine Strategie in den Daten"]
    return [smoke._portfolio(name, d)], []


def _bauen(portfolios, familie, d, nullen=()):
    """Baut IM SPEICHER und gibt (Praesentation, Meldungen) zurueck.

    Kein Schreiben auf die Platte: Die Broschueren sind 4-17 MB, und der Test
    will sie nur ansehen.

    Args:
        nullen: Indizes, deren "timeseries_df" auf None gesetzt wird. So wird
            der Ausfall injiziert, den die aktuellen Daten nicht mehr hergeben
            (19 von 19 Strategien sind vollstaendig).
    """
    pi = smoke._perf_inputs(portfolios, d, familie)
    for i in nullen:
        pi[i]["timeseries_df"] = None
    tpl, cfg = _vorlage_fuer_familie(familie)
    pptx_export.LAST_BUILD_ERRORS.clear()
    daten = generate_portfolioanalyse_pptx(
        portfolios, 0.0, performance_inputs=pi,
        template_path=tpl, template_config=cfg)
    return Presentation(io.BytesIO(daten)), list(pptx_export.LAST_BUILD_ERRORS)


# ─────────────────────────────────────────────────────────────────────────────
# Schritte
# ─────────────────────────────────────────────────────────────────────────────

def schritt1_vorlagen(vorlagen):
    print("Schritt 1 — die Fingerabdruecke der Vorlagen sind lesbar")
    f = 0
    for pfad in sorted(glob.glob(os.path.join("Vorlage", "*.pptx"))):
        try:
            fa = _fingerabdruecke(Presentation(pfad))
        except Exception as ex:
            print(f"    FEHLER — {os.path.basename(pfad)}: {ex}")
            f += 1
            continue
        vorlagen[os.path.basename(pfad)] = [_teile(x) for x in fa]
        if not fa:
            print(f"    FEHLER — {os.path.basename(pfad)}: keine Folie mit "
                  "Tabelle und beiden Diagrammen gefunden. Sind die "
                  "Shape-Namen umbenannt worden?")
            f += 1
            continue
        print(f"    OK — {os.path.basename(pfad):28s} "
              f"{len(fa)} Folie(n): {[x['nr'] for x in fa]}")
    if not vorlagen:
        print("    FEHLER — keine Vorlage gelesen")
        f += 1
    return f


def schritt2_regulaer(d, vorlagen):
    """Der Normalfall: vollstaendige Daten, alles muss befuellt sein.

    Zugleich die Gegenprobe zu Schritt 3: Ohne ihn wuesste man nicht, ob der
    Fingerabdruck-Vergleich ueberhaupt jemals "befuellt" erkennt (#64).
    """
    print("Schritt 2 — nach dem Befuellen traegt keine Folie ihre "
          "Vorlagenzahlen")
    f = 0
    for familie in sorted(VORLAGEN_FAMILIEN) + [""]:
        etikett = familie or "(Standard FFPB)"
        portfolios, fehlend = _portfolios_der_familie(familie, d)
        if portfolios is None:
            print(f"    UEBERSPRUNGEN — {etikett}: {', '.join(fehlend)}")
            continue
        tpl, _cfg = _vorlage_fuer_familie(familie)
        schluessel = os.path.basename(tpl) if tpl else "Vorlage_FFPB.pptx"
        teile = vorlagen.get(schluessel)
        if not teile:
            print(f"    UEBERSPRUNGEN — {etikett}: {schluessel} nicht gelesen")
            continue
        try:
            prs, meldungen = _bauen(portfolios, familie, d)
        except Exception as ex:
            print(f"    FEHLER — {etikett}: {type(ex).__name__}: {ex}")
            traceback.print_exc()
            f += 1
            continue

        unberuehrt, teilweise = [], []
        for fa in _fingerabdruecke(prs):
            treffer = _uebereinstimmung(fa, teile)
            if treffer == 3:
                unberuehrt.append(fa)
            elif treffer:
                teilweise.append((fa, treffer))

        if unberuehrt:
            f += 1
            print(f"    FEHLER — {etikett}: {len(unberuehrt)} Folie(n) tragen "
                  "unveraendert die Zahlen der Vorlage:")
            for fa in unberuehrt:
                print(f"             Folie {fa['nr']} — {fa['titel'][:60]!r}")
        elif meldungen:
            f += 1
            print(f"    FEHLER — {etikett}: Build-Meldungen, obwohl die Daten "
                  "vollstaendig sein sollten:")
            for m in meldungen:
                print(f"             {m[:110]}")
        else:
            zusatz = (f", {len(teilweise)} teilweise (erlaubt)"
                      if teilweise else "")
            print(f"    OK — {etikett:16s} {len(portfolios)} Strategie(n), "
                  f"alle Folien befuellt{zusatz}")
        for fa, treffer in teilweise:
            # Kein Fehler: So sieht der Leerjahr-Fall aus (Saeulen-Chart auf
            # Vorlagenstand, Tabelle befuellt) — er wird gesondert gemeldet
            # und ist per tests/test_kalenderjahre.py abgedeckt.
            print(f"        HINWEIS — Folie {fa['nr']} stimmt in {treffer} von "
                  f"3 Teilen mit der Vorlage ueberein")
    return f


# Was JE FAMILIE gemeldet werden muss, wenn eine Zeitreihe fehlt — abgeleitet
# aus den Rollen der jeweiligen Vorlage. Die Zahl ist die eigentliche Zusage:
# nicht mehr (keine Warnung ueber eine Folie, die es nicht gibt) und nicht
# weniger (kein stiller Ausfall).
ERWARTETE_MELDUNGEN = {
    "comdirect": ["Wertentwicklungs-Folie"],
    "": ["Wertentwicklungs-Folie", "Performance-Folie"],
    "Thema": ["Wertentwicklungs-Folie", "rollierenden Renditen"],
}


def schritt3_injektion(d, vorlagen):
    """DER EIGENTLICHE FALL — ohne die Korrektur ist dieser Schritt rot."""
    print("Schritt 3 — fehlt die Zeitreihe, bleibt die Folie auf "
          "Vorlagenstand UND es wird genau passend gemeldet")
    f = 0
    for familie, erwartet in ERWARTETE_MELDUNGEN.items():
        etikett = familie or "(Standard FFPB)"
        if familie == "Thema":
            namen = [n for n in smoke.THEMA_STRATEGIEN if n in d["d2c"]][:2]
            portfolios = [smoke._portfolio(n, d) for n in namen] or None
        else:
            portfolios, _fehlend = _portfolios_der_familie(familie, d)
        if not portfolios:
            print(f"    UEBERSPRUNGEN — {etikett}: Strategien fehlen")
            continue
        tpl, _cfg = _vorlage_fuer_familie(familie)
        schluessel = os.path.basename(tpl) if tpl else "Vorlage_FFPB.pptx"
        teile = vorlagen.get(schluessel)
        if not teile:
            print(f"    UEBERSPRUNGEN — {etikett}: {schluessel} nicht gelesen")
            continue

        # Bei mehreren Strategien die ZWEITE nullen: So faellt zugleich auf,
        # wenn die Meldung die falsche Strategie nennt.
        opfer = 1 if len(portfolios) > 1 else 0
        betroffen = portfolios[opfer][0]
        try:
            prs, meldungen = _bauen(portfolios, familie, d, nullen=(opfer,))
        except Exception as ex:
            print(f"    FEHLER — {etikett}: {type(ex).__name__}: {ex}")
            traceback.print_exc()
            f += 1
            continue

        unberuehrt = [fa for fa in _fingerabdruecke(prs)
                      if _uebereinstimmung(fa, teile) == 3]
        eigene = [m for m in meldungen if MELDUNG_KEINE_ZEITREIHE in m]
        fremde = [m for m in meldungen if MELDUNG_KEINE_ZEITREIHE not in m]

        fehlt = [s for s in erwartet
                 if not any(s in m and betroffen in m for m in eigene)]
        zuviel = len(eigene) - len(erwartet)

        if not unberuehrt:
            # Sollte der Export eines Tages die Vorlagenzahlen leeren, ist
            # dieser Test nicht falsch, sondern ueberholt — dann bitte HIER
            # nachziehen statt die Zusicherung aufzuweichen.
            print(f"    FEHLER — {etikett}: keine Folie behaelt die "
                  "Vorlagenzahlen; baut der Export inzwischen anders?")
            f += 1
        elif fehlt:
            print(f"    FEHLER — {etikett}: es fehlt die Meldung zu "
                  f"{fehlt} fuer {betroffen!r}. Genau dieser stille Ausfall "
                  "haette Platzhalterzahlen in ein Kundendokument gelassen.")
            for m in meldungen:
                print(f"             gemeldet wurde: {m[:100]}")
            f += 1
        elif zuviel > 0:
            print(f"    FEHLER — {etikett}: {len(eigene)} Meldungen statt "
                  f"{len(erwartet)}. Eine Warnung ueber eine Folie, die es in "
                  "dieser Vorlage nicht gibt, entwertet die uebrigen.")
            for m in eigene:
                print(f"             {m[:100]}")
            f += 1
        else:
            print(f"    OK — {etikett:16s} Folie(n) "
                  f"{[fa['nr'] for fa in unberuehrt]} auf Vorlagenstand, "
                  f"{len(eigene)} Meldung(en) zu {betroffen}: {erwartet}")
        for m in fremde:
            print(f"        HINWEIS — weitere Meldung: {m[:96]}")
    return f


def schritt4_abgrenzung(d):
    """Der Leerjahr-Fall darf NICHT als fehlende Zeitreihe gemeldet werden.

    Beide Faelle lassen Teile der Folie auf Vorlagenstand und sind fachlich
    verschieden. Liefen ihre Meldungen zusammen, waere entweder dieser Test
    oder `tests/test_kalenderjahre.py` nicht mehr aussagekraeftig.
    """
    print("Schritt 4 — der Leerjahr-Fall bleibt von der fehlenden Zeitreihe "
          "getrennt")
    csv_n = next(iter(d["ts"]), None)
    if csv_n is None:
        print("    UEBERSPRUNGEN — keine Zeitreihe in den Daten")
        return 0
    voll = d["ts"][csv_n]
    jahr = voll.index.max().year
    nur_laufendes = voll[voll.index >= f"{jahr}-01-01"]
    if nur_laufendes.empty:
        print("    UEBERSPRUNGEN — keine Zeilen im laufenden Jahr")
        return 0

    f = 0
    pptx_export.LAST_BUILD_ERRORS.clear()
    we = _build_we_data([{"timeseries_df": nur_laufendes, "fee_dec": 0.0}], 0,
                        "Testfall")
    meldungen = list(pptx_export.LAST_BUILD_ERRORS)
    if we is None:
        print("    FEHLER — ohne abgeschlossenes Kalenderjahr kam None; die "
              "Kennzahlen-Tabelle wuerde damit gar nicht befuellt")
        f += 1
    elif any(MELDUNG_KEINE_ZEITREIHE in m for m in meldungen):
        print("    FEHLER — der Leerjahr-Fall wird als fehlende Zeitreihe "
              "gemeldet; die beiden Faelle sind nicht mehr unterscheidbar")
        f += 1
    else:
        leerjahr = [m for m in meldungen if "Kalenderjahr" in m]
        print(f"    OK — Daten ab {nur_laufendes.index.min():%d.%m.%Y}: "
              f"Kennzahlen berechnet, {len(leerjahr)} Leerjahr-Meldung(en), "
              "keine zur fehlenden Zeitreihe")
    return f


def schritt5_grenzfaelle():
    """Leere Liste, None, Eintraege ohne Reihe — nichts davon darf werfen."""
    print("Schritt 5 — Grenzfaelle der performance_inputs")
    import pandas as pd
    f = 0
    faelle = [
        # (Etikett, performance_inputs, idx, erwartete Meldungen bei melden=True)
        ("None", None, 0, 0),
        ("leere Liste", [], 0, 0),
        ("idx ausserhalb", [{"timeseries_df": pd.DataFrame()}], 5, 0),
        ("Eintrag None", [None], 0, 1),
        ("leeres Dict", [{}], 0, 1),
        ("leerer DataFrame", [{"timeseries_df": pd.DataFrame()}], 0, 1),
    ]
    for etikett, pi, idx, soll in faelle:
        for name, funktion in (("we", _build_we_data),
                               ("perf", _build_perf_data),
                               ("roll", _build_rollierend_data)):
            try:
                pptx_export.LAST_BUILD_ERRORS.clear()
                ergebnis = funktion(pi, idx, "Testfall")
                ist = len(pptx_export.LAST_BUILD_ERRORS)
                if ergebnis is not None:
                    print(f"    FEHLER — {etikett}/{name}: Rueckgabe ist nicht "
                          "None")
                    f += 1
                elif ist != soll:
                    print(f"    FEHLER — {etikett}/{name}: {ist} Meldung(en) "
                          f"statt {soll}")
                    f += 1
            except Exception as ex:
                print(f"    FEHLER — {etikett}/{name}: {type(ex).__name__}: "
                      f"{ex}")
                f += 1
        # Und das Gegenstueck: melden=False schweigt in jedem dieser Faelle.
        try:
            pptx_export.LAST_BUILD_ERRORS.clear()
            _build_we_data(pi, idx, "Testfall", melden=False)
            if pptx_export.LAST_BUILD_ERRORS:
                print(f"    FEHLER — {etikett}: melden=False meldet trotzdem")
                f += 1
        except Exception as ex:
            print(f"    FEHLER — {etikett}/melden=False: "
                  f"{type(ex).__name__}: {ex}")
            f += 1
        if not f:
            print(f"    OK — {etikett:18s} alle drei Helfer liefern None, "
                  f"{soll} Meldung(en), melden=False schweigt")
    return f


def main():
    print("Pruefstein: Wertentwicklungs-Folien tragen keine Vorlagenzahlen\n")
    vorlagen = {}
    fehler = schritt1_vorlagen(vorlagen)
    print()
    fehler += schritt5_grenzfaelle()
    print()
    try:
        d = smoke._daten()
    except Exception as ex:
        print(f"UEBERSPRUNGEN — die echten Daten fehlen: "
              f"{type(ex).__name__}: {ex}")
        return 1 if fehler else 0
    print(f"Datenstand {d['tag']} — {len(d['pf_data'])} Bestaende, "
          f"{len(d['ts'])} Zeitreihen\n")
    fehler += schritt2_regulaer(d, vorlagen)
    print()
    fehler += schritt3_injektion(d, vorlagen)
    print()
    fehler += schritt4_abgrenzung(d)

    print()
    if fehler:
        print(f"FEHLGESCHLAGEN — {fehler} Punkt(e)")
        return 1
    print("BESTANDEN — keine Folie behaelt die Zahlen ihrer Vorlage, und wo "
          "die Zeitreihe fehlt, wird genau die betroffene Folie gemeldet")
    return 0


if __name__ == "__main__":
    sys.exit(main())

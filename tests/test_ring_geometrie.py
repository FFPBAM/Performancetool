"""Pruefstein fuer die GEOMETRIE der Ringdiagramme (NEU 25.08.2026).

Bis heute gab es dafuer keinen einzigen Test. `test_farben.py` prueft die
SEGMENTFARBEN, `test_export_smoke.py` baut die Broschueren durch — aber
Ringgroesse, Ringdicke, Label-Abstaende und Fuehrungslinien pruefte nichts.
Genau diese Werte sind das Ergebnis von fuenf Iterations-Wellen (Ende Juni,
09.07., 20.07., 27.-28.07., 10.08.2026) und wurden am 10.08.2026 als
Endzustand abgenommen ("wir sind am Zenit angekommen", Transferwissen #44).
Ein unbeabsichtigter Eingriff faellt sonst erst in einer Kundenbroschuere auf.

  1. Die VORLAGEN — Rahmen, `holeSize`, `plotArea` aller 22 Doughnut-Charts
  2. Die IST-Geometrie NACH `chart_dynamik.nachbearbeiten` (Durchmesser,
     Bandstaerke, Ringmitte) gegen eingefrorene Sollwerte
  3. Die ZUSICHERUNGEN, die unabhaengig von den konkreten Zahlen gelten
  4. Dieselben Messungen an ECHT GEBAUTEN Broschueren (nur mit Ordner-Argument)
     sowie: jede Fuehrungslinie traegt ihren Punkt am Label-Ende
  5. Der Familien-Look und wo Rueckkopplung und Seitentreue laufen
  6. Die Seite der Beschriftungen an ECHT GEBAUTEN Broschueren
  7. Die buendige Anbindung der Zahl an ihre Fuehrungslinie
  8. Steile Linien setzen an der Unterkante an, nicht seitlich

SCHRITT 1 HAENGT AM ARTEFAKT. Die Ringgeometrie kommt zu 100 Prozent aus der
.pptx-Vorlage — kein Code setzt jemals die Groesse eines Chart-Rahmens. Wer
eine Vorlage austauscht, verschiebt damit still jeden Ring.

SCHRITT 2 IST DER EIGENTLICHE WAECHTER. Die Ringgroesse entsteht in genau
einer Funktion: `ring_labels_aussen_dynamisch` (`modules/chart_dynamik.py`).
Seit dem 25.08.2026 sucht sie die groesste Groesse, bei der die Labels noch
kollisionsfrei liegen, statt pauschal Platz freizuhalten. Wer an `kopf_rand`,
am Deckel, an den Tabuflaechen oder an der Legende dreht, aendert JEDEN Ring —
dieser Schritt macht es sichtbar.

SCHRITT 3 UEBERLEBT EINE BEWUSSTE AENDERUNG. Die Zahlen aus Schritt 2 sind
dann neu einzufrieren; die Zusagen hier muessen weiter gelten.

Gemessen wird im Arbeitsspeicher — Schritt 1 bis 3 und 5 schreiben keine
Datei. Alle Laengen in Zoll (1 Zoll = 914400 EMU), Durchmesser zusaetzlich
in cm.

    python tests/test_ring_geometrie.py [ausgabeordner]

Ohne Ordner-Argument entfaellt Schritt 4. Die Optik selbst ist hiermit NICHT
geprueft — dafuer die Folien in ECHTEM PowerPoint ansehen (LibreOffice
ignoriert `holeSize` und zeichnet eigene Fuehrungslinien, #29).
"""

import math
import os
import sys

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)
os.chdir(WURZEL)

try:
    from pptx import Presentation
except ImportError:
    print("UEBERSPRUNGEN — python-pptx nicht installiert")
    sys.exit(0)

from modules import chart_dynamik as cd  # noqa: E402

NS_CHART = "http://schemas.openxmlformats.org/drawingml/2006/chart"
VORLAGEN_ORDNER = "Vorlage"

# Halbe Ausdehnung einer Prozent-Beschriftung in Zoll. Dieselben Werte, mit
# denen `chart_dynamik` die Labels setzt (dort HALB_W / HALB_H, #29).
HALB_BREITE, HALB_HOEHE = 0.33, 0.10

# Ueberlappungskriterium aus der Label-Entzerrung (Transferwissen #44)
UEBERLAPP_X, UEBERLAPP_Y = 0.55, 0.19

# Zoll. 0.005" = 0,13 mm — enger als jede sichtbare Verschiebung, weit genug
# fuer Rundungen beim Schreiben der f"{val:.5f}"-Bruchteile.
TOLERANZ = 0.005

# Kreuzende Fuehrungslinien auf den PLATZHALTERDATEN der Vorlagen. Kein
# Wunschwert, sondern der Ausgangswert: am 25.08.2026 gegen den Stand vom
# 24.08.2026 gemessen (LEGENDE_SPALTENWEISE=False) — zwei Stueck, in
# Vorlage_ETF F16 und Vorlage_comdirect F6. In echten Broschueren sind es
# null; deshalb steht hier eine Obergrenze und in Schritt 4 die Null.
#
# DIE ZAHL IST AM 26.08.2026 ZWEIMAL GEWANDERT, und beide Male aus gutem Grund:
#   Ausgangswert 25.08.  2 (Vorlage_ETF F16, Vorlage_comdirect F6)
#   mittags              4 — die Seitentreue kostete zwei zusaetzliche
#   abends               2 — Seitentreue zurueckgenommen, Preis mit zurueck
#   danach               1 — Seitentreue wieder ein, ABER mit der buendigen
#                        Anbindung aus Punkt 1: jetzt kreuzt sich nur noch EIN
#                        Paar (Vorlage_FFPB F9 C_Kennzahlen2).
#
# Die 1 steht hier, weil sie GEMESSEN ist und nicht gewuenscht. Die
# Vorlagendaten sind fest — der Wert schwankt nicht mit dem Bestand, also darf
# die Grenze eng sein. Wer eine Aenderung baut, die wieder auf 2 geht, soll
# das begruenden muessen.
#
# WER SIE ANHEBT, FUEHRT DEN NACHWEIS MIT: Schritt 4 verlangt die harte NULL
# an echt gebauten Broschueren (17 Ringe, 68 Linien) — sonst wird aus einer
# Obergrenze eine Ausrede.
KREUZUNGEN_VORLAGEN_MAX = 1

# Namentlich geduldete Beschriftungen auf der Legende, auf den PLATZHALTERDATEN
# der Vorlagen. Ein Eintrag, seit die Seitentreue laeuft (26.08.2026).
#
# Warum geduldet: Der Fall tritt mit den BEISPIELZAHLEN der Vorlage auf, nicht
# mit echten Daten — Schritt 4 misst dieselbe Zusage an echt gebauten
# Broschueren und meldet dort null Verletzungen. Und `Vorlage_FFPB.pptx` wird
# zu keiner Kundenbroschuere gebaut; die fuenf gebauten Familien benutzen
# andere Vorlagen (CVV etwa `Vorlage_cVV_Infoboard.pptx`).
#
# KORREKTUR EINER FRUEHEREN BEHAUPTUNG (26.08.2026): In STATUS.md stand, diese
# Vorlage sei "der Standard-Pfad, die Familie OHNE Eintrag in
# FAMILIE_RING_FORMAT". Das ist falsch — `_familie_aus_prs` erkennt sie an
# ihren Folientiteln als **CVV** und gibt ihr damit die volle CVV-Optik.
# Nachgemessen mit `cd._familie_aus_prs(Presentation("Vorlage/Vorlage_FFPB.pptx"))`.
# Der Fall ist also nicht deshalb harmlos, weil kein Familien-Format greift,
# sondern weil aus dieser Vorlage keine Broschuere entsteht.
#
# Warum als NAMENTLICHE Liste und nicht als abgeschaltete Pruefung: So bleibt
# die Zusage fuer JEDEN anderen Ring hart, der Fall bleibt sichtbar (er wird
# als HINWEIS gedruckt), und wer die Vorlage anfasst, bekommt sofort eine
# Meldung an anderer Stelle.
LEGENDE_GEDULDET = {("Vorlage_FFPB.pptx", 9, "C_Kennzahlen2")}

# SEITENTREUE (NEU 26.08.2026): Hoechstzahl an Fuehrungslinien, deren inneres
# Ende (am Segment) und aeusseres Ende (am Label) auf VERSCHIEDENEN Seiten der
# senkrechten Ringachse liegen. So eine Linie laeuft quer ueber den Ringkopf —
# das ist der gemeldete Eindruck "die Zahl steht neben ihrer Linie statt an
# ihr" (#44 Ansatzpunkt 2).
#
# Bewusst je Familie und als ZAHL, nicht als Quote: Eine Zahl je Familie zeigt,
# WO es schlechter wird; eine Gesamtquote laesst sich wegmitteln.
#
# DIE VIER FAMILIEN MIT ANLAGESTRATEGIE-FOLIEN HALTEN DIE HARTE NULL, Thema
# steht auf 3. Das ist der Stand nach der dritten Sichtpruefung (26.08.2026):
# Die Seitentreue laeuft dort, wo sie angesehen und fuer besser befunden wurde,
# und ist bei Thema ausdruecklich aus.
#
# Eine NULL ist hier mehr wert als eine Quote: Sie bricht sofort und sichtbar,
# waehrend sich eine Quote wegmitteln laesst. Bei Thema steht die 3, weil dort
# kein Schalter laeuft — sie ist eine Messung mit Deckel, kein Zielwert.
SEITENTREUE_MAX = {"CVV": 0, "ESG": 0, "ETF": 0, "comdirect": 0, "Thema": 3}

# Naeher als das an der Senkrechten ist "die Seite" keine sinnvolle Aussage.
SEITE_TOTZONE = 0.02

# Was in den sechs Vorlagen steht — Rahmen (Breite, Hoehe) in Zoll,
# `plotArea`-manualLayout (x, y, w, h) als Bruchteile des Rahmens und die
# `holeSize` der Vorlage. Maschinell aus den .pptx erzeugt, nicht abgetippt.
VORLAGEN_SOLL = {
    ("Vorlage_ESG.pptx", 16, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_ESG.pptx", 18, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_ESG.pptx", 20, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_ESG.pptx", 22, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_ETF.pptx", 16, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_ETF.pptx", 18, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 7, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 9, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 11, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 13, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 15, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_comdirect.pptx", 6, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_comdirect.pptx", 8, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_comdirect.pptx", 10, "C_Kennzahlen"): dict(rahmen=(4.3307, 3.3858), plot=(0.19333, 0.15984, 0.64409, 0.62135), hole=79),
    ("Vorlage_FFPB.pptx", 7, "C_Kennzahlen"): dict(rahmen=(4.4488, 4.5472), plot=(0.23959, 0.15984, 0.53916, 0.67915), hole=79),
    ("Vorlage_FFPB.pptx", 8, "C_Kennzahlen"): dict(rahmen=(4.4488, 4.5472), plot=(0.23959, 0.15984, 0.53916, 0.67915), hole=79),
    ("Vorlage_FFPB.pptx", 9, "C_Kennzahlen2"): dict(rahmen=(5.2362, 4.5472), plot=(0.49241, 0.22125, 0.34924, 0.44898), hole=79),
    ("Vorlage_FFPB.pptx", 9, "C_Kennzahlen1"): dict(rahmen=(5.2362, 4.5472), plot=(0.49531, 0.17508, 0.42582, 0.56419), hole=79),
    ("Vorlage_FFPB.pptx", 12, "C_Kennzahlen1"): dict(rahmen=(5.2362, 4.5472), plot=(0.40057, 0.21052, 0.41579, 0.48958), hole=79),
    ("Vorlage_Thema.pptx", 10, "C_Kennzahlen"): dict(rahmen=(5.2362, 4.5472), plot=(0.19333, 0.15984, 0.63988, 0.73742), hole=79),
    ("Vorlage_Thema.pptx", 11, "C_Kennzahlen2"): dict(rahmen=(5.2362, 4.5472), plot=(0.41742, 0.18658, 0.41579, 0.48958), hole=55),
    ("Vorlage_Thema.pptx", 11, "C_Kennzahlen1"): dict(rahmen=(5.2362, 4.5472), plot=(0.27607, 0.12133, 0.64409, 0.62135), hole=55),
}


# Was NACH `nachbearbeiten` herauskommt: Aussendurchmesser `d` und Ringmitte
# (cx, cy) in Zoll, dazu die gesetzte `holeSize`. Gemessen an den Vorlagen,
# damit der Schritt ohne Bestandsdaten laeuft; Schritt 4 wiederholt dieselbe
# Messung an echten Broschueren.
#
# ACHTUNG bei Vorlage_FFPB: Deren Platzhalter-Folientitel nennen eine
# cVV-Strategie, `_familie_aus_prs` liest daraus 'CVV'. Die FFPB-Ringe laufen
# hier also MIT Rueckkopplung, obwohl eine echte Standard-Broschuere ohne
# laeuft. Der Standard-Pfad wird von Schritt 4 nicht beruehrt (dort wird keine
# Broschuere ohne Familie gebaut) — festgenagelt ist er in Schritt 5.
NACH_SOLL = {
    ("Vorlage_ESG.pptx", 16, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_ESG.pptx", 18, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_ESG.pptx", 20, "C_Kennzahlen"): dict(d=1.9838, cx=2.2319, cy=1.5729, hole=79),
    ("Vorlage_ESG.pptx", 22, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_ETF.pptx", 16, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_ETF.pptx", 18, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 7, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 9, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 11, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 13, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_cVV_Infoboard.pptx", 15, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_comdirect.pptx", 6, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_comdirect.pptx", 8, "C_Kennzahlen"): dict(d=2.0638, cx=2.2319, cy=1.5729, hole=79),
    ("Vorlage_comdirect.pptx", 10, "C_Kennzahlen"): dict(d=2.1038, cx=2.2319, cy=1.6929, hole=79),
    ("Vorlage_FFPB.pptx", 7, "C_Kennzahlen"): dict(d=2.3986, cx=2.2652, cy=2.2736, hole=79),
    ("Vorlage_FFPB.pptx", 8, "C_Kennzahlen"): dict(d=2.3986, cx=2.2652, cy=2.2736, hole=79),
    ("Vorlage_FFPB.pptx", 9, "C_Kennzahlen2"): dict(d=0.8181, cx=3.4927, cy=1.1035, hole=79),
    ("Vorlage_FFPB.pptx", 9, "C_Kennzahlen1"): dict(d=1.3097, cx=3.7084, cy=2.2736, hole=79),
    ("Vorlage_FFPB.pptx", 12, "C_Kennzahlen1"): dict(d=2.1772, cx=3.1861, cy=2.2736, hole=79),
    ("Vorlage_Thema.pptx", 10, "C_Kennzahlen"): dict(d=3.3506, cx=2.6876, cy=2.2736, hole=79),
    ("Vorlage_Thema.pptx", 11, "C_Kennzahlen2"): dict(d=1.9771, cx=3.2743, cy=2.0936, hole=79),
    ("Vorlage_Thema.pptx", 11, "C_Kennzahlen1"): dict(d=2.5455, cx=3.1318, cy=2.2736, hole=79),
}


# ─────────────────────────────────────────────────────────────────────────────
# Messen

def _q(tag):
    return "{%s}%s" % (NS_CHART, tag)


def _ringe(prs):
    """(Folie 1-basiert, Folie, Shape) je Doughnut-Chart der Praesentation."""
    for nr, folie in enumerate(prs.slides, 1):
        for shape in folie.shapes:
            if shape.has_chart and "DOUGHNUT" in shape.chart.chart_type.name:
                yield nr, folie, shape


def _plot(shape):
    """Ring-Geometrie in Zoll aus dem Chart-XML. None, wenn kein manuelles
    plotArea-Layout dasteht (kommt in unseren Vorlagen nicht vor)."""
    wurzel = shape.chart._chartSpace
    pa = wurzel.find(".//" + _q("plotArea") + "/" + _q("layout")
                     + "/" + _q("manualLayout"))
    if pa is None:
        return None
    werte = {}
    for tag in ("x", "y", "w", "h"):
        el = pa.find(_q(tag))
        if el is None:
            return None
        werte[tag] = float(el.get("val"))
    loch = wurzel.find(".//" + _q("holeSize"))
    rahmen_b = shape.width / 914400.0
    rahmen_h = shape.height / 914400.0
    links = werte["x"] * rahmen_b
    rechts = (werte["x"] + werte["w"]) * rahmen_b
    oben = werte["y"] * rahmen_h
    unten = (werte["y"] + werte["h"]) * rahmen_h
    return {
        "rahmen_b": rahmen_b, "rahmen_h": rahmen_h,
        "roh": tuple(round(werte[t], 5) for t in ("x", "y", "w", "h")),
        "cx": (links + rechts) / 2.0, "cy": (oben + unten) / 2.0,
        "r": min(rechts - links, unten - oben) / 2.0,
        "hole": int(loch.get("val")) if loch is not None else None,
    }


def _labels(shape, rahmen_b, rahmen_h):
    """Mittelpunkte der Prozent-Beschriftungen in Zoll.

    `manualLayout` steht bei uns ABSOLUT (xMode/yMode="edge"), abgelegt wird
    die linke obere Ecke der Textbox: stored_x = (mx - 0.33) / rahmen_b.
    Hier wird zurueckgerechnet (#29).
    """
    raus = []
    for dl in shape.chart._chartSpace.iter(_q("dLbl")):
        ml = dl.find(".//" + _q("manualLayout"))
        if ml is None:
            continue
        ex, ey = ml.find(_q("x")), ml.find(_q("y"))
        if ex is None or ey is None:
            continue
        raus.append((float(ex.get("val")) * rahmen_b + HALB_BREITE,
                     float(ey.get("val")) * rahmen_h + HALB_HOEHE))
    return raus


def _legende_box(shape):
    """Die Legende als RECHTECK (links, oben, rechts, unten) in Zoll.

    Bewusst kein blosser Vergleich mit der Oberkante: Seit dem 25.08.2026
    (#71) darf ein Ring TIEFER reichen als die Legendenoberkante, solange er
    rechts neben der Legende steht. Geprueft gehoert deshalb die Flaeche, nicht
    die Kante — sonst prueft man eine Regel, die es nicht mehr gibt.
    """
    lg = shape.chart._chartSpace.find(".//" + _q("legend"))
    if lg is None:
        return None
    ml = lg.find(".//" + _q("manualLayout"))
    if ml is None or ml.find(_q("y")) is None:
        return None

    def wert(tag, ersatz):
        el = ml.find(_q(tag))
        return float(el.get("val")) if el is not None else ersatz

    breite = shape.width / 914400.0
    hoehe = shape.height / 914400.0
    lx, ly = wert("x", 0.0), wert("y", 0.0)
    lw, lh = wert("w", 1.0), wert("h", 1.0 - ly)
    return (lx * breite, ly * hoehe, (lx + lw) * breite, (ly + lh) * hoehe)


def _abstand_zu_rechteck(px, py, box):
    """Kuerzester Abstand des Punktes (px, py) zum Rechteck. 0 innerhalb."""
    links, oben, rechts, unten = box
    dx = max(links - px, 0.0, px - rechts)
    dy = max(oben - py, 0.0, py - unten)
    return math.hypot(dx, dy)


def _zeichnungsobjekte(shape):
    """Die Zeichnungsobjekte IM CHART als (links, oben, rechts, unten, Text),
    in Zoll ab der linken oberen Rahmenecke.

    Sie stehen nicht auf der Folie, sondern als `cdr:relSizeAnchor` in
    `ppt/drawings/drawingN.xml` des Chart-Teils — dort sitzen der
    Ueberschriftenbalken UND die Quellenangabe. `chart_dynamik` liest von hier
    bereits den Balken (`kopf_sperre_aus_usershapes`), aber nur den; fuer das
    untere Objekt gab es bis zum 25.08.2026 keine Entsprechung.

    Die Anker sind Bruchteile des Rahmens, nicht EMU.
    """
    CDR = "{http://schemas.openxmlformats.org/drawingml/2006/chartDrawing}"
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    breite = shape.width / 914400.0
    hoehe = shape.height / 914400.0
    raus = []
    try:
        from lxml import etree
        for rel in shape.chart.part.rels.values():
            if "chartUserShapes" not in rel.reltype:
                continue
            wurzel = etree.fromstring(rel.target_part.blob)
            for anker in wurzel:
                von, bis = anker.find(CDR + "from"), anker.find(CDR + "to")
                if von is None or bis is None:
                    continue
                fx = float(von.find(CDR + "x").text)
                fy = float(von.find(CDR + "y").text)
                tx = float(bis.find(CDR + "x").text)
                ty = float(bis.find(CDR + "y").text)
                text = " ".join(" ".join(e.text or "" for e in anker.iter(A + "t")).split())
                raus.append((fx * breite, fy * hoehe, tx * breite, ty * hoehe, text))
    except Exception:
        pass
    return raus


def _leader(folie, shape):
    """Die gezeichneten Fuehrungslinien dieses Rings als Strecken in Zoll,
    relativ zur linken oberen Ecke des Chart-Rahmens.

    `ring_leader_zeichnen` legt sie als eigene Connector-Shapes auf die Folie
    und benennt sie "RingLeader_<Shapename>_<i>" (#29) — daran haengen wir uns.
    """
    praefix = "RingLeader_%s_" % shape.name
    raus = []
    for kandidat in folie.shapes:
        if not (kandidat.name or "").startswith(praefix):
            continue
        try:
            ax, ay = kandidat.begin_x, kandidat.begin_y
            ex, ey = kandidat.end_x, kandidat.end_y
        except AttributeError:      # kein Connector (z.B. der Punkt am Ende)
            continue
        raus.append((((ax - shape.left) / 914400.0, (ay - shape.top) / 914400.0),
                     ((ex - shape.left) / 914400.0, (ey - shape.top) / 914400.0)))
    return raus


def _kreuzen(s1, s2):
    """Schneiden sich zwei Strecken? Beruehrungen an den Enden zaehlen nicht."""
    def richtung(a, b, c):
        return ((b[0] - a[0]) * (c[1] - a[1])
                - (b[1] - a[1]) * (c[0] - a[0]))
    (a, b), (c, d) = s1, s2
    d1, d2 = richtung(c, d, a), richtung(c, d, b)
    d3, d4 = richtung(a, b, c), richtung(a, b, d)
    # echtes Kreuzen: beide Strecken trennen die jeweils andere strikt
    return ((d1 > 0) != (d2 > 0)) and ((d3 > 0) != (d4 > 0)) \
        and min(abs(d1), abs(d2), abs(d3), abs(d4)) > 1e-9


def _abweichung(bezeichnung, ist, soll, toleranz=TOLERANZ):
    if abs(ist - soll) <= toleranz:
        return 0
    print(f"    FEHLER — {bezeichnung}: {ist:.4f} statt {soll:.4f} "
          f"(Abweichung {ist - soll:+.4f} Zoll)")
    return 1


def _vorlagen():
    """Die sechs Vorlagen, in fester Reihenfolge."""
    for name in sorted({schluessel[0] for schluessel in VORLAGEN_SOLL}):
        yield name, Presentation(os.path.join(VORLAGEN_ORDNER, name))


# ─────────────────────────────────────────────────────────────────────────────

def schritt1_vorlagen():
    print("Schritt 1 — die Vorlagen: Rahmen, holeSize, plotArea")
    fehler = 0
    gesehen = set()
    for name, prs in _vorlagen():
        for nr, _folie, shape in _ringe(prs):
            schluessel = (name, nr, shape.name)
            gesehen.add(schluessel)
            soll = VORLAGEN_SOLL.get(schluessel)
            if soll is None:
                print(f"    FEHLER — NEUER Ring ohne Sollwert: {schluessel}")
                fehler += 1
                continue
            geo = _plot(shape)
            if geo is None:
                print(f"    FEHLER — {schluessel}: kein manuelles "
                      f"plotArea-Layout (PowerPoint wuerde frei skalieren)")
                fehler += 1
                continue
            fehler += _abweichung(f"{name} F{nr} {shape.name} Rahmenbreite",
                                  geo["rahmen_b"], soll["rahmen"][0])
            fehler += _abweichung(f"{name} F{nr} {shape.name} Rahmenhoehe",
                                  geo["rahmen_h"], soll["rahmen"][1])
            if geo["roh"] != soll["plot"]:
                print(f"    FEHLER — {name} F{nr} {shape.name} plotArea: "
                      f"{geo['roh']} statt {soll['plot']}")
                fehler += 1
            if geo["hole"] != soll["hole"]:
                print(f"    FEHLER — {name} F{nr} {shape.name} holeSize: "
                      f"{geo['hole']} statt {soll['hole']}")
                fehler += 1
    for schluessel in sorted(set(VORLAGEN_SOLL) - gesehen):
        print(f"    FEHLER — Ring VERSCHWUNDEN: {schluessel}")
        fehler += 1
    if not fehler:
        anzahl_vorlagen = len({schluessel[0] for schluessel in VORLAGEN_SOLL})
        print(f"    OK — {len(VORLAGEN_SOLL)} Ringe in {anzahl_vorlagen} "
              f"Vorlagen unveraendert")
    return fehler


def schritt2_ist_geometrie():
    print("Schritt 2 — die Geometrie NACH chart_dynamik.nachbearbeiten")
    fehler = 0
    for name, prs in _vorlagen():
        cd.nachbearbeiten(prs)
        for nr, _folie, shape in _ringe(prs):
            schluessel = (name, nr, shape.name)
            soll = NACH_SOLL.get(schluessel)
            geo = _plot(shape)
            if soll is None or geo is None:
                print(f"    FEHLER — kein Sollwert oder kein Layout: "
                      f"{schluessel}")
                fehler += 1
                continue
            durchmesser = geo["r"] * 2.0
            bandstaerke = geo["r"] * (1.0 - geo["hole"] / 100.0)
            print(f"    {name:<26} F{nr:<3} {shape.name:<15} "
                  f"{durchmesser * 2.54:5.2f} cm breit, "
                  f"{bandstaerke * 2.54:4.2f} cm stark, hole {geo['hole']}")
            fehler += _abweichung(f"{name} F{nr} {shape.name} Durchmesser",
                                  durchmesser, soll["d"])
            fehler += _abweichung(f"{name} F{nr} {shape.name} Ringmitte x",
                                  geo["cx"], soll["cx"])
            fehler += _abweichung(f"{name} F{nr} {shape.name} Ringmitte y",
                                  geo["cy"], soll["cy"])
            if geo["hole"] != soll["hole"]:
                print(f"    FEHLER — {name} F{nr} {shape.name} holeSize: "
                      f"{geo['hole']} statt {soll['hole']}")
                fehler += 1
    if not fehler:
        print("    OK — jeder Ring liegt auf seinem eingefrorenen Mass")
    return fehler


def _zusicherungen(quelle, prs):
    """Die Zusagen, die unabhaengig von den konkreten Zahlen gelten muessen.

    Gibt (Verletzungen, sich kreuzende Fuehrungslinien) zurueck. Die
    Kreuzungen sind BEWUSST getrennt: Auf den PLATZHALTERDATEN der Vorlagen
    kreuzen sich auch im Stand vom 24.08.2026 zwei Linien (Vorlage_ETF F16 und
    Vorlage_comdirect F6) — eine Eigenschaft dieser kuenstlichen Daten, nicht
    des Codes. In ECHTEN Broschueren sind es vorher wie nachher null von 69
    Linien. Deshalb entscheidet der Aufrufer: Schritt 4 fordert die Null,
    Schritt 3 vergleicht gegen den gemessenen Ausgangswert.
    """
    fehler = 0
    kreuzungen = 0
    for nr, folie, shape in _ringe(prs):
        geo = _plot(shape)
        if geo is None:
            continue
        ort = f"{quelle} F{nr} {shape.name}"
        beschriftungen = _labels(shape, geo["rahmen_b"], geo["rahmen_h"])

        # (a) Keine Beschriftung ragt in den Ring. Der Text ist waagerecht:
        #     seitlich ragt die halbe BREITE zum Ring, oben/unten die halbe
        #     Hoehe (#26, Falle 6).
        #
        #     ACHTUNG BEIM ABSCHREIBEN (Fehler vom 25.08.2026, am selben Tag
        #     gefunden): #26 notiert 0,33*|sin| + 0,10*|cos| — das gilt fuer
        #     den Winkel, den `chart_dynamik` rechnet (`atan2(lvx, -lvy)`,
        #     gemessen ab der SENKRECHTEN). Hier wird ab der WAAGERECHTEN
        #     gemessen (`atan2(dy, dx)`), also gehoert die halbe Breite an den
        #     KOSINUS. Uebernommen wurden die Koeffizienten zunaechst
        #     unbesehen — dadurch verlangte der Test seitlich 0,23" zu wenig
        #     und oben/unten 0,23" zu viel.
        for mx, my in beschriftungen:
            abstand = math.hypot(mx - geo["cx"], my - geo["cy"])
            winkel = math.atan2(my - geo["cy"], mx - geo["cx"])
            ueberstand = (HALB_BREITE * abs(math.cos(winkel))
                          + HALB_HOEHE * abs(math.sin(winkel)))
            frei = abstand - ueberstand - geo["r"]
            if frei < -TOLERANZ:
                print(f"    FEHLER — {ort}: Beschriftung ragt {-frei:.3f} Zoll "
                      f"in den Ring")
                fehler += 1

        # (b) Keine zwei Beschriftungen ueberlappen
        for i in range(len(beschriftungen)):
            for j in range(i + 1, len(beschriftungen)):
                dx = abs(beschriftungen[i][0] - beschriftungen[j][0])
                dy = abs(beschriftungen[i][1] - beschriftungen[j][1])
                if dx < UEBERLAPP_X and dy < UEBERLAPP_Y:
                    print(f"    FEHLER — {ort}: zwei Beschriftungen "
                          f"ueberlappen (dx={dx:.3f}, dy={dy:.3f} Zoll)")
                    fehler += 1

        # (c) Jede Beschriftung liegt vollstaendig im Chart-Rahmen
        for mx, my in beschriftungen:
            if (mx - HALB_BREITE < -TOLERANZ
                    or mx + HALB_BREITE > geo["rahmen_b"] + TOLERANZ
                    or my - HALB_HOEHE < -TOLERANZ
                    or my + HALB_HOEHE > geo["rahmen_h"] + TOLERANZ):
                print(f"    FEHLER — {ort}: Beschriftung ({mx:.3f}, {my:.3f}) "
                      f"ragt aus dem Rahmen ({geo['rahmen_b']:.3f} x "
                      f"{geo['rahmen_h']:.3f})")
                fehler += 1

        # (d) Weder der Ring noch eine Beschriftung liegt AUF der Legende.
        #     Das ist die eigentliche Zusage — seit #71 darf der Ring tiefer
        #     reichen als die Legendenoberkante, aber niemals in ihre Flaeche.
        #     Der schmalste gemessene Abstand ist die Reserve, von der die
        #     Entscheidung fuer 4,64 cm statt 5,32 cm abhing.
        box = _legende_box(shape)
        if box is not None:
            abstand = _abstand_zu_rechteck(geo["cx"], geo["cy"], box)
            if abstand < geo["r"] - TOLERANZ:
                print(f"    FEHLER — {ort}: der Ring (Radius {geo['r']:.3f}) "
                      f"schneidet die Legende, Abstand nur {abstand:.3f} Zoll")
                fehler += 1
            for mx, my in beschriftungen:
                if (mx + HALB_BREITE > box[0] and mx - HALB_BREITE < box[2]
                        and my + HALB_HOEHE > box[1]
                        and my - HALB_HOEHE < box[3]):
                    if (quelle, nr, shape.name) in LEGENDE_GEDULDET:
                        # Namentlich geduldet, siehe LEGENDE_GEDULDET oben.
                        # Sichtbar bleiben soll es trotzdem.
                        print(f"    HINWEIS — {ort}: Beschriftung "
                              f"({mx:.3f}, {my:.3f}) liegt auf der Legende — "
                              f"namentlich geduldet (Platzhalterdaten, "
                              f"Standard-Pfad)")
                        continue
                    print(f"    FEHLER — {ort}: Beschriftung "
                          f"({mx:.3f}, {my:.3f}) liegt auf der Legende "
                          f"({box[0]:.2f}, {box[1]:.2f})-({box[2]:.2f}, "
                          f"{box[3]:.2f})")
                    fehler += 1

        # (f) Keine Beschriftung liegt auf der QUELLENANGABE (NEU 25.08.2026).
        #     Gemeldet von Philip an cVV Folie 7: "89,66 %" lag mitten auf
        #     "Quelle: Eigene Berechnung Stand: ...". Nachgemessen waren es
        #     vier Folien, nicht eine. Keine Pruefklasse deckte das ab, weil
        #     die Angabe kein Folien-Shape ist, sondern im Chart-Teil steckt —
        #     und weil unten bis dahin nur die Legende als Schranke galt. Die
        #     Quellenangabe liegt RECHTS daneben, also genau in dem Bereich,
        #     den die Spaltenregel vom 25.08.2026 freigegeben hat.
        for bl, bo, br, bu, text in _zeichnungsobjekte(shape):
            if "quelle" not in text.lower():
                continue
            for mx, my in beschriftungen:
                if (mx + HALB_BREITE > bl and mx - HALB_BREITE < br
                        and my + HALB_HOEHE > bo and my - HALB_HOEHE < bu):
                    print(f"    FEHLER — {ort}: Beschriftung "
                          f"({mx:.3f}, {my:.3f}) liegt auf der Quellenangabe "
                          f"({bl:.2f}, {bo:.2f})-({br:.2f}, {bu:.2f})")
                    fehler += 1

        # (e) Keine zwei Fuehrungslinien kreuzen sich. Am 10.08.2026 wurde
        #     dieser Wert ueber alle Broschueren mit NULL gemessen (#44) —
        #     er ist damit eine Zusage und keine Zufaelligkeit. Verschieben
        #     sich Labels, ist das die Klasse, die als erstes bricht.
        strecken = _leader(folie, shape)
        for i in range(len(strecken)):
            for j in range(i + 1, len(strecken)):
                if _kreuzen(strecken[i], strecken[j]):
                    kreuzungen += 1
                    print(f"    HINWEIS — {ort}: zwei Fuehrungslinien "
                          f"kreuzen sich")
    return fehler, kreuzungen


def schritt3_zusicherungen():
    print("Schritt 3 — die Zusicherungen (unabhaengig von den Sollwerten)")
    fehler = 0
    kreuzungen = 0
    for name, prs in _vorlagen():
        cd.nachbearbeiten(prs)
        f, k = _zusicherungen(name, prs)
        fehler += f
        kreuzungen += k
    if not fehler:
        print("    OK — keine Beschriftung im Ring, keine Ueberlappung, "
              "nichts aus dem Rahmen, kein Ring in der Legende")
    # Kreuzende Fuehrungslinien auf den PLATZHALTERDATEN der Vorlagen: der
    # Stand vom 24.08.2026 hat hier zwei (Vorlage_ETF F16, Vorlage_comdirect
    # F6), gemessen mit LEGENDE_SPALTENWEISE=False. Schlechter darf es nicht
    # werden; besser gern. Die harte Null steht in Schritt 4, wo echte Daten
    # gerechnet werden.
    if kreuzungen > KREUZUNGEN_VORLAGEN_MAX:
        print(f"    FEHLER — {kreuzungen} kreuzende Fuehrungslinien, erlaubt "
              f"sind {KREUZUNGEN_VORLAGEN_MAX} (Ausgangswert 24.08.2026)")
        fehler += kreuzungen - KREUZUNGEN_VORLAGEN_MAX
    else:
        print(f"    OK — {kreuzungen} kreuzende Fuehrungslinien auf den "
              f"Platzhalterdaten, Ausgangswert war {KREUZUNGEN_VORLAGEN_MAX}")
    return fehler


def schritt4_gebaute_broschueren(ausgabe):
    print("Schritt 4 — dieselben Zusicherungen an ECHT GEBAUTEN Broschueren")
    if not ausgabe:
        print("    UEBERSPRUNGEN — kein Ausgabeverzeichnis angegeben")
        return 0
    # Den Bau NICHT nachbauen: test_export_smoke kann das bereits.
    sys.path.insert(0, os.path.join(WURZEL, "tests"))
    try:
        import test_export_smoke as smoke
    except SystemExit:
        print("    UEBERSPRUNGEN — test_export_smoke hat abgebrochen "
              "(fehlende Abhaengigkeit oder Daten)")
        return 0
    except Exception as ex:
        print(f"    UEBERSPRUNGEN — test_export_smoke nicht ladbar: {ex}")
        return 0

    try:
        from modules.portfolioanalyse import (
            FAMILIE_ALLE_STRATEGIEN, VORLAGEN_FAMILIEN,
            _familie_fuer_strategie, _familien_portfolios,
            duration_info_aus_bestand,
        )
        daten = smoke._daten()
    except Exception as ex:
        print(f"    UEBERSPRUNGEN — Bestandsdaten nicht ladbar: {ex}")
        return 0

    os.makedirs(ausgabe, exist_ok=True)
    fehler = 0
    gemessen = 0
    for familie in sorted(VORLAGEN_FAMILIEN):
        strategie = next(
            (n for n in daten["namen"]
             if _familie_fuer_strategie(daten["nm"], n) == familie), None)
        if strategie is None:
            print(f"    UEBERSPRUNGEN — {familie}: keine Strategie in den Daten")
            continue
        alle = FAMILIE_ALLE_STRATEGIEN.get(familie)
        if alle:
            portfolios, fehlend = _familien_portfolios(
                alle, daten["namen"], daten["d2c"], daten["pf_data"],
                duration_info_aus_bestand)
            if fehlend:
                print(f"    UEBERSPRUNGEN — {familie}: fehlende Daten "
                      f"({', '.join(fehlend)})")
                continue
        else:
            portfolios = [smoke._portfolio(strategie, daten)]
        ziel, _groesse, meldungen = smoke._bauen(
            portfolios, familie, daten, ausgabe, f"ring_{familie}.pptx")
        for m in meldungen:
            print(f"    FEHLER — {familie} meldet beim Bauen: {m[:90]}")
            fehler += 1
        prs = Presentation(ziel)
        ringe = list(_ringe(prs))
        gemessen += len(ringe)
        masse = ", ".join(f"F{nr} {sh.name} {_plot(sh)['r'] * 2 * 2.54:.2f} cm"
                          for nr, _f, sh in ringe)
        print(f"    {familie:<12} {len(ringe)} Ring(e) — {masse}")
        # DIE ZWEI FAMILIENERKENNUNGEN MUESSEN SICH EINIG SEIN (25.08.2026).
        # `_familie_fuer_strategie` liest die MAPPING-TABELLE und waehlt die
        # VORLAGE; `_familie_aus_prs` liest den FOLIENTITEL und waehlt die
        # RING-OPTIK. Laufen sie auseinander, bekommt eine Broschuere die
        # richtige Vorlage und die falsche Optik — und beides sieht fuer sich
        # stimmig aus. Genau das war bei den SCHWEIZ-Strategien wochenlang so
        # (Vorlage Thema, Optik Standard), und es fiel erst auf, als sich die
        # beiden Optiken deutlich genug unterschieden.
        erkannt = cd._familie_aus_prs(prs)
        if erkannt != familie:
            print(f"    FEHLER — {familie}: die Ring-Optik erkennt "
                  f"{erkannt!r} statt {familie!r} — Vorlage und Optik "
                  f"laufen auseinander")
            fehler += 1
        f, k = _zusicherungen(familie, prs)
        # Hier gilt die harte Null: das sind die Zahlen, die beim Kunden
        # landen. Gemessen am 25.08.2026 ueber 69 Fuehrungslinien, vor und
        # nach der Vergroesserung jeweils null Kreuzungen.
        fehler += f + k

        # JEDE FUEHRUNGSLINIE TRAEGT EINEN PUNKT (NEU 26.08.2026).
        # Der Punkt am Label-Ende sagt, welche Zahl zu welcher Linie gehoert.
        # Bis zum 26.08.2026 war der REGIONEN-Ring der Themen-Broschuere als
        # einziger Ring im ganzen Produkt davon ausgenommen (Entscheidung
        # 20.07.2026) — gemessen: CVV 17/17, ESG 16/16, ETF 8/8,
        # comdirect 11/11, Thema nur 10 von 17. Philip meldete genau dort,
        # die Zuordnung Linie/Zahl sei "nicht gut ersichtlich".
        #
        # Die Zusage ist bewusst "alle" und keine Stueckzahl: Eine Stueckzahl
        # muesste bei jeder Datenaenderung nachgezogen werden und saehe dann
        # aus wie ein Messwert, waehrend sie nur eine Buchhaltung waere.
        n_lin, n_pkt = _punkte_je_linie(prs)
        if n_lin and n_pkt != n_lin:
            print(f"    FEHLER — {familie}: {n_pkt} von {n_lin} "
                  f"Fuehrungslinien tragen einen Punkt am Label-Ende")
            fehler += 1
    # Bei Thema steht der Strategiename IM Folientitel — dort kann die
    # Erkennung je Strategie anders ausfallen. Die vier anderen Familien
    # tragen statische Titel aus der Vorlage und sind mit einem Bau geprueft.
    thema_strategien = [n for n in daten["namen"]
                        if _familie_fuer_strategie(daten["nm"], n) == "Thema"]
    for name in thema_strategien:
        ziel, _g, _m = smoke._bauen([smoke._portfolio(name, daten)], "Thema",
                                    daten, ausgabe, f"famcheck_{name}.pptx")
        erkannt = cd._familie_aus_prs(Presentation(ziel))
        if erkannt != "Thema":
            print(f"    FEHLER — Thema-Strategie {name!r}: die Ring-Optik "
                  f"erkennt {erkannt!r} — sie bekaeme andere Ringe als der "
                  f"Rest ihrer Familie")
            fehler += 1
    if thema_strategien and not fehler:
        print(f"    OK — alle {len(thema_strategien)} Thema-Strategien werden "
              f"auch an ihrem Folientitel als Thema erkannt")

    if not fehler:
        print(f"    OK — {gemessen} Ringe in echten Broschueren halten "
              f"dieselben Zusagen, keine kreuzenden Fuehrungslinien, "
              f"jede Linie mit Punkt")
    return fehler


def _seitentreue(prs):
    """(Fuehrungslinien, davon auf der falschen Seite, Beispiele).

    Fuer jede gezeichnete Fuehrungslinie: Liegt ihr inneres Ende (das naeher
    an der Ringmitte, also beim Segment) auf derselben Seite der senkrechten
    Ringachse wie ihr aeusseres Ende (am Label)? Wenn nicht, laeuft sie quer
    ueber den Ringkopf.
    """
    gesamt = schief = 0
    beispiele = []
    for nr, folie, shape in _ringe(prs):
        geo = _plot(shape)
        if geo is None:
            continue
        cx, cy = geo["cx"], geo["cy"]
        for (ax, ay), (ex, ey) in _leader(folie, shape):
            innen, aussen = (((ax, ay), (ex, ey))
                             if (ax - cx) ** 2 + (ay - cy) ** 2
                             <= (ex - cx) ** 2 + (ey - cy) ** 2
                             else ((ex, ey), (ax, ay)))
            di, do = innen[0] - cx, aussen[0] - cx
            if abs(di) < SEITE_TOTZONE or abs(do) < SEITE_TOTZONE:
                continue
            gesamt += 1
            if (di >= 0) != (do >= 0):
                schief += 1
                beispiele.append(f"F{nr} {shape.name}")
    return gesamt, schief, beispiele


def _punkte_je_linie(prs):
    """(Fuehrungslinien, davon mit Punkt am Label-Ende).

    Gezaehlt werden die Shapes, die `ring_leader_zeichnen` selbst anlegt:
    `RingLeader_<chart>_<idx>_<teil>` fuer die Linien (eine Linie kann aus
    ZWEI Teilen bestehen, wenn sie geknickt ist — deshalb wird ueber
    (chart, idx) entdoppelt) und `RingLeaderDot_<chart>_<idx>` fuer den Punkt.

    DIE FOLIENNUMMER GEHOERT IN DEN SCHLUESSEL. Der Chart-Name ist nur
    INNERHALB einer Folie eindeutig: Eine CVV-Broschuere traegt auf fuenf
    Folien fuenf Ringe, die alle "C_Kennzahlen" heissen. Ohne die Folie faellt
    das auf vier Schluessel zusammen — die Zusage bliebe zufaellig richtig,
    die gemeldete Zahl waere aber falsch (4 statt 17), und ein Punkt, der auf
    genau einer Folie fehlt, verschwaende in der Vereinigung.
    """
    linien, punkte = set(), set()
    for nr, folie in enumerate(prs.slides, start=1):
        for sh in folie.shapes:
            n = sh.name or ""
            if n.startswith("RingLeaderDot_"):
                punkte.add((nr, n[len("RingLeaderDot_"):]))
            elif n.startswith("RingLeader_"):
                rest = n[len("RingLeader_"):]
                # den Teil-Index abwerfen: "<chart>_<idx>_<teil>"
                linien.add((nr, rest.rsplit("_", 1)[0]))
    return len(linien), len(linien & punkte)


def schritt6_seitentreue(ausgabe):
    print("Schritt 6 — die Beschriftungen stehen auf der Seite ihres Segments")
    if not ausgabe:
        print("    UEBERSPRUNGEN — kein Ausgabeverzeichnis angegeben")
        return 0
    fehler = 0
    ges_all = schief_all = 0
    geprueft = 0
    for familie, grenze in sorted(SEITENTREUE_MAX.items()):
        pfad = os.path.join(ausgabe, f"ring_{familie}.pptx")
        if not os.path.exists(pfad):
            print(f"    UEBERSPRUNGEN — {familie}: {os.path.basename(pfad)} "
                  f"fehlt (Schritt 4 baut sie)")
            continue
        geprueft += 1
        gesamt, schief, beispiele = _seitentreue(Presentation(pfad))
        ges_all += gesamt
        schief_all += schief
        if schief > grenze:
            print(f"    FEHLER — {familie}: {schief} von {gesamt} "
                  f"Fuehrungslinien laufen auf die falsche Seite, erlaubt "
                  f"sind {grenze} ({', '.join(sorted(set(beispiele))[:4])})")
            fehler += 1
        else:
            print(f"    {familie:<12} {schief} von {gesamt} falsch "
                  f"(erlaubt {grenze})")
    if geprueft and not fehler:
        anteil = (100.0 * schief_all / ges_all) if ges_all else 0.0
        # Die Zahl NICHT gegen die 22,4 % aus dem Sitzungsbericht halten: die
        # sind ueber ALLE gebauten Broschueren gemessen (143 Linien), hier
        # laufen nur die fuenf `ring_<Familie>.pptx` (69 Linien). Zwei
        # Grundgesamtheiten, eine Prozentzahl — genau so entstehen Vergleiche,
        # die stimmen wollen und nicht stimmen.
        print(f"    OK — {schief_all} von {ges_all} Fuehrungslinien "
              f"({anteil:.1f} %), jede Familie unter ihrer Obergrenze")
    return fehler


def _buendig_pruefen(pfad):
    """(geprueft, Meldungen) — die buendige Anbindung an einer Broschuere.

    WARUM DIESER SCHRITT UEBERHAUPT NOETIG IST: Der uebrige Pruefstein ist fuer
    Punkt 1 BLIND und war es auch vorher. Er rechnet mit Label-MITTEN, und die
    bewegen sich nicht. Der Fehler sass in der Breite der Box: Ohne `w`
    dimensioniert PowerPoint sie auf den Text, waehrend die Fuehrungslinie an
    einer gedachten Kante von 0,66 Zoll endete — links bis zu 8,6 mm hinter
    der letzten Ziffer. Ein gruener Lauf bedeutete hier also nichts.

    Geprueft wird deshalb am Artefakt, was die Anbindung traegt:
      1. `wMode`/`hMode` stehen auf "factor". Mit "edge" liest PowerPoint `w`
         als rechte KANTE, und alle Labels landen uebereinander.
      2. `w` mal Rahmenbreite ist 2 x HALB_BREITE. Damit ist die Textkante
         berechenbar, und `ring_leader_zeichnen` rechnet mit derselben Zahl.
      3. `h` ist da. Mit `w` ALLEIN rendert PowerPoint die Beschriftungen gar
         nicht mehr — im XML und ueber COM sind sie da, im Bild nicht.
      4. Die Ausrichtung passt zur Ringseite und steht an BEIDEN Stellen
         (`a:pPr` und `a:lstStyle/a:lvl1pPr`); der Absatz traegt keinen
         Textlauf, deshalb entscheidet der Listenstil.
      5. Die Innenabstaende stehen ausdruecklich im `bodyPr`.
    Die Punkte 1 und 3 hat der Renderer entschieden, nicht das Schema — sie
    stehen hier, damit das niemand ein zweites Mal herausfinden muss.
    """
    prs = Presentation(pfad)
    meldungen = []
    geprueft = 0
    for nr, _folie, shape in _ringe(prs):
        root = cd._root(shape.chart)
        pl = root.find(".//" + cd._q("plotArea") + "/" + cd._q("layout")
                       + "/" + cd._q("manualLayout"))
        if pl is None:
            continue
        mitte = (float(pl.find(cd._q("x")).get("val"))
                 + float(pl.find(cd._q("w")).get("val")) / 2.0)
        fw = shape.width / 914400.0
        for d in root.iter(cd._q("dLbl")):
            ml = d.find(".//" + cd._q("manualLayout"))
            if ml is None or ml.find(cd._q("x")) is None:
                continue
            geprueft += 1
            wo = f"F{nr} {shape.name} idx {d.find(cd._q('idx')).get('val')}"
            we = ml.find(cd._q("w"))
            he = ml.find(cd._q("h"))
            if we is None or he is None:
                meldungen.append(f"{wo}: manualLayout ohne w bzw. h — ohne "
                                 f"beide endet die Linie an einer gedachten "
                                 f"Kante (bzw. verschwindet die Zahl)")
                continue
            for tag in ("wMode", "hMode"):
                me = ml.find(cd._q(tag))
                if me is None or me.get("val") != "factor":
                    meldungen.append(f"{wo}: {tag} ist "
                                     f"{me.get('val') if me is not None else 'nicht gesetzt'}"
                                     f" statt 'factor' — PowerPoint liest den "
                                     f"Wert dann als Kante, nicht als Mass")
            breite = float(we.get("val")) * fw
            if abs(breite - 2 * HALB_BREITE) > TOLERANZ:
                meldungen.append(f"{wo}: Boxbreite {breite:.3f}\" statt "
                                 f"{2 * HALB_BREITE:.3f}\" — dann rechnet "
                                 f"ring_leader_zeichnen mit einer anderen "
                                 f"Kante als PowerPoint zeichnet")
            soll = "l" if float(ml.find(cd._q("x")).get("val")) \
                + float(we.get("val")) / 2.0 >= mitte else "r"
            txPr = d.find(cd._q("txPr"))
            if txPr is None:
                meldungen.append(f"{wo}: kein txPr — keine Ausrichtung")
                continue
            gefunden = [p.get("algn") for p in txPr.iter(cd._A + "pPr")]
            lvl = txPr.find(cd._A + "lstStyle/" + cd._A + "lvl1pPr")
            gefunden.append(lvl.get("algn") if lvl is not None else None)
            if any(g != soll for g in gefunden):
                meldungen.append(f"{wo}: Ausrichtung {gefunden} statt "
                                 f"{soll!r} an beiden Stellen")
            bodyPr = txPr.find(cd._A + "bodyPr")
            ins = str(int(round(cd.LABEL_INSET_IN * 914400)))
            if bodyPr is None or bodyPr.get("lIns") != ins \
                    or bodyPr.get("rIns") != ins:
                meldungen.append(f"{wo}: lIns/rIns nicht ausdruecklich auf "
                                 f"{ins} — dann ist die Textkante eine "
                                 f"Vorgabe und keine bekannte Zahl")
            # 6. UMBRUCH AUS. Das ist die Zeile, ohne die PowerPoint die
            #    laengeren Zahlen KUERZT ("37,1..." statt "37,13%"): Eine
            #    Prozentzahl ist ein unteilbares Wort, passt sie nicht in die
            #    feste Box, wird sie abgeschnitten. Am 26.08.2026 auf der
            #    Themen-Broschuere passiert, wo die Labels in 10 pt stehen
            #    und nicht in 9. Im XML und ueber COM stand der volle Text —
            #    gesehen hat es nur der PNG-Export.
            if bodyPr is None or bodyPr.get("wrap") != "none":
                meldungen.append(f"{wo}: wrap ist nicht 'none' — PowerPoint "
                                 f"kuerzt dann laengere Zahlen mit '...'")
            # 7. Oben/unten kein Innenabstand: die feste Boxhoehe ist knapp,
            #    und die Zeile soll auf der Boxmitte bleiben — dort rechnet
            #    `ring_leader_zeichnen` das Linienende hin.
            if bodyPr is None or bodyPr.get("tIns") != "0" \
                    or bodyPr.get("bIns") != "0":
                meldungen.append(f"{wo}: tIns/bIns nicht auf 0 — die Zeile "
                                 f"sitzt dann nicht auf der Boxmitte")
    return geprueft, meldungen


def _senkrecht_pruefen(pfad):
    """(geprueft, senkrecht, Meldungen) — die senkrechte Anbindung.

    Geprueft wird die REGEL, nicht eine Stueckzahl: Eine Fuehrungslinie setzt
    GENAU DANN senkrecht an (Ende mittig unter/ueber der Zahl), wenn sie
    mindestens `LEADER_SENKRECHT_STEIL` mal so steil wie breit ankommt.

    Warum die Regel und nicht die Zahl: Wie viele Labels betroffen sind, haengt
    am Bestand — heute sind es zwei (Thema F11, beide Ringe, je das Label ueber
    dem Ring). Eine feste Zahl muesste bei jeder Datenaenderung nachgezogen
    werden und saehe dann aus wie ein Messwert.

    Die Schwelle selbst ist gemessen: Median der Steilheit ueber alle 69
    Fuehrungslinien 0,23; die beiden beanstandeten Labels 13,1 und 11,3; der
    naechsthoehere Wert im Produkt 3,3. Ein erster Versuch mit "steiler als
    45 Grad" traf zwoelf Labels und machte diagonale wie CVV F13 SCHLECHTER.
    """
    prs = Presentation(pfad)
    meldungen = []
    geprueft = senkrecht = 0
    for nr, folie, shape in _ringe(prs):
        geo = _plot(shape)
        if geo is None:
            continue
        cx, cy = geo["cx"], geo["cy"]
        root = cd._root(shape.chart)
        fsa_el = root.find(".//" + cd._q("firstSliceAng"))
        fsa = float(fsa_el.get("val")) if fsa_el is not None else 0.0
        vals = [float(v.text) for v in root.findall(
            ".//" + cd._q("val") + "//" + cd._q("pt") + "/" + cd._q("v"))]
        if not vals:
            continue
        tot = sum(vals) or 1.0
        mids, kum = [], 0.0
        for v in vals:
            f = v / tot
            mids.append((fsa + (kum + f / 2) * 360) % 360)
            kum += f
        fw, fh = shape.width / 914400.0, shape.height / 914400.0
        boxen = {}
        for d in root.iter(cd._q("dLbl")):
            ml = d.find(".//" + cd._q("manualLayout"))
            if ml is None or ml.find(cd._q("x")) is None:
                continue
            boxen[int(d.find(cd._q("idx")).get("val"))] = (
                float(ml.find(cd._q("x")).get("val")) * fw + HALB_BREITE,
                float(ml.find(cd._q("y")).get("val")) * fh + HALB_HOEHE)
        # R_start wie in ring_leader_zeichnen: Mitte des Bandes
        R_out = geo["r"]
        R_in = R_out * 0.79
        R_start = R_out - 0.5 * (R_out - R_in)
        for sh in folie.shapes:
            n = sh.name or ""
            if not n.startswith("RingLeaderDot_" + shape.name + "_"):
                continue
            idx = int(n.rsplit("_", 1)[1])
            if idx not in boxen or idx >= len(mids):
                continue
            geprueft += 1
            mx, my = boxen[idx]
            sm = math.radians(mids[idx])
            sx = cx + R_start * math.sin(sm)
            sy = cy - R_start * math.cos(sm)
            steil = abs(my - sy) / max(abs(mx - sx), 1e-6)
            ist_senkrecht = abs(
                (sh.left + sh.width / 2) / 914400.0
                - shape.left / 914400.0 - mx) < 0.01
            soll_senkrecht = steil > cd.LEADER_SENKRECHT_STEIL
            if ist_senkrecht:
                senkrecht += 1
            if ist_senkrecht != soll_senkrecht:
                meldungen.append(
                    f"F{nr} {shape.name} idx {idx}: Steilheit {steil:.2f}, "
                    f"also {'senkrecht' if soll_senkrecht else 'seitlich'} "
                    f"erwartet — gezeichnet ist "
                    f"{'senkrecht' if ist_senkrecht else 'seitlich'}")
    return geprueft, senkrecht, meldungen


def schritt8_senkrechte_anbindung(ausgabe):
    print("Schritt 8 — steile Linien setzen unten an, nicht seitlich")
    if not ausgabe:
        print("    UEBERSPRUNGEN — kein Ausgabeverzeichnis angegeben")
        return 0
    fehler = 0
    for familie in sorted(SEITENTREUE_MAX):
        pfad = os.path.join(ausgabe, f"ring_{familie}.pptx")
        if not os.path.exists(pfad):
            print(f"    UEBERSPRUNGEN — {familie}: "
                  f"{os.path.basename(pfad)} fehlt (Schritt 4 baut sie)")
            continue
        n, senk, meldungen = _senkrecht_pruefen(pfad)
        for m in meldungen[:4]:
            print(f"    FEHLER — {familie}: {m}")
        if meldungen:
            fehler += 1
        else:
            print(f"    {familie:<12} {senk} von {n} Linien senkrecht "
                  f"(Schwelle {cd.LEADER_SENKRECHT_STEIL:.0f})")
    if not fehler:
        print("    OK — die Regel gilt fuer jede Linie: senkrecht genau dann, "
              "wenn sie steil genug ankommt")
    return fehler


def schritt7_buendige_anbindung(ausgabe):
    print("Schritt 7 — die Zahl klebt an ihrer Fuehrungslinie")
    if not ausgabe:
        print("    UEBERSPRUNGEN — kein Ausgabeverzeichnis angegeben")
        return 0
    fehler = 0
    for familie in sorted(SEITENTREUE_MAX):
        pfad = os.path.join(ausgabe, f"ring_{familie}.pptx")
        if not os.path.exists(pfad):
            print(f"    UEBERSPRUNGEN — {familie}: "
                  f"{os.path.basename(pfad)} fehlt (Schritt 4 baut sie)")
            continue
        n, meldungen = _buendig_pruefen(pfad)
        for m in meldungen[:4]:
            print(f"    FEHLER — {familie}: {m}")
        if meldungen:
            fehler += 1
        else:
            print(f"    {familie:<12} {n} Beschriftungen buendig angebunden")
    if not fehler:
        print("    OK — feste Boxbreite, Ausrichtung zur Ringseite, "
              "ausdrueckliche Innenabstaende")
        print("    HINWEIS — das beweist die OPTIK nicht. Gemessen am Bild "
              "(PNG-Export, CVV F7): Luecke links 5,14 -> 1,51 mm, "
              "rechts 0 -> 1,66 mm.")
    return fehler


def schritt5_familien_look():
    print("Schritt 5 — der Familien-Look bleibt getrennt")
    fehler = 0
    # Ohne erkannte Familie greift der duenne Original-Look aus der Vorlage.
    # Genau das unterscheidet die Standard-Broschuere von den fuenf Familien.
    standard = cd._ring_format(None, 79, 0.14)
    if standard["hole"] != 79:
        print(f"    FEHLER — Standard-Broschuere: hole {standard['hole']} "
              f"statt 79")
        fehler += 1
    else:
        print("    OK — ohne Familie bleibt es beim duennen Ring (hole 79)")
    # Seit 25.08.2026 tragen ALLE fuenf Familien die Bandstaerke der
    # Makro-Broschuere (hole 79). Fuer die vier mit Anlagestrategie-Folien
    # heisst das 0,56 cm bei 5,34 cm, fuer Thema 0,89 cm bei 8,51 cm auf F10 —
    # beides der Makro-Wert. Der kraeftige Ring vom 27.07.2026 (hole 68) ist
    # damit ueberall abgeloest; das war eine bewusste Entscheidung, keine
    # Nebenwirkung.
    SOLL_HOLE = {"CVV": 79, "ESG": 79, "ETF": 79, "comdirect": 79, "Thema": 79}
    for familie, soll in sorted(SOLL_HOLE.items()):
        fmt = cd._ring_format(familie, 79, 0.14)
        if fmt["hole"] != soll:
            print(f"    FEHLER — {familie}: hole {fmt['hole']} statt {soll}")
            fehler += 1
    if not fehler:
        print("    OK — alle fuenf Familien tragen den duennen Ring (hole 79)")

    # Und die Rueckkopplung laeuft NUR dort, wo sie beauftragt UND vermessen
    # ist. Ihre VOREINSTELLUNG ist AUS: eine Familie, die niemand gemessen hat,
    # bekommt sie nicht geschenkt. Ausdruecklich ein ist sie fuer die vier
    # Familien mit Anlagestrategie-Folien; Thema und "Standard" (Strategien
    # ohne Familienzuordnung, `_ring_format(None, ...)`) bleiben aus —
    # nachgewiesen ueber bytegleiche Chart-XML gegen den Stand vom 24.08.2026.
    #
    # Die None-Zeile ist der eigentliche Waechter: Bis zum 25.08.2026 stand die
    # Voreinstellung auf True, und die Standard-Broschuere erbte die
    # Rueckkopplung ungeprueft mit (5,73 -> 6,09 cm auf F7). Der Pfad wird von
    # test_export_smoke.py nicht gebaut und faellt sonst niemandem auf.
    SOLL_RK = {"CVV": True, "ESG": True, "ETF": True, "comdirect": True,
               "Thema": True, None: False}
    for familie, soll in sorted(SOLL_RK.items(), key=lambda p: str(p[0])):
        ist = cd._ring_format(familie, 79, 0.14)["rueckkopplung"]
        if ist != soll:
            bez = familie if familie is not None else "Standard (ohne Familie)"
            print(f"    FEHLER — {bez}: Rueckkopplung {ist} statt {soll}")
            fehler += 1
    if not fehler:
        print("    OK — die Rueckkopplung laeuft in den fuenf vermessenen "
              "Familien; Standard bleibt aus")

    # Und dieselbe Frage fuer die SEITENTREUE (Pass 6d, 26.08.2026).
    # Sie laeuft in den VIER Familien mit Anlagestrategie-Folien und ist bei
    # THEMA ausdruecklich aus — dort hat Philip sie zweimal angesehen und
    # zweimal abgelehnt. Die ganze Begruendung samt der drei Kehrtwenden eines
    # Tages steht im `seitentreue`-Kommentar in `chart_dynamik._RING_KRAEFTIG`.
    #
    # Die Zeile ist keine Formalie, sondern der Waechter ueber eine
    # Entscheidung nach Augenschein: Wer den Schalter umlegt, ohne dass jemand
    # das Ergebnis angesehen hat, bricht hier. Und die None-Zeile ist der
    # eigentliche Waechter — bei der Rueckkopplung hat genau sie am 25.08.2026
    # gefehlt, und die Standard-Broschuere erbte eine Aenderung ungeprueft mit.
    st_fehler = 0
    SOLL_ST = {"CVV": True, "ESG": True, "ETF": True, "comdirect": True,
               "Thema": False, None: False}
    for familie, soll in sorted(SOLL_ST.items(), key=lambda p: str(p[0])):
        ist = cd._ring_format(familie, 79, 0.14)["seitentreue"]
        if ist != soll:
            bez = familie if familie is not None else "Standard (ohne Familie)"
            print(f"    FEHLER — {bez}: Seitentreue {ist} statt {soll}")
            st_fehler += 1
    if not st_fehler:
        print("    OK — die Seitentreue laeuft in den vier "
              "Anlagestrategie-Familien; Thema und Standard bleiben aus")
    fehler += st_fehler

    # Und dieselbe Frage fuer die BUENDIGE ANBINDUNG (Punkt 1, 26.08.2026).
    # Sie ist fuer alle fuenf Familien beauftragt und in echtem PowerPoint
    # angesehen; der Standard-Pfad (None) bleibt aus. Die None-Zeile ist auch
    # hier die eigentliche Waechterin: Sie ist bei der Rueckkopplung am
    # 25.08.2026 einmal gefehlt, und die Standard-Broschuere erbte eine
    # Aenderung ungeprueft mit.
    bu_fehler = 0
    SOLL_BU = {"CVV": True, "ESG": True, "ETF": True, "comdirect": True,
               "Thema": True, None: False}
    for familie, soll in sorted(SOLL_BU.items(), key=lambda p: str(p[0])):
        ist = cd._ring_format(familie, 79, 0.14)["label_buendig"]
        if ist != soll:
            bez = familie if familie is not None else "Standard (ohne Familie)"
            print(f"    FEHLER — {bez}: label_buendig {ist} statt {soll}")
            bu_fehler += 1
    if not bu_fehler:
        print("    OK — die buendige Anbindung laeuft in den fuenf Familien; "
              "Standard bleibt aus")
    fehler += bu_fehler

    # Und fuer die SENKRECHTE ANBINDUNG (26.08.2026): Steht die Zahl fast
    # senkrecht ueber ihrem Segment, setzt die Linie unten an statt seitlich.
    sk_fehler = 0
    SOLL_SK = {"CVV": True, "ESG": True, "ETF": True, "comdirect": True,
               "Thema": True, None: False}
    for familie, soll in sorted(SOLL_SK.items(), key=lambda p: str(p[0])):
        ist = cd._ring_format(familie, 79, 0.14)["leader_senkrecht"]
        if ist != soll:
            bez = familie if familie is not None else "Standard (ohne Familie)"
            print(f"    FEHLER — {bez}: leader_senkrecht {ist} statt {soll}")
            sk_fehler += 1
    if not sk_fehler:
        print("    OK — die senkrechte Anbindung laeuft in den fuenf "
              "Familien; Standard bleibt aus")
    fehler += sk_fehler
    return fehler


def main():
    ausgabe = sys.argv[1] if len(sys.argv) > 1 else None
    print("Pruefstein: Geometrie der Ringdiagramme\n")
    fehler = 0
    for schritt in (schritt1_vorlagen, schritt2_ist_geometrie,
                    schritt3_zusicherungen):
        fehler += schritt()
        print()
    fehler += schritt4_gebaute_broschueren(ausgabe)
    print()
    fehler += schritt5_familien_look()
    print()
    fehler += schritt6_seitentreue(ausgabe)
    print()
    fehler += schritt7_buendige_anbindung(ausgabe)
    print()
    fehler += schritt8_senkrechte_anbindung(ausgabe)
    print()
    if fehler:
        print(f"FEHLGESCHLAGEN — {fehler} Abweichung(en)")
        return 1
    print("BESTANDEN — Ringgroesse, Ringdicke und Beschriftungen stehen auf "
          "ihren eingefrorenen Massen (Stand 25.08.2026)")
    print("Hinweis: die OPTIK beweist das nicht — dafuer in ECHTEM PowerPoint "
          "ansehen (#29).")
    return 0


if __name__ == "__main__":
    sys.exit(main())

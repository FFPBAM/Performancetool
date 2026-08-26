# -*- coding: utf-8 -*-
"""Welche Strategien landen in der Broschuere? (NEU 26.08.2026)

Entscheidung Philip, 26.08.2026: **Die Broschuere fuehrt immer nur die oben
gewaehlte Strategie.** Das Vergleichsportfolio ist eine Sache des Bildschirms.
Damit gilt fuer die Familie "Thema" dieselbe Regel wie fuer CVV/ESG/ETF/
comdirect, die die Vergleichsauswahl seit jeher ignorieren.

Warum es dafuer einen Pruefstein braucht: Der Unterschied zwischen dem, was am
Bildschirm steht, und dem, was im Dokument landet, ist unsichtbar. Genau solche
stillen Unterschiede sind in diesem Projekt mehrfach als Fehler zurueckgekommen
(#46, #51, #59). Und die Regel steht in der Oberflaeche, nicht in einer
Funktion, die man einzeln aufrufen koennte — geprueft wird deshalb an der
LAUFENDEN App (AppTest), nicht am Quelltext.

Schritte:
  1. Haken gesetzt -> der Export bekommt GENAU EIN Portfolio, das obere.
  2. Der Hinweis dazu steht in der Oberflaeche (und nur dann).
  3. Gegenprobe: Zwei Portfolios ergaeben nachweislich eine ANDERE Broschuere
     (25 statt 21 Folien). Ohne diesen Schritt koennte Schritt 1 gruen sein,
     ohne etwas zu bedeuten.

    python tests/test_broschuere_auswahl.py     (braucht streamlit, python-pptx)
"""

import os
import re
import sys

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)
sys.path.insert(0, os.path.join(WURZEL, "tests"))
os.chdir(WURZEL)

try:
    from streamlit.testing.v1 import AppTest
except ImportError as ex:
    print("UEBERSPRUNGEN — streamlit fehlt: %s" % ex)
    sys.exit(0)

try:
    from modules import pptx_export
except ImportError as ex:
    print("UEBERSPRUNGEN — Abhaengigkeit fehlt: %s" % ex)
    sys.exit(0)

# Eine Strategie der Familie "Thema" — die einzige Familie, bei der der Haken
# ueberhaupt etwas bewirken koennte (Dupliziermodus). Bei CVV/ESG/ETF/comdirect
# ueberschreibt der Export die Auswahl ohnehin mit allen Familienstrategien.
THEMA_ERSTE = "Pro"
THEMA_ZWEITE = "Offensiv"


def _view_pf_beschriftung():
    """Beschriftung der Portfolioanalyse-Ansicht aus streamlit_app lesen —
    damit sie nur an EINER Stelle steht (wie in test_anlagekriterien.py)."""
    with open(os.path.join(WURZEL, "streamlit_app.py"), encoding="utf-8") as fh:
        m = re.search(r'^_VIEW_PF\s*=\s*"([^"]*)"', fh.read(), re.M)
    return m.group(1) if m else "Portfolioanalyse"


def _app_in_der_portfolioanalyse(mit_haken):
    """App hochfahren, in die Portfolioanalyse wechseln, Auswahl setzen."""
    at = AppTest.from_file(os.path.join(WURZEL, "streamlit_app.py"),
                           default_timeout=400)
    at.secrets["passwords"] = {"t": "t"}
    at.session_state["logged_in"] = True
    at.session_state["username"] = "t"
    at.session_state["nav_view"] = _view_pf_beschriftung()
    at.session_state["pf_sel_1"] = THEMA_ERSTE
    at.session_state["pf_compare"] = bool(mit_haken)
    if mit_haken:
        at.session_state["pf_sel_2"] = THEMA_ZWEITE
    at.run()
    return at


def _ss(at, key, default=None):
    """session_state von AppTest kennt kein .get() — daher dieser Zugriff
    (gleiche Hilfe wie in test_bedienung.py)."""
    try:
        return at.session_state[key]
    except (KeyError, AttributeError):
        return default


def _knopf(at, key):
    for b in at.button:
        if b.key == key:
            return b
    return None


def schritt1_export_bekommt_eine_strategie():
    print("1. Haken gesetzt -> der Export bekommt genau EIN Portfolio")
    aufzeichnung = []
    original = pptx_export.generate_portfolioanalyse_pptx

    def aufzeichnen(portfolios, *a, **kw):
        aufzeichnung.append(list(portfolios))
        # Echte Bytes waeren hier Verschwendung — geprueft wird die AUSWAHL.
        return b"PK\x03\x04-attrappe"

    pptx_export.generate_portfolioanalyse_pptx = aufzeichnen
    try:
        at = _app_in_der_portfolioanalyse(mit_haken=True)
        if at.exception:
            for ex in at.exception:
                print("   FEHLER — die App warf: %s" % str(ex.value)[:200])
            return 1
        if _ss(at, "pf_sel_1") != THEMA_ERSTE:
            print("   UEBERSPRUNGEN — '%s' steht nicht in den Daten "
                  "(gewaehlt: %r)" % (THEMA_ERSTE, _ss(at, "pf_sel_1")))
            return 0
        if not _ss(at, "pf_compare"):
            print("   FEHLER — der Haken 'Vergleichsportfolio' liess sich "
                  "nicht setzen; der Test prueft sonst den falschen Zustand")
            return 1
        knopf = _knopf(at, "pf_pptx_btn")
        if knopf is None:
            print("   UEBERSPRUNGEN — Schaltflaeche 'PowerPoint erstellen' "
                  "nicht gefunden (liegt evtl. schon eine Datei im Cache)")
            return 0
        knopf.click().run()
        if at.exception:
            for ex in at.exception:
                print("   FEHLER — der Export warf: %s" % str(ex.value)[:200])
            return 1
    finally:
        pptx_export.generate_portfolioanalyse_pptx = original

    if not aufzeichnung:
        print("   FEHLER — der Export wurde gar nicht aufgerufen")
        return 1
    portfolios = aufzeichnung[-1]
    namen = [p[0] for p in portfolios]
    if len(portfolios) != 1:
        print("   FEHLER — %d Portfolios statt 1: %s" % (len(portfolios), namen))
        return 1
    if namen[0] != THEMA_ERSTE:
        print("   FEHLER — uebergeben wurde %r statt der oberen Auswahl %r"
              % (namen[0], THEMA_ERSTE))
        return 1
    print("   OK — trotz gesetztem Haken genau 1 Portfolio: %r" % namen[0])
    return 0


def schritt2_hinweis_steht_da():
    print("\n2. Der Hinweis steht in der Oberflaeche — und nur mit Haken")
    fehler = 0
    for mit_haken in (True, False):
        at = _app_in_der_portfolioanalyse(mit_haken=mit_haken)
        if at.exception:
            for ex in at.exception:
                print("   FEHLER — die App warf: %s" % str(ex.value)[:200])
            return 1
        # st.info landet bei AppTest unter .info; zur Sicherheit auch markdown.
        texte = [i.value for i in at.info] + [m.value for m in at.markdown]
        gefunden = any("nicht** in die PowerPoint" in t
                       or "nicht in die PowerPoint" in t for t in texte)
        if mit_haken and not gefunden:
            print("   FEHLER — mit Haken fehlt der Hinweis, dass das "
                  "Vergleichsportfolio NICHT in die Broschuere kommt")
            fehler += 1
        elif not mit_haken and gefunden:
            print("   FEHLER — ohne Haken steht der Hinweis trotzdem da")
            fehler += 1
    if not fehler:
        print("   OK — Hinweis erscheint mit Haken und fehlt ohne")
    return fehler


def schritt3_gegenprobe():
    """Zwei Portfolios ergaeben eine ANDERE Broschuere — sonst prueft 1 nichts."""
    print("\n3. Gegenprobe: zwei Portfolios ergaeben 25 statt 21 Folien")
    try:
        import test_export_smoke as SMOKE
        from pptx import Presentation
    except ImportError as ex:
        print("   UEBERSPRUNGEN — %s" % ex)
        return 0
    import tempfile
    d = SMOKE._daten()
    if THEMA_ERSTE not in d["d2c"] or THEMA_ZWEITE not in d["d2c"]:
        print("   UEBERSPRUNGEN — '%s'/'%s' nicht in den Daten"
              % (THEMA_ERSTE, THEMA_ZWEITE))
        return 0
    ausgabe = tempfile.mkdtemp(prefix="ffpb_auswahl_")
    einzeln = [SMOKE._portfolio(THEMA_ERSTE, d)]
    zu_zweit = einzeln + [SMOKE._portfolio(THEMA_ZWEITE, d)]
    z1, _g, _m = SMOKE._bauen(einzeln, "Thema", d, ausgabe, "eine.pptx")
    z2, _g, _m = SMOKE._bauen(zu_zweit, "Thema", d, ausgabe, "zwei.pptx")
    n1, n2 = len(Presentation(z1).slides), len(Presentation(z2).slides)
    if n1 == n2:
        print("   FEHLER — beide Broschueren haben %d Folien; die Auswahl "
              "wirkt gar nicht, Schritt 1 waere bedeutungslos" % n1)
        return 1
    print("   OK — 1 Strategie: %d Folien, 2 Strategien: %d Folien "
          "(der Unterschied ist messbar)" % (n1, n2))
    return 0


def main():
    fehler = 0
    fehler += schritt1_export_bekommt_eine_strategie()
    fehler += schritt2_hinweis_steht_da()
    fehler += schritt3_gegenprobe()
    print()
    if fehler:
        print("FEHLGESCHLAGEN — %d Abweichung(en)" % fehler)
        return 1
    print("BESTANDEN — die Broschuere fuehrt nur die obere Strategie")
    return 0


if __name__ == "__main__":
    sys.exit(main())

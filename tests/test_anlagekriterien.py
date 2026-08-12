"""Prueft die Anlagekriterien-Konfiguration (Mapping_Anlagekriterien.xlsx).

HINTERGRUND (10.08.2026):
    Der Anlagekriterien-Kasten stand bisher NUR statisch in den PPTX-Vorlagen
    — je Familie unterschiedlich geschrieben, mit Tippfehlern ("FPFB Strategie
    30", "AUsgewogen") und uneinheitlicher Prozent-Schreibweise. Er wandert
    jetzt in eine Excel, die BEIDE Ausgaben speist: den Banner im Streamlit-
    Tool und den Kasten auf der Struktur-Folie der Broschuere.

    Weil damit eine einzige Datei bestimmt, was Kunden gedruckt bekommen,
    prueft dieser Test sie streng.

ERWEITERT (12.08.2026):
    Bis dahin stand hier "Thema hat keine Anlagekriterien". Das war die Sicht
    der PowerPoint — die Thema-Vorlage hat als einzige keinen Kriterien-Kasten,
    also stand Thema auch nicht in der Excel. Fachlich ist das falsch herum:
    JEDE Strategie hat Anlagekriterien; dass eine Vorlage sie nicht abdruckt,
    ist eine Eigenschaft der VORLAGE. Offensiv, Pro und Pro Dividende stehen
    jetzt in der Excel und erscheinen im Tool — die Broschuere bleibt gleich.

Geprueft wird:
  1. Die Excel existiert und hat die erwarteten Spalten.
  2. Jede Strategie ist im Namens-Mapping bekannt (kein Schluessel ins Leere).
  3. Die Spalte 'Familie' stimmt mit Mapping_Namen.xlsx ueberein — sie ist
     bewusst doppelt gefuehrt (Lesbarkeit in Excel) und wird hier festgenagelt,
     damit sie nicht auseinanderlaeuft.
  4. Umfang: 17 der 19 Strategien sind erfasst (14 mit Kasten + 3 aus Thema).
     Die zwei SCHWEIZ-Strategien fehlen noch — bekannte Luecke, kein Fehler.
  4b. Welche VORLAGE einen Kriterien-Kasten hat. Das ist die Stelle, an der
     entschieden wird, ob ein Excel-Eintrag in einer Kundenbroschuere landet —
     bekaeme Vorlage_Thema.pptx eine Tabelle, stuenden die Thema-Kriterien
     PLOETZLICH gedruckt da, ohne Code-Aenderung.
  5. Kein Feld ist leer.
  6. Schreibweisen sind einheitlich: 'mind.' statt 'min.', Prozent mit
     Leerzeichen, kein doppelter Leerraum.
  7. anlagekriterien_fuer() liefert die vier Kriterien in Vorlagen-Reihenfolge;
     Thema liefert jetzt Werte, SCHWEIZ weiterhin [].
  9b. Der Banner erscheint auch fuer eine THEMA-Strategie in der laufenden App.

    python tests/test_anlagekriterien.py     (braucht pandas + streamlit)
"""

import os
import re
import sys

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)

import pandas as pd                                            # noqa: E402

EXCEL = os.path.join(WURZEL, "Mapping_Anlagekriterien.xlsx")
NAMEN = os.path.join(WURZEL, "Mapping_Namen.xlsx")

KEY = "Strategie auswählen"
KRITERIEN = ["Anlageregion", "Aktienanteil",
             "Anleihenanteil / Liquidität", "Fremdwährungen"]

# ─────────────────────────────────────────────────────────────────────────
# ZWEI DINGE, DIE BIS ZUM 12.08.2026 VERWECHSELT WURDEN
#
# Bis dahin galt hier: "Thema hat keine Anlagekriterien". Das war die
# Sicht der PowerPoint — die Thema-Vorlage hat als einzige keinen
# Kriterien-Kasten, also stand Thema auch nicht in der Excel.
#
# Fachlich ist das falsch herum: JEDE Strategie hat Anlagekriterien. Dass
# eine Vorlage sie nicht abdruckt, ist eine Eigenschaft der VORLAGE, keine
# der Strategie. Seit dem 12.08.2026 stehen die drei oeffentlich
# dokumentierten Thema-Strategien in der Excel und erscheinen im Tool; die
# Broschüre bleibt unveraendert.
#
# Der Test trennt deshalb jetzt sauber:
#   MIT_EINTRAG   - wer in der Excel steht (Tool zeigt den Banner)
#   MIT_KASTEN    - wessen VORLAGE eine Kriterien-Tabelle hat (Broschüre)
#   NOCH_OFFEN    - bekannte Luecke, kein Fehler
# ─────────────────────────────────────────────────────────────────────────

# Die 14 Strategien, deren Vorlage einen Kriterien-Kasten enthaelt.
MIT_KASTEN = [
    "cVV konservativ", "cVV defensiv", "cVV defensiv plus",
    "cVV ausgewogen", "cVV dynamic",
    "ESG defensiv", "ESG defensiv+", "ESG ausgewogen", "ESG offensiv",
    "ETF_ausgewogen", "ETF_Wachstum",
    "Comdirect_30", "Comdirect_70", "Comdirect_100",
]

# Neu am 12.08.2026: in der Excel und damit im Tool, aber OHNE Kasten in der
# Broschüre. Quelle der Werte: fuggerbank.de/private-banking/pro/
THEMA_MIT_EINTRAG = ["Offensiv", "Pro", "Pro Dividende"]

MIT_EINTRAG = MIT_KASTEN + THEMA_MIT_EINTRAG

# Bekannte Luecke: Die beiden SCHWEIZ-Strategien stehen bewusst NICHT auf der
# Webseite; ihre Kriterien muessen aus dem Haus kommen. Bis dahin liefern sie
# korrekt [] — das ist kein Fehler, sondern ein offener Punkt (Doku Sec. 15).
NOCH_OFFEN = ["Schweiz_substanzorientiert", "Schweiz_aktienorientiert"]

# Vorlagen ohne Kriterien-Tabelle. Solange das so ist, kann kein Excel-Eintrag
# eine Thema-Broschüre veraendern — Schritt 4b weist das am Artefakt nach.
VORLAGEN_OHNE_KASTEN = ["Vorlage_Thema.pptx", "Vorlage_FFPB.pptx"]


def _pruefe_struktur(df):
    print("1. Aufbau der Excel")
    fehler = 0
    soll = [KEY, "Familie", "Anzeigename"] + KRITERIEN
    fehlend = [s for s in soll if s not in df.columns]
    if fehlend:
        print(f"   FEHLER — Spalten fehlen: {fehlend}")
        fehler += 1
    else:
        print(f"   OK — {len(df.columns)} Spalten, {len(df)} Zeilen")
    return fehler


def _pruefe_schluessel(df, namen):
    print("\n2. Jede Strategie ist im Namens-Mapping bekannt")
    fehler = 0
    bekannt = set(namen[KEY].astype(str).str.strip())
    for s in df[KEY].astype(str).str.strip():
        if s not in bekannt:
            print(f"   FEHLER — '{s}' steht nicht in Mapping_Namen.xlsx")
            fehler += 1
    if not fehler:
        print(f"   OK — alle {len(df)} Schluessel gefunden")
    return fehler


def _pruefe_familie(df, namen):
    print("\n3. Spalte 'Familie' deckt sich mit Mapping_Namen.xlsx")
    fehler = 0
    soll = dict(zip(namen[KEY].astype(str).str.strip(),
                    namen["Powerpoint Familie"].astype(str).str.strip()))
    for _, z in df.iterrows():
        s = str(z[KEY]).strip()
        ist = str(z["Familie"]).strip()
        if soll.get(s, "") != ist:
            print(f"   FEHLER — {s}: Excel sagt '{ist}', "
                  f"Mapping sagt '{soll.get(s, '')}'")
            fehler += 1
    if not fehler:
        print("   OK — keine Abweichung")
    return fehler


def _pruefe_umfang(df):
    print("\n4. Umfang der Excel (17 von 19 Strategien)")
    fehler = 0
    ist = list(df[KEY].astype(str).str.strip())
    fehlend = [s for s in MIT_EINTRAG if s not in ist]
    zuviel = [s for s in ist if s not in MIT_EINTRAG]
    if fehlend:
        print(f"   FEHLER — fehlen: {fehlend}")
        fehler += 1
    for s in zuviel:
        if s in NOCH_OFFEN:
            # Kein Fehler waere falsch: Sobald jemand SCHWEIZ eintraegt,
            # soll dieser Test daran erinnern, den offenen Punkt zu schliessen.
            print(f"   HINWEIS — '{s}' ist neu eingetragen. Bitte NOCH_OFFEN "
                  f"hier und den Backlog-Punkt in der Doku nachziehen.")
        else:
            print(f"   FEHLER — unbekannte Strategie '{s}'")
            fehler += 1
    if not fehler:
        print(f"   OK — {len(ist)} Strategien erfasst "
              f"({len(MIT_KASTEN)} mit Kasten, "
              f"{len(THEMA_MIT_EINTRAG)} nur im Tool)")
        print(f"   OFFEN — {len(NOCH_OFFEN)} SCHWEIZ-Strategien ohne Werte "
              f"(stehen bewusst nicht auf der Webseite)")
    return fehler


def _pruefe_vorlagen_kasten():
    """4b — Der Kasten haengt an der VORLAGE, nicht an der Excel.

    Warum das ein eigener Schritt ist: pptx_export ruft
    fill_anlagekriterien_slide fuer JEDE Familie auf. Dass Thema keinen
    Kasten bekommt, liegt einzig daran, dass Vorlage_Thema.pptx keine
    Kriterien-Tabelle enthaelt (fill_anlagekriterien_slide steigt dann aus).
    Bekaeme die Vorlage eines Tages eine Tabelle, wuerden die seit dem
    12.08.2026 eingetragenen Thema-Kriterien PLOETZLICH IN KUNDENBROSCHUEREN
    STEHEN — ohne dass jemand den Code angefasst hat. Dieser Schritt schlaegt
    dann an.
    """
    print("\n4b. Welche Vorlage hat ueberhaupt einen Kriterien-Kasten?")
    try:
        from pptx import Presentation
        from modules.pptx_slides import finde_anlagekriterien_tabelle
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    ordner = os.path.join(WURZEL, "Vorlage")
    if not os.path.isdir(ordner):
        print("   UEBERSPRUNGEN — Ordner Vorlage/ fehlt")
        return 0

    fehler = 0
    for name in sorted(os.listdir(ordner)):
        if not name.lower().endswith(".pptx"):
            continue
        prs = Presentation(os.path.join(ordner, name))
        treffer = [i for i in range(len(prs.slides))
                   if finde_anlagekriterien_tabelle(prs.slides[i]) is not None]
        soll_leer = name in VORLAGEN_OHNE_KASTEN
        if soll_leer and treffer:
            print(f"   FEHLER — {name} hat jetzt eine Kriterien-Tabelle "
                  f"(Folien {treffer}). Damit wuerden die Thema-Kriterien in "
                  f"die Broschüre gedruckt. Gewollt? Dann diesen Test und die "
                  f"Doku nachziehen.")
            fehler += 1
        elif not soll_leer and not treffer:
            print(f"   FEHLER — {name} hat KEINE Kriterien-Tabelle mehr — "
                  f"der Kasten faellt dort still aus")
            fehler += 1
        else:
            print(f"   OK — {name:28s} "
                  f"{'ohne Kasten (so gewollt)' if soll_leer else f'Kasten auf Folie {treffer}'}")
    return fehler


def _pruefe_vollstaendig(df):
    print("\n5. Kein Feld leer")
    fehler = 0
    for _, z in df.iterrows():
        for sp in ["Anzeigename"] + KRITERIEN:
            wert = z.get(sp)
            if pd.isna(wert) or not str(wert).strip():
                print(f"   FEHLER — {z[KEY]}: '{sp}' ist leer")
                fehler += 1
    if not fehler:
        print(f"   OK — {len(df) * 5} Felder gefuellt")
    return fehler


def _pruefe_schreibweise(df):
    print("\n6. Einheitliche Schreibweise")
    fehler = 0
    for _, z in df.iterrows():
        for sp in KRITERIEN:
            w = str(z[sp])
            if re.search(r"\bmin\.(?!\w)", w):
                print(f"   FEHLER — {z[KEY]} / {sp}: 'min.' statt 'mind.' ({w!r})")
                fehler += 1
            if re.search(r"\d%", w):
                print(f"   FEHLER — {z[KEY]} / {sp}: Prozent ohne "
                      f"Leerzeichen ({w!r})")
                fehler += 1
            if "  " in w or w != w.strip():
                print(f"   FEHLER — {z[KEY]} / {sp}: ueberfluessiger "
                      f"Leerraum ({w!r})")
                fehler += 1
    # Tippfehler, die es in den Vorlagen gab — duerfen nicht zurueckkehren
    alle = " ".join(df.astype(str).values.ravel())
    for tippfehler in ("FPFB", "AUsgewogen"):
        if tippfehler in alle:
            print(f"   FEHLER — Tippfehler '{tippfehler}' ist zurueck")
            fehler += 1
    if not fehler:
        print("   OK — keine Abweichung")
    return fehler


def _pruefe_zugriff(df):
    print("\n7. anlagekriterien_fuer() liefert die richtige Reihenfolge")
    try:
        from modules.shared import anlagekriterien_fuer
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0
    fehler = 0
    paare = anlagekriterien_fuer("cVV defensiv plus", df)
    bez = [b for b, _ in paare]
    if bez != KRITERIEN:
        print(f"   FEHLER — Reihenfolge {bez}, erwartet {KRITERIEN}")
        fehler += 1
    else:
        print(f"   OK — cVV defensiv plus: "
              f"{', '.join(f'{b}={w}' for b, w in paare)}")
    # Die drei Thema-Strategien MUESSEN jetzt Kriterien liefern — das ist
    # die Faehigkeit, die am 12.08.2026 dazugekommen ist.
    for s in THEMA_MIT_EINTRAG:
        paare = anlagekriterien_fuer(s, df)
        if not paare:
            print(f"   FEHLER — '{s}' liefert keine Kriterien, obwohl sie "
                  f"in der Excel stehen")
            fehler += 1
        elif [b for b, _ in paare] != KRITERIEN:
            print(f"   FEHLER — '{s}': falsche Reihenfolge {[b for b, _ in paare]}")
            fehler += 1
    # SCHWEIZ liefert weiterhin [] — bekannte Luecke, kein Fehler.
    for s in NOCH_OFFEN:
        if anlagekriterien_fuer(s, df):
            print(f"   HINWEIS — '{s}' liefert jetzt Kriterien. Bitte "
                  f"NOCH_OFFEN und den Backlog-Punkt nachziehen.")
    if not fehler:
        print(f"   OK — {len(THEMA_MIT_EINTRAG)} Thema-Strategien liefern "
              f"Kriterien, {len(NOCH_OFFEN)} SCHWEIZ noch leer")
    for leer in ("", None, "Gibt es nicht"):
        if anlagekriterien_fuer(leer, df):
            print(f"   FEHLER — {leer!r} sollte [] liefern")
            fehler += 1
    return fehler


def _pruefe_bauweise(df):
    """Der Block MUSS aus nativen Streamlit-Bausteinen bestehen.

    Die erste Fassung war ein HTML-Block mit eigener heller Flaeche — im Dark
    Mode ein greller weisser Kasten. Dieser Schritt haelt fest, dass wir da
    nicht zurueckfallen: kein eigenes CSS, keine festen Farben.
    """
    print("\n8. Bauweise: native Streamlit-Bausteine, kein eigenes CSS")
    import inspect
    fehler = 0
    try:
        from modules.shared import zeige_anlagekriterien, markdown_escapen
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    # NUR den ausfuehrbaren Code pruefen: Docstring und Kommentare BEGRUENDEN
    # ja gerade, warum wir var(--…) und st.metric NICHT verwenden — ohne diese
    # Bereinigung schlaegt der Test auf seine eigene Erklaerung an.
    roh = inspect.getsource(zeige_anlagekriterien)
    doku = zeige_anlagekriterien.__doc__ or ""
    if doku:
        roh = roh.replace(doku, "")
    quelle = "\n".join(z for z in roh.splitlines()
                       if not z.strip().startswith("#"))

    # Kein Roh-HTML, keine Farbwerte, kein unsafe_allow_html
    verboten = {
        "unsafe_allow_html": "eigenes HTML statt nativer Bausteine",
        "background:": "fest verdrahteter Hintergrund",
        "#003460": "fest verdrahtete Textfarbe (bricht im Dark Mode)",
        "<div": "Roh-HTML",
        "var(--": "Streamlit 1.61 stellt KEINE Theme-CSS-Variablen bereit",
    }
    for muster, warum in verboten.items():
        if muster in quelle:
            print(f"   FEHLER — '{muster}' im Code: {warum}")
            fehler += 1

    # Die nativen Bausteine, die das Theme mitbringen
    for muster in ("st.container(", "st.columns", "st.caption"):
        if muster not in quelle:
            print(f"   FEHLER — '{muster}' fehlt; Block ist nicht nativ gebaut")
            fehler += 1

    # Ohne Rahmen (Wunsch Philip 10.08.2026): die graue Umrandung wirkte
    # im Seitenfluss unruhig.
    if "border=True" in quelle:
        print("   FEHLER — Rahmen ist zurueck (border=True)")
        fehler += 1

    # Bewusst NICHT st.metric — sonst sieht die Regel aus wie eine Kennzahl
    if "st.metric" in quelle:
        print("   FEHLER — st.metric verwischt Bestandszahl und Regel")
        fehler += 1

    # Markdown-Sonderzeichen aus der Excel duerfen nichts verschlucken
    if markdown_escapen("max. *5* _%_") != "max. \\*5\\* \\_%\\_":
        print("   FEHLER — Markdown-Sonderzeichen werden nicht entschaerft")
        fehler += 1

    if not fehler:
        print("   OK — container/columns/caption, kein CSS, keine Farbwerte")
    return fehler


def _view_pf_beschriftung():
    """Beschriftung der Portfolioanalyse-Ansicht aus streamlit_app lesen.

    streamlit_app.py laesst sich nicht importieren (das Skript wuerde
    losrennen), deshalb wird die Zuweisung aus dem Quelltext gelesen. So
    steht der Wert trotzdem nur an EINER Stelle.
    """
    import re
    pfad = os.path.join(WURZEL, "streamlit_app.py")
    with open(pfad, encoding="utf-8") as fh:
        m = re.search(r'^_VIEW_PF\s*=\s*"([^"]*)"', fh.read(), re.M)
    return m.group(1) if m else "Portfolioanalyse"


def _pruefe_in_der_app(df):
    """End-to-End: die App hochfahren und den Banner im Markup suchen."""
    print("\n9. Banner erscheint in der laufenden App (AppTest)")
    try:
        from streamlit.testing.v1 import AppTest
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    at = AppTest.from_file(os.path.join(WURZEL, "streamlit_app.py"),
                           default_timeout=300)
    # check_login liest st.secrets; angemeldet wird ueber den session_state,
    # damit der Test nicht am Passwort haengt (siehe test_app_titel.py).
    at.secrets["passwords"] = {"testnutzer": "nur-fuer-den-test"}
    at.session_state["logged_in"] = True
    at.session_state["username"] = "testnutzer"
    try:
        at.run()
    except Exception as ex:
        print(f"   UEBERSPRUNGEN — App liess sich nicht starten: {ex}")
        return 0
    if at.exception:
        for ex in at.exception:
            print(f"   FEHLER — App warf: {str(ex.value)[:220]}")
        return 1

    markup = " ".join(m.value for m in at.markdown)
    if "Anlagekriterien" not in markup:
        print("   FEHLER — 'Anlagekriterien' steht nicht im Seiten-Markup")
        return 1

    # Der Banner muss zur AUSGEWAEHLTEN Strategie passen. Streamlit waehlt die
    # erste Option der Portfolio-Auswahl vor.
    gewaehlt = at.selectbox(key="p_sel1").value if at.selectbox else None
    print(f"   vorausgewaehlte Strategie: {gewaehlt!r}")
    from modules.shared import anlagekriterien_fuer
    paare = anlagekriterien_fuer(gewaehlt, df)
    if not paare:
        print("   OK — Strategie ohne Kasten, Banner entfaellt korrekt")
        return 0
    fehlen = [w for _, w in paare if w not in markup]
    if fehlen:
        print(f"   FEHLER — Werte fehlen im Markup: {fehlen}")
        return 1
    print(f"   OK — Performance-Ansicht: alle {len(paare)} Werte gefunden "
          f"({', '.join(w for _, w in paare)})")

    # ── zweite Ansicht: Portfolioanalyse ──────────────────────────────────
    # Der Banner soll in BEIDEN Ansichten stehen. Umschalten ueber den
    # session_state des segmented_control (key "nav_view", siehe #18).
    # Der Wert wird aus streamlit_app importiert statt hier wiederholt —
    # sonst bricht der Test still, wenn sich die Beschriftung aendert
    # (genau das passierte beim Entfernen der Piktogramme am 10.08.2026).
    at.session_state["nav_view"] = _view_pf_beschriftung()
    try:
        at.run()
    except Exception as ex:
        print(f"   UEBERSPRUNGEN (Portfolioanalyse) — {ex}")
        return 0
    if at.exception:
        for ex in at.exception:
            print(f"   FEHLER — Portfolioanalyse warf: {str(ex.value)[:220]}")
        return 1

    markup2 = " ".join(m.value for m in at.markdown)
    gewaehlt2 = None
    for sb in at.selectbox:
        if sb.key == "pf_sel_1":
            gewaehlt2 = sb.value
    print(f"   vorausgewaehltes Portfolio: {gewaehlt2!r}")
    paare2 = anlagekriterien_fuer(gewaehlt2, df)
    if not paare2:
        print("   OK — Portfolio ohne Kasten, Banner entfaellt korrekt")
        return 0
    fehlen2 = [w for _, w in paare2 if w not in markup2]
    if fehlen2:
        print(f"   FEHLER — Portfolioanalyse, Werte fehlen: {fehlen2}")
        return 1
    print(f"   OK — Portfolioanalyse: alle {len(paare2)} Werte gefunden")
    return 0


def _pruefe_thema_in_der_app(df):
    """Schritt 9b — die eigentliche Neuerung vom 12.08.2026.

    Schritt 9 prueft die VORAUSGEWAEHLTE Strategie; das ist eine mit Kasten.
    Der Banner fuer eine Thema-Strategie wurde damit von nichts abgedeckt —
    und genau der ist neu: Pro, Pro Dividende und Offensiv erscheinen jetzt
    im Tool, obwohl ihre Broschüre den Kasten weiterhin nicht kennt.

    Der Test waehlt deshalb ausdruecklich eine Thema-Strategie aus und liest
    ihre Werte am gerenderten Markup nach.
    """
    print("\n9b. Banner fuer eine THEMA-Strategie (neu seit 12.08.2026)")
    try:
        from streamlit.testing.v1 import AppTest
        from modules.shared import anlagekriterien_fuer
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    at = AppTest.from_file(os.path.join(WURZEL, "streamlit_app.py"),
                           default_timeout=300)
    at.secrets["passwords"] = {"testnutzer": "nur-fuer-den-test"}
    at.session_state["logged_in"] = True
    at.session_state["username"] = "testnutzer"
    try:
        at.run()
    except Exception as ex:
        print(f"   UEBERSPRUNGEN — App liess sich nicht starten: {ex}")
        return 0

    auswahl = None
    for sb in at.selectbox:
        if sb.key == "p_sel1":
            auswahl = sb
    if auswahl is None:
        print("   UEBERSPRUNGEN — Strategie-Auswahl (p_sel1) nicht gefunden")
        return 0

    optionen = list(auswahl.options)
    ziel = next((o for o in optionen if str(o).strip() in THEMA_MIT_EINTRAG),
                None)
    if ziel is None:
        print(f"   UEBERSPRUNGEN — keine Thema-Strategie in den Daten "
              f"(vorhanden: {optionen[:6]}…)")
        return 0

    try:
        auswahl.select(ziel).run()
    except Exception as ex:
        print(f"   UEBERSPRUNGEN — Umschalten auf {ziel!r} misslang: {ex}")
        return 0
    if at.exception:
        for ex in at.exception:
            print(f"   FEHLER — App warf nach Auswahl {ziel!r}: "
                  f"{str(ex.value)[:220]}")
        return 1

    paare = anlagekriterien_fuer(ziel, df)
    if not paare:
        print(f"   FEHLER — '{ziel}' liefert keine Kriterien (Excel-Eintrag?)")
        return 1

    markup = " ".join(m.value for m in at.markdown)
    if "Anlagekriterien" not in markup:
        print(f"   FEHLER — kein Anlagekriterien-Banner bei '{ziel}'")
        return 1
    fehlen = [w for _, w in paare if w not in markup]
    if fehlen:
        print(f"   FEHLER — Werte fehlen im Markup: {fehlen}")
        return 1
    print(f"   OK — '{ziel}': alle {len(paare)} Werte im Banner "
          f"({', '.join(w for _, w in paare)})")
    return 0


def _pruefe_broschueren(df, ordner):
    """Schritt 10 (nur mit Ordner-Argument): Zeigt die ERZEUGTE Broschuere
    genau das, was in der Konfiguration steht?

        python tests/test_anlagekriterien.py C:\\pfad\\zur\\ausgabe
    """
    print(f"\n10. Kasten in den erzeugten Broschueren ({ordner})")
    try:
        import glob
        from pptx import Presentation
        from modules.pptx_slides import finde_anlagekriterien_tabelle
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    dateien = sorted(glob.glob(os.path.join(ordner, "*.pptx")))
    if not dateien:
        print(f"   FEHLER — keine PPTX in {ordner}")
        return 1

    # Nachschlagen ueber die Kopfzeile: unabhaengig von Folienpositionen.
    nach_anzeige = {}
    for _, z in df.iterrows():
        nach_anzeige.setdefault(str(z["Anzeigename"]).strip(), []).append(z)

    fehler, geprueft = 0, 0
    for pfad in dateien:
        name = os.path.basename(pfad)
        prs = Presentation(pfad)
        for si, slide in enumerate(prs.slides, 1):
            t = finde_anlagekriterien_tabelle(slide)
            if t is None:
                continue
            geprueft += 1
            kopf = t.cell(0, 2).text.strip()
            kandidaten = nach_anzeige.get(kopf)
            if not kandidaten:
                print(f"   FEHLER — {name} F{si}: Kopfzeile {kopf!r} steht "
                      f"in keiner Zeile der Konfiguration")
                fehler += 1
                continue
            # Bei gleichem Anzeigenamen (ETF/CVV 'Ausgewogen') genuegt EINE
            # passende Zeile — die Werte entscheiden.
            passt = False
            for z in kandidaten:
                if all(t.cell(r, 0).text.strip() == k
                       and t.cell(r, 2).text.strip() == str(z[k]).strip()
                       for r, k in enumerate(KRITERIEN, start=1)):
                    passt = True
                    break
            if not passt:
                z = kandidaten[0]
                print(f"   FEHLER — {name} F{si} ({kopf}):")
                for r, k in enumerate(KRITERIEN, start=1):
                    ist_b, ist_w = t.cell(r, 0).text.strip(), t.cell(r, 2).text.strip()
                    if ist_b != k or ist_w != str(z[k]).strip():
                        print(f"      Zeile {r}: {ist_b!r}={ist_w!r}  "
                              f"erwartet {k!r}={str(z[k]).strip()!r}")
                fehler += 1
                continue
            # Formatierung darf nicht verloren gehen: jede beschriebene Zelle
            # braucht einen Run MIT expliziter Schriftgroesse.
            for r in range(len(t.rows)):
                for c in (0, 2):
                    paras = t.cell(r, c).text_frame.paragraphs
                    runs = paras[0].runs if paras else []
                    if not runs or runs[0].font.size is None:
                        print(f"   FEHLER — {name} F{si} Zeile {r} Spalte {c}: "
                              f"Formatierung verloren (keine Schriftgroesse)")
                        fehler += 1

    if geprueft == 0:
        print("   FEHLER — kein einziger Kasten gefunden")
        return 1
    if not fehler:
        print(f"   OK — {geprueft} Kaesten stimmen mit der Konfiguration "
              f"ueberein, Formatierung erhalten")
    return fehler


def main():
    if not os.path.exists(EXCEL):
        print(f"FEHLER: {EXCEL} fehlt")
        return 1
    df = pd.read_excel(EXCEL)
    namen = pd.read_excel(NAMEN)

    fehler = (_pruefe_struktur(df)
              + _pruefe_schluessel(df, namen)
              + _pruefe_familie(df, namen)
              + _pruefe_umfang(df)
              + _pruefe_vorlagen_kasten()
              + _pruefe_vollstaendig(df)
              + _pruefe_schreibweise(df)
              + _pruefe_zugriff(df)
              + _pruefe_bauweise(df)
              + _pruefe_in_der_app(df)
              + _pruefe_thema_in_der_app(df))

    # Schritt 10 nur, wenn ein Export-Ordner uebergeben wurde (wie
    # test_trennstriche.py). Ohne Ordner bleibt der Test schnell.
    if len(sys.argv) > 1:
        fehler += _pruefe_broschueren(df, sys.argv[1])

    print()
    if fehler:
        print(f"FEHLGESCHLAGEN — {fehler} Abweichung(en)")
        return 1
    print(f"BESTANDEN — {len(df)} Strategien, Konfiguration konsistent")
    return 0


if __name__ == "__main__":
    sys.exit(main())

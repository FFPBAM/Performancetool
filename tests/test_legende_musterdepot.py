"""Prueft die Legende der Wertentwicklungs-Folie: "Musterdepot", nicht
"Referenzportfolio".

REGEL (festgelegt mit Philip am 10.08.2026):
    Die Legenden-Textbox "Legende Diagramm links" der Wertentwicklungs-Folie
    zeigt den Begriff, der in der VORLAGE steht — und dort steht ueberall
    "Musterdepot". Der Code faesst sie nicht mehr an.

    Vorgeschichte: Am 02.07.2026 (Punkt 3) schrieb fill_wertentwicklung_slide
    beim Befuellen 'Musterdepot ' -> 'Referenzportfolio ' um und kuerzte zum
    Ausgleich den 5-Leerzeichen-Luecken-Run auf 3. Begruendung damals: die
    performance-Folie fuehre den Begriff "Referenzportfolio", beide sollten
    gleich heissen. Das Argument traegt nur fuer die FFPB-Standardvorlage —
    die Rolle "performance" kommt in KEINER Familien-Konfiguration vor. In
    den Familien-Broschueren gab es also gar keine Folie, zu der die
    Angleichung gepasst haette. Zurueckgenommen am 10.08.2026.

    BEWUSSTE AUSNAHME: Vorlage_FFPB.pptx Folie 10 ist die performance-Folie.
    Sie fuehrt "Referenzportfolio" statisch in der Vorlage (Legenden-Box UND
    Folientitel) und behaelt den Begriff. Jede Folie folgt ihrer Vorlage.

Geprueft wird:
  1. Vorlagen-Invariante: jede Legenden-Box sagt "Musterdepot" — ausser der
     bekannten performance-Folie. Faengt ab, dass jemand eine Vorlage
     austauscht und der Begriff unbemerkt wandert. Braucht NUR die
     Standardbibliothek (PPTX ist ein ZIP).
  2. Der Serienname des Saeulen-Charts (WE_SERIES_PORTFOLIO) folgt der
     Vorlage.
  3. Wirkung am echten Artefakt: fill_wertentwicklung_slide auf die echte
     ETF-Vorlage angewendet laesst die Legende von F19 unveraendert.
     Dieser Teil braucht python-pptx und wird sonst uebersprungen.

Auf dem Stand vom 07.08.2026 sind Schritt 2 und 3 ROT.

    python tests/test_legende_musterdepot.py
"""

import os
import sys
import zipfile
import xml.etree.ElementTree as ET

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)

VORLAGEN_ORDNER = os.path.join(WURZEL, "Vorlage")

NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

LEGENDE_SHAPE = "Legende Diagramm links"
SOLL_BEGRIFF = "Musterdepot"
ALT_BEGRIFF = "Referenzportfolio"

# Die eine Folie, die "Referenzportfolio" behaelt: die performance-Folie der
# Standardvorlage (Position 10 laut pptx_export.DEFAULT_TEMPLATE_CONFIG).
AUSNAHME = ("Vorlage_FFPB.pptx", "slide10.xml")


def _legenden_aus_vorlage(pfad):
    """Liest alle Legenden-Textboxen einer PPTX direkt aus dem ZIP.

    Gibt {slide-dateiname: text} zurueck. Bewusst ohne python-pptx, damit
    Schritt 1 auch ohne Paketinstallation laeuft.

    WICHTIG: Der Text wird aus den <a:t>-Knoten zusammengesetzt, OHNE die
    Whitespace-Knoten zu verwerfen — der Luecken-Run zwischen den beiden
    Legendeneintraegen besteht aus 5 Leerzeichen und ist Teil des Ergebnisses.
    """
    gefunden = {}
    with zipfile.ZipFile(pfad) as z:
        namen = [n for n in z.namelist() if n.startswith("ppt/slides/slide")
                 and n.endswith(".xml")]
        for name in sorted(namen):
            wurzel = ET.fromstring(z.read(name))
            for sp in wurzel.iter(f"{{{NS['p']}}}sp"):
                cNvPr = sp.find(f".//{{{NS['p']}}}cNvPr")
                if cNvPr is None or cNvPr.get("name") != LEGENDE_SHAPE:
                    continue
                text = "".join(t.text or "" for t in sp.iter(f"{{{NS['a']}}}t"))
                gefunden[os.path.basename(name)] = text
    return gefunden


def _pruefe_vorlagen():
    """Schritt 1 — Vorlagen-Invariante."""
    print("1. Vorlagen: Legenden-Box sagt 'Musterdepot'")
    if not os.path.isdir(VORLAGEN_ORDNER):
        print(f"   FEHLER: Ordner fehlt: {VORLAGEN_ORDNER}")
        return 1

    fehler = 0
    geprueft = 0
    for datei in sorted(os.listdir(VORLAGEN_ORDNER)):
        if not datei.lower().endswith(".pptx"):
            continue
        pfad = os.path.join(VORLAGEN_ORDNER, datei)
        for slide, text in sorted(_legenden_aus_vorlage(pfad).items()):
            geprueft += 1
            ist_ausnahme = (datei, slide) == AUSNAHME
            soll = ALT_BEGRIFF if ist_ausnahme else SOLL_BEGRIFF
            ok = soll in text
            if not ok:
                fehler += 1
            marke = "  (performance-Folie, Ausnahme)" if ist_ausnahme else ""
            # repr(), damit der Luecken-Run sichtbar bleibt
            print(f"   {datei[:26]:26s} {slide:12s} {'OK' if ok else 'FEHLER':6s} "
                  f"{text.strip()[:34]!r}{marke}")

    if geprueft == 0:
        print(f"   FEHLER: keine Shape '{LEGENDE_SHAPE}' gefunden — "
              f"wurde sie umbenannt?")
        return 1
    print(f"   {geprueft} Legenden geprueft")
    return fehler


def _pruefe_konstante():
    """Schritt 2 — Serienname folgt der Vorlage."""
    print("\n2. Serienname des Saeulen-Charts")
    try:
        # pptx_slides importiert python-pptx auf Modulebene
        from modules.pptx_slides import WE_SERIES_PORTFOLIO
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0
    ok = WE_SERIES_PORTFOLIO == SOLL_BEGRIFF
    print(f"   WE_SERIES_PORTFOLIO = {WE_SERIES_PORTFOLIO!r}  "
          f"{'OK' if ok else 'FEHLER — erwartet ' + repr(SOLL_BEGRIFF)}")
    return 0 if ok else 1


def _wertentwicklungs_folien(familie):
    """1-indexierte Positionen der wertentwicklung-Folien einer Familie."""
    from modules.vorlagen_config import VORLAGEN_FAMILIEN
    _, cfg = VORLAGEN_FAMILIEN[familie]
    if "feste_bloecke" in cfg:
        return [b["wertentwicklung"] for b in cfg["feste_bloecke"]
                if "wertentwicklung" in b]
    return [cfg["block_positionen"]["wertentwicklung"]]


def _pruefe_wirkung():
    """Schritt 3 — echte Vorlage befuellen, Legende muss stehen bleiben."""
    print("\n3. Wirkung: fill_wertentwicklung_slide auf die echte ETF-Vorlage")
    try:
        from pptx import Presentation
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0
    try:
        from modules.pptx_slides import fill_wertentwicklung_slide
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    pfad = os.path.join(VORLAGEN_ORDNER, "Vorlage_ETF.pptx")
    if not os.path.exists(pfad):
        print(f"   FEHLER: {pfad} fehlt")
        return 1

    # Minimale we_data: Kennzahlen ja, Charts leer. Die Chart-Zweige haengen
    # an pa["jahre"] bzw. we["dates"] — leer heisst, sie werden uebersprungen.
    # Das haelt den Test schnell und unabhaengig von den CSV-Daten; die
    # Legende liegt hinter diesen Zweigen und wird trotzdem erreicht.
    we_data = {
        "auflage_jahr": 2019, "laufendes_jahr": 2026,
        "kum_nach_kosten": 0.42, "pa_nach_kosten": 0.05,
        "ytd": 0.03, "duration": None,
        "benchmark_text": "Testbenchmark",
        "performance_pa": {}, "wertentwicklung": {},
    }

    fehler = 0
    positionen = _wertentwicklungs_folien("ETF")
    for pos in positionen:
        prs = Presentation(pfad)          # je Folie frisch: keine Seiteneffekte
        slide = prs.slides[pos - 1]

        vorher = None
        for shape in slide.shapes:
            if shape.name == LEGENDE_SHAPE:
                vorher = shape.text_frame.text
        if vorher is None:
            print(f"   F{pos}: FEHLER — Shape '{LEGENDE_SHAPE}' nicht gefunden")
            fehler += 1
            continue

        fill_wertentwicklung_slide(prs, pos - 1, "ETF Wachstum", we_data)

        nachher = None
        for shape in slide.shapes:
            if shape.name == LEGENDE_SHAPE:
                nachher = shape.text_frame.text

        unveraendert = nachher == vorher
        sagt_soll = SOLL_BEGRIFF in (nachher or "")
        kein_alt = ALT_BEGRIFF not in (nachher or "")
        ok = unveraendert and sagt_soll and kein_alt
        if not ok:
            fehler += 1
        print(f"   F{pos}: {'OK' if ok else 'FEHLER'}  "
              f"vorher={vorher.strip()[:30]!r} nachher={(nachher or '').strip()[:30]!r}")
        if not ok and not unveraendert:
            print("        Die Legende wurde umgeschrieben — genau das soll "
                  "seit dem 10.08.2026 nicht mehr passieren.")
    return fehler


def main():
    fehler = 0
    fehler += _pruefe_vorlagen()
    fehler += _pruefe_konstante()
    fehler += _pruefe_wirkung()

    print()
    if fehler:
        print(f"FEHLGESCHLAGEN — {fehler} Abweichung(en)")
        return 1
    print("BESTANDEN — Legende zeigt 'Musterdepot', Vorlage bleibt unangetastet")
    return 0


if __name__ == "__main__":
    sys.exit(main())

"""Regressionstest: die Datumsachse zeigt den letzten Datenzeitraum (12.08.2026).

Gemeldet wurde: In der ETF-Broschuere laufen die Daten bis Juli 2026, die
Datumsachse endet aber bei "Dez/25" — 2026 kommt auf der Achse nicht vor.

URSACHE. PowerPoint setzt die Ticks einer Datumsachse beim ACHSEN-MINIMUM an
und zaehlt in Schritten von majorUnit x majorTimeUnit weiter — NICHT an
Kalendergrenzen. chart_dynamik.datumsachse_an_daten setzte das Minimum auf den
Monatsanfang des ersten Datenpunkts (gegen Leerraum vor der Kurve). Damit
liegen alle Jahres-Ticks auf dem ANFANGSMONAT der Reihe, und der letzte faellt
vor das laufende Jahr. Am gebauten Artefakt gemessen (Datenstand 260721,
21 Datumsachsen — es war KEINE in Ordnung):

    ETF (2)           Achse ab 01.11.2015   letzter Tick Nov/25
    ESG (4)           Achse ab 01.09.2020   letzter Tick Sep/25
    cVV klassisch (4) Achse ab 01.12.2008   letzter Tick Dez/25
    cVV Dynamic       Achse ab 01.10.2018   letzter Tick Okt/25
    cVV Vergleich     Achse ab 01.01.2009   letzter Tick Jan/21 (!)
    Thema             Monatsticks: 37 bzw. 23 Beschriftungen
    comdirect (3)     majorTimeUnit fehlt ganz

Zwei weitere Fehler derselben Stelle, die dieser Test mit abdeckt:

  * majorUnit wurde nie mitgezogen. Die cVV-Vergleichsfolie traegt in der
    Vorlage majorUnit=12 + majorTimeUnit=months; das Umstellen auf "years"
    machte daraus einen Tick alle ZWOELF JAHRE.
  * majorTimeUnit fehlt in Vorlage_comdirect.pptx ganz. Wegen
    "if mtu is not None" lief die Anpassung dort ueberhaupt nicht.

Geprueft wird in zwei Schritten:

  Schritt 1 (ohne Installation) — der Rechenkern achsen_raster gegen von Hand
                     nachgerechnete Faelle, dazu die Grenzfaelle (ein
                     Datenpunkt, Beginn exakt im Januar, Spanne unter einem
                     Monat).
  Schritt 2 (+ python-pptx, streamlit) — am ECHTEN Artefakt: je Familie eine
                     gebaute Broschuere, jedes Chart mit Datumsachse
                     nachgerechnet. Der letzte Tick muss im Jahr des letzten
                     Datenpunkts liegen.

Ohne pptx/streamlit wird Schritt 2 sauber uebersprungen.

    python tests/test_datumsachse.py [ausgabeordner]

Rueckgabewert 0 = bestanden, 1 = fehlgeschlagen.
"""

import datetime as dt
import importlib.util
import os
import sys
import tempfile

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)
sys.path.insert(0, os.path.join(WURZEL, "tests"))
os.chdir(WURZEL)   # Vorlage/ und Daten/ werden relativ geladen

from modules.chart_dynamik import achsen_raster  # noqa: E402

NS_C = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
EPOCHE = dt.date(1899, 12, 30)      # Excel-Seriennummer 0

# Erwartete Beschriftungszahl je Chart — darunter wird die Achse nichtssagend,
# darueber unleserlich (7 pt, um 90 Grad gedreht).
MIN_TICKS, MAX_TICKS = 5, 20


def _d(seriennummer):
    """Excel-Seriennummer -> Datum."""
    return EPOCHE + dt.timedelta(days=int(float(seriennummer)))


def _monate(a, b):
    """Anzahl Monate von a bis b (beide Monatsanfaenge)."""
    return (b.year - a.year) * 12 + (b.month - a.month)


# ───────────────────────────── Schritt 1 ──────────────────────────────────

# Die echten Chart-Spannen (Datenstand 260721, aus den gebauten Broschueren
# ausgelesen) plus Grenzfaelle. Die Sollwerte sind von Hand nachgerechnet,
# nicht aus dem Code uebernommen.
#
# ACHTUNG: Der erste Chart-Punkt liegt einen Handelstag VOR dem ersten Tag der
# Zeitreihe — das ist der Indexstand 100 % (bei den cVV-Reihen der 31.12.2008,
# obwohl die Historie erst am 01.01.2009 beginnt, siehe HISTORIE_AB und
# Transferwissen #43).
#
#   (Bezeichnung, erster Tag, letzter Tag,
#    Soll-Minimum, Soll-Maximum, Soll-majorUnit, Soll-majorTimeUnit)
FAELLE = [
    # Lange Historien -> Jahresschritt.
    ("ETF (beide)",        (2015, 11, 30), (2026, 7, 21),
     (2015, 7, 1),  (2026, 8, 1), 1, "years"),
    ("cVV klassisch (4)",  (2008, 12, 31), (2026, 7, 21),
     (2008, 7, 1),  (2026, 8, 1), 1, "years"),
    ("cVV Dynamic",        (2018, 10, 5),  (2026, 7, 21),
     (2018, 7, 1),  (2026, 8, 1), 1, "years"),
    # Vergleichsfolie: Monatsreihe, beginnt Ende Januar -> hier gewinnt der
    # KALENDARISCHE Anker (Vorlauf 0 gegen 6 Monate). In der Vorlage steht an
    # dieser Achse majorUnit=12 + majorTimeUnit=months; wird majorUnit nicht
    # mitgezogen, entsteht daraus ein Tick alle zwoelf Jahre.
    ("cVV Vergleich",      (2009, 1, 31),  (2026, 7, 21),
     (2009, 1, 1),  (2026, 8, 1), 1, "years"),
    # Mittlere Historien -> Halbjahresschritt.
    ("ESG (alle vier)",    (2020, 9, 30),  (2026, 7, 21),
     (2020, 7, 1),  (2026, 8, 1), 6, "months"),
    ("SCHWEIZ (beide)",    (2022, 9, 12),  (2026, 7, 21),
     (2022, 7, 1),  (2026, 8, 1), 6, "months"),
    # Kurze Historien -> Quartalsschritt (vorher: Monatsticks, bis zu 37 Stueck).
    ("Pro",                (2023, 8, 31),  (2026, 7, 21),
     (2023, 7, 1),  (2026, 8, 1), 3, "months"),
    ("Pro Dividende",      (2024, 10, 21), (2026, 7, 21),
     (2024, 10, 1), (2026, 8, 1), 3, "months"),
    ("comdirect (alle 3)", (2024, 3, 11),  (2026, 7, 21),
     (2024, 1, 1),  (2026, 8, 1), 3, "months"),
    # ── Grenzfaelle ────────────────────────────────────────────────────────
    ("ein einziger Tag",   (2026, 7, 21),  (2026, 7, 21),
     (2026, 7, 1),  (2026, 8, 1), 3, "months"),
    ("Spanne unter 1 Monat", (2026, 7, 2),  (2026, 7, 21),
     (2026, 7, 1),  (2026, 8, 1), 3, "months"),
    ("Beginn im Januar, kurz", (2026, 1, 5), (2026, 7, 21),
     (2026, 1, 1),  (2026, 8, 1), 3, "months"),
    # Genau auf der Stufengrenze: 36 Monate Spanne -> noch Quartalsschritt.
    ("genau 36 Monate",    (2023, 8, 15),  (2026, 7, 21),
     (2023, 7, 1),  (2026, 8, 1), 3, "months"),
]


def _ticks(start, ende, schritt_monate):
    """Die Tick-Positionen, die PowerPoint zeichnet: ab dem Achsen-Minimum
    in festen Schritten, solange sie das Maximum nicht ueberschreiten."""
    raus, m = [], 0
    while True:
        mon = start.year * 12 + (start.month - 1) + m
        pos = dt.date(mon // 12, mon % 12 + 1, 1)
        if pos > ende:
            return raus
        raus.append(pos)
        m += schritt_monate


def _pruefe_rechenkern():
    print("1. Rechenkern achsen_raster (ohne PowerPoint)")
    kopf = (f"   {'Fall':24s} {'Achse ab':>10s} {'Vorl.':>5s} {'Schritt':>9s} "
            f"{'letzter Tick':>12s} {'n':>3s}  Ergebnis")
    print(kopf)
    print("   " + "-" * (len(kopf) - 3))

    fehler = 0
    for (name, erst, letzt, s_min, s_max, s_mu, s_mtu) in FAELLE:
        d_erst, d_letzt = dt.date(*erst), dt.date(*letzt)
        ist_min, ist_max, ist_mu, ist_mtu = achsen_raster(d_erst, d_letzt)

        schritt = ist_mu * (12 if ist_mtu == "years" else 1)
        ticks = _ticks(ist_min, ist_max, schritt)
        vorlauf = _monate(ist_min, dt.date(d_erst.year, d_erst.month, 1))

        ok = (ist_min == dt.date(*s_min) and ist_max == dt.date(*s_max)
              and ist_mu == s_mu and ist_mtu == s_mtu)
        # Der eigentliche Zweck: der letzte Tick liegt im Jahr des letzten
        # Datenpunkts — sonst fehlt dem Leser das aktuelle Jahr.
        if not ticks or ticks[-1].year != d_letzt.year:
            ok = False
        # Und die Achse darf keinen Datenpunkt abschneiden.
        if ist_min > d_erst or ist_max < d_letzt:
            ok = False
        fehler += 0 if ok else 1

        schritt_txt = (f"{ist_mu} Jahr(e)" if ist_mtu == "years"
                       else f"{ist_mu} Monate")
        print(f"   {name[:24]:24s} {ist_min.isoformat():>10s} "
              f"{vorlauf:4d}M {schritt_txt:>9s} "
              f"{(ticks[-1].isoformat() if ticks else '-'):>12s} "
              f"{len(ticks):3d}  {'OK' if ok else 'FEHLER'}")
        if not ok:
            print(f"        erwartet: ab {dt.date(*s_min)} bis {dt.date(*s_max)}, "
                  f"majorUnit={s_mu}, majorTimeUnit={s_mtu}")
    return fehler


# ───────────────────────────── Schritt 2 ──────────────────────────────────

def _achsen_befund(chart_shape):
    """Liest je Datumsachse des Charts, was tatsaechlich in der XML steht."""
    cs = chart_shape.chart._chartSpace
    raus = []
    for ax in cs.findall(".//c:dateAx", NS_C):
        werte = []
        for cat in cs.findall(".//c:cat", NS_C):
            werte += [float(v.text) for v in cat.iter(f"{{{NS_C['c']}}}v")
                      if v.text]
        if not werte:
            continue
        sc = ax.find("c:scaling", NS_C)
        mn = sc.find("c:min", NS_C) if sc is not None else None
        mx = sc.find("c:max", NS_C) if sc is not None else None
        mu = ax.find("c:majorUnit", NS_C)
        mtu = ax.find("c:majorTimeUnit", NS_C)
        raus.append({
            "erst": _d(min(werte)), "letzt": _d(max(werte)),
            "min": _d(mn.get("val")) if mn is not None else None,
            "max": _d(mx.get("val")) if mx is not None else None,
            "major_unit": mu.get("val") if mu is not None else None,
            "major_time_unit": mtu.get("val") if mtu is not None else None,
        })
    return raus


def _pruefe_datei(pfad, etikett):
    """Prueft jede Datumsachse einer gebauten Broschuere.
    Gibt (Anzahl geprueft, Anzahl Abweichungen) zurueck."""
    from pptx import Presentation

    geprueft = fehler = 0
    for nr, slide in enumerate(Presentation(pfad).slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            for b in _achsen_befund(shape):
                geprueft += 1
                maengel = []
                if b["min"] is None or b["max"] is None:
                    maengel.append("keine Achsengrenzen")
                if b["major_unit"] is None:
                    maengel.append("majorUnit fehlt")
                if b["major_time_unit"] is None:
                    maengel.append("majorTimeUnit fehlt")
                if maengel:
                    fehler += 1
                    print(f"   {etikett:12s} {nr:3d} "
                          f"{b['letzt'].isoformat():>10s}  FEHLER: "
                          f"{', '.join(maengel)}")
                    continue

                schritt = int(b["major_unit"]) * (
                    12 if b["major_time_unit"] == "years" else 1)
                ticks = _ticks(b["min"], b["max"], schritt)
                vorlauf = _monate(b["min"],
                                  dt.date(b["erst"].year, b["erst"].month, 1))

                # Der eigentliche Zweck der Uebung.
                if not ticks or ticks[-1].year != b["letzt"].year:
                    maengel.append("letzter Tick nicht im Jahr des letzten "
                                   "Datenpunkts")
                # Ein Achsendatum NACH den Daten waere eine falsche Aussage.
                if ticks and ticks[-1] > b["letzt"]:
                    maengel.append("Beschriftung liegt hinter den Daten")
                if b["min"] > b["erst"] or b["max"] < b["letzt"]:
                    maengel.append("Achse schneidet Daten ab")
                if not MIN_TICKS <= len(ticks) <= MAX_TICKS:
                    maengel.append(f"{len(ticks)} Beschriftungen")
                if vorlauf >= schritt:
                    maengel.append(f"{vorlauf} Monate Vorlauf")

                ok = not maengel
                fehler += 0 if ok else 1
                schritt_txt = (f"{b['major_unit']} Jahr(e)"
                               if b["major_time_unit"] == "years"
                               else f"{b['major_unit']} Monate")
                print(f"   {etikett:12s} {nr:3d} "
                      f"{b['letzt'].isoformat():>10s} "
                      f"{b['min'].isoformat():>10s} {vorlauf:4d}M "
                      f"{schritt_txt:>9s} "
                      f"{(ticks[-1].isoformat() if ticks else '-'):>12s} "
                      f"{len(ticks):3d}  {'OK' if ok else 'FEHLER'}")
                if maengel:
                    print(f"        {'; '.join(maengel)}")
    return geprueft, fehler


# Zusaetzliche Faelle ueber die Familien-Broschueren hinaus. Beide waren VOR
# dem 12.08.2026 die dichtesten Achsen ueberhaupt: die Themen-Duplikation
# brachte 37 bzw. 23 Monatsbeschriftungen, SCHWEIZ 47.
THEMA_ZUSATZ = [
    ("Thema x3", ["Offensiv", "Pro", "Pro Dividende"]),
    ("SCHWEIZ", ["Schweiz_substanzorientiert"]),
]


def _pruefe_artefakt(ausgabe):
    print("\n2. Wirkung am echten Artefakt (gebaute Broschueren)")
    # find_spec statt import: pyflakes kennt kein noqa, und ein importierter,
    # aber hier nicht benutzter Name waere eine Meldung (_pruefe_datei
    # importiert Presentation selbst).
    if importlib.util.find_spec("pptx") is None:
        print("   UEBERSPRUNGEN — python-pptx nicht installiert")
        return 0
    try:
        from modules.portfolioanalyse import (
            VORLAGEN_FAMILIEN, FAMILIE_ALLE_STRATEGIEN, _familien_portfolios,
            _familie_fuer_strategie, duration_info_aus_bestand,
        )
        # Datenbeschaffung und Bau kommen aus dem Export-Smoketest — beide
        # Suiten laufen damit ueber denselben Pfad wie die Oberflaeche.
        from test_export_smoke import _daten, _portfolio, _bauen
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    d = _daten()
    kopf = (f"   {'Broschuere':12s} {'Fo':>3s} {'Daten bis':>10s} "
            f"{'Achse ab':>10s} {'Vorl.':>5s} {'Schritt':>9s} "
            f"{'letzter Tick':>12s} {'n':>3s}  Ergebnis")
    print(kopf)
    print("   " + "-" * (len(kopf) - 3))

    fehler = 0
    geprueft = 0

    # ── je Familie eine Broschuere ─────────────────────────────────────────
    for familie in sorted(VORLAGEN_FAMILIEN):
        strategie = next((n for n in d["namen"]
                          if _familie_fuer_strategie(d["nm"], n) == familie), None)
        if strategie is None:
            print(f"   {familie:12s} UEBERSPRUNGEN (keine Strategie in den Daten)")
            continue
        alle = FAMILIE_ALLE_STRATEGIEN.get(familie)
        if alle:
            portfolios, fehlend = _familien_portfolios(
                alle, d["namen"], d["d2c"], d["pf_data"], duration_info_aus_bestand)
            if fehlend:
                print(f"   {familie:12s} UEBERSPRUNGEN (fehlende Daten: "
                      f"{', '.join(fehlend)})")
                continue
        else:
            portfolios = [_portfolio(strategie, d)]

        ziel, _gr, meldungen = _bauen(portfolios, familie, d, ausgabe,
                                      f"{familie}.pptx")
        for m in meldungen:
            print(f"   ! BUILD-FEHLER {familie}: {m[:90]}")
            fehler += 1
        n, f = _pruefe_datei(ziel, familie)
        geprueft += n
        fehler += f

    # ── Themen-Sonderfaelle (Duplikation und SCHWEIZ) ──────────────────────
    for etikett, namen in THEMA_ZUSATZ:
        fehlend = [n for n in namen if n not in d["d2c"]]
        if fehlend:
            print(f"   {etikett:12s} UEBERSPRUNGEN (nicht in den Daten: "
                  f"{', '.join(fehlend)})")
            continue
        portfolios = [_portfolio(n, d) for n in namen]
        ziel, _gr, meldungen = _bauen(
            portfolios, "Thema", d, ausgabe,
            f"Thema_{etikett.replace(' ', '_')}.pptx")
        for m in meldungen:
            print(f"   ! BUILD-FEHLER {etikett}: {m[:90]}")
            fehler += 1
        n, f = _pruefe_datei(ziel, etikett)
        geprueft += n
        fehler += f

    if not geprueft:
        print("   UEBERSPRUNGEN — kein Chart mit Datumsachse gefunden")
    else:
        print(f"\n   {geprueft} Datumsachsen geprueft")
    return fehler


def main():
    ausgabe = (sys.argv[1] if len(sys.argv) > 1
               else tempfile.mkdtemp(prefix="ffpb_datumsachse_"))
    os.makedirs(ausgabe, exist_ok=True)

    fehler = _pruefe_rechenkern() + _pruefe_artefakt(ausgabe)
    print()
    if fehler:
        print(f"FEHLGESCHLAGEN — {fehler} Abweichung(en)")
        return 1
    print("BESTANDEN — jede Datumsachse beschriftet den letzten Datenzeitraum,")
    print("            schneidet nichts ab und bleibt lesbar.")
    return 0


if __name__ == "__main__":
    sys.exit(main())

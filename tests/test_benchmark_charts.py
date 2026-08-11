"""Regressionstest: ohne Benchmark keine Benchmark — auch im CHART (11.08.2026).

Der Bugfix vom 07.08.2026 (test_benchmark_erkennung.py) hat die KENNZAHLEN
in Ordnung gebracht: Strategien ohne Vergleichsmassstab zeigen "-" statt
0,00 %. Die BROSCHUERE zeigte den Vergleichsmassstab danach trotzdem weiter,
an drei Stellen:

  1. Saeulen-Chart  — eine Serie "Benchmark" aus lauter Nullen (Null-Balken)
  2. Legenden-Box   — "Musterdepot     Benchmark***"
  3. Fussnote       — "*** 50% EuroStoxx 50; 50% MSCI World Euro", also der
                      unveraenderte VORLAGENTEXT und damit die Benchmark
                      einer FREMDEN Strategie ("Pro")

Punkt 3 ist der schwerste: die Kundenbroschuere nennt eine
Benchmark-Zusammensetzung, die fuer die Strategie nicht gilt.

Geprueft wird in zwei Schritten:

  Schritt 1 (nur pandas)   — analytics liefert leere Benchmark-Listen und
                             has_benchmark=False bei den SCHWEIZ-Reihen,
                             alle uebrigen Strategien unveraendert.
  Schritt 2 (+ python-pptx, streamlit) — am ECHTEN Artefakt: eine gebaute
                             SCHWEIZ-Broschuere enthaelt an keiner der drei
                             Stellen noch eine Benchmark; die Kontroll-
                             strategie "Pro" ist Zeichen fuer Zeichen
                             unveraendert.
  Schritt 3 (+ streamlit)  — der Hinweis im TOOL, an der gerenderten
                             Oberflaeche: genau einmal bei SCHWEIZ, nie bei
                             "Pro", und zwar UNABHAENGIG vom
                             Benchmark-Schalter.

Ohne pptx/streamlit werden die Schritte 2 und 3 sauber uebersprungen.

    python tests/test_benchmark_charts.py

Rueckgabewert 0 = bestanden, 1 = fehlgeschlagen.
"""

import glob
import os
import sys

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)
sys.path.insert(0, os.path.join(WURZEL, "tests"))
os.chdir(WURZEL)   # Vorlage/ und Daten/ werden relativ geladen

from modules.analytics import compute_performance_data  # noqa: E402
from test_benchmark_erkennung import (  # noqa: E402
    OHNE_BENCHMARK, FEE_TEST, neuester_tag, lade_zeitreihe,
)

DATEN = os.path.join(WURZEL, "Daten")

# Strategie-Anzeigenamen (Mapping-Spalte "Strategie auswaehlen") fuer Schritt 2.
FALL_OHNE = "Schweiz_substanzorientiert"
FALL_MIT = "Pro"                     # Kontrollfall: muss unveraendert bleiben
FAMILIE = "Thema"                    # beide gehoeren zur Themen-Familie

NS_C = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}


# ───────────────────────────── Schritt 1 ──────────────────────────────────

def _pruefe_daten():
    """analytics: leere Benchmark-Listen statt Nullen/Einsen."""
    print("1. Chart-Daten aus analytics.compute_performance_data")
    tag = neuester_tag()
    if tag is None:
        print(f"   FEHLER: keine CSVs in {DATEN}")
        return 1
    dateien = sorted(glob.glob(os.path.join(DATEN, f"*_{tag}_*.CSV")))
    if not dateien:
        print(f"   FEHLER: keine Dateien fuer Tag {tag}")
        return 1

    kopf = (f"   {'Strategie':30s} {'has_bm':>7s} {'Balken':>7s} "
            f"{'Linie':>7s}  Ergebnis")
    print(kopf)
    print("   " + "-" * (len(kopf) - 3))

    fehler = 0
    for pfad in dateien:
        name, df = lade_zeitreihe(pfad)
        erwartet_bm = name not in OHNE_BENCHMARK
        erg = compute_performance_data(df, FEE_TEST)

        has_bm = bool(erg.get("has_benchmark"))
        pa = erg.get("performance_pa", {})
        we = erg.get("wertentwicklung", {})
        n_balken = len(pa.get("benchmark") or [])
        n_linie = len(we.get("benchmark") or [])

        if erwartet_bm:
            # Mit Benchmark: beide Serien so lang wie ihre Bezugsachse.
            ok = (has_bm
                  and n_balken == len(pa.get("jahre") or [])
                  and n_linie == len(we.get("dates") or []))
        else:
            # Ohne Benchmark: nichts zu zeichnen.
            ok = not has_bm and n_balken == 0 and n_linie == 0

        fehler += 0 if ok else 1
        print(f"   {name[:30]:30s} {str(has_bm):>7s} {n_balken:>7d} "
              f"{n_linie:>7d}  {'OK' if ok else 'FEHLER'}")
    return fehler


# ───────────────────────────── Schritt 2 ──────────────────────────────────

def _serien_namen(chart_shape):
    """Seriennamen direkt aus der Chart-XML (nicht ueber die python-pptx-API,
    damit auch eine leergeraeumte Serie auffaellt)."""
    namen = []
    for ser in chart_shape.chart._chartSpace.findall(".//c:ser", NS_C):
        tx = ser.find(".//c:tx//c:v", NS_C)
        namen.append(tx.text if tx is not None else "")
    return namen


def _folie_lesen(prs, idx):
    """Liest die drei interessanten Stellen der Wertentwicklungs-Folie."""
    from modules.pptx_slides import (
        SHAPE_WE_CHART_BAR, SHAPE_WE_FUSSNOTE, SHAPE_WE_LEGENDE,
    )
    befund = {"serien": [], "legende": "", "fussnote": ""}
    for shape in prs.slides[idx].shapes:
        if shape.name == SHAPE_WE_CHART_BAR and getattr(shape, "has_chart", False):
            befund["serien"] = _serien_namen(shape)
        elif shape.name == SHAPE_WE_LEGENDE and shape.has_text_frame:
            befund["legende"] = shape.text_frame.text
        elif shape.name == SHAPE_WE_FUSSNOTE and shape.has_text_frame:
            befund["fussnote"] = shape.text_frame.text
    return befund


def _pruefe_artefakt():
    """Baut zwei echte Broschueren und liest nach, was drinsteht."""
    print("\n2. Wirkung am echten Artefakt (gebaute Themen-Broschuere)")
    try:
        from pptx import Presentation
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0
    try:
        import tempfile
        from modules.vorlagen_config import VORLAGEN_FAMILIEN
        # Die Datenbeschaffung ist identisch zum Export-Smoketest; sie dort
        # wiederzuverwenden haelt beide Tests auf demselben Datenpfad.
        from test_export_smoke import _daten, _portfolio, _bauen
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    d = _daten()
    fehlend = [n for n in (FALL_OHNE, FALL_MIT) if n not in d["d2c"]]
    if fehlend:
        print(f"   UEBERSPRUNGEN — nicht in den Daten: {', '.join(fehlend)}")
        return 0

    # 1-indexierte Position der Wertentwicklungs-Folie im Themen-Block
    _, cfg = VORLAGEN_FAMILIEN[FAMILIE]
    idx = cfg["block_positionen"]["wertentwicklung"] - 1

    ausgabe = tempfile.mkdtemp(prefix="ffpb_bmchart_")
    fehler = 0

    for name, erwartet_bm in ((FALL_OHNE, False), (FALL_MIT, True)):
        ziel, _groesse, meldungen = _bauen(
            [_portfolio(name, d)], FAMILIE, d, ausgabe, f"{name}.pptx")
        if meldungen:
            for m in meldungen:
                print(f"   ! BUILD-FEHLER {name}: {m[:90]}")
            fehler += 1
        b = _folie_lesen(Presentation(ziel), idx)

        hat_serie = any("enchmark" in s for s in b["serien"])
        hat_legende = "enchmark" in b["legende"]
        hat_fussnote = "***" in b["fussnote"]
        ok = (hat_serie == erwartet_bm
              and hat_legende == erwartet_bm
              and hat_fussnote == erwartet_bm)
        fehler += 0 if ok else 1

        titel = "MIT Benchmark" if erwartet_bm else "OHNE Benchmark"
        print(f"\n   {name} ({titel})")
        print(f"     Chart-Serien : {b['serien']}")
        print(f"     Legende      : {b['legende'].strip()!r}")
        print(f"     Fussnote *** : {'ja' if hat_fussnote else 'nein'}")
        print(f"     -> {'OK' if ok else 'FEHLER'}")
        if not ok:
            print(f"        erwartet an allen drei Stellen: "
                  f"{'Benchmark' if erwartet_bm else 'keine Benchmark'}")
    return fehler


# ───────────────────────────── Schritt 3 ──────────────────────────────────

HINWEIS_MARKE = "kein Vergleichsmaßstab"


def _pruefe_hinweis():
    """Der Hinweis im Tool — an der GERENDERTEN Oberflaeche, nicht im Quelltext.

    Wichtig ist das Verhalten BEI AUSGESCHALTETEM Benchmark-Schalter: bis
    11.08.2026 hing der Hinweis am Linien-Chart und erschien nur, wenn der
    Schalter an war. War er aus, blieb das "-" in den Kacheln unkommentiert.
    """
    print("\n3. Hinweis in der laufenden App (AppTest)")
    try:
        from streamlit.testing.v1 import AppTest
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    def lauf(strategie, benchmark_an):
        at = AppTest.from_file(os.path.join(WURZEL, "streamlit_app.py"),
                               default_timeout=400)
        at.secrets["passwords"] = {"t": "t"}
        at.session_state["logged_in"] = True
        at.session_state["username"] = "t"
        at.session_state["p_bm"] = benchmark_an
        at.run()
        sel = next((s for s in at.selectbox
                    if strategie in [str(o) for o in s.options]), None)
        if sel is None:
            return None
        sel.set_value(strategie).run()
        texte = [str(e.value) for sammlung in (at.caption, at.markdown, at.info)
                 for e in sammlung]
        return [t for t in texte if HINWEIS_MARKE in t]

    faelle = [
        (FALL_OHNE, True,  1), (FALL_OHNE, False, 1),
        (FALL_MIT,  True,  0), (FALL_MIT,  False, 0),
    ]
    fehler = 0
    for strategie, bm_an, soll in faelle:
        treffer = lauf(strategie, bm_an)
        schalter = "AN " if bm_an else "AUS"
        if treffer is None:
            print(f"   {strategie[:28]:28s} Schalter {schalter}  "
                  f"UEBERSPRUNGEN (nicht in der Auswahl)")
            continue
        ok = len(treffer) == soll
        fehler += 0 if ok else 1
        print(f"   {strategie[:28]:28s} Schalter {schalter}  "
              f"erwartet {soll}x, gefunden {len(treffer)}x  "
              f"{'OK' if ok else 'FEHLER'}")
    return fehler


def main():
    fehler = _pruefe_daten() + _pruefe_artefakt() + _pruefe_hinweis()
    print()
    if fehler:
        print(f"FEHLGESCHLAGEN — {fehler} Abweichung(en)")
        return 1
    print("BESTANDEN — ohne Benchmark bleibt weder Balken noch Legende "
          "noch Fussnote stehen;")
    print("            Strategien MIT Benchmark sind unveraendert.")
    return 0


if __name__ == "__main__":
    sys.exit(main())

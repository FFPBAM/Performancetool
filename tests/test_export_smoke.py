"""End-to-End-Test: erzeugt fuer jede Familie eine echte Broschuere.

Bildet nach, was render_portfolioanalyse beim Klick auf "PowerPoint
erstellen" tut — ohne Oberflaeche. Geprueft wird je Datei:
  - der Export laeuft ohne Exception durch
  - die Folienzahl stimmt (bei Thema inkl. Duplikation fuer N Strategien)
  - LAST_BUILD_ERRORS ist leer
  - die erzeugte Datei laesst sich wieder oeffnen

Braucht python-pptx, streamlit und die echten Daten. Fehlt etwas davon,
wird sauber uebersprungen statt zu scheitern.

    python tests/test_export_smoke.py [ausgabeordner]

Die erzeugten PPTX bitte STICHPROBENARTIG IN ECHTEM PowerPoint oeffnen —
LibreOffice reicht nicht (Transferwissen #16/#28).
"""

import os
import sys
import tempfile
import traceback

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)

# Vorlage/ und Daten/ werden im Code relativ zum Arbeitsverzeichnis geladen
os.chdir(WURZEL)

try:
    from pptx import Presentation
except ImportError:
    print("UEBERSPRUNGEN — python-pptx nicht installiert")
    sys.exit(0)

try:
    from modules.shared import (
        DATA_FOLDER, DATA_FOLDER_PF, EXCLUDE_SUBSTRINGS,
        detect_newest_date_tag, load_all_csvs, load_mapping,
        load_name_mapping, build_portfolio_timeseries,
    )
    from modules.portfolioanalyse import (
        load_pf_csvs, build_pf_data, duration_info_aus_bestand,
        VORLAGEN_FAMILIEN, FAMILIE_ALLE_STRATEGIEN, _familien_portfolios,
        _vorlage_fuer_familie, _familie_fuer_strategie,
    )
    from modules import pptx_export
    from modules.pptx_export import generate_portfolioanalyse_pptx
except ImportError as ex:
    print(f"UEBERSPRUNGEN — Abhaengigkeit fehlt: {ex}")
    sys.exit(0)

# Familie "Thema" laeuft als einzige im Dupliziermodus: pro zusaetzlicher
# Strategie waechst die Broschuere um den Block F10-13 (vier Folien).
THEMA_STRATEGIEN = ["Offensiv", "Pro", "Pro Dividende"]
THEMA_BLOCK = 4
THEMA_BASIS = 21


def _daten():
    tag = detect_newest_date_tag(DATA_FOLDER_PF, EXCLUDE_SUBSTRINGS)
    pf_data = build_pf_data(load_pf_csvs(DATA_FOLDER_PF, tag))
    nm = load_name_mapping()
    sp = nm.columns
    gefiltert = nm[nm[sp[1]].isin(set(pf_data.keys()))]
    mapping = load_mapping()
    ts = build_portfolio_timeseries(
        load_all_csvs(DATA_FOLDER,
                      detect_newest_date_tag(DATA_FOLDER, EXCLUDE_SUBSTRINGS),
                      EXCLUDE_SUBSTRINGS),
        mapping)
    return {
        "tag": tag,
        "pf_data": pf_data,
        "nm": nm,
        "namen": gefiltert[sp[0]].tolist(),
        "d2c": dict(zip(gefiltert[sp[0]], gefiltert[sp[1]])),
        "d2b": dict(zip(gefiltert[sp[0]], gefiltert[sp[3]])),
        "mapping": mapping,
        "ts": ts,
    }


def _portfolio(name, d):
    csv_n = d["d2c"][name]
    df = d["pf_data"][csv_n]
    ad = (df["Auswertungsdatum"].iloc[0]
          if "Auswertungsdatum" in df.columns else None)
    return (name, df, ad, duration_info_aus_bestand(df))


def _perf_inputs(portfolios, d):
    raus = []
    for name, _df, _ad, dur in portfolios:
        csv_n = d["d2c"].get(name)
        treffer = d["mapping"].loc[d["mapping"]["Inhaber"] == csv_n,
                                   "Honorarsatz Standard"]
        bm = d["d2b"].get(name)
        if bm is None or str(bm).strip().lower() in (
                "", "nan", "none", "haben keine benchmark"):
            bm = None
        else:
            bm = str(bm).strip()
        raus.append({
            "timeseries_df": d["ts"].get(csv_n) if csv_n else None,
            "fee_dec": float(treffer.values[0]) if len(treffer) else 0.0,
            "duration": dur.get("duration") if isinstance(dur, dict) else None,
            "benchmark_text": bm,
        })
    return raus


def _bauen(portfolios, familie, d, ausgabe, dateiname):
    tpl, cfg = _vorlage_fuer_familie(familie)
    daten = generate_portfolioanalyse_pptx(
        portfolios, 0.0, performance_inputs=_perf_inputs(portfolios, d),
        template_path=tpl, template_config=cfg)
    ziel = os.path.join(ausgabe, dateiname)
    with open(ziel, "wb") as f:
        f.write(daten)
    return ziel, len(daten), list(pptx_export.LAST_BUILD_ERRORS)


def main():
    ausgabe = sys.argv[1] if len(sys.argv) > 1 else tempfile.mkdtemp(prefix="ffpb_export_")
    os.makedirs(ausgabe, exist_ok=True)

    d = _daten()
    print(f"Datenstand {d['tag']} — {len(d['pf_data'])} Portfolios, "
          f"{len(d['ts'])} Zeitreihen")
    print(f"Ausgabe: {ausgabe}\n")
    print(f"{'Fall':28s} {'Folien':>6s} {'soll':>5s} {'MB':>6s}  Status")
    print("-" * 68)

    fehler = 0

    # ── Teil 1: je Familie eine Broschuere ──────────────────────────────
    for familie in sorted(VORLAGEN_FAMILIEN):
        strategie = next((n for n in d["namen"]
                          if _familie_fuer_strategie(d["nm"], n) == familie), None)
        if strategie is None:
            print(f"{familie:28s} {'-':>6s} {'-':>5s} {'-':>6s}  "
                  f"UEBERSPRUNGEN (keine Strategie in den Daten)")
            continue
        try:
            alle = FAMILIE_ALLE_STRATEGIEN.get(familie)
            if alle:
                portfolios, fehlend = _familien_portfolios(
                    alle, d["namen"], d["d2c"], d["pf_data"],
                    duration_info_aus_bestand)
                if fehlend:
                    print(f"{familie:28s} {'-':>6s} {'-':>5s} {'-':>6s}  "
                          f"UEBERSPRUNGEN (fehlende Daten: {', '.join(fehlend)})")
                    continue
            else:
                portfolios = [_portfolio(strategie, d)]

            ziel, groesse, meldungen = _bauen(portfolios, familie, d, ausgabe,
                                              f"{familie}.pptx")
            n = len(Presentation(ziel).slides)
            soll = VORLAGEN_FAMILIEN[familie][1].get("erwartete_folien")
            # Im Dupliziermodus waechst die Folienzahl mit den Strategien
            if VORLAGEN_FAMILIEN[familie][1].get("block_positionen"):
                soll += THEMA_BLOCK * (len(portfolios) - 1)
            ok = n == soll and not meldungen
            fehler += 0 if ok else 1
            print(f"{familie:28s} {n:6d} {soll:5d} {groesse/1048576:6.2f}  "
                  f"{'OK' if ok else 'ABWEICHUNG'}")
            for m in meldungen:
                print(f"    ! {m[:96]}")
        except Exception as ex:
            fehler += 1
            print(f"{familie:28s} {'-':>6s} {'-':>5s} {'-':>6s}  "
                  f"FEHLER: {type(ex).__name__}: {ex}")
            traceback.print_exc()

    # ── Teil 2: Thema mit mehreren Strategien (Dupliziermodus) ──────────
    # Der einzige Pfad, auf dem _vervielfaeltige_block laeuft. Waere
    # _THEMA_CONFIG faelschlich auf modus="fest" gestellt, blieben es
    # immer 21 Folien und die Zusatzstrategien haetten stillschweigend keine.
    for anzahl in (2, 3):
        namen = [n for n in THEMA_STRATEGIEN if n in d["d2c"]][:anzahl]
        if len(namen) < anzahl:
            print(f"{'Thema x' + str(anzahl):28s} {'-':>6s} {'-':>5s} {'-':>6s}  "
                  f"UEBERSPRUNGEN (nur {len(namen)} Themen-Strategien in den Daten)")
            continue
        try:
            portfolios = [_portfolio(n, d) for n in namen]
            ziel, groesse, meldungen = _bauen(portfolios, "Thema", d, ausgabe,
                                              f"Thema_{anzahl}.pptx")
            n = len(Presentation(ziel).slides)
            soll = THEMA_BASIS + THEMA_BLOCK * (anzahl - 1)
            ok = n == soll and not meldungen
            fehler += 0 if ok else 1
            print(f"{'Thema x' + str(anzahl) + ' (Duplikation)':28s} {n:6d} {soll:5d} "
                  f"{groesse/1048576:6.2f}  {'OK' if ok else 'ABWEICHUNG'}")
            for m in meldungen:
                print(f"    ! {m[:96]}")
        except Exception as ex:
            fehler += 1
            print(f"{'Thema x' + str(anzahl):28s} {'-':>6s} {'-':>5s} {'-':>6s}  "
                  f"FEHLER: {type(ex).__name__}: {ex}")
            traceback.print_exc()

    print()
    if fehler:
        print(f"FEHLGESCHLAGEN — {fehler} Fall/Faelle")
        return 1
    print("BESTANDEN — alle Broschueren erzeugt und wieder lesbar")
    print("Hinweis: stichprobenartig in ECHTEM PowerPoint oeffnen (#16/#28).")
    return 0


if __name__ == "__main__":
    sys.exit(main())

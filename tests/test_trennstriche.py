"""Prueft die Trennstriche der Positionstabellen in erzeugten Broschueren.

REGEL (festgelegt mit Philip am 07.08.2026):
    Der DICKE Trennstrich steht genau unter einer Kategorie-Ueberschrift
    (RENTEN, AKTIEN, EDELMETALLE, LIQUIDITAET). Zwischen den einzelnen
    Positionen und vor der naechsten Ueberschrift steht die duenne Linie.

        Wuerth Finance IHS 3 %     <- duenn darunter
        AKTIEN                     <- DICK darunter
        Future of Defence ETF      <- duenn darunter

HINTERGRUND: fill_table_with_positions schrieb frueher nur Text und
Fettung; die Rahmenlinien blieben dort stehen, wo die VORLAGE ihre Gruppen
hatte. Bei "CVV Defensiv" lief der dicke Strich dadurch mitten durch die
Rentenliste (zwischen Fraport und Fresenius), waehrend der Uebergang
Wuerth -> AKTIEN keinen bekam. Ueber alle Familien waren es 80 falsch
platzierte Striche. Behoben durch tabelle_kategorie_trennlinien.

Aufruf mit einem Ordner voller erzeugter PPTX:

    python tests/test_export_smoke.py /pfad/zur/ausgabe
    python tests/test_trennstriche.py /pfad/zur/ausgabe

Ohne Argument wird ein temporaerer Ordner erzeugt und vorher exportiert.
"""

import glob
import os
import subprocess
import sys
import tempfile

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
except ImportError:
    print("UEBERSPRUNGEN — python-pptx nicht installiert")
    sys.exit(0)

DICK_AB_PT = 0.5      # alles darueber gilt als dicker Trennstrich
TABELLEN_SHAPES = ("T_Kennzahlen", "Tabelle Einzeltitel")


def _breite_pt(cell, tag):
    tcPr = cell._tc.find(qn("a:tcPr"))
    if tcPr is None:
        return None
    el = tcPr.find(qn("a:" + tag))
    if el is None or el.find(qn("a:noFill")) is not None:
        return None
    w = el.get("w")
    return int(w) / 12700.0 if w else None


def _ist_gruppenzeile(cell):
    if not " ".join(cell.text_frame.text.split()):
        return False
    return any(run.font.bold
               for p in cell.text_frame.paragraphs for run in p.runs)


def pruefe_datei(pfad):
    """Returns Liste von Befunden (leer = alles korrekt)."""
    befunde = []
    for s_idx, slide in enumerate(Presentation(pfad).slides, start=1):
        for sh in slide.shapes:
            if not getattr(sh, "has_table", False) or sh.name not in TABELLEN_SHAPES:
                continue
            zeilen = [list(r.cells)[0] for r in sh.table.rows]
            # Kopfzeile (0) und die beiden Schlusszeilen (Abstand + Summe)
            # gehoeren nicht zur Regel — die regelt eigene Logik.
            erste, letzte = 1, len(zeilen) - 3
            for i in range(erste, letzte):
                dick = (_breite_pt(zeilen[i], "lnB") or 0) >= DICK_AB_PT
                soll_dick = _ist_gruppenzeile(zeilen[i])
                if dick != soll_dick:
                    befunde.append((
                        s_idx, i,
                        "DICK, obwohl keine Ueberschrift" if dick
                        else "DUENN unter einer Ueberschrift",
                        " ".join(zeilen[i].text_frame.text.split())[:28]))
    return befunde


def main():
    if len(sys.argv) > 1:
        ordner = sys.argv[1]
    else:
        ordner = tempfile.mkdtemp(prefix="ffpb_trenn_")
        print(f"Kein Ordner angegeben — exportiere zuerst nach {ordner}\n")
        r = subprocess.run([sys.executable,
                            os.path.join(WURZEL, "tests", "test_export_smoke.py"),
                            ordner], capture_output=True, text=True)
        if r.returncode != 0:
            print("Export fehlgeschlagen:")
            print(r.stdout[-2000:])
            return 1

    dateien = sorted(glob.glob(os.path.join(ordner, "*.pptx")))
    if not dateien:
        print(f"Keine PPTX in {ordner}")
        return 1

    gesamt = 0
    for pfad in dateien:
        befunde = pruefe_datei(pfad)
        name = os.path.basename(pfad)
        if befunde:
            gesamt += len(befunde)
            print(f"{name:18s} {len(befunde):3d} Abweichung(en)")
            for s, i, art, text in befunde[:10]:
                print(f"    Folie {s:2d} Zeile {i:2d}  {art:32s} {text}")
            if len(befunde) > 10:
                print(f"    ... und {len(befunde)-10} weitere")
        else:
            print(f"{name:18s}   OK")

    print()
    if gesamt:
        print(f"FEHLGESCHLAGEN — {gesamt} falsch platzierte Trennstriche")
        return 1
    print("BESTANDEN — dicker Trennstrich steht ueberall genau unter der "
          "Kategorie-Ueberschrift")
    return 0


if __name__ == "__main__":
    sys.exit(main())

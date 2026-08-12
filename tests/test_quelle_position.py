"""Regressionstest: die Quellenangabe darf nicht im Disclaimer liegen (12.08.2026).

GEMELDET wurde eine Folie — die Wertentwicklungs-Folie der Offensiv-Broschuere:
"Quelle: Eigene Berechnung, Stand 20.07.2026" wird vom Disclaimer-Fliesstext
ueberdruckt. NACHGEMESSEN am PowerPoint-Rendering waren es alle: In allen sechs
Vorlagen liegt die Textbox "Quelle" auf den Emu identisch bei

    Quelle   23,30-28,10 x 13,89-14,19 cm
    Fussnote 12,50-28,10 x 11,16-16,20 cm     <- die Quelle liegt DARIN

und der Disclaimer reicht gerendert bis 14,47 cm. Betroffen sind 16 Folien:

    Vorlage_cVV_Infoboard.pptx   F8, F10, F12, F14, F16
    Vorlage_ESG.pptx             F17, F19, F21, F23
    Vorlage_comdirect.pptx       F7, F9, F11
    Vorlage_ETF.pptx             F17, F19
    Vorlage_Thema.pptx           F12          (Offensiv, Pro, Pro Dividende)
    Vorlage_FFPB.pptx            F11

Alle 16 tragen die Rolle "wertentwicklung" und laufen durch EINE Funktion,
pptx_slides.fill_wertentwicklung_slide — deshalb genuegt dort eine Korrektur
(WE_QUELLE_TOP_CM).

ZWEI URSACHEN, die sich addieren:

  1. Die Vorlagen-Position der Quelle liegt im Fussnotenfeld. Das faellt nicht
     auf, solange der Disclaimer kurz genug bleibt — er ist es nie gewesen.
  2. Der Disclaimer ist in der Vorlage HART umbrochen: Die Absaetze sind von
     Hand auf die Boxbreite gebrochen (laengste Vorlagenzeile 149 Zeichen bei
     6 pt). WE_DISCLAIMER_REPLACEMENTS ersetzte einen davon durch 189 Zeichen
     — der Absatz bricht um, ALLES darunter rutscht eine Zeile tiefer. Der
     Kommentar ueber der Konstanten versprach seit Juli 2026, die Ersatztexte
     seien "auf aehnliche Laenge kalibriert, damit das Layout haelt".
     Gemessen hat das nie jemand. Schritt 2 misst es.

Geprueft wird in drei Schritten:

  Schritt 1 — an allen sechs VORLAGEN: der Fussnoten-Textblock wird mit den
                     echten Ersatztexten nachgerechnet, und WE_QUELLE_TOP_CM
                     muss darunter liegen. Eine neue Vorlage mit laengerem
                     Disclaimer faellt hier auf.
  Schritt 2 — jeder Ersatztext passt in eine Zeile (WE_FUSSNOTE_ZEILE_MAX).
                     Das ist der Test, der den Fehler verhindert haette.
  Schritt 3 (+ streamlit) — am ECHTEN Artefakt: je Familie eine gebaute
                     Broschuere; auf jeder Wertentwicklungs-Folie liegt die
                     Quelle vollstaendig unter dem Disclaimer, ist nicht leer
                     und traegt das Stand-Datum.

Die Konstanten liegen in pptx_slides, das pandas und python-pptx hereinzieht —
ohne die beiden ueberspringt sich die Suite ganz. Fehlt nur streamlit, laufen
Schritt 1 und 2 und Schritt 3 wird uebersprungen.

    python tests/test_quelle_position.py [ausgabeordner]

Rueckgabewert 0 = bestanden, 1 = fehlgeschlagen.
"""

import importlib.util
import math
import os
import sys
import tempfile

WURZEL = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, WURZEL)
sys.path.insert(0, os.path.join(WURZEL, "tests"))
os.chdir(WURZEL)   # Vorlage/ und Daten/ werden relativ geladen

try:
    from modules.pptx_slides import (                              # noqa: E402
        SHAPE_WE_FUSSNOTE, SHAPE_WE_QUELLE,
        WE_DISCLAIMER_REPLACEMENTS, WE_FUSSNOTE_ZEILE_MAX, WE_QUELLE_TOP_CM,
        WE_FOOTNOTE_STAR1_NEW, WE_FOOTNOTE_STAR1_PREFIX,
        WE_FOOTNOTE_STAR2_NEW, WE_FOOTNOTE_STAR2_PREFIX,
    )
except ImportError as ex:
    # pptx_slides zieht pandas und python-pptx herein. Ohne die beiden ist
    # hier nichts zu messen — ueberspringen statt scheitern.
    print(f"UEBERSPRUNGEN — Abhaengigkeit fehlt: {ex}")
    sys.exit(0)

EMU_PRO_CM = 360000

# ── Rendering-Kennwerte, am Artefakt gemessen (nicht geschaetzt) ────────────
# Quelle: Thema.pptx F12, von PowerPoint 16.0 als PNG mit 1920x1225 px
# ausgegeben (Folie 29,0 x 18,5 cm) und zeilenweise vermessen. Die erste
# Textzeile beginnt 0,06 cm unter der Boxoberkante; von der ersten bis zur
# dreizehnten Zeile sind es 3,04 cm, also 0,2533 cm Zeilenabstand (6 pt,
# Zeilenabstand 100 %). Die Unterlaenge der Glyphen kommt unten drauf.
ZEILENHOEHE_CM = 0.2533
ERSTE_ZEILE_CM = 0.06
GLYPHENHOEHE_CM = 0.21

# Mindestabstand zwischen Disclaimer-Unterkante und Quellenzeile. Weniger
# waere zwar kollisionsfrei, sieht aber nach Versehen aus.
MINDESTABSTAND_CM = 0.15

VORLAGEN = [
    "Vorlage_cVV_Infoboard.pptx", "Vorlage_ESG.pptx", "Vorlage_comdirect.pptx",
    "Vorlage_ETF.pptx", "Vorlage_Thema.pptx", "Vorlage_FFPB.pptx",
]

# So viele Folien mit Quelle UND Disclaimer traegt jede Vorlage. Sinkt eine
# Zahl, ist eine Folie aus der Vorlage verschwunden — dann greift die
# Korrektur dort nicht mehr, ohne dass irgendetwas scheitern wuerde.
#
# Vorlage_FFPB.pptx hat ZWEI: F11 ist die Wertentwicklungs-Folie (betroffen),
# F10 die Performance-Folie. Deren Disclaimer ist kuerzer (9 Zeilen bis
# 13,45 cm), die Quelle steht dort bei 13,89 cm also frei — sie wird bewusst
# NICHT verschoben. Geprueft wird die Zusage ("Quelle unter dem Disclaimer"),
# nicht eine feste Koordinate.
FOLIEN_SOLL = {
    "Vorlage_cVV_Infoboard.pptx": 5, "Vorlage_ESG.pptx": 4,
    "Vorlage_comdirect.pptx": 3, "Vorlage_ETF.pptx": 2,
    "Vorlage_Thema.pptx": 1, "Vorlage_FFPB.pptx": 2,
}


def _cm(emu):
    return None if emu is None else emu / EMU_PRO_CM


def _we_folien(prs):
    """Alle Folien, die eine Quellenangabe UND einen Disclaimer tragen.

    Erkannt wird an den beiden Shapes statt an einer Folienliste: So findet
    der Test sie auch in einer gebauten Broschuere, in der die Themen-Folien
    vervielfaeltigt wurden und die Positionen sich verschoben haben — und er
    findet eine neue Folie, die dieselbe Konstellation erbt, ohne dass jemand
    daran denkt, hier eine Nummer nachzutragen.
    """
    raus = []
    for nr, slide in enumerate(prs.slides, start=1):
        namen = {sh.name: sh for sh in slide.shapes}
        if SHAPE_WE_FUSSNOTE in namen and SHAPE_WE_QUELLE in namen:
            raus.append((nr, namen[SHAPE_WE_FUSSNOTE], namen[SHAPE_WE_QUELLE]))
    return raus


def _fussnote_unterkante(fussnote):
    """Untere Kante des gerenderten Fussnotentextes in cm.

    Der Disclaimer ist hart umbrochen; jeder Absatz belegt eine Zeile, ein zu
    langer Absatz bricht zusaetzlich um. Leere Absaetze zaehlen mit — sie
    erzeugen eine Leerzeile.
    """
    zeilen = 0
    for para in fussnote.text_frame.paragraphs:
        zeilen += max(1, math.ceil(len(para.text) / WE_FUSSNOTE_ZEILE_MAX))
    return (_cm(fussnote.top) + ERSTE_ZEILE_CM
            + (zeilen - 1) * ZEILENHOEHE_CM + GLYPHENHOEHE_CM), zeilen


# ───────────────────────────── Schritt 1 ──────────────────────────────────

def _ersatztexte_anwenden(fussnote):
    """Schreibt die Vorlagen-Fussnote so um, wie der Export es tut.

    Bewusst ueber die PRODUKTIVE Funktion und die produktiven Konstanten —
    eine Nachbildung wuerde genau das nicht messen, worauf es ankommt.
    Die ***-Zeile bleibt Vorlagentext: Ihre echte Laenge haengt an der
    Strategie und wird in Schritt 3 am Artefakt geprueft.
    """
    from modules.pptx_slides import replace_paragraph_text_by_prefix
    tf = fussnote.text_frame
    replace_paragraph_text_by_prefix(tf, WE_FOOTNOTE_STAR1_PREFIX,
                                     WE_FOOTNOTE_STAR1_NEW)
    replace_paragraph_text_by_prefix(tf, WE_FOOTNOTE_STAR2_PREFIX,
                                     WE_FOOTNOTE_STAR2_NEW)
    for prefix, neu in WE_DISCLAIMER_REPLACEMENTS:
        replace_paragraph_text_by_prefix(tf, prefix, neu)


def _pruefe_vorlagen():
    print("1. Vorlagen: liegt WE_QUELLE_TOP_CM unter dem Disclaimer?")
    if importlib.util.find_spec("pptx") is None:
        print("   UEBERSPRUNGEN — python-pptx nicht installiert")
        return 0
    from pptx import Presentation

    kopf = (f"   {'Vorlage':28s} {'Fo':>3s} {'Zeilen':>6s} {'Unterkante':>10s} "
            f"{'Quelle neu':>10s} {'Luft':>6s}  Ergebnis")
    print(kopf)
    print("   " + "-" * (len(kopf) - 3))

    fehler = 0
    for datei in VORLAGEN:
        pfad = os.path.join("Vorlage", datei)
        if not os.path.exists(pfad):
            print(f"   {datei:28s} UEBERSPRUNGEN (nicht vorhanden)")
            continue
        folien = _we_folien(Presentation(pfad))
        soll = FOLIEN_SOLL.get(datei)
        if soll is not None and len(folien) != soll:
            fehler += 1
            print(f"   {datei:28s} FEHLER: {len(folien)} statt {soll} "
                  f"Wertentwicklungs-Folien")
        for nr, fussnote, quelle in folien:
            _ersatztexte_anwenden(fussnote)
            unten, zeilen = _fussnote_unterkante(fussnote)
            luft = WE_QUELLE_TOP_CM - unten
            maengel = []
            if luft < MINDESTABSTAND_CM:
                maengel.append(f"nur {luft:.2f} cm Abstand zum Disclaimer "
                               f"(mindestens {MINDESTABSTAND_CM:.2f})")
            # Die Box darf nicht in die Foliennummer oder aus der Folie laufen.
            unterkante_quelle = WE_QUELLE_TOP_CM + _cm(quelle.height)
            if unterkante_quelle > 17.0:
                maengel.append(f"Quelle endet bei {unterkante_quelle:.2f} cm "
                               f"— zu nah an der Foliennummer")
            ok = not maengel
            fehler += 0 if ok else 1
            print(f"   {datei:28s} {nr:3d} {zeilen:6d} {unten:10.2f} "
                  f"{WE_QUELLE_TOP_CM:10.2f} {luft:6.2f}  "
                  f"{'OK' if ok else 'FEHLER'}")
            for m in maengel:
                print(f"        {m}")
    return fehler


# ───────────────────────────── Schritt 2 ──────────────────────────────────

def _pruefe_zeilenlaengen():
    print("\n2. Ersatztexte: passt jeder in eine Zeile?")
    faelle = [("WE_FOOTNOTE_STAR1_NEW", WE_FOOTNOTE_STAR1_NEW),
              ("WE_FOOTNOTE_STAR2_NEW", WE_FOOTNOTE_STAR2_NEW)]
    faelle += [(f"WE_DISCLAIMER_REPLACEMENTS[{i}]", neu)
               for i, (_p, neu) in enumerate(WE_DISCLAIMER_REPLACEMENTS)]

    fehler = 0
    for name, text in faelle:
        ok = len(text) <= WE_FUSSNOTE_ZEILE_MAX
        fehler += 0 if ok else 1
        print(f"   {name:30s} {len(text):4d} / {WE_FUSSNOTE_ZEILE_MAX}  "
              f"{'OK' if ok else 'ZU LANG'}")
        if not ok:
            print(f"        bricht um — alles darunter rutscht eine Zeile "
                  f"tiefer: {text[:70]}...")
    return fehler


# ───────────────────────────── Schritt 3 ──────────────────────────────────

def _pruefe_datei(pfad, etikett, stand):
    """Misst jede Wertentwicklungs-Folie einer gebauten Broschuere."""
    from pptx import Presentation

    geprueft = fehler = 0
    for nr, fussnote, quelle in _we_folien(Presentation(pfad)):
        geprueft += 1
        unten, zeilen = _fussnote_unterkante(fussnote)
        oben_q = _cm(quelle.top)
        text = quelle.text_frame.text.strip()
        maengel = []
        if oben_q < unten + MINDESTABSTAND_CM:
            maengel.append(f"Quelle bei {oben_q:.2f} cm, Disclaimer reicht "
                           f"bis {unten:.2f} cm")
        if not text:
            maengel.append("Quellenangabe ist leer")
        elif stand and stand not in text:
            maengel.append(f"Stand-Datum fehlt: {text!r}")
        # Absicht: Hier wird NICHT die Absatzlaenge geprueft. Die Vorlagen
        # tragen selbst Absaetze ueber der Zeilenbreite (237 bzw. 254
        # Zeichen) — das ist gewollter Vorlagentext, der sauber umbricht,
        # und _fussnote_unterkante rechnet den Umbruch mit. Die Laengen der
        # vom CODE geschriebenen Zeilen misst Schritt 2.
        ok = not maengel
        fehler += 0 if ok else 1
        print(f"   {etikett:12s} {nr:3d} {zeilen:6d} {unten:10.2f} "
              f"{oben_q:10.2f} {oben_q - unten:6.2f}  "
              f"{'OK' if ok else 'FEHLER'}")
        for m in maengel:
            print(f"        {m}")
    return geprueft, fehler


# Die Themen-Duplikation gehoert dazu: Dort entstehen die Wertentwicklungs-
# Folien erst beim Bauen (Offensiv ist die gemeldete), und SCHWEIZ ist der
# Fall OHNE Benchmark — dort faellt die ***-Zeile weg, die Fussnote wird also
# kuerzer als in allen anderen Broschueren.
THEMA_ZUSATZ = [
    ("Thema x3", ["Offensiv", "Pro", "Pro Dividende"]),
    ("SCHWEIZ", ["Schweiz_substanzorientiert"]),
]


def _pruefe_artefakt(ausgabe):
    print("\n3. Wirkung am echten Artefakt (gebaute Broschueren)")
    if importlib.util.find_spec("pptx") is None:
        print("   UEBERSPRUNGEN — python-pptx nicht installiert")
        return 0
    try:
        from modules.portfolioanalyse import (
            VORLAGEN_FAMILIEN, FAMILIE_ALLE_STRATEGIEN, _familien_portfolios,
            _familie_fuer_strategie, duration_info_aus_bestand,
        )
        # Datenbeschaffung und Bau kommen aus dem Export-Smoketest — damit
        # laeuft diese Suite ueber denselben Pfad wie die Oberflaeche.
        from test_export_smoke import _daten, _portfolio, _bauen
    except ImportError as ex:
        print(f"   UEBERSPRUNGEN — {ex}")
        return 0

    d = _daten()
    kopf = (f"   {'Broschuere':12s} {'Fo':>3s} {'Zeilen':>6s} {'Unterkante':>10s} "
            f"{'Quelle':>10s} {'Luft':>6s}  Ergebnis")
    print(kopf)
    print("   " + "-" * (len(kopf) - 3))

    fehler = geprueft = 0

    def _lauf(portfolios, familie, dateiname, etikett):
        nonlocal fehler, geprueft
        ziel, _gr, meldungen = _bauen(portfolios, familie, d, ausgabe, dateiname)
        for m in meldungen:
            print(f"   ! BUILD-FEHLER {etikett}: {m[:90]}")
            fehler += 1
        # Der erwartete Datenstand: das Auswertungsdatum des ersten Portfolios,
        # genau wie pptx_export._stand_str ihn bildet.
        ad = portfolios[0][2]
        stand = ad.strftime("%d.%m.%Y") if hasattr(ad, "strftime") else None
        n, f = _pruefe_datei(ziel, etikett, stand)
        geprueft += n
        fehler += f

    for familie in sorted(VORLAGEN_FAMILIEN):
        strategie = next((n for n in d["namen"]
                          if _familie_fuer_strategie(d["nm"], n) == familie), None)
        if strategie is None:
            print(f"   {familie:12s} UEBERSPRUNGEN (keine Strategie in den Daten)")
            continue
        alle = FAMILIE_ALLE_STRATEGIEN.get(familie)
        if alle:
            portfolios, fehlend = _familien_portfolios(
                alle, d["namen"], d["d2c"], d["pf_data"], duration_info_aus_bestand)
            if fehlend:
                print(f"   {familie:12s} UEBERSPRUNGEN (fehlende Daten: "
                      f"{', '.join(fehlend)})")
                continue
        else:
            portfolios = [_portfolio(strategie, d)]
        _lauf(portfolios, familie, f"{familie}.pptx", familie)

    for etikett, namen in THEMA_ZUSATZ:
        fehlend = [n for n in namen if n not in d["d2c"]]
        if fehlend:
            print(f"   {etikett:12s} UEBERSPRUNGEN (nicht in den Daten: "
                  f"{', '.join(fehlend)})")
            continue
        _lauf([_portfolio(n, d) for n in namen], "Thema",
              f"Thema_{etikett.replace(' ', '_')}.pptx", etikett)

    if not geprueft:
        print("   UEBERSPRUNGEN — keine Wertentwicklungs-Folie gefunden")
    else:
        print(f"\n   {geprueft} Wertentwicklungs-Folien geprueft")
    return fehler


def main():
    ausgabe = (sys.argv[1] if len(sys.argv) > 1
               else tempfile.mkdtemp(prefix="ffpb_quelle_"))
    os.makedirs(ausgabe, exist_ok=True)

    fehler = _pruefe_vorlagen() + _pruefe_zeilenlaengen() + _pruefe_artefakt(ausgabe)
    print()
    if fehler:
        print(f"FEHLGESCHLAGEN — {fehler} Abweichung(en)")
        return 1
    print("BESTANDEN — die Quellenangabe steht auf jeder Wertentwicklungs-Folie")
    print("            unter dem Disclaimer, und jeder Ersatztext passt in")
    print("            eine Zeile.")
    return 0


if __name__ == "__main__":
    sys.exit(main())

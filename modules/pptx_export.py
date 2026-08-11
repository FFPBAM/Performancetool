# modules/pptx_export.py
"""
PowerPoint-Export für Performance und Portfolioanalyse.

Nutzt eine Corporate-Vorlage (Vorlage/Vorlage_FFPB.pptx) und befüllt sie
mit den Daten aus dem Streamlit-Tool.

Öffentliche API:
- generate_portfolioanalyse_pptx(...) -> bytes
- generate_performance_pptx(...) -> bytes  (später)

Wichtigste Mechanik:
- Die Vorlage enthält benannte Shapes (C_Kennzahlen, T_Kennzahlen, C_Kennzahlen1, C_Kennzahlen2)
- Wir finden diese Shapes per Name und tauschen Inhalte aus
- Für Vergleichsportfolios werden die 4 dynamischen Slides dupliziert

⚠️ WICHTIG (Juli 2026): Diese Version setzt die NEUE Vorlage mit 26 Slides
voraus (Wertentwicklungs-Folie aus der alten cVV-Broschüre an Template-
Position 11). Mit der alten 25-Slide-Vorlage crasht der Export bewusst
früh mit einer klaren Fehlermeldung (siehe _EXPECTED_TEMPLATE_SLIDES).
"""

import os
import io
import datetime as dt
from typing import Optional

import pandas as pd
from pptx import Presentation

# Datenbasierte Chart-Nachbearbeitung (Achse, Ring-Labels, holeSize) — siehe
# modules/chart_dynamik.py. Läuft am Ende über alle Charts.
try:
    from modules.chart_dynamik import nachbearbeiten as _charts_nachbearbeiten
except ImportError:
    from chart_dynamik import nachbearbeiten as _charts_nachbearbeiten

# Generische PPTX-Helpers (Shape-Lookup, Vorlage, Slide-Manipulation)
try:
    from modules.pptx_helpers import (
        find_shape_by_name, load_template,
        duplicate_slide, remove_slide, save_and_reload,
        update_quelle_datum, update_slide_numbers,
    )
except ImportError:
    from pptx_helpers import (
        find_shape_by_name, load_template,
        duplicate_slide, remove_slide, save_and_reload,
        update_quelle_datum, update_slide_numbers,
    )

# Chart-Nachbearbeitung der CVV-Vergleichsfolie (Folie 19)
try:
    from modules.pptx_charts import set_line_series_sparse, set_series_line_colors
except ImportError:
    from pptx_charts import set_line_series_sparse, set_series_line_colors

# Slide-Befüllung (Domain: Anlagevorschlag, Wertentwicklung, Performance,
# Portfoliozusammenstellung). Dieses Modul ORCHESTRIERT nur: es ruft je Rolle
# die passende fill_*-Funktion auf, die Folien-Details liegen in pptx_slides.
#
# Der Import war bis 11.08.2026 rund viermal so lang — er zog Konstanten und
# Helfer herein, die ausschließlich die Durchreich-Wrapper unten brauchten
# (siehe das Aufräum-Band weiter unten). Mit den Wrappern sind sie entfallen.
try:
    from modules.pptx_slides import (
        clean_strategy_name,
        fill_anlagevorschlag_slides, fill_performance_slide,
        fill_wertentwicklung_slide, fill_zusammenstellung_slide,
        fill_rollierend_slide, fill_einzeltitel_themen_slide,
        fill_uebersicht_slide, fill_anlagekriterien_slide,
    )
except ImportError:
    from pptx_slides import (
        clean_strategy_name,
        fill_anlagevorschlag_slides, fill_performance_slide,
        fill_wertentwicklung_slide, fill_zusammenstellung_slide,
        fill_rollierend_slide, fill_einzeltitel_themen_slide,
        fill_uebersicht_slide, fill_anlagekriterien_slide,
    )

# Anlagekriterien-Konfiguration. BEWUSST das UI-freie Modul, nicht shared.py:
# dieser Export laeuft ohne Streamlit (Batch-Faehigkeit, Doku Abschnitt 13).
try:
    from modules import anlagekriterien as _anlagekriterien
except ImportError:
    import anlagekriterien as _anlagekriterien


# ---------------------------------------------------------------------------
# Konstanten
# ---------------------------------------------------------------------------
TEMPLATE_PATH = os.path.join("Vorlage", "Vorlage_FFPB.pptx")

_EXPECTED_TEMPLATE_SLIDES = 26
"""Erwartete Slide-Anzahl der Standard-Vorlage (Streamlit-Pfad). Schutz
gegen den klassischen Deploy-Fehler 'Code neu, Vorlage alt'."""

# ── Vorlagen-Konfiguration (NEU 03.07.2026, für Vorlagen-Familie) ──────────
# Der dynamische 4-Folien-Block kann in jeder Vorlage (ESG/CVV/ETF/Themen/…)
# an anderer Position liegen. Die Konfiguration beschreibt die Vorlage;
# _normalisiere_vorlage bringt den Block in die kanonische Reihenfolge.
BLOCK_REIHENFOLGE = ["anlagevorschlag", "wertentwicklung", "performance",
                     "zusammenstellung"]
"""Kanonische Reihenfolge des dynamischen Blocks im fertigen Export
(F: Strategieentwurf → Wertentwicklung → Performance mit BM → Zusammenst.)."""

DEFAULT_TEMPLATE_CONFIG = {
    # Erwartete Folienzahl der Vorlage (Guard gegen Code/Vorlage-Mismatch)
    "erwartete_folien": _EXPECTED_TEMPLATE_SLIDES,
    # 1-indexierte POSITIONEN der Block-Folien in der Vorlage:
    "block_positionen": {
        "anlagevorschlag": 7,
        "zusammenstellung": 9,
        "performance": 10,
        "wertentwicklung": 11,
    },
    # 1-indexierte Positionen, die beim Export ENTFERNT werden
    # (Standard-Vorlage: 8 = Anlagevorschlag-Teil-2, 12 = Währungen):
    "entfernen": [8, 12],
}
"""Konfiguration der Standard-Vorlage (Vorlage/Vorlage_FFPB.pptx).
Für andere Vorlagen der Familie eine analoge Konfiguration an
generate_portfolioanalyse_pptx(template_config=...) übergeben."""

# Slide-Positionen in der Vorlage (1-indexed, Stand Juli 2026 / 26 Slides):
#   Position 7  = Anlagevorschlag (Aktien + Allokations-Ring)
#   Position 8  = Anlagevorschlag-Teil-2       ← wird beim Export ENTFERNT
#   Position 9  = Portfoliozusammenstellung (Regionen + Branchen)
#   Position 10 = Performance/Wertentwicklung (mit Benchmark)
#   Position 11 = Wertentwicklung/Kurzübersicht (NEU — alte cVV-Folie)
#   Position 12 = Währungen                     ← wird beim Export ENTFERNT
SLIDE_ANLAGEVORSCHLAG_1 = 7
SLIDE_ANLAGEVORSCHLAG_2 = 8
SLIDE_ZUSAMMENSTELLUNG_1 = 9
SLIDE_PERFORMANCE = 10
SLIDE_WERTENTWICKLUNG = 11
SLIDE_WAEHRUNGEN = 12

# Foliennummer-Shape-Namen (uneinheitlich in der Vorlage):
# - Die meisten Slides nutzen "Foliennummer"
# - Slides 7, 8, 10 nutzen "Foliennummernplatzhalter 1" (anderer Layout-Master)
# - Die neue Wertentwicklungs-Folie (cVV-Import) nutzt "Foliennummer"
# Slides ohne eines dieser Shapes (Cover, Sub-Cover, Impressum) behalten keine Seitenzahl.
SHAPE_FOLIENNUMMER_NAMES = (
    "Foliennummer",
    "Foliennummernplatzhalter 1",
    "Slide Number",
    "Folienzahl",
)


# ---------------------------------------------------------------------------
# Strategienamen-Konvertierung
# ---------------------------------------------------------------------------
# clean_strategy_name kommt aus modules/pptx_slides.py (via Import oben).
# Die Funktion ist Public und wird intern in generate_portfolioanalyse_pptx aufgerufen.


# ---------------------------------------------------------------------------
# Vorlage laden
# ---------------------------------------------------------------------------
def _load_template(template_path: Optional[str] = None,
                   erwartete_folien: Optional[int] = None) -> Presentation:
    """Lädt die Vorlage über pptx_helpers.load_template und prüft dabei
    die Folienzahl.

    Ein Folienzahl-Mismatch bedeutet fast immer 'Code/Konfig und Vorlage
    passen nicht zusammen' (klassischer Deploy-Fehler) und würde sonst
    später still die falschen Folien treffen.

    Args:
        template_path: Pfad zur Vorlage. None = modul-lokales TEMPLATE_PATH.
        erwartete_folien: erwartete Folienzahl. None = Standard-Vorlage (26).
    """
    pfad = template_path or TEMPLATE_PATH
    erwartet = erwartete_folien or _EXPECTED_TEMPLATE_SLIDES
    prs = load_template(pfad)
    n = len(prs.slides)
    if n != erwartet:
        raise ValueError(
            f"Vorlage '{pfad}' hat {n} Folien, erwartet werden {erwartet}. "
            f"Vermutlich passen Vorlage und Konfiguration/Code nicht zusammen "
            f"(Vorlage aktualisieren oder 'erwartete_folien' der "
            f"Vorlagen-Konfiguration korrigieren)."
        )
    return prs


# ---------------------------------------------------------------------------
# AUFGERÄUMT 11.08.2026: hier standen 40 Durchreich-Funktionen
# ---------------------------------------------------------------------------
# Nach der Modul-Aufteilung Ende Juni 2026 (pptx_helpers / pptx_charts /
# pptx_slides / pptx_export) blieben in dieser Datei 40 Funktionen der Form
#
#     def _find_shape_by_name(slide, name):
#         """Wrapper für pptx_helpers.find_shape_by_name."""
#         return find_shape_by_name(slide, name)
#
# stehen — rund 290 Zeilen, ein Fünftel der Datei. 27 davon wurden nirgends
# aufgerufen, die übrigen 13 an genau einer bis drei Stellen, alle in dieser
# Datei. Schaden haben sie keinen angerichtet, aber Substanz vorgetäuscht:
# wer die Logik suchte, landete hier statt im zuständigen Modul.
#
# Die Aufrufe gehen jetzt direkt an die importierten Funktionen (Import-Block
# oben). Wer eine alte `_name`-Schreibweise sucht: Unterstrich weglassen —
# die Funktion liegt unverändert in ihrem Modul.
#
# _load_template ist als einzige geblieben. Sie ist kein Durchreicher,
# sondern ergänzt load_template um den Folienzahl-Guard.
#
# Ebenfalls entfallen: der Alias _NS_CHART auf pptx_charts.NS_CHART, den
# seit der Modultrennung niemand mehr referenziert hat.
# ---------------------------------------------------------------------------


# ─────────────────────────────────────────────────────────────────────────
# Performance-Berechnungs-Helfer
#
# AUFGERÄUMT 07.08.2026: Hier standen elf Kopien der analytics-Mathematik,
# entstanden als Duplikate aus streamlit_app.py. Acht davon wurden nie
# aufgerufen (_calc_cagr, _calc_vola, _drawdown_from_index,
# _calc_max_drawdown, _calc_sharpe_excess, _calc_period_return,
# _calc_period_return_after_fee, _calc_daily_returns_after_fee,
# _make_index_from_returns) — sie luden nur dazu ein, versehentlich
# weiterverwendet zu werden statt modules/analytics.py.
#
# Geblieben sind die zwei tatsächlich genutzten Helfer. Sie sind identisch
# zu ihren analytics-Gegenstücken; perspektivisch sollten auch sie von dort
# kommen, das berührt aber compute_wertentwicklung_data und wird deshalb
# getrennt gemacht.
# ─────────────────────────────────────────────────────────────────────────
import numpy as _np


def _annual_fee_to_daily_drag(fee_pa_decimal):
    """Jährlicher Honorarsatz → äquivalente tägliche Belastung."""
    return (1.0 + fee_pa_decimal) ** (1 / 365) - 1


def _make_index_after_fee(d_returns_decimal, fee_pa_decimal, startwert=100.0):
    """Index aus Tagesrenditen nach taggenauem Honorarabzug."""
    e = _annual_fee_to_daily_drag(fee_pa_decimal)
    idx = _np.empty(len(d_returns_decimal) + 1, dtype=float)
    idx[0] = startwert
    for i, d in enumerate(d_returns_decimal, start=1):
        idx[i] = idx[i-1] * (1.0 + (d - e))
    return idx


def compute_performance_data(timeseries_df: pd.DataFrame, fee_dec: float,
                             n_years_bar_chart: int = 5) -> dict:
    """
    Berechnet alle Performance-Daten für die Performance-Folie.

    Diese Funktion ist seit Juni 2026 ein Wrapper — die eigentliche Berechnungs-Logik
    lebt jetzt zentral in `modules/analytics.py` und wird auch vom Streamlit-Performance-Tab
    genutzt. So gibt es nur EINE Stelle für die Mathematik.

    Args:
        timeseries_df: DataFrame mit Spalten 'ret_port', 'ret_bm', 'rf'. Index = Datum.
        fee_dec: Honorarsatz p.a. als Dezimal (z.B. 0.01023 für 1,023% inkl. MwSt).
        n_years_bar_chart: Anzahl Jahre für Säulen-Chart (Default: 5).

    Returns: Dict mit Keys 'kennzahlen', 'performance_pa', 'wertentwicklung'.
             Siehe modules/analytics.py für Details.
    """
    # Lazy import: bricht modules/analytics.py weg, schlägt erst hier auf.
    from modules.analytics import compute_performance_data as _ac
    return _ac(timeseries_df, fee_dec, n_years_bar_chart)


def compute_rollierend_data(timeseries_df: pd.DataFrame, fee_dec: float) -> dict:
    """Rollierende Perioden-Renditen NACH KOSTEN (NEU 06.07.2026, für die
    Themen-Broschüren-Folie "Wertentwicklung rollierend", Tabelle 8x7).

    Identische Logik wie streamlit_app.build_rolling_table (Konsistenz-
    Doktrin): Index nach Kosten (taggenauer Drag), Perioden
      YTD (ab 31.12. Vorjahr — Jahresgrenze!),
      1 / 3 / 5 / 10 Jahre (Punkt-zu-Punkt: Enddatum minus n Jahre),
    jeweils als asof-Verhältnis. Fehlt die Historie (z.B. Strategie < 3
    Jahre), ist der Wert None → in der Folie "-".

    Returns:
        {"ytd": float|None, "1J": float|None, "3J": float|None,
         "5J": float|None, "10J": float|None} — Dezimalwerte (0.1589 = 15,89%).
    """
    ts = timeseries_df.sort_index()
    dates = pd.to_datetime(ts.index)
    r = pd.to_numeric(ts["ret_port"], errors="coerce").fillna(0.0).to_numpy(dtype=float)
    if len(r) == 0:
        return {k: None for k in ("ytd", "1J", "3J", "5J", "10J")}

    # Index nach Kosten mit synthetischem Startpunkt am Vortag (wie Tool:
    # xd = [erstes_datum - 1 Tag] + dates), damit asof(Startdatum) den
    # Schlussstand VOR der ersten Rendite trifft. _make_index_after_fee
    # liefert bereits len(r)+1 Werte inkl. Startwert 100 an Position 0.
    idx_vals = _make_index_after_fee(r, fee_dec, startwert=100.0)
    start_stamp = dates[0] - pd.Timedelta(days=1)
    full_index = pd.DatetimeIndex([start_stamp]).append(pd.DatetimeIndex(dates))
    index_ser = pd.Series(list(idx_vals), index=full_index).sort_index()

    end_ts = index_ser.index.max()

    def _asof(ts_target):
        s = index_ser.dropna()
        if s.empty or ts_target < s.index.min():
            return None
        return float(s.asof(ts_target))

    def _period_return(start_ts):
        v_end = _asof(end_ts)
        v_start = _asof(start_ts)
        if v_end is None or v_start is None or v_start == 0:
            return None
        return v_end / v_start - 1.0

    perioden = {
        "ytd": pd.Timestamp(end_ts.year - 1, 12, 31),
        "1J": end_ts - pd.DateOffset(years=1),
        "3J": end_ts - pd.DateOffset(years=3),
        "5J": end_ts - pd.DateOffset(years=5),
        "10J": end_ts - pd.DateOffset(years=10),
    }
    return {k: _period_return(v) for k, v in perioden.items()}


def compute_wertentwicklung_data(timeseries_df: pd.DataFrame, fee_dec: float,
                                 duration: Optional[float] = None,
                                 benchmark_text: Optional[str] = None) -> dict:
    """
    Berechnet die Daten für die Wertentwicklungs-/Kurzübersichts-Folie
    (NEU Juli 2026 — die aus dem alten VBA-Tool übernommene cVV-Folie).

    HINWEIS ARCHITEKTUR: Die Kennzahlen-Mathematik gehört perspektivisch nach
    modules/analytics.py (zentrale Berechnungs-Stelle). Sie liegt vorerst
    hier, weil analytics.py in dieser Session nicht angefasst wurde — beim
    nächsten analytics-Update bitte umziehen.

    Kennzahlen-Definitionen:

    1. "Wertentwicklung seit {Auflagejahr} kumuliert*"
       → Kumulierte Rendite NACH Kosten über die GESAMTE Historie
         (= Endstand des Linien-Charts relativ zu dessen Start — was der
         Kunde auf der Kurve sieht, ist exakt die Kennzahl).
         GEÄNDERT 02.07.2026 (Punkt 7): vorher bis
         31.12. des Vorjahres (Original-cVV-Konvention) — das erzeugte auf
         F8/F9 zwei verschiedene p.a.-Werte nebeneinander (z.B. 6,16% vs
         6,19%). Jetzt identische Basis wie F9 und Tool-UI; die *-Fußnote
         wird von fill_wertentwicklung_slide entsprechend umgeschrieben.

    2. "Rendite p.a. seit {Auflagejahr} nach Kosten"
       → Annualisierung (365-Tage-Basis, identische Mathematik wie
         analytics.calc_cagr) der Kennzahl 1 über die gesamte Historie —
         ergibt EXAKT den Wert, der auf F9 als "Performance p.a. (Referenz)"
         steht (Konsistenz-Kriterium, im Test verifiziert).

    3. "Wertentwicklung seit 01.01.{Jahr}**"
       → YTD-Rendite NACH Kosten mit taggenauem Honorarabzug — identische
         Konvention wie alle übrigen Tool-Kennzahlen (Juli 2026, Philip:
         "die Performance aus dem Tool wird nach Kosten angezeigt und
         täglich abgezogen"). Die alte VBA-Regel "vor Kosten, ab 30.06.
         abzüglich halbjährigen Honorarsatz" wurde damit VERWORFEN; die
         zugehörigen statischen Texte der Vorlage (Fußnote ** und
         Disclaimer-Satz) werden von fill_wertentwicklung_slide dynamisch
         auf die Tool-Konvention umgeschrieben.

    4. "Duration" — wird NICHT hier berechnet, sondern durchgereicht
       (kommt aus den Duration-CSVs / dur_info im Portfolioanalyse-Tab).
       None → Folie zeigt "–" (z.B. reine Aktien-Strategien).

    Chart-Daten (Balken + Linie) kommen 1:1 aus compute_performance_data
    (modules/analytics.py) — identische Datenbasis wie die Performance-Folie:
    Balken = volle Kalenderjahre nach Kosten (max. 5), Linie = gesamte
    Historie als Index (Start 1.0) nach Kosten.

    Args:
        timeseries_df: DataFrame mit Spalten 'ret_port', 'ret_bm', 'rf'
            (Tagesrenditen dezimal), Index = Datum.
        fee_dec: Effektiver Honorarsatz p.a. dezimal (ggf. inkl. MwSt —
            gleiche Konvention wie compute_performance_data).
        duration: Gewichtete Portfolio-Duration (oder None).
        benchmark_text: Benchmark-Zusammensetzung für die ***-Fußnote
            (aus Mapping_Namen.xlsx Spalte D), oder None.

    Returns:
        Dict mit Keys: auflage_jahr, laufendes_jahr, kum_nach_kosten,
        pa_nach_kosten, ytd, duration, benchmark_text, has_benchmark,
        performance_pa, wertentwicklung.
    """
    ts = timeseries_df.sort_index()
    dates = pd.to_datetime(ts.index)
    r = pd.to_numeric(ts["ret_port"], errors="coerce").fillna(0.0).to_numpy(dtype=float)

    first_date = dates[0]
    last_date = dates[-1]
    auflage_jahr = int(first_date.year)
    laufendes_jahr = int(last_date.year)

    drag = _annual_fee_to_daily_drag(fee_dec)

    # ── Chart-Daten: identische Basis wie Performance-Folie ──
    charts = compute_performance_data(timeseries_df, fee_dec)

    # ── Kennzahl 1 + 2: kumuliert / p.a. über die GESAMTE Historie ──
    # (02.07.2026, Punkt 7 — vorher bis 31.12. Vorjahr, s. Docstring)
    # BEWUSST direkt aus dem analytics-Ergebnis abgeleitet statt parallel
    # nachgerechnet: Kennzahl 2 ist damit KONSTRUKTIV identisch mit dem
    # "Performance p.a. (Referenz)"-Wert der F9 (gleiche Quelle, gleiche
    # Rundung), Kennzahl 1 ist der Endstand des F9-/F8-Linien-Index.
    kum_nach_kosten = None
    pa_nach_kosten = charts.get("kennzahlen", {}).get("performance_pa_ref")
    _idx_ref = charts.get("wertentwicklung", {}).get("referenz") or []
    if len(_idx_ref) > 0 and float(_idx_ref[0]) != 0.0:
        kum_nach_kosten = float(_idx_ref[-1]) / float(_idx_ref[0]) - 1.0

    # ── Kennzahl 3: YTD nach Kosten (taggenauer Honorarabzug, Tool-Konvention) ──
    ytd = None
    ytd_start = pd.Timestamp(laufendes_jahr, 1, 1)
    mask_ytd = dates >= ytd_start
    if mask_ytd.any():
        r_af_ytd = r[mask_ytd] - drag
        ytd = float(_np.prod(1.0 + r_af_ytd) - 1.0)

    return {
        "auflage_jahr": auflage_jahr,
        "laufendes_jahr": laufendes_jahr,
        "kum_nach_kosten": kum_nach_kosten,
        "pa_nach_kosten": pa_nach_kosten,
        "ytd": ytd,
        "duration": duration,
        "benchmark_text": benchmark_text,
        # Durchgereicht (11.08.2026): die Folie muss Chart-Serie, Legende und
        # ***-Fußnote weglassen, wenn es keinen Vergleichsmaßstab gibt.
        "has_benchmark": bool(charts.get("has_benchmark")),
        "performance_pa": charts.get("performance_pa", {}),
        "wertentwicklung": charts.get("wertentwicklung", {}),
    }


# ---------------------------------------------------------------------------
# Portfolioanalyse-Export (Hauptfunktion)
# ---------------------------------------------------------------------------
LAST_BUILD_ERRORS: list = []
"""Diagnose (NEU Juli 2026): Fehler, die beim Berechnen der Folien-Daten
aufgetreten sind. Wird von generate_portfolioanalyse_pptx zu Beginn geleert.
Hintergrund: _build_perf_data/_build_we_data fangen Exceptions bewusst ab
(eine kaputte Kennzahl soll nicht den ganzen Export crashen — die Folie
zeigt dann Platzhalter). Vorher passierte das STILL und war im UI nicht
diagnostizierbar ("Folie wird nicht befüllt, aber kein Fehler"). Jetzt
sammelt diese Liste die Fehlermeldungen; der Aufrufer (portfolioanalyse.py)
zeigt sie nach dem Export als Warnung an."""


def _record_build_error(context: str, exc: Exception):
    import traceback
    LAST_BUILD_ERRORS.append(
        f"{context}: {type(exc).__name__}: {exc}"
    )
    # Voller Traceback zusätzlich in die Server-Logs (streamlit-Konsole)
    traceback.print_exc()


F9_BAR_INCLUDE_CURRENT_YEAR = True
"""02.07.2026 (Punkt 4, Redundanz F8/F9-Balken): Der F9-Balken-Chart zeigt
zusätzlich das LAUFENDE Jahr (YTD nach Kosten vs. Benchmark) als weiteren
Balken — damit unterscheiden sich F8 (nur volle Kalenderjahre, gemäß
*-Fußnoten-Logik) und F9 (inkl. aktuellem Jahresverlauf) inhaltlich.
Entspricht der Darstellung des Streamlit-Tools selbst (Balken-Chart
"Kalenderjahre" zeigt dort ebenfalls das laufende Jahr mit an).
False = altes Verhalten (beide Charts identisch)."""


def _append_current_year_bar(perf: dict, ts: pd.DataFrame, fee: float) -> dict:
    """Hängt das laufende Jahr an performance_pa an (siehe
    F9_BAR_INCLUDE_CURRENT_YEAR). Portfolio nach Kosten (taggenauer Drag,
    identische Mathematik wie Kennzahl 3 der F8), Benchmark brutto —
    gleiche Konvention wie die vollen Jahre aus analytics.

    Non-destruktiv: arbeitet auf einer Kopie des performance_pa-Dicts.
    No-op wenn das laufende Jahr bereits enthalten ist oder keine Daten hat.
    """
    pa = perf.get("performance_pa") or {}
    jahre = list(pa.get("jahre") or [])
    if not jahre:
        return perf
    ts_sorted = ts.sort_index()
    dates = pd.to_datetime(ts_sorted.index)
    cur_year = int(dates[-1].year)
    if any(int(y) == cur_year for y in jahre):
        return perf  # laufendes Jahr ist (z.B. bei Jahreswechsel) schon drin
    mask = _np.array([d.year == cur_year for d in dates])
    if not mask.any():
        return perf
    r = pd.to_numeric(ts_sorted["ret_port"], errors="coerce").fillna(0.0).to_numpy(dtype=float)
    drag = _annual_fee_to_daily_drag(fee)
    ref_ytd = float(_np.prod(1.0 + (r[mask] - drag)) - 1.0)
    # Ohne Benchmark bleibt die Benchmark-Liste leer (11.08.2026) — sonst
    # bekäme eine ansonsten leere Liste hier doch noch einen (Null-)Wert und
    # die Serie wäre um genau einen Balken länger als die Jahresliste.
    bench = list(pa.get("benchmark") or [])
    if perf.get("has_benchmark"):
        rb = pd.to_numeric(ts_sorted.get("ret_bm"),
                           errors="coerce").fillna(0.0).to_numpy(dtype=float)
        bench = bench + [float(_np.prod(1.0 + rb[mask]) - 1.0)]
    new_pa = {
        "jahre": jahre + [cur_year],
        "referenz": list(pa.get("referenz") or []) + [ref_ytd],
        "benchmark": bench,
    }
    out = dict(perf)
    out["performance_pa"] = new_pa
    return out


def _build_perf_data(performance_inputs, idx: int) -> Optional[dict]:
    """
    Helfer: Berechnet performance_data Dict aus performance_inputs[idx].

    Returns None wenn keine Daten oder ungültige Eingabe — dann zeigt die
    Performance-Folie die Vorlagen-Platzhalter. Berechnungsfehler landen
    in LAST_BUILD_ERRORS (siehe oben) statt still verschluckt zu werden.

    02.07.2026 (Punkt 4): hängt optional das laufende Jahr an den
    Balken-Chart an (F9_BAR_INCLUDE_CURRENT_YEAR).
    """
    if not performance_inputs or idx >= len(performance_inputs):
        return None
    pi = performance_inputs[idx]
    if pi is None:
        return None
    ts = pi.get("timeseries_df")
    fee = pi.get("fee_dec", 0.0)
    if ts is None or len(ts) == 0:
        return None
    try:
        perf = compute_performance_data(ts, fee)
        if F9_BAR_INCLUDE_CURRENT_YEAR:
            perf = _append_current_year_bar(perf, ts, fee)
        return perf
    except Exception as exc:
        _record_build_error(f"Performance-Folie, Portfolio {idx + 1}", exc)
        return None


def _build_we_data(performance_inputs, idx: int) -> Optional[dict]:
    """
    Helfer (NEU Juli 2026): Berechnet das we_data-Dict für die
    Wertentwicklungs-Folie aus performance_inputs[idx].

    Nutzt dieselbe Zeitreihe/fee wie die Performance-Folie, zusätzlich die
    OPTIONALEN Keys "duration" und "benchmark_text" im performance_inputs-
    Eintrag (siehe Snippet für portfolioanalyse.py). Fehlen die Keys, zeigt
    die Folie an den Stellen "–" bzw. behält die Vorlagen-Fußnote —
    vollständig rückwärtskompatibel zu Alt-Aufrufern.

    Returns None wenn keine Daten — dann zeigt die Folie Vorlagen-Platzhalter
    (nur der Titel wird gesetzt).
    """
    if not performance_inputs or idx >= len(performance_inputs):
        return None
    pi = performance_inputs[idx]
    if pi is None:
        return None
    ts = pi.get("timeseries_df")
    fee = pi.get("fee_dec", 0.0)
    if ts is None or len(ts) == 0:
        return None
    try:
        return compute_wertentwicklung_data(
            ts, fee,
            duration=pi.get("duration"),
            benchmark_text=pi.get("benchmark_text"),
        )
    except Exception as exc:
        _record_build_error(f"Wertentwicklungs-Folie, Portfolio {idx + 1}", exc)
        return None



# ── CVV Folie 19: Vergleichs-Chart der fünf Strategien ────────────────────
# Farben laut Corporate Design (aus den Vorlagen extrahiert), NAMENSBASIERT —
# die Vorlage nutzt sonst Office-Standard-Akzente (Blau/Orange/Grau/Gelb).
VERGLEICH_FARBEN = {
    "Konservativ":   "9FD0EF",   # sehr helles Blau
    "Defensiv":      "66A4CE",   # helles Blau
    "Defensiv Plus": "5F8CA1",   # mittleres Blau
    "Ausgewogen":    "14355C",   # dunkles Marineblau
    "Dynamic":       "BB9256",   # Gold/Braun
}
VERGLEICH_LINIENBREITE_EMU = 15875   # 1.25 pt (wie Vorlage)

_EXCEL_EPOCHE = dt.date(1899, 12, 30)


def _excel_serial(datum) -> int:
    """Datum → Excel-Seriennummer (für die Datums-Achse des Charts)."""
    d = datum.date() if hasattr(datum, "date") else datum
    return (d - _EXCEL_EPOCHE).days


def _monats_index_nach_kosten(timeseries_df, fee_dec: float):
    """Monatsend-Indexreihe NACH KOSTEN, normiert auf 1.0 am ERSTEN eigenen
    Datenpunkt (1.0 = 100 %).

    Returns: dict {date: float} oder {} wenn keine Daten.

    Wichtig: Die Normierung erfolgt je Strategie auf ihren EIGENEN Start —
    eine später startende Strategie (z.B. Dynamic) beginnt ebenfalls bei 100 %,
    nicht bei dem Wert, den die anderen zu diesem Zeitpunkt schon haben.
    """
    if timeseries_df is None or len(timeseries_df) == 0:
        return {}
    ts = timeseries_df.sort_index()
    r = pd.to_numeric(ts["ret_port"], errors="coerce").fillna(0.0).to_numpy(dtype=float)
    if len(r) == 0:
        return {}
    idx_werte = _make_index_after_fee(r, fee_dec, startwert=1.0)   # len = n+1
    daten = pd.to_datetime(ts.index)
    reihe = pd.Series(idx_werte[1:], index=daten)
    # letzter Wert je Monat = Monatsende
    monats = reihe.groupby([daten.year, daten.month]).last()
    letzte_daten = pd.Series(daten).groupby([daten.year, daten.month]).last()
    ergebnis = {}
    basis = None
    for schluessel, wert in monats.items():
        tag = letzte_daten.loc[schluessel]
        if basis is None:
            basis = wert
        ergebnis[tag.date()] = float(wert / basis)
    return ergebnis


def _build_vergleich_data(performance_inputs, n_strategien: int):
    """Baut Kategorien (Excel-Serials) + Serien mit LÜCKEN für Folie 19.

    Returns: (kategorien, serien) — serien = [(None, {idx: wert}), …] in der
    Reihenfolge der Strategien. (None = Serienname der Vorlage behalten.)
    Gibt (None, None) zurück, wenn keine verwertbaren Zeitreihen da sind.
    """
    reihen = []
    for k in range(n_strategien):
        pi = (performance_inputs[k]
              if performance_inputs and k < len(performance_inputs) else None)
        if not pi:
            reihen.append({})
            continue
        reihen.append(_monats_index_nach_kosten(pi.get("timeseries_df"),
                                                pi.get("fee_dec", 0.0)))
    alle_daten = sorted({d for r in reihen for d in r})
    if not alle_daten:
        return None, None
    pos = {d: i for i, d in enumerate(alle_daten)}
    kategorien = [_excel_serial(d) for d in alle_daten]
    serien = [(None, {pos[d]: w for d, w in r.items()}) for r in reihen]
    return kategorien, serien

def _build_rollierend_data(performance_inputs, idx: int) -> Optional[dict]:
    """Helfer (NEU 06.07.2026): rollierende Perioden-Renditen für die
    Themen-Broschüren-Tabelle aus performance_inputs[idx].

    Nutzt dieselbe Zeitreihe/fee wie die übrigen Folien (Konsistenz).
    Returns None wenn keine Daten → Folie behält Vorlagen-Platzhalter.
    """
    if not performance_inputs or idx >= len(performance_inputs):
        return None
    pi = performance_inputs[idx]
    if pi is None:
        return None
    ts = pi.get("timeseries_df")
    fee = pi.get("fee_dec", 0.0)
    if ts is None or len(ts) == 0:
        return None
    try:
        return compute_rollierend_data(ts, fee)
    except Exception as exc:
        _record_build_error(f"Rollierende Tabelle, Portfolio {idx + 1}", exc)
        return None


def _reorder_slides(prs, new_order: list):
    """Sortiert die Folien gemäß new_order (Liste aktueller Indizes) um —
    EIN atomarer Reorder der sldIdLst statt fehleranfälliger Move-Ketten."""
    xml_slides = prs.slides._sldIdLst
    elems = list(xml_slides)
    assert sorted(new_order) == list(range(len(elems))), \
        "new_order muss eine Permutation aller Folien-Indizes sein"
    for el in elems:
        xml_slides.remove(el)
    for i in new_order:
        xml_slides.append(elems[i])


def _normalisiere_vorlage(prs, template_config: dict) -> int:
    """Bringt eine Vorlage in die kanonische Export-Grundform
    (NEU 03.07.2026 — ersetzt die frühere hartcodierte remove/move-Kette):

    1. Entfernt die in template_config["entfernen"] gelisteten Folien
       (1-indexierte Vorlagen-Positionen).
    2. Sortiert die vier Block-Folien (template_config["block_positionen"])
       in die kanonische Reihenfolge BLOCK_REIHENFOLGE und rückt sie als
       zusammenhängenden Block an die Position der vordersten Block-Folie.
       Alle übrigen (statischen) Folien behalten ihre relative Reihenfolge.

    Returns:
        0-indexierter Start-Index des Blocks im normalisierten Deck.
    """
    entfernen_idx = sorted([p - 1 for p in template_config.get("entfernen", [])],
                           reverse=True)
    for i in entfernen_idx:
        remove_slide(prs, i)

    # Block-Positionen um die Entfernungen korrigieren
    def _nach_entfernung(pos1: int) -> int:
        i = pos1 - 1
        return i - sum(1 for r in entfernen_idx if r < i)

    bp = template_config["block_positionen"]
    reihenfolge = template_config.get("block_reihenfolge", BLOCK_REIHENFOLGE)
    fehlend = [t for t in reihenfolge if t not in bp]
    if fehlend:
        raise ValueError(f"template_config['block_positionen'] unvollständig, "
                         f"fehlt: {fehlend}")
    block_in_kanon = [_nach_entfernung(bp[t]) for t in reihenfolge]

    n = len(prs.slides)
    block_set = set(block_in_kanon)
    if len(block_set) != len(reihenfolge):
        raise ValueError("block_positionen zeigen auf dieselbe Folie")
    nicht_block = [i for i in range(n) if i not in block_set]
    start = min(block_in_kanon)
    vor = sum(1 for i in nicht_block if i < start)
    new_order = nicht_block[:vor] + block_in_kanon + nicht_block[vor:]
    _reorder_slides(prs, new_order)
    return vor  # = Block-Start im neuen Deck


def _vervielfaeltige_block(prs, block_start: int, n_strategien: int,
                           block_laenge: Optional[int] = None):
    """Dupliziert den 4-Folien-Block für n_strategien (NEU 03.07.2026 —
    generalisiert die frühere 2-Portfolio-Sonderlogik auf beliebiges N).

    Vorgehen (indexstabil):
    - Jede Block-Folie wird — von der LETZTEN zur ERSTEN — (N−1)-mal
      dupliziert (duplicate_slide fügt Kopien direkt hinter der Quelle ein;
      die Reihenfolge unter identischen Kopien ist egal).
      → Zwischenstand: Folien nach TYP gruppiert ([AV×N][WE×N][Perf×N][Zus×N]).
    - Ein einziger Reorder gruppiert danach nach STRATEGIE:
      [AV1,WE1,Perf1,Zus1, AV2,WE2, …].
    """
    if n_strategien <= 1:
        return
    B = block_laenge if block_laenge is not None else len(BLOCK_REIHENFOLGE)
    for offset in range(B - 1, -1, -1):
        src = block_start + offset
        for _ in range(n_strategien - 1):
            duplicate_slide(prs, src)
    n = len(prs.slides._sldIdLst)
    new_order = list(range(block_start))
    for k in range(n_strategien):
        for t in range(B):
            new_order.append(block_start + t * n_strategien + k)
    new_order += list(range(block_start + B * n_strategien, n))
    _reorder_slides(prs, new_order)


def _stand_str(eval_date) -> Optional[str]:
    """Formatiert das Auswertungsdatum für die statischen 'Quelle'-Zeilen
    (02.07.2026, Punkt 6). None/ungültig → None (Quelle bleibt unangetastet)."""
    try:
        if eval_date is not None and hasattr(eval_date, "strftime"):
            return eval_date.strftime("%d.%m.%Y")
    except Exception:
        pass
    return None


def generate_portfolioanalyse_pptx(
    portfolios: list,   # Liste von (display_name, df, auswertungsdatum, dur_info)
    anlagevolumen: float = 0.0,
    performance_inputs: Optional[list] = None,
    template_path: Optional[str] = None,
    template_config: Optional[dict] = None,
) -> bytes:
    """
    Erstellt eine PPTX aus einer Corporate-Vorlage und befüllt den dynamischen
    4-Folien-Block für BELIEBIG VIELE Strategien (generalisiert 03.07.2026 —
    vorher nur 1 oder 2; der Streamlit-Aufruf bleibt unverändert kompatibel).

    Block pro Strategie (kanonische Reihenfolge BLOCK_REIHENFOLGE):
      Strategieentwurf → Wertentwicklung → Performance (mit BM) → Zusammenst.
    Bei N Strategien stehen die Blöcke hintereinander (Strategie 1 = F7-10,
    Strategie 2 = F11-14, …), statische Folien davor/danach unverändert.

    Args:
        portfolios: Liste von Tupeln (display_name, df, auswertungsdatum,
            duration_info) — EIN Eintrag pro Strategie, beliebige Anzahl ≥ 1.
        anlagevolumen: aktuell ungenutzt (Zukunftsfeature).
        performance_inputs: Liste in gleicher Reihenfolge wie `portfolios`
            (oder None). Format pro Eintrag:
            {"timeseries_df": df(ret_port,ret_bm,rf), "fee_dec": 0.0119,
             "duration": 4.24|None, "benchmark_text": "50% iBoxx …"|None}
            None-Einträge → betroffene Folien zeigen Vorlagen-Platzhalter.
        template_path: Pfad zur Vorlage (NEU 03.07. — für die Vorlagen-
            Familie ESG/CVV/ETF/Themen). None = Standard TEMPLATE_PATH.
        template_config: Vorlagen-Beschreibung (erwartete_folien,
            block_positionen, entfernen) — siehe DEFAULT_TEMPLATE_CONFIG.
            None = Standard-Vorlage.

    Returns:
        PPTX-Bytes
    """
    if not portfolios:
        raise ValueError("Mindestens ein Portfolio erforderlich.")
    cfg = template_config or DEFAULT_TEMPLATE_CONFIG
    n_strategien = len(portfolios)

    LAST_BUILD_ERRORS.clear()  # Diagnose-Liste pro Export frisch
    prs = _load_template(template_path, cfg.get("erwartete_folien"))

    # ── FESTE BLÖCKE (NEU 09.07.2026, für die CVV-Vorlage) ──────────────────
    # Manche Vorlagen enthalten die Folien ALLER Strategien bereits fertig
    # vorgebaut (CVV: fünf Strategie-Paare, jedes mit einem eigenen, starren
    # Anlagekriterien-Kasten). Sie dürfen deshalb NICHT dupliziert werden —
    # sonst bekämen alle Strategien den Kasten der ersten. Stattdessen wird
    # an FESTEN Vorlagen-Positionen befüllt.
    #
    #   cfg["feste_bloecke"] = [ {rolle: 1-indexierte Vorlagenposition}, … ]
    #                          ein Eintrag je Strategie, in Reihenfolge
    #   cfg["rollen_optionen"] = {rolle: {kwarg: wert}}   (optional)
    #
    # Fehlt der Schlüssel, läuft exakt der bisherige Pfad
    # (normalisieren → vervielfältigen). Blast-Radius für Standard/Themen: 0.
    feste_bloecke = cfg.get("feste_bloecke")
    rollen_optionen = cfg.get("rollen_optionen", {})

    if feste_bloecke:
        # Nur Entfernungen anwenden; keine Umsortierung, keine Duplikation.
        for i in sorted([p - 1 for p in cfg.get("entfernen", [])], reverse=True):
            remove_slide(prs, i)
        if cfg.get("entfernen"):
            prs = save_and_reload(prs)
        if n_strategien > len(feste_bloecke):
            _record_build_error(
                "Vorlage",
                ValueError(f"{n_strategien} Strategien übergeben, die Vorlage "
                           f"hat aber nur {len(feste_bloecke)} feste Blöcke — "
                           f"überzählige Strategien werden ignoriert."))
        block_start, B, reihenfolge = None, None, None
    else:
        # ── Schritt 1: Vorlage normalisieren (Entfernungen + Block kanonisch) ──
        block_start = _normalisiere_vorlage(prs, cfg)
        prs = save_and_reload(prs)

        # ── Schritt 2: Block auf N Strategien vervielfältigen ──
        reihenfolge = cfg.get("block_reihenfolge", BLOCK_REIHENFOLGE)
        B = len(reihenfolge)
        _vervielfaeltige_block(prs, block_start, n_strategien, block_laenge=B)
        if n_strategien > 1:
            prs = save_and_reload(prs)

    # ── Schritt 3: jede Strategie in ihren Block füllen ──
    # Dispatch: Rolle → Fill-Aufruf. Der Offset innerhalb des Blocks ergibt
    # sich aus der Position der Rolle in `reihenfolge` (so kann jede Vorlage
    # eine andere Folien-Zusammenstellung/-Reihenfolge haben, z.B. die
    # Themen-Broschüren mit "rollierend" statt "performance").
    # Anlagekriterien EINMAL laden (nicht je Strategie) — die Datei ist klein,
    # aber der Export laeuft ueber bis zu fuenf Strategien.
    kriterien_cfg = _anlagekriterien.lade()

    for k, (display_name, df, eval_date, _dur) in enumerate(portfolios):
        strategy_name = clean_strategy_name(display_name)
        perf_data = _build_perf_data(performance_inputs, k)
        we_data = _build_we_data(performance_inputs, k)
        roll_data = _build_rollierend_data(performance_inputs, k)
        stand = _stand_str(eval_date)

        # Folienindizes dieser Strategie ermitteln:
        #  - feste Blöcke: direkt aus der Vorlagen-Konfiguration (1-indexiert)
        #  - sonst: wie bisher über base + offset im vervielfältigten Block
        if feste_bloecke:
            if k >= len(feste_bloecke):
                continue
            ziele = [(pos - 1, rolle)
                     for rolle, pos in feste_bloecke[k].items()]
        else:
            base = block_start + B * k
            ziele = [(base + offset, rolle)
                     for offset, rolle in enumerate(reihenfolge)]

        for idx, rolle in ziele:
            opt = dict(rollen_optionen.get(rolle, {}))
            if rolle == "anlagevorschlag":
                fill_anlagevorschlag_slides(prs, idx, df, strategy_name,
                                             eval_date=eval_date, **opt)
                # Anlagekriterien-Kasten aus der Konfiguration (NEU 10.08.2026).
                # Schluessel ist der UNGEKUERZTE Anzeigename aus dem Mapping
                # ("cVV konservativ"), nicht strategy_name — der hat den
                # Praefix bereits verloren. Strategien ohne Eintrag (Familie
                # Thema) lassen die Vorlage unberuehrt.
                fill_anlagekriterien_slide(
                    prs, idx,
                    _anlagekriterien.fuer(display_name, kriterien_cfg),
                    _anlagekriterien.anzeigename(display_name, kriterien_cfg))
            elif rolle == "wertentwicklung":
                fill_wertentwicklung_slide(prs, idx, strategy_name,
                                            we_data=we_data, stand_date_str=stand)
            elif rolle == "performance":
                fill_performance_slide(prs, idx, strategy_name,
                                        performance_data=perf_data,
                                        stand_date_str=stand)
            elif rolle == "zusammenstellung":
                fill_zusammenstellung_slide(prs, idx, df, strategy_name,
                                             eval_date=eval_date)
            elif rolle == "rollierend":
                fill_rollierend_slide(prs, idx, strategy_name,
                                       rollierend_data=roll_data,
                                       stand_date_str=stand)
            elif rolle == "einzeltitel_themen":
                fill_einzeltitel_themen_slide(prs, idx, df, strategy_name,
                                               eval_date=eval_date)
            else:
                _record_build_error(
                    f"Portfolio {k + 1}",
                    ValueError(f"Unbekannte Block-Rolle '{rolle}' in "
                               f"block_reihenfolge"))

    # Quelle-Datum aktualisieren auf das Auswertungsdatum des ersten Portfolios.
    # Steht statisch in den Chart-Annotationen (drawing*.xml) der Vorlage als
    # "Quelle: Eigene Berechnung Stand: 12.02.2026" — muss auf das echte
    # Auswertungsdatum aktualisiert werden. Deckt auch die "Quelle"-Textbox
    # der neuen Wertentwicklungs-Folie ab (gleiches "Stand ..."-Muster).
    if portfolios and portfolios[0][2] is not None:
        datum_obj = portfolios[0][2]
        try:
            if hasattr(datum_obj, 'strftime'):
                datum_str = datum_obj.strftime("%d.%m.%Y")
                update_quelle_datum(prs, datum_str)
        except Exception:
            pass

    # ── Schritt 3b: Folien, die EINMAL für ALLE Strategien laufen ──────────
    # (NEU 10.07.2026, für CVV Folie 17: Wertentwicklung aller fünf Strategien
    # nebeneinander.) Konfiguration:
    #     cfg["einmal_folien"] = {"uebersicht": 17}   # 1-indexiert
    # Fehlt der Schlüssel, passiert nichts — Standard/Themen unberührt.
    for rolle, pos in (cfg.get("einmal_folien") or {}).items():
        idx = pos - 1
        try:
            if rolle == "uebersicht":
                roll_liste = [_build_rollierend_data(performance_inputs, k)
                              for k in range(len(portfolios))]
                if not any(roll_liste):
                    _record_build_error(
                        f"Folie {pos}",
                        ValueError("Keine Performance-Zeitreihen übergeben — "
                                   "Übersichtstabelle behält die Vorlagen-Werte."))
                    continue
                _stand = _stand_str(portfolios[0][2]) if portfolios else None
                fill_uebersicht_slide(prs, idx, roll_liste,
                                       stand_date_str=_stand,
                                       **(rollen_optionen.get(rolle, {})))
            elif rolle == "vergleich":
                # Linien-Chart aller Strategien, Index 100 je eigenem Start.
                kat, serien = _build_vergleich_data(performance_inputs,
                                                    len(portfolios))
                if not kat:
                    _record_build_error(
                        f"Folie {pos}",
                        ValueError("Keine Performance-Zeitreihen übergeben — "
                                   "Vergleichs-Chart behält die Vorlagen-Werte."))
                    continue
                shape = find_shape_by_name(prs.slides[idx], "Diagramm")
                if shape is None or not getattr(shape, "has_chart", False):
                    _record_build_error(
                        f"Folie {pos}",
                        ValueError("Chart-Shape 'Diagramm' nicht gefunden."))
                    continue
                set_line_series_sparse(shape, kat, serien)
                set_series_line_colors(shape, VERGLEICH_FARBEN,
                                       breite_emu=VERGLEICH_LINIENBREITE_EMU)
            else:
                _record_build_error(
                    f"Folie {pos}",
                    ValueError(f"Unbekannte Einmal-Rolle '{rolle}'"))
        except Exception as _ex:
            _record_build_error(f"Folie {pos} ({rolle})", _ex)

    # Foliennummern dynamisch setzen (NACH allen Add/Remove/Duplicate-Operationen,
    # VOR dem Speichern). Korrigiert die statischen Werte aus der Vorlage
    # (Slide 7 hat z.B. "13", soll aber "7" sein nach Renumber).
    update_slide_numbers(prs)

    # Charts datenbasiert nachziehen (NEU 07.07.2026): Wertentwicklungs-Linie
    # auf die echte Datenspanne skalieren (kein Leerraum vor/nach der Kurve),
    # Doughnut-Ringe auf dünnen Original-Look (holeSize=79) + Außen-Labels
    # radial aus dem tatsächlichen Segmentwinkel. Läuft NACH allen Fill-/
    # Slide-Operationen, VOR dem Speichern. Balken (catAx) bleiben unberührt.
    # Rührt die Download-Logik NICHT an — arbeitet ausschließlich an Chart-XML.
    try:
        _charts_nachbearbeiten(prs)
    except Exception as _ex:
        # Chart-Kosmetik darf den Export nie abbrechen.
        LAST_BUILD_ERRORS.append(f"Chart-Nachbearbeitung übersprungen: {_ex}")

    # Speichern
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

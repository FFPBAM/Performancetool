"""
modules/pptx_helpers.py — Generische PPTX-Helper-Funktionen.

Pure PPTX-Mechanik OHNE Domain-Wissen (keine Performance/Anlagevorschlag/
Portfolio-Logik). Diese Helper sind die Bausteine, mit denen die höheren
Module (pptx_charts, pptx_slides, pptx_export) arbeiten.

Was hier hingehört:
- Shape-Lookup nach Namen
- Text-Manipulation in Shapes und Tabellen-Zellen
- Tabellen leeren / Header behalten
- PPTX-Vorlage laden
- Pure Daten-Konvertierungs-Helper

Was hier NICHT hingehört:
- Chart-Manipulation (→ pptx_charts.py, Schritt 7)
- Slide-Befüllung mit Domain-Daten (→ pptx_slides.py, Schritt 8)
- Slide-Duplikation / Reorder (→ separates Modul, Schritt 6c)

Diese Datei hat KEINE Imports von Streamlit.
Sie kann unverändert in lokalen Python-Skripten genutzt werden.
"""

import io
import os
import re
from copy import deepcopy
from typing import Optional, Sequence

import pandas as pd
from pptx import Presentation
from pptx.opc.packuri import PackURI


# ─────────────────────────────────────────────────────────────────────────────
# Default-Vorlage
# ─────────────────────────────────────────────────────────────────────────────

DEFAULT_TEMPLATE_PATH = os.path.join("Vorlage", "Vorlage_FFPB.pptx")
"""Pfad zur Standard-Vorlage. Kann pro Aufruf von load_template() überschrieben werden."""


# ─────────────────────────────────────────────────────────────────────────────
# Shape-Lookup
# ─────────────────────────────────────────────────────────────────────────────

def find_shape_by_name(slide, name: str):
    """Findet ein Shape auf einer Slide anhand seines Namens.

    Args:
        slide: pptx.slide.Slide
        name: Shape-Name wie er in PowerPoint vergeben wurde (z.B. "Tabelle", "Diagramm links")

    Returns:
        Das Shape-Objekt oder None falls nicht gefunden.
    """
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


# ─────────────────────────────────────────────────────────────────────────────
# Text-Manipulation in Shapes
# ─────────────────────────────────────────────────────────────────────────────

def replace_text_in_shape(shape, new_text: str):
    """Ersetzt den Text in einem Text-Shape (z.B. Placeholder).

    Behält die Formatierung des ersten Runs bei (Schriftart, Größe, Farbe).
    Alle weiteren Paragraphen und Runs werden entfernt.

    Args:
        shape: Shape mit has_text_frame=True
        new_text: Neuer Text

    No-op wenn das Shape kein text_frame hat.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Alle Paragraphen außer dem ersten löschen
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    # Im ersten Paragraphen: ersten Run behalten, Text ersetzen, Rest löschen
    p = tf.paragraphs[0]
    if len(p.runs) == 0:
        p.text = new_text
    else:
        p.runs[0].text = new_text
        for run in p.runs[1:]:
            r = run._r
            r.getparent().remove(r)


# ─────────────────────────────────────────────────────────────────────────────
# Tabellen-Zellen-Text
# ─────────────────────────────────────────────────────────────────────────────

def set_cell_text(cell, text: str, is_bold: Optional[bool] = None):
    """Setzt den Text einer Tabellenzelle.

    WICHTIG: Leere Strings werden zu NBSP (U+00A0) konvertiert. Grund:
    Die Vorlage verwendet in nicht-befüllten Zellen ebenfalls NBSP als
    Platzhalter. Lässt man die Zelle mit leerem <a:t/> zurück, rendert
    LibreOffice sie mit Default-Font-Metriken (größere Zeilenhöhe), was
    die gesamte Tabelle vertikal streckt und zu Überlauf führen kann.

    Args:
        cell: Die Zelle
        text: Der neue Text (leer → NBSP)
        is_bold: Wenn explizit True/False: setzt Bold-Formatierung.
                 Wenn None: behält vorherige Formatierung bei.
    """
    # Leere Zellen auf NBSP setzen
    if text == "":
        text = "\u00A0"

    tf = cell.text_frame
    # Alle Paragraphen außer dem ersten löschen
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    p = tf.paragraphs[0]
    if len(p.runs) == 0:
        p.text = text
        # Bold explizit setzen wenn gewünscht
        if is_bold is not None and p.runs:
            p.runs[0].font.bold = is_bold
    else:
        p.runs[0].text = text
        # Bold explizit setzen wenn gewünscht
        if is_bold is not None:
            p.runs[0].font.bold = is_bold
        for run in p.runs[1:]:
            r = run._r
            r.getparent().remove(r)


def set_cell_text_preserve_format(cell, text: str):
    """Setzt Zellen-Text und erhält das Format des ersten Runs.

    Im Gegensatz zu `set_cell_text` (das alle Runs durch einen neuen leeren Run
    ersetzt) bleibt hier die Font-Formatierung (Größe, Farbe, Bold) erhalten —
    wichtig für gestylte Tabellen-Zellen (z.B. fett, weiß auf blauem Header).

    Args:
        cell: Die Zelle
        text: Der neue Text
    """
    if not cell.text_frame.paragraphs:
        # Fallback: kein Paragraph → normales set_cell_text Verhalten
        set_cell_text(cell, text)
        return
    para = cell.text_frame.paragraphs[0]
    if not para.runs:
        set_cell_text(cell, text)
        return
    # Erste Run behält ihr Format, alle weiteren Runs löschen
    runs = list(para.runs)
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)
    # Weitere Paragraphs löschen
    for p in cell.text_frame.paragraphs[1:]:
        p._p.getparent().remove(p._p)


# ─────────────────────────────────────────────────────────────────────────────
# Tabellen-Operationen
# ─────────────────────────────────────────────────────────────────────────────

def clear_table(table, keep_header_rows: int = 1):
    """Leert alle Zellen einer Tabelle ab der angegebenen Start-Zeile.

    Args:
        table: pptx.table.Table
        keep_header_rows: Anzahl der Header-Zeilen die behalten werden (default: 1)
    """
    for row_idx in range(keep_header_rows, len(table.rows)):
        for cell in table.rows[row_idx].cells:
            set_cell_text(cell, "")


# ─────────────────────────────────────────────────────────────────────────────
# Daten-Konvertierung
# ─────────────────────────────────────────────────────────────────────────────

def safe_float(value, default: float = 0.0) -> float:
    """Konvertiert einen Wert zu float. NaN, NaT, None, ungültige Werte → default.

    Wichtig: Verhindert TypeError beim Vergleich/Sortieren gemischter Typen.

    Args:
        value: Beliebiger Wert
        default: Rückgabewert bei ungültigem Input (default: 0.0)

    Returns:
        Float-Wert oder default.
    """
    if value is None:
        return default
    try:
        if pd.isna(value):
            return default
    except (TypeError, ValueError):
        pass
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


# ─────────────────────────────────────────────────────────────────────────────
# Template laden
# ─────────────────────────────────────────────────────────────────────────────

def load_template(path: Optional[str] = None) -> Presentation:
    """Lädt eine PPTX-Vorlage.

    Args:
        path: Pfad zur PPTX-Datei. None = DEFAULT_TEMPLATE_PATH.

    Returns:
        pptx.Presentation Objekt.

    Raises:
        FileNotFoundError: wenn die Datei nicht existiert.
    """
    template_path = path if path is not None else DEFAULT_TEMPLATE_PATH
    if not os.path.exists(template_path):
        raise FileNotFoundError(
            f"Vorlage nicht gefunden: {template_path}\n"
            f"Bitte 'Vorlage_FFPB.pptx' im Ordner 'Vorlage/' ablegen."
        )
    return Presentation(template_path)


# ─────────────────────────────────────────────────────────────────────────────
# Slide-Manipulation: Duplikation, Reorder, Remove, Reload
# ─────────────────────────────────────────────────────────────────────────────

def clone_chart_part(prs, source_chart_part):
    """Erstellt eine tiefe Kopie eines Chart-Parts mit eigener URI.

    Kopiert auch die Sub-Relationships (z.B. embeddings zu XLSX).
    Wird von `duplicate_slide` aufgerufen, kann auch standalone genutzt werden.

    Args:
        prs: Presentation-Objekt
        source_chart_part: Der zu klonende Chart-Part

    Returns:
        Der neue Chart-Part (eigene URI, eigene Sub-Relationships).
    """
    package = source_chart_part.package

    # Neue URI finden (nächste freie chartN.xml)
    existing_nums = set()
    for part in package.iter_parts():
        partname_str = str(part.partname)
        m = re.search(r'/ppt/charts/chart(\d+)\.xml$', partname_str)
        if m:
            existing_nums.add(int(m.group(1)))
    n = 1
    while n in existing_nums:
        n += 1
    new_partname = PackURI(f"/ppt/charts/chart{n}.xml")

    # Chart-Part-Klasse nutzen und neuen Part mit neuer URI erstellen
    chart_part_cls = type(source_chart_part)
    new_chart_part = chart_part_cls.load(
        new_partname,
        source_chart_part.content_type,
        package,
        source_chart_part.blob,
    )

    # Sub-Relationships kopieren (z.B. Excel-Embedding)
    for rel in source_chart_part.rels.values():
        if rel.is_external:
            new_chart_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_chart_part.relate_to(rel.target_part, rel.reltype)

    return new_chart_part


def duplicate_slide(prs, source_idx: int):
    """Dupliziert eine Slide samt Chart-Teilen und Image-Referenzen.

    Die Charts werden so kopiert, dass Änderungen am Duplikat NICHT das
    Original überschreiben (eigene Chart-Parts).
    Image-Parts werden geteilt (Bilder werden nicht modifiziert).

    Fügt die neue Slide direkt hinter die Quelle ein.

    Args:
        prs: Presentation-Objekt
        source_idx: 0-basierter Index der Quell-Slide

    Returns:
        Die neue Slide (an Position source_idx + 1).
    """
    source = prs.slides[source_idx]

    # Neue Slide mit gleichem Layout anlegen
    new_slide = prs.slides.add_slide(source.slide_layout)

    # Alle Shapes entfernen die beim Layout-Add automatisch kamen
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Mapping alte rId → neue rId für Duplikate
    rid_map = {}

    for rel in list(source.part.rels.values()):
        if "chart" in rel.reltype:
            # Chart-Part duplizieren (eigener Part mit neuer Datei-URI)
            new_chart_part = clone_chart_part(prs, rel.target_part)
            new_rel_id = new_slide.part.relate_to(new_chart_part, rel.reltype)
            rid_map[rel.rId] = new_rel_id
        elif "image" in rel.reltype:
            # Image-Part wiederverwenden – Bilder werden nicht modifiziert
            new_rel_id = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rel.rId] = new_rel_id

    # Shapes kopieren und rId-Referenzen patchen
    R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    for shape in source.shapes:
        el = shape.element
        new_el = deepcopy(el)

        # Alle r:id / r:embed / r:link Attribute im XML durchgehen und mappen
        for attr in ["id", "embed", "link"]:
            attr_name = R_NS + attr
            for el_with_rid in new_el.iter():
                if attr_name in el_with_rid.attrib:
                    old_rid = el_with_rid.attrib[attr_name]
                    if old_rid in rid_map:
                        el_with_rid.attrib[attr_name] = rid_map[old_rid]

        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Die neue Slide ist aktuell am Ende – verschieben hinter die Quelle
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_slide_element = slides[-1]
    xml_slides.remove(new_slide_element)
    xml_slides.insert(source_idx + 1, new_slide_element)

    return new_slide


def move_slide(prs, from_idx: int, to_idx: int):
    """Verschiebt eine Slide innerhalb der Präsentation.

    Args:
        prs: Presentation-Objekt
        from_idx: 0-basierte Position der Slide, die verschoben werden soll
        to_idx: 0-basierte Ziel-Position (vor Reordering)

    Realisiert via direkter Manipulation von <p:sldIdLst> in presentation.xml.

    Raises:
        IndexError: bei ungültigen Indizes.
    """
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    if from_idx < 0 or from_idx >= len(slides):
        raise IndexError(f"move_slide: from_idx {from_idx} out of range (max {len(slides)-1})")
    if to_idx < 0 or to_idx >= len(slides):
        raise IndexError(f"move_slide: to_idx {to_idx} out of range (max {len(slides)-1})")
    target = slides[from_idx]
    xml_slides.remove(target)
    xml_slides.insert(to_idx, target)


def remove_slide(prs, slide_idx: int):
    """Entfernt eine Slide an gegebener Position (0-indexed).

    Räumt sowohl die <p:sldId>-Referenz als auch die Slide-Part-Relationship auf.
    """
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    slide_elem = slides_list[slide_idx]
    rId = slide_elem.rId
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_elem)


def save_and_reload(prs) -> Presentation:
    """Speichert die Präsentation in Memory und lädt sie neu.

    Das räumt interne Slide-IDs auf (wichtig nach remove/duplicate Operationen),
    sonst können "Duplicate name"-Warnungen beim späteren Speichern auftreten.

    Returns:
        Frisch geladene Presentation.
    """
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


# ─────────────────────────────────────────────────────────────────────────────
# Drawing-XML-Patches (Quelle-Datum, Foliennummern)
# ─────────────────────────────────────────────────────────────────────────────

def update_quelle_datum(prs, datum_str: str):
    """Aktualisiert das 'Quelle: Eigene Berechnung Stand: XX.XX.XXXX' Datum.

    Sucht in ALLEN drawing*.xml Parts der PPTX nach dem Quelle-Pattern und
    ersetzt das Datum. Diese Quelle-Zeile steht statisch in den drawing-Parts
    der Vorlage (im Chart-Annotation-Layer der Ring-Charts) und kann nicht
    über die normale python-pptx-Slide-API erreicht werden.

    Args:
        prs: Presentation-Objekt
        datum_str: Datum im Format 'DD.MM.YYYY' (z.B. '17.06.2026')

    No-op wenn datum_str leer ist.
    """
    if not datum_str:
        return
    package = prs.part.package
    for part in package.iter_parts():
        pn = str(part.partname)
        if not pn.startswith("/ppt/drawings/drawing"):
            continue
        try:
            xml = part.blob.decode('utf-8')
        except Exception:
            continue
        if 'Quelle: Eigene Berechnung Stand:' not in xml:
            continue
        new_xml = re.sub(
            r'(Quelle: Eigene Berechnung Stand: )\d{2}\.\d{2}\.\d{4}',
            f'\\g<1>{datum_str}',
            xml
        )
        if new_xml != xml:
            part._blob = new_xml.encode('utf-8')


# Default-Namen für Foliennummer-Shapes (PowerPoint generiert verschiedene
# Namen je nach Sprache und Vorlagen-Herkunft)
DEFAULT_FOLIENNUMMER_NAMES = (
    "Foliennummer",
    "Foliennummernplatzhalter 1",
    "Slide Number",
    "Folienzahl",
)


def update_slide_numbers(prs, foliennummer_names: Optional[Sequence[str]] = None):
    """Setzt die Foliennummer auf jeder Slide auf die korrekte 1-indexed Position.

    HINTERGRUND:
    Die Vorlage hat statische Seitenzahlen (z.B. Slides 7-9 zeigen "13"-"15",
    weil der Designer eine Lücke für dynamische Folien angenommen hat). Nach
    Add/Remove/Duplicate-Operationen stimmen diese Werte nicht mehr — daher
    nach allen Slide-Manipulationen einmal alle Foliennummern auf die korrekte
    Position überschreiben.

    Slides ohne Foliennummer-Shape (Cover, Sub-Cover, Impressum) bleiben
    unverändert — das ist gewollt: solche Slides sollen keine Seitenzahl tragen.

    Args:
        prs: Presentation-Objekt
        foliennummer_names: Sequenz der Shape-Namen die als Foliennummer
            erkannt werden sollen. None = DEFAULT_FOLIENNUMMER_NAMES.

    Sollte als LETZTER Schritt vor `prs.save()` aufgerufen werden.
    """
    names = foliennummer_names if foliennummer_names is not None else DEFAULT_FOLIENNUMMER_NAMES
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.name in names and shape.has_text_frame:
                replace_text_in_shape(shape, str(idx))
                break

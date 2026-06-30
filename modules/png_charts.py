"""PNG-basierte Ring-Charts mit Labels außen.

Dieses Modul ersetzt PowerPoint-native Ring-Charts (Donuts) durch eine
Kombination aus:
- Banner-Rectangle (oben) mit Titel
- PNG-gerenderter Donut-Chart (matplotlib) mit Labels außen + Leader-Lines
- Legende als python-pptx Shapes (Rectangle + Text pro Item)
- Source-Annotation unten rechts

Hintergrund: PowerPoint's native Ring-Charts platzieren Datenbeschriftungen
bei großen Segmenten (>50%) trotz `dLblPos="outEnd"` innen — Microsoft's
"smart auto-placement" lässt sich nicht per Property abschalten. matplotlib
gibt uns volle Kontrolle über die Label-Positionen.

Hauptfunktion: `replace_donut_chart(slide, shape_name, ...)`.
"""
from __future__ import annotations

import io
from typing import Sequence

import matplotlib
matplotlib.use("Agg")  # Headless rendering (Streamlit Cloud)
import matplotlib.pyplot as plt
import numpy as np

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt


# ────────────────────────────────────────────────────────────────────────────
# FFPB-Farbpalette (aus der Vorlage extrahiert)
# ────────────────────────────────────────────────────────────────────────────
FFPB_COLORS = [
    "#1F3A5F",  # Dunkelblau (primär)
    "#7FA8C8",  # Hellblau
    "#C8A45C",  # Gold
    "#D5E0EB",  # Hellblau-grau (pale)
    "#A8A8A8",  # Grau (für "Sonstige")
    "#5E7A99",  # Mittelblau
    "#E5B97F",  # Hellgold
    "#3D5778",  # Dunkelgrau-blau
]
FFPB_DARKBLUE = "#1F3A5F"  # Banner-Hintergrund


# ────────────────────────────────────────────────────────────────────────────
# PNG-Rendering
# ────────────────────────────────────────────────────────────────────────────
def render_donut_png(
    values: Sequence[float],
    labels: Sequence[str],
    colors: Sequence[str] | None = None,
    dpi: int = 200,
    figsize: tuple[float, float] = (4.5, 4.0),
    label_distance: float = 1.30,
    label_fontsize: int = 10,
) -> bytes:
    """Rendert einen Donut-Chart als PNG mit Labels außen + Leader-Lines.

    Args:
        values: Werte pro Segment (Summe egal — wird relativ berechnet).
        labels: Anzeige-Strings pro Segment (z.B. "57,03%").
        colors: Hex-Farben pro Segment. Default: FFPB_COLORS.
        dpi: Auflösung. 200 = scharfes Rendering bei moderater Dateigröße
            (~45 KB pro PNG). Höher als der bisherige Default 150, weil die
            PNG auf der Folie mit ~2,9" Breite dargestellt wird und 150 dpi
            dort sichtbar weich wirkte.
        figsize: matplotlib-Figur-Größe in inches.
        label_distance: Distanz der Labels vom Mittelpunkt (1.0 = Ring-Außenrand).
        label_fontsize: Schriftgröße der Prozent-Labels.

    Returns:
        PNG-Bytes (transparenter Hintergrund).

    Hinweis zum Seitenverhältnis: Die zurückgegebene PNG ist breiter als hoch
    (die Labels brauchen horizontalen Platz links/rechts). Der Donut SELBST ist
    durch `ax.set_aspect("equal")` exakt kreisrund. Beim Einfügen auf die Folie
    MUSS das native Pixel-Seitenverhältnis erhalten bleiben (siehe
    `replace_donut_chart`), sonst wird der Kreis zum Oval verzerrt.
    """
    if colors is None:
        colors = FFPB_COLORS[: len(values)]
    colors = list(colors)[: len(values)]  # zur Sicherheit kürzen

    fig, ax = plt.subplots(figsize=figsize, dpi=dpi)
    wedges, _ = ax.pie(
        list(values),
        colors=colors,
        startangle=90,
        counterclock=False,
        wedgeprops=dict(width=0.45, edgecolor="white", linewidth=1.5),
    )
    for wedge, lbl in zip(wedges, labels):
        # Mittel-Winkel des Segments
        ang = (wedge.theta2 + wedge.theta1) / 2
        x_ring = np.cos(np.deg2rad(ang))
        y_ring = np.sin(np.deg2rad(ang))
        x_label = label_distance * x_ring
        y_label = label_distance * y_ring
        ha = "left" if x_ring > 0 else "right"
        ax.annotate(
            lbl,
            xy=(x_ring, y_ring),       # Leader-Line-Spitze: am Ring-Rand
            xytext=(x_label, y_label),  # Label-Position: weiter außen
            ha=ha, va="center",
            fontsize=label_fontsize, fontweight="bold",
            arrowprops=dict(arrowstyle="-", color="#333333", lw=0.7,
                            connectionstyle="arc3,rad=0"),
        )
    ax.set_aspect("equal")
    # Padding für Labels (sonst werden sie abgeschnitten)
    ax.set_xlim(-1.65, 1.65)
    ax.set_ylim(-1.50, 1.50)
    ax.axis("off")
    fig.patch.set_alpha(0)  # transparenter Hintergrund

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, transparent=True,
                bbox_inches="tight", pad_inches=0.02)
    plt.close(fig)
    return buf.getvalue()


def _png_pixel_size(png_bytes: bytes) -> tuple[int, int]:
    """Liest die Pixel-Dimensionen (Breite, Höhe) aus dem IHDR-Chunk einer PNG.

    Vermeidet eine zusätzliche PIL-Abhängigkeit an dieser Stelle: Die
    PNG-Signatur ist 8 Byte lang, danach folgt der IHDR-Chunk
    (4 Byte Länge + 'IHDR' + 4 Byte Breite + 4 Byte Höhe, jeweils Big-Endian).
    Breite liegt damit ab Byte 16, Höhe ab Byte 20.
    """
    if len(png_bytes) < 24:
        return (1, 1)
    width = int.from_bytes(png_bytes[16:20], "big")
    height = int.from_bytes(png_bytes[20:24], "big")
    if width <= 0 or height <= 0:
        return (1, 1)
    return (width, height)


# ────────────────────────────────────────────────────────────────────────────
# Slide-Manipulation
# ────────────────────────────────────────────────────────────────────────────
def _remove_shape(slide, shape) -> None:
    """Entfernt eine Shape vom Slide."""
    sp = shape._element
    sp.getparent().remove(sp)


def _add_banner(slide, left, top, width, height, *, text, bg_color=FFPB_DARKBLUE):
    """Banner-Rectangle mit dunkelblauer Füllung + weißer Schrift (wie nativ)."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor.from_string(bg_color.lstrip("#"))
    rect.line.fill.background()  # Kein Rahmen
    tf = rect.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.text = text
    for run in p.runs:
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = "Calibri"


def _add_legend(slide, left, top, width, items, *, font_size_pt=10, line_h_pt=22):
    """Vertikale Legende: Farb-Rectangle + Text pro Item.

    Args:
        items: Liste von (label, hex_color) Tupeln.
    """
    rect_size = Pt(10)
    text_indent = Pt(16)
    line_h_emu = Pt(line_h_pt)
    for i, (label, color) in enumerate(items):
        item_top = top + i * line_h_emu
        # Farb-Indikator (Rectangle)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left,
            item_top + (line_h_emu - rect_size) // 2,  # vertikal zentriert
            rect_size, rect_size,
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
        rect.line.fill.background()
        # Text-Box rechts neben Indikator
        tb = slide.shapes.add_textbox(
            left + text_indent, item_top, width - text_indent, line_h_emu
        )
        tf = tb.text_frame
        tf.margin_left = 0; tf.margin_right = 0
        tf.margin_top = 0; tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = label
        for run in p.runs:
            run.font.size = Pt(font_size_pt)
            run.font.bold = True
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def _add_source_text(slide, left, top, width, height, *, text):
    """Kleine Quelle-Annotation rechts unten."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = text
    for run in p.runs:
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        run.font.name = "Calibri"


# ────────────────────────────────────────────────────────────────────────────
# Hauptfunktion
# ────────────────────────────────────────────────────────────────────────────
def replace_donut_chart(
    slide,
    chart_shape_name: str,
    *,
    values: Sequence[float],
    item_labels: Sequence[str],
    percent_labels: Sequence[str],
    colors: Sequence[str] | None = None,
    banner_text: str,
    source_text: str,
) -> None:
    """Ersetzt einen nativen Donut-Chart komplett durch PNG + Banner + Legende + Source.

    Diese Funktion ist der Drop-In-Ersatz für `replace_chart_data()` bei
    Ring-Charts. Sie:
    1. Findet die Chart-Shape, merkt sich Position + Größe
    2. Entfernt die Chart-Shape (entfernt damit auch native Banner/Source aus drawing*.xml)
    3. Erzeugt das Layout-Equivalent neu:
       - Banner oben (z.B. "REGIONEN", "Branchen")
       - PNG-Ring rechts (mit Labels außen + Leader-Lines), SEITENVERHÄLTNIS-ERHALTEND
       - Legende links (vertikal, Rectangle + Text pro Item)
       - Source-Annotation unten rechts

    Args:
        slide: Die Slide
        chart_shape_name: Name der zu ersetzenden Chart-Shape
            (z.B. "C_Kennzahlen", "C_Kennzahlen1", "C_Kennzahlen2")
        values: Numerische Werte pro Segment
        item_labels: Labels für die Legende (z.B. "Global", "Index")
        percent_labels: Vorformatierte Prozent-Strings (z.B. "57,03%")
        colors: Hex-Farben pro Segment. Default: FFPB_COLORS.
        banner_text: Text im Banner oben
        source_text: Quellen-Annotation unten rechts

    Raises:
        ValueError: Wenn Shape nicht gefunden
    """
    # Chart-Shape finden
    chart_shape = None
    for sh in slide.shapes:
        if sh.name == chart_shape_name:
            chart_shape = sh
            break
    if chart_shape is None:
        raise ValueError(f"Shape '{chart_shape_name}' nicht gefunden auf Slide")

    # Bounding-Box merken
    cx, cy = chart_shape.left, chart_shape.top
    cw, ch = chart_shape.width, chart_shape.height

    # Layout-Berechnung (alles relativ zur Chart-Shape):
    BANNER_H = int(ch * 0.06)            # 6%: dünner Banner oben
    SOURCE_H = int(ch * 0.05)            # 5%: Source-Text unten
    INNER_TOP = cy + BANNER_H + Pt(4)
    INNER_BOTTOM = cy + ch - SOURCE_H
    INNER_H = INNER_BOTTOM - INNER_TOP

    # Innerhalb des Inner-Bereichs: links Legend (35%), rechts Ring (65%)
    LEGEND_W = int(cw * 0.35)
    RING_W = cw - LEGEND_W

    # 0. Original entfernen
    _remove_shape(slide, chart_shape)

    # 1. Banner oben
    _add_banner(slide, cx, cy, cw, BANNER_H, text=banner_text)

    # 2. PNG-Ring (rechts) — SEITENVERHÄLTNIS-ERHALTEND einpassen + zentrieren
    #
    # WICHTIG: Früher wurde die PNG mit fester Breite UND Höhe (RING_W × INNER_H)
    # eingefügt. Da die Ring-Box hochkant ist (≈0,72) und die PNG quer (≈1,15),
    # wurde der kreisrunde Donut vertikal gestreckt → sichtbares Oval.
    # Jetzt: native Pixel-Geometrie auslesen, in die verfügbare Box einpassen
    # (das kleinere Limit gewinnt) und im Box-Bereich zentrieren.
    png_bytes = render_donut_png(
        values=values, labels=percent_labels, colors=colors,
    )
    ring_left = cx + LEGEND_W

    img_w_px, img_h_px = _png_pixel_size(png_bytes)
    img_aspect = img_w_px / img_h_px          # Breite / Höhe (nativ)
    box_aspect = RING_W / INNER_H

    if img_aspect >= box_aspect:
        # Bild relativ breiter als die Box → Breite ist limitierend
        draw_w = RING_W
        draw_h = int(round(RING_W / img_aspect))
    else:
        # Bild relativ höher als die Box → Höhe ist limitierend
        draw_h = INNER_H
        draw_w = int(round(INNER_H * img_aspect))

    # Im verfügbaren Ring-Bereich horizontal + vertikal zentrieren
    pic_left = ring_left + (RING_W - draw_w) // 2
    pic_top = INNER_TOP + (INNER_H - draw_h) // 2

    slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        pic_left, pic_top, draw_w, draw_h,
    )

    # 3. Legend links, vertikal mittig
    if colors is None:
        colors = FFPB_COLORS[: len(item_labels)]
    items = list(zip(item_labels, colors))
    n = len(items)
    line_h_emu = Pt(22)
    legend_total_h = n * line_h_emu
    legend_top = INNER_TOP + (INNER_H - legend_total_h) // 2
    _add_legend(slide, cx + Pt(12), legend_top, LEGEND_W - Pt(24), items)

    # 4. Source-Annotation unten rechts
    _add_source_text(
        slide, cx, cy + ch - SOURCE_H, cw, SOURCE_H, text=source_text,
    )

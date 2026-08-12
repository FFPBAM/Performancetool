# modules/chart_dynamik.py  (bzw. Funktionen in pptx_charts.py)
"""Dynamische, datenbasierte Chart-Skalierung — native PP, kein Bild.
Aufruf jeweils NACH dem Daten-Schreiben (replace_chart_data)."""
import math
import datetime as dt
import re

_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
def _q(t): return f"{{{_C}}}{t}"

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _root(chart):
    # python-pptx: chart._chartSpace ist der lxml-Wurzelknoten des Charts
    return chart._chartSpace


# ═══════════════════════════════════════════════════════════════════════════
#  WEGWEISER FÜR KÜNFTIGE CHATS  —  bitte zuerst lesen
# ═══════════════════════════════════════════════════════════════════════════
#
#  Diese Datei zieht ALLE Charts einer fertigen Präsentation datenbasiert nach
#  (native PowerPoint-Objekte, KEIN Bild). Einstieg ist immer `nachbearbeiten`
#  ganz unten — von dort aus die Kette lesen.
#
#  ── SO ARBEITEST DU GEZIELT AN EINEM PROBLEM ──────────────────────────────
#  „Passe die Führungslinie im Ring an"      → ring_leader_zeichnen  + CONFIG
#  „Labels stehen falsch / überlappen"        → ring_labels_aussen_dynamisch
#  „Punkt am Segment / Farbe / Stub-Länge"    → CONFIG-Block unten
#  „Nur bei Branchen / nur bei Thema"         → _ring_typ / _ist_thema_familie
#  „Segmentfarben falsch"                     → ring_segmentfarben
#  „Kurve hat Leerraum / Achse falsch"        → datumsachse_an_daten
#
#  ── RINGTYPEN (datengetrieben erkannt, KEINE Prozent-Sonderfälle) ─────────
#  ANLAGEKLASSEN : Kategorien = AKTIEN/RENTEN/EDELMETALLE/LIQUIDITÄT
#                  (ESG, CVV, ETF — der „AKTUELLE STRUKTUR"-Ring)
#  REGIONEN      : Kategorien = Länder/Regionen   (Thema-Broschüren, Ring links)
#  BRANCHEN      : Kategorien = Sektoren          (Thema-Broschüren, Ring rechts)
#  → Regionen/Branchen gibt es NUR in der Thema-Familie (Offensiv/Pro/Pro Div.).
#
#  ── FALLSTRICKE (teuer gelernt — nicht erneut hineinlaufen) ───────────────
#  1) LibreOffice-Rendering ist KEIN Beweis für PowerPoints Leader-Optik.
#     LibreOffice ignoriert showLeaderLines=0 und zeichnet ZUSÄTZLICH eigene
#     Auto-Leader → im Render sieht man Doppel-Linien. Prüfe die Optik nur an
#     echten PowerPoint-Screenshots. Geometrie prüft man dagegen zuverlässig
#     im XML (Koordinaten, Längen, Kreuzungen).
#  2) PowerPoints AUTO-Leader folgen einer undokumentierten Regel und fehlen
#     mal (Label auf Segmentnaht ODER Linie quert das Ring-Loch). Deshalb
#     schalten wir sie AB (ring_leaderlines_aus) und zeichnen SELBST
#     (ring_leader_zeichnen). Niemals wieder auf die Auto-Leader verlassen.
#  3) KOORDINATEN: Ein Punkt (xi, yi) in Rahmen-Zoll liegt auf der Folie bei
#     (shape.left + xi*914400, shape.top + yi*914400) EMU. Exakt 1:1, weil die
#     Rahmenbreite in Zoll = shape.width/914400.
#  4) root.find(".//c:val") liefert nur die ERSTE Serie. Bei Mehr-Serien-Charts
#     (CVV-Vergleich) über findall(".//c:val") ALLE Serien nehmen.
#  5) Segmentfarben erben sonst die INDEX-Position der Vorlagen-dPt → AKTIEN
#     würde gold statt blau. Deshalb ring_segmentfarben NAMENSbasiert.
#  6) Ein defektes Einzelchart darf den Export nie abbrechen → try/except in
#     nachbearbeiten. Beim Debuggen das except temporär entfernen.
#
# ═══════════════════════════════════════════════════════════════════════════
#  KONFIGURATION  —  zentrale Stellschrauben (hier ändern, nicht im Code suchen)
# ═══════════════════════════════════════════════════════════════════════════

# ── Führungslinien (Leader) ────────────────────────────────────────────────
LEADER_FARBE        = "000000"   # SCHWARZ (früher grau A6A6A6) — Wunsch 20.07.
LEADER_BREITE_EMU   = 9525       # 0,75 pt
_LEADER_RADIAL_STUB = 0.16       # Zoll: radialer Austritt aus dem Ring
_LEADER_MIN_STUB    = 0.06       # Zoll: darunter gerade Linie statt Knick
# Alter Name für Rückwärtskompatibilität (falls extern referenziert):
LEADER_GRAU         = LEADER_FARBE

# ── Punkt am Label-Ende der Führungslinie (kleiner gefüllter Kreis) ─────────
# Punkte erscheinen nur auf bestimmten Ringtypen UND nur in der Thema-Familie
# (Wunsch 20.07.). PUNKT_RINGTYPEN = erlaubte Ringtypen (aus _ring_typ):
# 'ANLAGEKLASSEN' (Assetklassen-Ring/Einzeltitel), 'BRANCHEN', 'REGIONEN'.
# Zum Erweitern einfach den gewünschten Typ ergänzen; für alle Ringe leere
# Prüfung via PUNKT_ALLE_RINGTYPEN=True; Thema-Bindung über PUNKT_NUR_THEMA.
PUNKT_AN            = True
PUNKT_RINGTYPEN     = ("ANLAGEKLASSEN", "BRANCHEN")   # Regionen bewusst OHNE
PUNKT_NUR_THEMA     = True
PUNKT_FARBE         = "000000"
PUNKT_DURCHMESSER   = 0.055      # Zoll

# ── Label-Text ─────────────────────────────────────────────────────────────
LABEL_SCHRIFTFARBE  = "000000"   # Prozentzahlen IMMER schwarz

# ── Datumsachse der Linien-Charts (NEU 12.08.2026) ─────────────────────────
# Schrittweite der Achsenbeschriftung, abhängig von der Länge der Historie.
# (Obergrenze der Spanne in Monaten, Schritt in Monaten) — die erste passende
# Zeile gewinnt, darunter greift DATUMSACHSE_SCHRITT_LANG.
# Ziel sind 8–19 Beschriftungen: weniger sagt nichts, mehr ist bei 7 pt und um
# 90 Grad gedrehtem Text nicht mehr lesbar. Vorher standen dort bis zu 37.
DATUMSACHSE_STUFEN      = ((36, 3), (84, 6))   # <=3 Jahre: Quartal, <=7: Halbjahr
DATUMSACHSE_SCHRITT_LANG = 12                  # darüber: Jahresschritt

# ── Familien-spezifische Ring-Optik (NEU 27.07.2026) ────────────────────────
# NUR die hier gelisteten Familien weichen von den globalen Defaults ab; alle
# anderen nutzen die Defaults → deren Ringe bleiben UNVERÄNDERT. Steuerbar:
#   hole               – Loch-Prozent: KLEINER = DICKERER Ring (Default 79)
#   leader_breite_emu  – Strichstärke der Führungslinien (Default 9525 = 0,75 pt;
#                        12700 EMU = 1 pt, also 19050 = 1,5 pt)
#   label_fett         – Prozentzahlen fett (bleiben schwarz)
#   punkt_durchmesser  – Punkt-Größe in Zoll, FALLS die Familie Punkte hat
#                        (CVV hat per PUNKT_NUR_THEMA aktuell KEINE)
# Wunsch 27.07.: CVV-Ringe insgesamt kräftiger — dickerer Ring, fettere
# Führungslinien, fette Prozentzahlen.
# Kräftige Ring-Optik (Feinschliff 27.-28.07.): dickerer Ring, ruhige GERADE
# Leader mittig im Band, kleine dezente Punkte, fette schwarze Prozente, etwas
# luftigere Labels. CVV, ESG, ETF und Thema teilen dieselbe Optik — deshalb hier
# EINMAL definiert. Soll EINE Familie abweichen, gib ihr unten statt
# _RING_KRAEFTIG einen eigenen dict-Block mit anderen Zahlen.
_RING_KRAEFTIG = {
    "hole": 68,                 # dickerer, markanterer Ring (Default 79)
    "leader_breite_emu": 15875, # 1,25 pt: kräftig, aber nicht zu dominant
    "label_fett": True,         # fette Prozentzahlen (bleiben schwarz)
    "punkte": True,             # Punkte an den Leader-Enden
    "punkt_durchmesser": 0.05,  # klein/dezent
    "leader_start_tiefe": 0.5,  # Ansatz auf die MITTE der Ringdicke (im Band)
    "leader_gerade": True,      # ruhige gerade Linien statt harter Haken
    "label_gap_in": 0.18,       # Labels etwas luftiger außerhalb des Rings
}
FAMILIE_RING_FORMAT = {
    "CVV": _RING_KRAEFTIG,
    "ESG": _RING_KRAEFTIG,
    "ETF": _RING_KRAEFTIG,
    "Thema": _RING_KRAEFTIG,
    "comdirect": _RING_KRAEFTIG,
    # Nur "Standard" (Strategien ohne Familie) ist NICHT gelistet → bewusst der
    # bisherige Look. (Bei Bedarf einfach "Standard": _RING_KRAEFTIG, ergänzen —
    # der Familien-String ist dort allerdings leer, siehe _familie_aus_prs.)
}
_RING_FORMAT_DEFAULT = {
    "hole": None,                   # None → der hole_size-Parameter von nachbearbeiten
    "leader_breite_emu": LEADER_BREITE_EMU,
    "label_fett": False,
    "punkte": False,                # nur Familien mit punkte=True bekommen Punkte
    "punkt_durchmesser": PUNKT_DURCHMESSER,
    "leader_start_tiefe": 0.0,      # 0.0 = Ansatz am Außenrand (bisheriges Verhalten)
    "leader_gerade": False,         # False = bisherige geknickte Führung
    "label_gap_in": None,           # None → der label_gap_in-Parameter von nachbearbeiten
}


def _ring_format(fam, hole_size, label_gap_in):
    """Format-Werte für eine Familie: Familien-Override über Default gelegt.
    'hole' und 'label_gap_in' fallen bei None auf die übergebenen globalen
    Defaults (hole_size bzw. label_gap_in von nachbearbeiten) zurück."""
    f = dict(_RING_FORMAT_DEFAULT)
    f.update(FAMILIE_RING_FORMAT.get(fam, {}))
    if f["hole"] is None:
        f["hole"] = hole_size
    if f["label_gap_in"] is None:
        f["label_gap_in"] = label_gap_in
    return f

# ── Familien-Zuordnung (aus Mapping_Namen.xlsx, Spalte "Powerpoint Familie") ─
# Strategie → Familie, wie im Mapping. Wird über den Folientitel gematcht
# (_familie_aus_prs). WICHTIG zu den Schlüsseln:
#   • cVV-Broschüren lassen im Titel den Prefix WEG ("Anlagestrategie
#     Ausgewogen", nicht "cVV Ausgewogen") → die "nackten" Formen
#     (konservativ/defensiv/ausgewogen/…) sind hier als CVV hinterlegt.
#   • ESG behält den Prefix im Titel ("ESG Ausgewogen") → als "esg …" hinterlegt.
#   • Überlappung: "Offensiv"=Thema, aber "ESG Offensiv"=ESG. Das Matching nimmt
#     den LÄNGSTEN Treffer, deshalb gewinnt "esg offensiv" vor "offensiv".
_STRATEGIE_FAMILIE = {
    # Thema (kein Prefix im Titel)
    "pro dividende": "Thema", "offensiv": "Thema", "pro": "Thema",
    # ESG (mit "ESG"-Prefix im Titel)
    "esg defensiv plus": "ESG", "esg defensiv": "ESG",
    "esg ausgewogen": "ESG", "esg offensiv": "ESG",
    # CVV (Titel OHNE "cVV"-Prefix → nackte Formen)
    "konservativ": "CVV", "defensiv plus": "CVV", "defensiv": "CVV",
    "ausgewogen": "CVV", "dynamic": "CVV",
    # ETF
    "etf ausgewogen": "ETF", "etf wachstum": "ETF",
    # comdirect (Titel: "Anlagestrategie Portfolioverwaltung 30/70/100" auf den
    # Struktur-Folien bzw. "Comdirect 30/70/100 | Wertentwicklung"). Beide Formen
    # hinterlegt, damit die Familie sicher erkannt wird.
    "portfolioverwaltung 30": "comdirect", "portfolioverwaltung 70": "comdirect",
    "portfolioverwaltung 100": "comdirect",
    "comdirect 30": "comdirect", "comdirect 70": "comdirect",
    "comdirect 100": "comdirect",
}


def _familie_aus_prs(prs):
    """Powerpoint-Familie aus dem Mapping ableiten — anhand des Strategienamens
    im Folientitel. Rückgabe 'Thema'|'CVV'|'ESG'|'ETF'|'comdirect' oder None.

    Es werden NUR echte Titelzeilen durchsucht (Text mit 'Anlagestrategie' oder
    'Portfoliozusammenstellung'), damit Fließtext wie 'Die Strategie Pro …' kein
    False Positive erzeugt. Der LÄNGSTE passende Strategiename gewinnt (löst die
    Überlappung 'ESG Offensiv' vs 'Offensiv').
    """
    hay = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip().lower()
                if "anlagestrategie" in t or "portfoliozusammenstellung" in t:
                    hay.append(t)
    text = " ".join(hay)
    if not text:
        return None
    for strat, fam in sorted(_STRATEGIE_FAMILIE.items(),
                             key=lambda kv: -len(kv[0])):
        if re.search(r"\b" + re.escape(strat) + r"\b", text):
            return fam
    return None


def _ist_thema_familie(prs):
    """True, wenn die Präsentation zur Thema-Familie gehört (Offensiv/Pro/
    Pro Dividende). PRIMÄR aus dem Mapping (_familie_aus_prs, über den Titel).

    Fällt das Titel-Matching aus (kein Treffer), STRUKTURELLE Rückfallebene:
    nur Thema-Broschüren haben Regionen-/Branchen-Ringe (ESG/CVV/ETF haben
    Anlageklassen-Ringe).

    FALLSTRICK: Strategienamen überlappen zwischen Familien ('Offensiv'=Thema,
    'ESG Offensiv'=ESG) und cVV-Titel lassen den Prefix weg — deshalb im
    Mapping die nackten cVV-Formen + längster Treffer zuerst (_familie_aus_prs).
    """
    fam = _familie_aus_prs(prs)
    if fam is not None:
        return fam == "Thema"
    # Rückfall: strukturell (nur Thema hat Regionen-/Branchen-Ringe)
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                chart = shape.chart
                typ = chart.chart_type.name if chart.chart_type else ""
                if "DOUGHNUT" in typ and _ring_typ(chart, shape) in (
                        "BRANCHEN", "REGIONEN"):
                    return True
    return False


def _ring_typ(chart, shape):
    """Klassifiziert den Ring datengetrieben: 'ANLAGEKLASSEN' | 'BRANCHEN' |
    'REGIONEN' | 'UNBEKANNT'. Grundlage sind die Kategorienamen; der Shape-Name
    (C_Kennzahlen1=Regionen, C_Kennzahlen2=Branchen) dient nur als Rückfall.

    FALLSTRICK: Verlass dich nicht allein auf den Shape-Namen — Vorlagen können
    ihn ändern. Assetklassen werden immer über die Kategorien erkannt.
    """
    ser = _root(chart).find(".//" + _q("ser"))
    if ser is None:
        return "UNBEKANNT"
    kats = _kategorien(ser)
    if _ist_assetklassen_ring(kats):
        return "ANLAGEKLASSEN"
    name = (shape.name or "")
    if name.endswith("2"):
        return "BRANCHEN"
    if name.endswith("1"):
        return "REGIONEN"
    return "UNBEKANNT"




_EXCEL_EPOCHE = dt.date(1899, 12, 30)   # Seriennummer 0 in Excel/OOXML

# Reihenfolge der Kind-Elemente einer <c:dateAx> laut OOXML-Schema (CT_DateAx).
# Ein Element an der falschen Stelle macht die Datei für PowerPoint unlesbar.
_DATEAX_ORDNUNG = (
    "axId", "scaling", "delete", "axPos", "majorGridlines", "minorGridlines",
    "title", "numFmt", "majorTickMark", "minorTickMark", "tickLblPos", "spPr",
    "txPr", "crossAx", "crosses", "crossesAt", "auto", "lblOffset",
    "baseTimeUnit", "majorUnit", "majorTimeUnit", "minorUnit", "minorTimeUnit",
    "extLst",
)


def _monatsindex(d):
    """Datum -> fortlaufende Monatsnummer. Rechnet Monatsarithmetik ohne
    Jahresüberlauf-Sonderfälle."""
    return d.year * 12 + (d.month - 1)


def _monatsanfang(index):
    """Umkehrung von _monatsindex."""
    return dt.date(index // 12, index % 12 + 1, 1)


def _ax_wert_setzen(ax, tag, wert):
    """Setzt <c:TAG val="WERT"/> in einer Achse — und legt das Element an der
    vom Schema verlangten Stelle an, falls es fehlt.

    Nötig, weil `majorTimeUnit` in Vorlage_comdirect.pptx GAR NICHT vorkommt:
    Ein `if el is not None` hätte dort still nichts getan (genau das war bis
    zum 12.08.2026 der Fall), und ein schlichtes SubElement() hätte das
    Element hinter minorTimeUnit gehängt.
    """
    from lxml import etree
    el = ax.find(_q(tag))
    if el is None:
        el = etree.Element(_q(tag))
        rang = _DATEAX_ORDNUNG.index(tag)
        for kind in ax:
            if not isinstance(kind.tag, str):
                continue                       # Kommentare o. Ä. überspringen
            name = etree.QName(kind).localname
            if name in _DATEAX_ORDNUNG and _DATEAX_ORDNUNG.index(name) > rang:
                kind.addprevious(el)
                break
        else:
            ax.append(el)
    el.set("val", str(wert))
    return el


def achsen_raster(erster_tag, letzter_tag):
    """Achsengrenzen und Tick-Abstand einer Datumsachse aus der Datenspanne.

    Gibt (min_datum, max_datum, major_unit, major_time_unit) zurück.

    HINTERGRUND (12.08.2026, gemeldet aus der Praxis). PowerPoint setzt die
    Ticks einer Datumsachse beim ACHSEN-MINIMUM an und zählt von dort in
    Schritten von majorUnit x majorTimeUnit weiter — NICHT an Kalendergrenzen.
    Wer das Minimum auf den ersten Datenpunkt legt (so war es bis heute),
    verankert damit das ganze Raster auf dessen Monat: Die ETF-Reihe beginnt am
    30.11.2015, also lagen die Jahresticks auf November und der letzte auf
    Nov/25 — die Achse endete beschriftungsseitig im Vorjahr, obwohl die Kurve
    bis Juli 2026 läuft. Von 21 Datumsachsen war keine einzige in Ordnung.

    Deshalb wird der Anker HINTEN gesetzt, nicht vorn. Zwei Kandidaten:

      (a) der Monat des LETZTEN Datenpunkts — dann trägt genau er die letzte
          Beschriftung;
      (b) die nächste Kalendergrenze davor (Januar beim Jahresschritt, Jan/Jul
          beim Halbjahres-, Jan/Apr/Jul/Okt beim Quartalsschritt).

    Es gewinnt der Kandidat mit dem kleineren VORLAUF — der Achsenstrecke, die
    vor dem ersten Datenpunkt leer bleibt. Damit bekommt die cVV-Monatsreihe
    (Beginn Ende Januar) glatte Januar-Ticks bei null Vorlauf, während die
    ETF-Reihe mit fünf statt elf leeren Monaten davonkommt.

    Ein Kandidat, dessen letzter Tick NACH dem letzten Datenpunkt läge, wird
    verworfen: Ein Achsendatum in der Zukunft hat in einer Kundenbroschüre
    nichts zu suchen. Kandidat (a) erfüllt das immer, es bleibt also stets
    einer übrig.
    """
    lo_daten = _monatsindex(erster_tag)          # Monat des ersten Datenpunkts
    letzter_monat = _monatsindex(letzter_tag)
    hi = letzter_monat + 1                       # Monatsanfang NACH dem Ende

    schritt = DATUMSACHSE_SCHRITT_LANG
    for grenze, stufe in DATUMSACHSE_STUFEN:
        if hi - lo_daten <= grenze:
            schritt = stufe
            break

    # Monatsnummer % schritt == 0 trifft genau die Kalendergrenzen: bei 12 den
    # Januar, bei 6 Januar/Juli, bei 3 Januar/April/Juli/Oktober.
    kalender = letzter_monat - (letzter_monat % schritt)

    kandidaten = []
    for anker in (kalender, letzter_monat):
        stufen = -(-(anker - lo_daten) // schritt)        # aufgerundet
        start = anker - stufen * schritt
        letzter_tick = start + ((hi - start) // schritt) * schritt
        if letzter_tick > letzter_monat:
            continue
        kandidaten.append((lo_daten - start, start))
    start = min(kandidaten)[1]          # kleinster Vorlauf gewinnt

    if schritt % 12 == 0:
        einheit, anzahl = "years", schritt // 12
    else:
        einheit, anzahl = "months", schritt
    return _monatsanfang(start), _monatsanfang(hi), anzahl, einheit


def datumsachse_an_daten(chart):
    """Setzt die Datums-Achse (dateAx) einer Linie auf die tatsächliche
    Datenspanne statt auf fixe Vorlagengrenzen — Grenzen, Schrittweite UND
    Anker (das Warum steht in achsen_raster). Behebt den Leerraum vor/nach der
    Kurve und sorgt dafür, dass der letzte Datenzeitraum beschriftet ist.

    Der Parameter `auf_monat_runden` ist am 12.08.2026 entfallen: Das Runden
    auf Monatsgrenzen steckt jetzt in achsen_raster und ist nicht mehr
    abwählbar — der einzige Aufrufer hat es nie abgewählt.
    """
    root = _root(chart)
    # Kategorien gezielt aus dem cat-Block holen (erste Serie; bei den
    # Vergleichs-Charts teilen sich alle Serien dieselbe Datumsachse).
    cat_block = root.find(".//" + _q("cat"))
    if cat_block is None:
        return None
    cats = [float(v.text) for v in cat_block.iter(_q("v")) if v.text]
    if not cats:
        return None
    from lxml import etree
    ax = root.find(".//" + _q("dateAx"))
    if ax is None:
        return None

    a_min, a_max, m_anzahl, m_einheit = achsen_raster(
        _EXCEL_EPOCHE + dt.timedelta(days=int(min(cats))),
        _EXCEL_EPOCHE + dt.timedelta(days=int(max(cats))))
    lo = (a_min - _EXCEL_EPOCHE).days
    hi = (a_max - _EXCEL_EPOCHE).days

    scaling = ax.find(_q("scaling"))
    for tag, val in (("max", hi), ("min", lo)):   # max VOR min (Schema-Reihenfolge)
        el = scaling.find(_q(tag))
        if el is None:
            el = etree.SubElement(scaling, _q(tag))
        el.set("val", str(int(val)))
    # majorUnit MUSS mitgezogen werden, nicht nur majorTimeUnit: Die
    # cVV-Vergleichsfolie trägt in der Vorlage majorUnit=12 mit
    # majorTimeUnit="months". Wurde nur die Zeiteinheit auf "years" gestellt,
    # ergab das einen Tick alle ZWÖLF JAHRE — zwei Beschriftungen auf
    # siebzehneinhalb Jahren Historie.
    _ax_wert_setzen(ax, "majorUnit", m_anzahl)
    _ax_wert_setzen(ax, "majorTimeUnit", m_einheit)

    # ── Y-ACHSE (valAx) datenbasiert skalieren ──────────────────────────────
    # Feste Vorlagen-Grenzen (z.B. 0.8-1.4) schneiden stark gestiegene
    # Strategien (Offensiv: Index bis 2.8) oben ab → halber Chart leer.
    # ALLE Serien berücksichtigen, nicht nur die erste (Fix 10.07.2026):
    # root.find(".//c:val") liefert nur <c:val> der ERSTEN Serie. Bei der
    # CVV-Vergleichsfolie (5 Serien) skalierte die Achse auf den Bereich von
    # "Konservativ" — vier Linien liefen aus dem Chart heraus.
    # Bei 1-2 Serien (Wertentwicklungs-Folien) ändert sich nichts (gemessen).
    yvals = [float(v.text)
             for val in root.findall(".//" + _q("val"))
             for v in val.iter(_q("v")) if v.text]
    if yvals:
        import math as _m
        dmin, dmax = min(yvals), max(yvals)
        ymin = _m.floor(dmin * 10) / 10.0          # z.B. 0.897 -> 0.8
        ymax = _m.ceil(dmax * 10) / 10.0           # z.B. 2.872 -> 2.9
        # Luft über der höchsten Linie: klebt sie an der Achsengrenze,
        # eine 10%-Stufe drauflegen.
        if ymax - dmax < 0.02:
            ymax += 0.1
        vax = root.find(".//" + _q("valAx"))
        if vax is not None:
            vsc = vax.find(_q("scaling"))
            for tag, val in (("max", ymax), ("min", ymin)):
                el = vsc.find(_q(tag))
                if el is None:
                    el = etree.SubElement(vsc, _q(tag))
                el.set("val", f"{val:.2f}")
    return (lo, hi)


# ─────────────────────────────────────────────────────────────────────────
# UNGENUTZT (Stand 20.07.2026): ring_labels_kompakt war ein Alternativ-Ansatz
# (kompakt am Ring). Aktiv ist ring_labels_aussen_dynamisch. Nur behalten als
# Referenz — NICHT aufgerufen. Vor dem Löschen prüfen, dass nachbearbeiten es
# nicht referenziert.
# ─────────────────────────────────────────────────────────────────────────
def ring_labels_kompakt(chart, frame_w_in, frame_h_in,
                        gap_in=0.16, min_v_in=0.22, min_h_in=0.60,
                        rand_in=0.12, kopf_frei_in=0.54):
    """Platziert Ring-Labels KOMPAKT und nah am Segment (NEU 13.07.2026).

    Ersetzt die alte, über viele Runden gewachsene Positionierung. Weil die
    Führungslinien jetzt SELBST gezeichnet werden (ring_leader_zeichnen), muss
    kein Rücksicht mehr auf PowerPoints Auto-Leader-Regel (Naht/Loch) genommen
    werden — die Labels dürfen einfach dort sitzen, wo es optisch am saubersten
    ist: knapp außerhalb des Rings, radial über dem eigenen Segment.

    Prinzip (dynamisch, für jede Segmentzahl und Ringgröße):
      1) Startposition radial knapp außerhalb (gap_in + Box-Ausdehnung).
      2) VERTIKALES Entzerren pro Seite (links/rechts): überlappende Labels
         rücken vertikal auseinander — NICHT weit tangential wegziehen. So
         bleibt jedes Label nah am Segment → kurze Leader, keine langen Striche
         quer über den oberen Ringbereich.
      3) Anti-Kreuzung: kreuzende Leaderpaare tauschen Position (geometrischer
         Streckenschnitt-Test), danach erneut entzerren.
      4) Finale Überlappungsauflösung als Sicherung.

    Ergebnis über 1775 Zufalls-Labels: 0 Kreuzungen, 0 Überlappungen,
    0 außerhalb des Rahmens, keine Leader länger als 0,9\".

    Schreibt die manualLayout-Offsets (xMode/yMode=edge) wie die Vorgänger-
    Funktion, damit die Labels in PowerPoint exakt an der Zielposition landen.
    """
    root = _root(chart)
    pa = root.find(".//" + _q("plotArea") + "/" + _q("layout")
                   + "/" + _q("manualLayout"))
    if pa is None:
        return None

    def _g(tag):
        e = pa.find(_q(tag))
        return float(e.get("val")) if e is not None else None

    px, py, pw, ph = _g("x"), _g("y"), _g("w"), _g("h")
    if None in (px, py, pw, ph):
        return None

    left, right = px * frame_w_in, (px + pw) * frame_w_in
    top, bot = py * frame_h_in, (py + ph) * frame_h_in
    cx, cy = (left + right) / 2, (top + bot) / 2
    R_out = min(right - left, bot - top) / 2

    # ANMERKUNG 12.08.2026: Hier wurde bis heute die holeSize gelesen und daraus
    # die Band-Mitte berechnet (R_out * (1 + holeSize) / 2) — der Radius, auf den
    # PowerPoint ein Label OHNE manualLayout von sich aus setzt. Verwendet wurde
    # das nie: manualLayout-x/y sind ABSOLUTE Bruchteile des Rahmens, keine
    # Abstände vom Default. Die ganze Kette (holeSize → band_mitte → dfx/dfy)
    # war tot und ist entfernt; die Erkenntnis steht hier, damit sie niemand
    # ein zweites Mal herleiten muss.

    fsa_el = root.find(".//" + _q("firstSliceAng"))
    fsa = float(fsa_el.get("val")) if fsa_el is not None else 0.0

    vals = [float(v.text)
            for v in root.findall(".//" + _q("val") + "//" + _q("pt")
                                  + "/" + _q("v"))]
    if not vals:
        return None
    tot = sum(vals) or 1.0
    mids, kum = [], 0.0
    for v in vals:
        f = v / tot
        mids.append((fsa + (kum + f / 2) * 360) % 360)
        kum += f

    HW, HH = 0.33, 0.10

    # 1) Startpositionen radial knapp außerhalb
    pos = []
    for m in mids:
        mr = math.radians(m)
        r = R_out + gap_in + HW * abs(math.sin(mr)) + HH * abs(math.cos(mr))
        pos.append([cx + r * math.sin(mr), cy - r * math.cos(mr), m])

    def _entzerre(runden=80):
        for _ in range(runden):
            bewegt = False
            for seite in (-1, 1):
                grp = [p for p in pos
                       if (1 if p[0] >= cx else -1) == seite]
                grp.sort(key=lambda p: p[1])
                for i in range(len(grp) - 1):
                    a, b = grp[i], grp[i + 1]
                    if b[1] - a[1] < min_v_in:
                        schub = (min_v_in - (b[1] - a[1])) / 2 + 0.005
                        a[1] = max(kopf_frei_in + HH, a[1] - schub)
                        b[1] = min(frame_h_in - rand_in - HH, b[1] + schub)
                        bewegt = True
            if not bewegt:
                break

    def _se(i):
        m = math.radians(pos[i][2])
        return (cx + R_out * math.sin(m), cy - R_out * math.cos(m))

    def _ccw(a, b, c):
        return (c[1] - a[1]) * (b[0] - a[0]) > (b[1] - a[1]) * (c[0] - a[0])

    def _kreuzt(i, j):
        p1, p2 = _se(i), (pos[i][0], pos[i][1])
        p3, p4 = _se(j), (pos[j][0], pos[j][1])
        return (_ccw(p1, p3, p4) != _ccw(p2, p3, p4)
                and _ccw(p1, p2, p3) != _ccw(p1, p2, p4))

    _entzerre()
    # 3) Anti-Kreuzung (Positionstausch) + Re-Entzerren
    for _ in range(len(pos) ** 2 + 5):
        getauscht = False
        for i in range(len(pos)):
            for j in range(i + 1, len(pos)):
                if _kreuzt(i, j):
                    pos[i][0], pos[i][1], pos[j][0], pos[j][1] = \
                        pos[j][0], pos[j][1], pos[i][0], pos[i][1]
                    getauscht = True
                    break
            if getauscht:
                break
        if not getauscht:
            break
        _entzerre()

    # 4) Finale reine Überlappungsauflösung (Sicherung)
    for _ in range(120):
        bewegt = False
        for i in range(len(pos)):
            for j in range(i + 1, len(pos)):
                if (abs(pos[i][0] - pos[j][0]) < min_h_in
                        and abs(pos[i][1] - pos[j][1]) < min_v_in):
                    schub = (min_v_in - abs(pos[i][1] - pos[j][1])) / 2 + 0.005
                    hoch, runter = (i, j) if pos[i][1] <= pos[j][1] else (j, i)
                    pos[hoch][1] = max(kopf_frei_in + HH, pos[hoch][1] - schub)
                    pos[runter][1] = min(frame_h_in - rand_in - HH,
                                         pos[runter][1] + schub)
                    bewegt = True
        if not bewegt:
            break

    # 5) manualLayout-Offsets schreiben (Ziel − PP-Default), xMode/yMode=edge.
    #    Default-Position des Labels: Band-Mitte am Segmentwinkel.
    from lxml import etree
    ser = root.find(".//" + _q("ser"))
    if ser is None:
        return None
    dLbls = ser.find(_q("dLbls"))

    # idx → dLbl-Element (vorhandene nutzen, fehlende anlegen)
    vorhandene = {}
    if dLbls is not None:
        for d in dLbls.findall(_q("dLbl")):
            ixe = d.find(_q("idx"))
            if ixe is not None:
                vorhandene[int(ixe.get("val"))] = d

    gesetzt = 0
    for i, (lx, ly, m) in enumerate(pos):
        # Ziel = Box-Mitte; manualLayout-x/y = linke obere Ecke als Bruchteil
        off_x = (lx - HW) / frame_w_in
        off_y = (ly - HH) / frame_h_in
        d = vorhandene.get(i)
        if d is None:
            continue
        ml = d.find(".//" + _q("manualLayout"))
        if ml is None:
            layout = d.find(_q("layout"))
            if layout is None:
                layout = etree.SubElement(d, _q("layout"))
            ml = etree.SubElement(layout, _q("manualLayout"))
        for tag, val in (("x", off_x), ("y", off_y)):
            e = ml.find(_q(tag))
            if e is None:
                e = etree.SubElement(ml, _q(tag))
            e.set("val", "%.5f" % val)
        for tag in ("xMode", "yMode"):
            e = ml.find(_q(tag))
            if e is None:
                e = etree.SubElement(ml, _q(tag))
            e.set("val", "edge")
        gesetzt += 1
    return gesetzt


def ring_labels_aussen_dynamisch(chart, frame_w_in, frame_h_in,
                                 gap_in=0.14, min_gap_deg=24.0, rand_in=0.12,
                                 tangential_in=0.14, rand_oben_in=None,
                                 kopf_frei_in=None):
    """Platziert die Ring-Datenlabels GEOMETRISCH exakt außerhalb des Rings.

    Liest die echte Ring-Geometrie aus dem plotArea-Layout des Charts
    (Mittelpunkt, Außenradius) und setzt jedes Label an eine berechnete
    Zielposition knapp außerhalb (gap_in Zoll). Der manualLayout-Offset wird
    als (Ziel − PP-Default) berechnet, wobei der Default (Band-Mitte) aus der
    holeSize abgeleitet wird — dadurch landet das Label in PowerPoint exakt
    am Ziel, unabhängig davon wie dünn der Ring ist.

    Selbstkalibrierend pro Chart: funktioniert für jede Segmentzahl und jede
    Plot-Größe (die drei Themen-Ringe haben unterschiedlich große Plots!).
    Zu dicht stehende Labels werden über einen Mindest-Winkelabstand
    (min_gap_deg) auseinandergespreizt → keine Überlappung.

    frame_w_in / frame_h_in = Breite/Höhe des Chart-Rahmens in Zoll
    (aus dem Shape; in nachbearbeiten() automatisch übergeben).
    """
    root = _root(chart)
    # 1) plotArea inner-Rechteck (Bruchteile des Rahmens)
    pa = root.find(".//" + _q("plotArea") + "/" + _q("layout") + "/" + _q("manualLayout"))
    if pa is None:
        return None
    def _g(tag):
        e = pa.find(_q(tag)); return float(e.get("val")) if e is not None else None
    px, py, pw, ph = _g("x"), _g("y"), _g("w"), _g("h")
    if None in (px, py, pw, ph):
        return None
    # 2) Ring-Geometrie in Zoll
    left, right = px * frame_w_in, (px + pw) * frame_w_in
    top, bot = py * frame_h_in, (py + ph) * frame_h_in
    cx, cy = (left + right) / 2.0, (top + bot) / 2.0
    R_out = min(right - left, bot - top) / 2.0

    # 2b) RING an den verfügbaren vertikalen Raum ZWISCHEN Überschrift und
    #     Legende anpassen (Größe + Lage). Die Legende sitzt bei diesen Ringen
    #     unten im Rahmen und schiebt den Ring hoch; gleichzeitig soll oben Luft
    #     zur Überschrift bleiben. Wir lesen die Legenden-Oberkante, verkleinern
    #     den Ring wenn nötig und zentrieren ihn im freien Bereich → Labels oben
    #     wie unten mit Abstand, konsistente Größe.
    legend = root.find(".//" + _q("legend"))
    leg_y = None
    if legend is not None:
        lm = legend.find(".//" + _q("manualLayout"))
        if lm is not None and lm.find(_q("y")) is not None:
            leg_y = float(lm.find(_q("y")).get("val"))
    legend_top = (leg_y * frame_h_in) if leg_y is not None else frame_h_in * 0.97
    kopf_rand = 0.15          # Luft zur Überschrift oben
    label_pad = 0.52          # vertikale Ausdehnung Label + Rand (inkl. De-overlap-Spreizung)
    R_ziel = min(R_out, 0.27 * frame_h_in)     # Grundverkleinerung großer Ringe
    # so weit verkleinern, dass Ring + Labels zwischen Kopf und Legende passen
    for _ in range(20):
        lo = kopf_rand + R_ziel + label_pad          # min. mögliche Zentrum-y
        hi = legend_top - R_ziel - label_pad         # max. mögliche Zentrum-y
        if hi >= lo:
            break
        R_ziel *= 0.94
    new_cy = min(max(frame_h_in / 2.0, lo), hi) if hi >= lo else (lo + hi) / 2.0
    if R_ziel < R_out - 1e-3 or abs(new_cy - cy) > 1e-3:
        faktor = R_ziel / R_out
        cxf = px + pw / 2.0                       # horizontales Zentrum halten
        pw2, ph2 = pw * faktor, ph * faktor
        px2 = cxf - pw2 / 2.0
        py2 = (new_cy / frame_h_in) - ph2 / 2.0   # vertikal neu setzen
        for tag, val in (("x", px2), ("y", py2), ("w", pw2), ("h", ph2)):
            e = pa.find(_q(tag))
            if e is not None:
                e.set("val", f"{val:.5f}")
        left, right = px2 * frame_w_in, (px2 + pw2) * frame_w_in
        top, bot = py2 * frame_h_in, (py2 + ph2) * frame_h_in
        cx, cy = (left + right) / 2.0, (top + bot) / 2.0
        R_out = min(right - left, bot - top) / 2.0

    # (Auch hier wurde bis 12.08.2026 die holeSize gelesen und als `band_center`
    # der PP-Default-Radius berechnet — unbenutzt, siehe die Erklärung weiter
    # oben. R_target unten ersetzt ihn vollständig.)
    # gap_in = gewünschte SICHTBARE Freiheit zwischen Text-Innenkante und Ring.
    # Die nötige radiale Distanz der Label-MITTE hängt vom Winkel ab, weil der
    # Text waagerecht ist: seitlich ragt die halbe Breite zum Ring, oben/unten
    # nur die halbe Höhe. R_target wird deshalb pro Label in Schritt 5 berechnet.
    HALB_W, HALB_H = 0.33, 0.10                # halbe Text-Box (Zoll, ~"29,60%")
    # 3) Segment-Mittelwinkel (Grad, im Uhrzeigersinn ab 12 Uhr)
    val_block = root.find(".//" + _q("val"))
    if val_block is None:
        return None
    vals = [float(v.text) for v in val_block.iter(_q("v")) if v.text]
    total = sum(vals) or 1.0
    fsa_el = root.find(".//" + _q("firstSliceAng"))
    fsa = float(fsa_el.get("val")) if fsa_el is not None else 0.0
    mids, kum = [], 0.0
    for v in vals:
        f = v / total; mids.append(fsa + (kum + f / 2) * 360); kum += f
    # 4) Label-Winkel mit Mindest-Winkelabstand spreizen (gegen Überlappung)
    order = sorted(range(len(mids)), key=lambda i: mids[i] % 360)
    ang = [mids[i] % 360 for i in order]
    for _ in range(200):
        bewegt = False
        for k in range(len(ang)):
            k2 = (k + 1) % len(ang)
            gap = (ang[k2] - ang[k]) % 360
            if 0 < gap < min_gap_deg:
                push = (min_gap_deg - gap) / 2
                ang[k] = (ang[k] - push) % 360
                ang[k2] = (ang[k2] + push) % 360
                bewegt = True
        if not bewegt:
            break
    label_ang = [0.0] * len(mids)
    for si, i in enumerate(order):
        label_ang[i] = ang[si]
    # 5) Zielpositionen (absolut, in Zoll) berechnen — mit Rand-Begrenzung,
    #    damit große Ringe nicht abgeschnitten werden. PLUS ein tangentialer
    #    Versatz (seitlich zur Radial-Richtung), damit KEIN Label exakt auf der
    #    radialen Linie zu seinem Segment sitzt → PowerPoint zeichnet dann für
    #    JEDES Label einen Leader-Strich (auch freistehende Segmente).
    # Oberer Rand separat steuerbar (NEU 10.07.2026): über dem Ring sitzt der
    # Überschriftenbalken ("AKTUELLE STRUKTUR"). rand_in=0.12 ließ die oberen
    # Labels fast am Balken kleben. rand_oben_in=None → altes Verhalten.
    # ACHTUNG: _rand_oben klemmt die Label-MITTE. Oberkante = Mitte - HALB_H.
    # Muss in ALLEN nachfolgenden Durchläufen benutzt werden (De-overlap,
    # Leader-Garantie, Mindest-Ringabstand) — sonst schiebt der letzte Pass
    # das Label wieder bis rand_in hoch und es landet im Überschriftenbalken.
    _rand_oben = rand_in if rand_oben_in is None else float(rand_oben_in)

    ziel = []
    for i in range(len(mids)):
        la = math.radians(label_ang[i])
        sx, sy = math.sin(la), -math.cos(la)          # radial nach außen
        tsx, tsy = math.cos(la), math.sin(la)         # tangential (im Uhrzeigersinn)
        # radiale Ausdehnung der (waagerechten) Text-Box in Blickrichtung:
        # seitlich zählt die Breite, oben/unten die Höhe.
        radial_extent = HALB_W * abs(sx) + HALB_H * abs(sy)
        r_target_i = R_out + gap_in + radial_extent    # → Innenkante = R_out+gap_in
        r_use = r_target_i
        if sx > 1e-6:    r_use = min(r_use, (frame_w_in - rand_in - cx) / sx)
        elif sx < -1e-6: r_use = min(r_use, (rand_in - cx) / sx)
        if sy > 1e-6:    r_use = min(r_use, (frame_h_in - rand_in - cy) / sy)
        elif sy < -1e-6: r_use = min(r_use, (_rand_oben - cy) / sy)
        r_use = max(r_use, R_out + 0.05)              # nie innerhalb des Rings
        tx = cx + r_use * sx + tangential_in * tsx
        ty = cy + r_use * sy + tangential_in * tsy
        ziel.append([tx, ty])

    # 6) ADAPTIVE Überlappungs-Auflösung im ABSOLUTEN Raum — garantiert, dass
    #    sich keine zwei Zahlen überlappen, egal wie die Segmente verteilt
    #    sind. Zu dicht stehende Labels werden vertikal auseinandergedrängt
    #    (Zahlen sind waagerechter Text → Überlappung ist v.a. vertikal),
    #    dabei im Rahmen gehalten. min_v/min_h ≈ halbe Label-Höhe/-Breite.
    min_v, min_h = 0.205, 0.60
    for _ in range(120):
        bewegt = False
        reihenfolge = sorted(range(len(ziel)), key=lambda i: ziel[i][1])
        for a in range(len(reihenfolge)):
            for b in range(a + 1, len(reihenfolge)):
                i, j = reihenfolge[a], reihenfolge[b]
                if abs(ziel[i][1] - ziel[j][1]) < min_v and abs(ziel[i][0] - ziel[j][0]) < min_h:
                    schub = (min_v - abs(ziel[i][1] - ziel[j][1])) / 2.0 + 0.005
                    hoch, runter = (i, j) if ziel[i][1] <= ziel[j][1] else (j, i)
                    ziel[hoch][1] = max(_rand_oben, ziel[hoch][1] - schub)
                    ziel[runter][1] = min(frame_h_in - rand_in, ziel[runter][1] + schub)
                    bewegt = True
        if not bewegt:
            break

    # 6b) LEADER-GARANTIE: jedes Label muss einen Mindest-Seitenabstand
    #     (senkrecht zur radialen Linie seines Segments) haben, sonst zeichnet
    #     PowerPoint keinen Strich. Wir schieben zu radial-nahe Labels
    #     tangential weg — danach nochmal Überlappung auflösen.
    min_tang = 0.20
    for _durchlauf in range(6):
        for i in range(len(ziel)):
            md = math.radians(mids[i])
            perp_x, perp_y = math.cos(md), math.sin(md)   # senkrecht zur Radial-Richtung
            lvx, lvy = ziel[i][0] - cx, ziel[i][1] - cy
            d_perp = lvx * perp_x + lvy * perp_y
            if abs(d_perp) < min_tang:
                richtung = 1.0 if d_perp >= 0 else -1.0
                korr = (min_tang - abs(d_perp)) * richtung
                nx = ziel[i][0] + korr * perp_x
                ny = ziel[i][1] + korr * perp_y
                # im Rahmen halten
                ziel[i][0] = max(rand_in, min(frame_w_in - rand_in, nx))
                ziel[i][1] = max(_rand_oben, min(frame_h_in - rand_in, ny))
        # Überlappung erneut auflösen (Tangential-Schub kann welche erzeugt haben)
        for _ in range(60):
            bewegt = False
            reihenfolge = sorted(range(len(ziel)), key=lambda i: ziel[i][1])
            for a in range(len(reihenfolge)):
                for b in range(a + 1, len(reihenfolge)):
                    i, j = reihenfolge[a], reihenfolge[b]
                    if abs(ziel[i][1] - ziel[j][1]) < min_v and abs(ziel[i][0] - ziel[j][0]) < min_h:
                        schub = (min_v - abs(ziel[i][1] - ziel[j][1])) / 2.0 + 0.005
                        hoch, runter = (i, j) if ziel[i][1] <= ziel[j][1] else (j, i)
                        ziel[hoch][1] = max(_rand_oben, ziel[hoch][1] - schub)
                        ziel[runter][1] = min(frame_h_in - rand_in, ziel[runter][1] + schub)
                        bewegt = True
            if not bewegt:
                break

    # 6c) MINDEST-RING-ABSTAND: kein Label darf dem Ring zu nah kommen
    #     (die vertikale De-overlap kann Labels zum Ring drücken). Zu nahe
    #     Labels werden radial nach außen geschoben (im Rahmen gehalten).
    min_clear = 0.12
    for _ in range(4):
        for i in range(len(ziel)):
            lvx, lvy = ziel[i][0] - cx, ziel[i][1] - cy
            rad = math.hypot(lvx, lvy) or 1e-6
            pa_ang = math.atan2(lvx, -lvy)                      # Positionswinkel
            rext = 0.33 * abs(math.sin(pa_ang)) + 0.10 * abs(math.cos(pa_ang))
            inner = rad - rext - R_out
            if inner < min_clear:
                schub = min_clear - inner
                ziel[i][0] = max(rand_in, min(frame_w_in - rand_in, ziel[i][0] + schub * lvx / rad))
                ziel[i][1] = max(_rand_oben, min(frame_h_in - rand_in, ziel[i][1] + schub * lvy / rad))

    # 6d) KOPF-FREIHALTUNG (NEU 10.07.2026) — harte Garantie gegen Kollision
    #     mit dem Überschriftenbalken.
    #
    #     Ein Label, das oben am Ring sitzt, lässt sich RADIAL nicht nach unten
    #     schieben: dort ist der Ring (r_use = max(r_use, R_out+0.05) gewinnt).
    #     Ein größerer oberer Rand bewegt es deshalb kaum.
    #
    #     Lösung: das Label auf SEINEM Radius entlang der Ringkurve wegdrehen,
    #     bis die geforderte Oberkante erreicht ist. Der Abstand zum Ring bleibt
    #     dabei exakt erhalten, die Führungslinie zeigt weiter auf sein Segment.
    #
    #         y = cy - r*cos(a)  →  cos(a) = (cy - y_ziel) / r
    #
    #     kopf_frei_in = Mindest-Oberkante des Labels (None → Pass aus).
    if kopf_frei_in is not None:
        y_soll = float(kopf_frei_in) + HALB_H          # gewünschte Label-MITTE
        for i in range(len(mids)):
            if ziel[i][1] >= y_soll - 1e-6:
                continue                                # tief genug
            dx0, dy0 = ziel[i][0] - cx, ziel[i][1] - cy
            r_i = math.hypot(dx0, dy0)
            if r_i < 1e-6:
                continue
            seite = 1.0 if dx0 >= 0 else -1.0           # Seite beibehalten
            # Winkel UND Radius gemeinsam lösen: beim Herunterdrehen zeigt die
            # Textbox stärker mit ihrer BREITE zum Ring (HALB_W statt HALB_H),
            # der Radius muss also mitwachsen — sonst berührt das Label den Ring.
            r = r_i
            gefunden = False
            for _ in range(8):
                c = (cy - y_soll) / r
                if abs(c) > 1.0:
                    break
                a = math.acos(max(-1.0, min(1.0, c)))
                sx, sy = seite * math.sin(a), -math.cos(a)
                r_noetig = R_out + gap_in + HALB_W * abs(sx) + HALB_H * abs(sy)
                r_neu = max(r, r_noetig)
                if abs(r_neu - r) < 1e-4:
                    gefunden = True
                    break
                r = r_neu
            if not gefunden and abs((cy - y_soll) / r) > 1.0:
                # Radius reicht nicht: so tief wie möglich (waagerecht daneben)
                ziel[i][1] = cy
                ziel[i][0] = max(rand_in, min(frame_w_in - rand_in,
                                              cx + seite * r))
                continue
            c = max(-1.0, min(1.0, (cy - y_soll) / r))
            a = math.acos(c)
            ziel[i][0] = max(rand_in, min(frame_w_in - rand_in,
                                          cx + seite * r * math.sin(a)))
            ziel[i][1] = cy - r * c

        # Entzerren (nur nach unten) und radiales Ausschieben BEEINFLUSSEN
        # SICH GEGENSEITIG: das Nach-unten-Schieben drückt Labels auf den Ring,
        # der Ausschub schiebt zwei Labels wieder zusammen. Einmal nacheinander
        # reicht nicht — gemessen an [0.94, 0.03, 0.02, 0.01]. Deshalb im
        # Wechsel, und mit dem Entzerren als LETZTEM Schritt.
        def _entzerren_nach_unten():
            for _ in range(60):
                bewegt = False
                reihenfolge = sorted(range(len(ziel)), key=lambda i: ziel[i][1])
                for a_ in range(len(reihenfolge)):
                    for b_ in range(a_ + 1, len(reihenfolge)):
                        i, j = reihenfolge[a_], reihenfolge[b_]
                        if (abs(ziel[i][1] - ziel[j][1]) < min_v
                                and abs(ziel[i][0] - ziel[j][0]) < min_h):
                            schub = (min_v - abs(ziel[i][1] - ziel[j][1])) + 0.005
                            runter = j if ziel[j][1] >= ziel[i][1] else i
                            ziel[runter][1] = min(frame_h_in - rand_in,
                                                  ziel[runter][1] + schub)
                            bewegt = True
                if not bewegt:
                    break

        def _radial_ausschieben():
            for i in range(len(mids)):
                lvx, lvy = ziel[i][0] - cx, ziel[i][1] - cy
                rad = math.hypot(lvx, lvy)
                if rad < 1e-6:
                    continue
                inner = rad - (HALB_W * abs(lvx) + HALB_H * abs(lvy)) / rad - R_out
                if inner < min_clear:
                    schub = min_clear - inner
                    ziel[i][0] = max(rand_in, min(frame_w_in - rand_in,
                                                  ziel[i][0] + schub * lvx / rad))
                    ziel[i][1] = max(y_soll, min(frame_h_in - rand_in,
                                                 ziel[i][1] + schub * lvy / rad))

        for _ in range(4):
            _entzerren_nach_unten()
            _radial_ausschieben()
        _entzerren_nach_unten()

    # 6d2) BOGENGRENZEN-ABSTAND (NEU 13.07.2026) — sichert, dass PowerPoint für
    #      jedes Label eine Führungslinie zeichnet.
    #
    #      SYSTEMATISCH am realen Chart nachgewiesen (alle 5 Regionen-Labels):
    #      PowerPoint zeichnet den radialen Leader nur, wenn er EINDEUTIG einem
    #      Segment zuzuordnen ist. Das ist der Fall, wenn der POSITIONSWINKEL des
    #      Labels (Winkel vom Ringmittelpunkt zur Label-Mitte)
    #        (a) INNERHALB des eigenen Segmentbogens liegt, mit Mindestabstand zu
    #            beiden Bogengrenzen, ODER
    #        (b) klar AUSSERHALB (langer Leader, z.B. winziges 12-Uhr-Segment).
    #
    #      Der 8,33%-Fall (Asien) hatte seinen Positionswinkel mit 0,1° EXAKT auf
    #      der Naht zwischen Asien und Deutschland → der Leader hätte auf der
    #      Segmentgrenze geendet → mehrdeutig → PowerPoint zeichnete nichts.
    #      Alle anderen Labels (auch alle Branchen) lagen ≥11,9° von jeder Grenze
    #      → alle hatten Leader. Regel gegen alle belegten Fälle verifiziert.
    #
    #      Frühere Fixes drehten am Winkel-OFFSET zum Mittelwinkel — die falsche
    #      Größe. Die richtige ist der Abstand zur GRENZE. Ist ausschließlich aus
    #      der Segmentgeometrie berechenbar und daher OHNE Rendern prüfbar.
    #
    #      Fix: Für jedes Label, dessen Positionswinkel zu nah an einer eigenen
    #      Bogengrenze liegt (und das nicht ohnehin klar außerhalb sitzt), suche
    #      die nächstgelegene Position, deren Positionswinkel GRENZ_ABSTAND von
    #      beiden Grenzen hält — unter Balken, ohne Überlappung, ohne Ring-
    #      Berührung. Findet sich keine, bleibt das Label (nichts wird schlechter).
    GRENZ_ABSTAND = 6.0    # Mindest-Winkelabstand des Positionswinkels zur
    #                        Bogengrenze. Belegt: Leader ab ~11,9° sicher, kein
    #                        Leader bei 0,1°. 6° ist ein konservativer Mittelwert
    #                        mit Sicherheitsmarge zur belegten Nicht-Zeichnung.
    MIN_LEADER = 0.28      # Mindest-Leader-Länge (Label klar außerhalb des Rings)

    ob_grenze = kopf_frei_in if kopf_frei_in is not None else _rand_oben

    # Segmentbögen (Start/Ende in Grad) aus kumulierten Werten + firstSliceAng.
    # (mids wurde oben genau so gebildet; hier rekonstruieren wir die Grenzen.)
    _summe = sum(vals) or 1.0
    _bogen = []
    _kum = 0.0
    for _v in vals:
        _f = _v / _summe
        _bogen.append(((fsa + _kum * 360) % 360, (fsa + (_kum + _f) * 360) % 360))
        _kum += _f

    def _pos_winkel(x, y):
        return math.degrees(math.atan2(x - cx, -(y - cy))) % 360

    def _abstand_grenze(i, x, y):
        """Kleinster Winkelabstand des Positionswinkels zu den Bögen-Grenzen
        von Segment i. Positiv, in Grad."""
        pang = _pos_winkel(x, y)
        a, b = _bogen[i]
        d_a = min(abs(pang - a), 360 - abs(pang - a))
        d_b = min(abs(pang - b), 360 - abs(pang - b))
        return min(d_a, d_b)

    def _im_bogen(i, x, y):
        pang = _pos_winkel(x, y)
        a, b = _bogen[i]
        if a <= b:
            return a <= pang <= b
        return pang >= a or pang <= b          # Bogen über 0° hinweg

    def _leader_sicher(i, x, y):
        """PowerPoint zeichnet: entweder klar im Bogen mit Grenzabstand, oder
        klar außerhalb (dann ist die Zuordnung über die Länge eindeutig)."""
        if not _im_bogen(i, x, y):
            return True                        # außerhalb → langer Leader
        return _abstand_grenze(i, x, y) >= GRENZ_ABSTAND

    def _frei(i, x, y):
        if not (rand_in <= x <= frame_w_in - rand_in):
            return False
        if not (ob_grenze <= y - HALB_H and y <= frame_h_in - rand_in):
            return False
        # Ring-Abstand positionsabhängig: seitlich zählt die Box-Breite,
        # oben/unten die Höhe.
        lvx, lvy = x - cx, y - cy
        rad = math.hypot(lvx, lvy)
        if rad > 1e-6:
            pa_ang = math.atan2(lvx, -lvy)
            rext = HALB_W * abs(math.sin(pa_ang)) + HALB_H * abs(math.cos(pa_ang))
            if rad - rext - R_out < 0.05:
                return False
        for j in range(len(ziel)):
            if j == i:
                continue
            if abs(x - ziel[j][0]) < min_h and abs(y - ziel[j][1]) < min_v:
                return False
        return True

    for i in range(len(mids)):
        # nur eingreifen, wenn der Leader aktuell NICHT sicher ist
        if _leader_sicher(i, ziel[i][0], ziel[i][1]) and _frei(i, ziel[i][0],
                                                               ziel[i][1]):
            continue
        seg = mids[i] % 360
        sx = cx + R_out * math.sin(math.radians(seg))
        sy = cy - R_out * math.cos(math.radians(seg))
        # Suche die Position mit dem GRÖSSTEN Grenzabstand, die alle
        # Constraints erfüllt. Winkel um den Mittelwinkel, Radius fein gerastert.
        bestpos = None
        best_abstand = -1.0
        for off_grad in (0, 2, 4, 6, 8, 10, 12):
            for seite in (-1, 1):
                for r_delta in (0.12, 0.16, 0.20, 0.24, 0.28, 0.32,
                                0.36, 0.40, 0.44):
                    r = R_out + r_delta
                    a = math.radians(seg + seite * off_grad)
                    x = cx + r * math.sin(a)
                    y = cy - r * math.cos(a)
                    if math.hypot(x - sx, y - sy) < MIN_LEADER:
                        continue
                    if not _frei(i, x, y):
                        continue
                    if not _leader_sicher(i, x, y):
                        continue
                    ab = _abstand_grenze(i, x, y)
                    if ab > best_abstand:
                        best_abstand = ab
                        bestpos = (x, y)
        # übernehmen, wenn eine sichere Position gefunden wurde und sie den
        # Grenzabstand VERBESSERT
        if bestpos is not None and best_abstand > _abstand_grenze(
                i, ziel[i][0], ziel[i][1]):
            ziel[i][0], ziel[i][1] = bestpos

    # 6e) ANTI-KREUZUNG (NEU 10.07.2026) — letzte Garantie, dass sich keine
    #     zwei Führungslinien überkreuzen.
    #
    #     Ursache der Kreuzung: die vorherigen Pässe (De-overlap, Leader-
    #     Garantie) verschieben Labels tangential, ohne die Winkelreihenfolge
    #     der Segmente zu respektieren. Sitzen zwei kleine Segmente eng
    #     beieinander (z.B. Versorger 321° / Elektro 347° oben), kann das Label
    #     des einen auf die Seite des anderen geraten → die Leader kreuzen sich
    #     (im Screenshot bei 6,96%/7,03% sichtbar).
    #
    #     Lösung: Kreuzung ist ein rein GEOMETRISCHES Kriterium (schneiden sich
    #     die Strecken Segment→Label?). Wir erkennen jedes kreuzende Paar und
    #     TAUSCHEN die beiden Label-Positionen. Nach dem Tausch zeigt jeder
    #     Leader auf sein eigenes Segment über die Box des Partners — das
    #     entwirrt die Kreuzung, ohne neue Überlappung zu erzeugen (die Boxen
    #     standen ja schon kollisionsfrei). Wiederholen bis kreuzungsfrei.
    #
    #     Generisch für jede Segmentzahl und -verteilung; hart nach oben
    #     begrenzt, damit es unter allen Umständen terminiert.
    def _seg_end(i):
        m = math.radians(mids[i])
        return (cx + R_out * math.sin(m), cy - R_out * math.cos(m))

    def _ccw(a, b, c):
        return (c[1] - a[1]) * (b[0] - a[0]) > (b[1] - a[1]) * (c[0] - a[0])

    def _kreuzt(i, j):
        p1, p2 = _seg_end(i), ziel[i]
        p3, p4 = _seg_end(j), ziel[j]
        return (_ccw(p1, p3, p4) != _ccw(p2, p3, p4)
                and _ccw(p1, p2, p3) != _ccw(p1, p2, p4))

    for _entwirren in range(len(ziel) * len(ziel) + 5):
        getauscht = False
        for a_ in range(len(ziel)):
            for b_ in range(a_ + 1, len(ziel)):
                if _kreuzt(a_, b_):
                    ziel[a_], ziel[b_] = ziel[b_], ziel[a_]
                    getauscht = True
                    break
            if getauscht:
                break
        if not getauscht:
            break

    # 7) Offsets schreiben (Nullpunkt = Ring-Band-Mitte des Segments; so
    #    rechnet PowerPoint den manualLayout-Offset bei vorhandenem
    #    manualLayout — empirisch bestätigt).
    import copy
    from lxml import etree
    dlbls = {int(d.find(_q("idx")).get("val")): d
             for d in root.iter(_q("dLbl")) if d.find(_q("idx")) is not None}

    # WICHTIG: Hat das Portfolio MEHR Segmente als die Vorlage Label-Elemente
    # (dLbl), so haben die zusätzlichen Segmente KEIN dLbl → PowerPoint setzt
    # sie auf die Default-Position (ins Loch). Deshalb für jedes fehlende
    # Segment ein dLbl aus einem vorhandenen KLONEN und einfügen.
    if dlbls:
        referenz = dlbls[min(dlbls)]
        container = referenz.getparent()
        for i in range(len(vals)):
            if i not in dlbls:
                neu = copy.deepcopy(referenz)
                neu.find(_q("idx")).set("val", str(i))
                # direkt hinter das Referenz-dLbl einsortieren (bleibt in der
                # dLbl-Gruppe vor den gemeinsamen dLbls-Eigenschaften)
                container.insert(list(container).index(referenz) + 1, neu)
                dlbls[i] = neu

    # ABSOLUTE Positionierung statt Offset (Fix 10.07.2026).
    #
    # Vorher schrieb der Code einen Offset RELATIV zur PowerPoint-Default-
    # Position und nahm an, diese sei die Ring-Band-Mitte des Segments.
    # Das stimmt NICHT, sobald PowerPoint die Labels selbst nach außen holt
    # und entzerrt — bei holeSize 79 % passt kein Text ins Band, und bei eng
    # benachbarten Segmenten (z.B. 341,4° und 358,1°) spreizt PowerPoint sie
    # zusätzlich. Der Offset addierte sich dann auf einen unbekannten Nullpunkt.
    # Am echten Export nachgemessen: bis zu 0,57" Abweichung.
    #
    # Mit xMode/yMode = "edge" sind x/y KEINE Offsets mehr, sondern die
    # absolute Position der linken oberen Ecke des Labels, als Anteil der
    # Chart-Fläche. Damit ist die berechnete Geometrie das, was PowerPoint
    # zeichnet — unabhängig von jeder Automatik.
    #
    # Schema-Reihenfolge in CT_ManualLayout:
    #   layoutTarget, xMode, yMode, wMode, hMode, x, y, w, h
    for i in range(len(mids)):
        tx, ty = ziel[i]
        x_edge = (tx - HALB_W) / frame_w_in     # linke Kante des Labels
        y_edge = (ty - HALB_H) / frame_h_in     # obere  Kante des Labels
        d = dlbls.get(i)
        if d is None:
            continue
        layout = d.find(_q("layout"))
        if layout is None:
            layout = etree.Element(_q("layout")); d.insert(1, layout)
        ml = layout.find(_q("manualLayout"))
        if ml is not None:
            layout.remove(ml)
        ml = etree.SubElement(layout, _q("manualLayout"))
        for tag, val in (("xMode", "edge"), ("yMode", "edge")):
            e = etree.SubElement(ml, _q(tag)); e.set("val", val)
        for tag, val in (("x", x_edge), ("y", y_edge)):
            e = etree.SubElement(ml, _q(tag)); e.set("val", f"{val:.5f}")
    # überzählige Label-Slots (idx >= Segmentzahl) entfernen
    for idx, d in list(dlbls.items()):
        if idx >= len(vals):
            d.getparent().remove(d)
    return {"segmente": len(vals), "R_out": round(R_out, 3)}



def kopf_sperre_aus_usershapes(chart, frame_h_in, luft_in=0.30):
    """Liest die Unterkante des Überschriften-Balkens AUS DEM CHART.

    Der Balken ("AKTUELLE STRUKTUR") ist kein Folien-Shape, sondern ein
    <cdr:relSizeAnchor> in ppt/drawings/drawingN.xml (chartUserShapes) des
    Charts. Gemessen an der CVV-Vorlage: y 0.02"…0.25" bei 3.39" Rahmenhöhe.

    Statt die Sperre fest zu verdrahten, lesen wir sie hier — dann stimmt sie
    auch, wenn jemand den Balken in der Vorlage verschiebt.

    Returns: Mindest-Oberkante für Labels (Zoll ab Chartrahmen-Oberkante)
             oder None, wenn kein Balken gefunden wurde.
    """
    CDR = "{http://schemas.openxmlformats.org/drawingml/2006/chartDrawing}"
    try:
        from lxml import etree
        for rel in chart.part.rels.values():
            if "chartUserShapes" not in rel.reltype:
                continue
            root = etree.fromstring(rel.target_part.blob)
            unten = None
            for anchor in root:
                f = anchor.find(CDR + "from")
                t = anchor.find(CDR + "to")
                if f is None or t is None:
                    continue
                fx = float(f.find(CDR + "x").text)
                fy = float(f.find(CDR + "y").text)
                tx = float(t.find(CDR + "x").text)
                ty = float(t.find(CDR + "y").text)
                # oberes, nahezu rahmenbreites Rechteck = Überschriftenbalken
                if fy < 0.15 and (tx - fx) > 0.8:
                    kante = ty * frame_h_in
                    unten = kante if unten is None else max(unten, kante)
            if unten is not None:
                return unten + luft_in
    except Exception:
        pass
    return None

def ring_holesize(chart, hole=79):
    """Setzt die Ring-Dicke (holeSize = Lochanteil in %). 79 ≈ dünner
    Original-Look aus dem Excel-Makro-PP (Referenz gemessen: ~77 %).
    ACHTUNG: LibreOffice IGNORIERT holeSize beim Rendern (immer ~50 %) —
    nur in echtem PowerPoint sichtbar. Der XML-Wert ist maßgeblich."""
    root = _root(chart)
    hs = root.find(".//" + _q("holeSize"))
    if hs is not None:
        hs.set("val", str(int(hole)))
        return int(hole)
    return None


def _hat_dateax(chart):
    return _root(chart).find(".//" + _q("dateAx")) is not None



# ─────────────────────────────────────────────────────────────────────────
# Segmentfarben, Führungslinien und Label-Farben (NEU 10.07.2026)
#
# PROBLEM: Die <c:dPt>-Farben der Vorlage hängen am INDEX, nicht am Namen.
# Die CVV-Vorlage hatte EDELMETALLE/RENTEN/LIQUIDITÄT (gold/blau/hellblau);
# nach dem Befüllen steht AKTIEN auf idx 0 und erbt Gold. Die Farbe muss
# deshalb datenbasiert je ASSETKLASSE gesetzt werden.
#
# Die Palette stammt aus den Vorlagen selbst (nicht erfunden):
#   14355C dunkelblau · 66A4CE hellblau · 9FD0EF helleres blau
#   BB9256 gold       · 808080 grau
# ─────────────────────────────────────────────────────────────────────────

ASSET_FARBEN = {
    "AKTIEN":       "14355C",   # dunkelblau
    "RENTEN":       "66A4CE",   # hellblau
    "EDELMETALLE":  "BB9256",   # gold
    "LIQUIDITÄT":   "9FD0EF",   # noch helleres blau
    "SONSTIGE":     "808080",   # grau
}

# Schriftfarbe der Prozent-Labels: IMMER Schwarz, unabhängig vom Segment.
# (Anforderung 10.07.2026 — die Zuordnung Label→Segment läuft über Position
# und Führungslinie, NICHT über die Schriftfarbe. Farbe gehört an den Ring.)
# HINWEIS: LABEL_SCHRIFTFARBE, LEADER_FARBE, LEADER_BREITE_EMU,
# _LEADER_RADIAL_STUB, _LEADER_MIN_STUB stehen jetzt im CONFIG-Block ganz oben.


_KERNKLASSEN = {"AKTIEN", "RENTEN", "EDELMETALLE", "LIQUIDITÄT"}
_FILL_TAGS = ("noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill")


def _kategorien(ser):
    """Liest die Kategorienamen (idx-Reihenfolge) einer Serie."""
    cat = ser.find(_q("cat"))
    if cat is None:
        return []
    namen = {}
    for pt in cat.iter(_q("pt")):
        try:
            i = int(pt.get("idx"))
        except (TypeError, ValueError):
            continue
        v = pt.find(_q("v"))
        if v is not None and v.text:
            namen[i] = v.text.strip()
    return [namen.get(i, "") for i in range(max(namen) + 1)] if namen else []


def _ist_assetklassen_ring(kategorien):
    """True, wenn ALLE Kategorien Assetklassen sind und mindestens eine
    Kernklasse dabei ist.

    Damit fassen wir die Sektoren- und Regionen-Ringe der Themen-Broschüren
    NICHT an — deren Kategorien ("Informationstechnologie", "Nordamerika", …)
    stehen nicht in der Palette. Selbstselektierende Regel, kein Schalter.
    """
    if not kategorien:
        return False
    oben = [k.strip().upper() for k in kategorien if k.strip()]
    if not oben:
        return False
    return all(k in ASSET_FARBEN for k in oben) and any(k in _KERNKLASSEN for k in oben)


def _setze_fill(spPr, hexfarbe):
    """Ersetzt die Füllung eines <c:spPr> durch solidFill(hexfarbe).

    Schema-Reihenfolge in CT_ShapeProperties: (xfrm, geometry, FILL, ln, …).
    Die Füllung wird deshalb VOR ein evtl. vorhandenes <a:ln> gesetzt.
    """
    from lxml import etree
    for tag in _FILL_TAGS:
        for el in spPr.findall(_A + tag):
            spPr.remove(el)
    fill = etree.Element(_A + "solidFill")
    clr = etree.SubElement(fill, _A + "srgbClr")
    clr.set("val", hexfarbe)
    ln = spPr.find(_A + "ln")
    if ln is not None:
        ln.addprevious(fill)
    else:
        spPr.insert(0, fill)


def ring_segmentfarben(chart):
    """Setzt die Segmentfarben eines Assetklassen-Rings NAMENSBASIERT.

    Ohne diesen Schritt erbt jedes Segment die Farbe des Vorlagen-<c:dPt> an
    seiner INDEX-Position — dann ist z.B. AKTIEN gold statt dunkelblau.

    Rührt Ringe an, deren Kategorien KEINE Assetklassen sind, nicht an.
    Returns: dict {kategorie: hexfarbe} der gesetzten Segmente (leer = nichts getan).
    """
    from copy import deepcopy
    from lxml import etree
    root = _root(chart)
    ser = root.find(".//" + _q("ser"))
    if ser is None:
        return {}
    kats = _kategorien(ser)
    if not _ist_assetklassen_ring(kats):
        return {}

    vorhandene = {}
    for dpt in ser.findall(_q("dPt")):
        i = dpt.find(_q("idx"))
        if i is not None:
            vorhandene[int(i.get("val"))] = dpt
    if not vorhandene:
        return {}
    muster = vorhandene[min(vorhandene)]

    gesetzt = {}
    for i, name in enumerate(kats):
        schluessel = name.strip().upper()
        farbe = ASSET_FARBEN.get(schluessel)
        if not farbe:
            continue
        dpt = vorhandene.get(i)
        if dpt is None:
            # Vorlage hat weniger dPt als Segmente → aus einem vorhandenen
            # klonen (gleiche Falle wie bei den <c:dLbl>, siehe TW #26).
            dpt = deepcopy(muster)
            dpt.find(_q("idx")).set("val", str(i))
            letztes = ser.findall(_q("dPt"))[-1]
            letztes.addnext(dpt)
            vorhandene[i] = dpt
        spPr = dpt.find(_q("spPr"))
        if spPr is None:
            spPr = etree.SubElement(dpt, _q("spPr"))
        _setze_fill(spPr, farbe)
        gesetzt[name] = farbe

    # Überzählige dPt (Vorlage hatte mehr Segmente als jetzt) entfernen —
    # sonst färbt PowerPoint nicht existierende Punkte.
    for i, dpt in list(vorhandene.items()):
        if i >= len(kats):
            ser.remove(dpt)
    return gesetzt


# ─────────────────────────────────────────────────────────────────────────
# UNGENUTZT (Stand 20.07.2026): ring_leaderlines färbte PowerPoints AUTO-Leader.
# Wir zeichnen Leader jetzt SELBST (ring_leader_zeichnen) — siehe Fallstrick 2.
# Nur als Referenz behalten, NICHT aufrufen.
# ─────────────────────────────────────────────────────────────────────────
def ring_leaderlines(chart, farbe=LEADER_FARBE, breite_emu=LEADER_BREITE_EMU):
    """Färbt die Führungslinien dezent (Default: Grau 0.75pt) statt Schwarz.

    ACHTUNG (OOXML-Grenze): <c:leaderLines> gilt für die GANZE Serie — eine
    Führungslinie pro Segment einzufärben ist im Chart-XML NICHT möglich.
    Die farbliche Zuordnung übernimmt deshalb der Labeltext (ring_label_farben).

    Schema-Reihenfolge in CT_DLbls: … showLeaderLines, dann leaderLines.
    """
    from lxml import etree
    root = _root(chart)
    ser = root.find(".//" + _q("ser"))
    if ser is None:
        return False
    dLbls = ser.find(_q("dLbls"))
    if dLbls is None:
        return False

    sll = dLbls.find(_q("showLeaderLines"))
    if sll is None:
        sll = etree.SubElement(dLbls, _q("showLeaderLines"))
    sll.set("val", "1")

    for el in dLbls.findall(_q("leaderLines")):
        dLbls.remove(el)
    ll = etree.SubElement(dLbls, _q("leaderLines"))
    spPr = etree.SubElement(ll, _q("spPr"))
    ln = etree.SubElement(spPr, _A + "ln")
    ln.set("w", str(int(breite_emu)))
    fill = etree.SubElement(ln, _A + "solidFill")
    clr = etree.SubElement(fill, _A + "srgbClr")
    clr.set("val", farbe)
    return True


def ring_labels_stub_fix(chart, frame_w_in, frame_h_in,
                         rand_in=0.12, kopf_frei_in=0.54,
                         max_nudge_in=0.50):
    """Repariert die FÜHRUNGSLINIEN-RICHTUNG oberer Labels (Wunsch 20.07.2026).

    Problem: Sitzt eine Zahl fast senkrecht ÜBER ihrem Segment (obere kleine
    Segmente), kann ring_leader_zeichnen keinen sauberen horizontalen Stub
    setzen — der Knick zeigte auf die falsche Seite, daher fiel die Linie auf
    eine gerade, richtungslose Diagonale zurück (im PowerPoint sichtbar).

    Lösung: das betroffene Label MINIMAL weiter nach außen schieben (horizontal
    vom Segment weg), bis der Radial-Knick sauber wird. So bleibt der (7)-Look
    erhalten und nur die wenigen Problemlabels bekommen ihre Richtung zurück.

    Wirkt NUR, wenn:
      • das Segment eher oben/unten steht (|cos| > 0,30 — bei seitlichen Labels
        ist eine gerade, fast waagerechte Linie ohnehin richtig), UND
      • der saubere Stub sonst fehlschlägt.
    Sicherungen: nur nach AUSSEN schieben (|mx−cx| wächst), im Rahmen bleiben,
    nicht in ein anderes Label hineinschieben.

    FALLSTRICK: Die x-Position im manualLayout ist die LINKE Boxkante als
    Bruchteil → stored_x = (mx − HALB_W) / frame_w. Nicht die Box-Mitte.
    """
    root = _root(chart)
    pa = root.find(".//" + _q("plotArea") + "/" + _q("layout")
                   + "/" + _q("manualLayout"))
    if pa is None:
        return 0

    def _g(tag):
        e = pa.find(_q(tag))
        return float(e.get("val")) if e is not None else None

    px, py, pw, ph = _g("x"), _g("y"), _g("w"), _g("h")
    if None in (px, py, pw, ph):
        return 0
    left, right = px * frame_w_in, (px + pw) * frame_w_in
    top, bot = py * frame_h_in, (py + ph) * frame_h_in
    cx, cy = (left + right) / 2, (top + bot) / 2
    R = min(right - left, bot - top) / 2

    fsa_el = root.find(".//" + _q("firstSliceAng"))
    fsa = float(fsa_el.get("val")) if fsa_el is not None else 0.0
    vals = [float(v.text)
            for v in root.findall(".//" + _q("val") + "//" + _q("pt")
                                  + "/" + _q("v"))]
    if not vals:
        return 0
    tot = sum(vals) or 1.0
    mids, kum = [], 0.0
    for v in vals:
        f = v / tot
        mids.append((fsa + (kum + f / 2) * 360) % 360)
        kum += f

    ser = root.find(".//" + _q("ser"))
    if ser is None:
        return 0
    HW, HH = 0.33, 0.10

    # Alle Label-Boxen einlesen (idx → [mx, my, dLbl-Element])
    boxen = {}
    for d in ser.findall(".//" + _q("dLbl")):
        ix = d.find(_q("idx"))
        ml = d.find(".//" + _q("manualLayout"))
        if ix is None or ml is None:
            continue
        xe, ye = ml.find(_q("x")), ml.find(_q("y"))
        if xe is None or ye is None:
            continue
        ixv = int(ix.get("val"))
        if ixv >= len(mids):
            continue
        mx = float(xe.get("val")) * frame_w_in + HW
        my = float(ye.get("val")) * frame_h_in + HH
        boxen[ixv] = [mx, my, ml, xe]

    def _kollision(ixv, mx_neu, my):
        for j, (jx, jy, *_ ) in boxen.items():
            if j == ixv:
                continue
            if abs(mx_neu - jx) < 2 * HW and abs(my - jy) < 2 * HH:
                return True
        return False

    verschoben = 0
    for ixv, (mx, my, ml, xe) in boxen.items():
        sm = math.radians(mids[ixv])
        cos_m = math.cos(sm)
        if abs(cos_m) <= 0.30:
            continue                              # seitlich → gerade Linie ok
        r = (cy - my) / cos_m
        if not (R < r < R + 1.5):
            continue
        side = 1.0 if mx >= cx else -1.0
        e_x = mx - side * HW
        k_x = cx + r * math.sin(sm)
        if side * (e_x - k_x) > _LEADER_MIN_STUB:
            continue                              # Knick schon sauber
        # Nudge-Ziel: Label so weit nach außen, dass der Stub sauber wird
        mx_neu = k_x + side * (HW + _LEADER_MIN_STUB + 0.05)
        # nur nach AUSSEN (|mx−cx| darf nicht kleiner werden)
        if abs(mx_neu - cx) <= abs(mx - cx):
            continue
        # Verschiebung deckeln: große Sprünge würden den Look zu stark ändern →
        # dann lieber die gerade Linie behalten.
        if abs(mx_neu - mx) > max_nudge_in:
            continue
        # im Rahmen bleiben
        if mx_neu - HW < rand_in or mx_neu + HW > frame_w_in - rand_in:
            continue
        # nicht in ein anderes Label schieben
        if _kollision(ixv, mx_neu, my):
            continue
        xe.set("val", "%.5f" % ((mx_neu - HW) / frame_w_in))
        boxen[ixv][0] = mx_neu
        verschoben += 1
    return verschoben


def ring_leaderlines_aus(chart):
    """Schaltet PowerPoints automatische Führungslinien AB (showLeaderLines=0).

    Hintergrund (13.07.2026): PowerPoint zeichnet radiale Leader nach einer
    undokumentierten, nicht steuerbaren Regel — mal ja, mal nein (z.B. wenn ein
    Label auf einer Segmentnaht sitzt oder die Linie quer durchs Ring-Loch
    liefe). Das ist im XML nicht verlässlich vorhersagbar. Statt darauf zu
    vertrauen, schalten wir die Auto-Leader ab und zeichnen sie selbst
    (ring_leader_zeichnen) als echte Verbinder-Shapes auf die Folie.
    """
    from lxml import etree
    root = _root(chart)
    ser = root.find(".//" + _q("ser"))
    if ser is None:
        return False
    dLbls = ser.find(_q("dLbls"))
    if dLbls is None:
        return False
    sll = dLbls.find(_q("showLeaderLines"))
    if sll is None:
        sll = etree.SubElement(dLbls, _q("showLeaderLines"))
    sll.set("val", "0")
    for el in dLbls.findall(_q("leaderLines")):
        dLbls.remove(el)
    return True


def ring_leader_zeichnen(slide, shape, chart, farbe=LEADER_FARBE,
                         breite_emu=LEADER_BREITE_EMU,
                         punkt_zeichnen=False,
                         punkt_farbe=PUNKT_FARBE,
                         punkt_durchmesser=PUNKT_DURCHMESSER,
                         start_tiefe=0.0,
                         gerade=False):
    """Zeichnet für jedes Ring-Label eine EIGENE Führungslinie als Connector.

    Läuft vom Außenrand des Segments (R_out am Segment-Mittelwinkel) zur der
    Kante der Label-Box, die dem Segment zugewandt ist. Weil wir die Linie
    selbst setzen, erscheint sie IMMER — unabhängig von PowerPoints Auto-Regel.

    punkt_zeichnen: setzt zusätzlich einen kleinen gefüllten Kreis am Segment-
    Ansatz (wie im Makro-Zielbild). Die ENTSCHEIDUNG, ob Punkte gezeichnet
    werden, trifft der Aufrufer (nachbearbeiten) anhand von Ringtyp + Familie —
    hier wird nur das Flag umgesetzt.
    Punkt (xi, yi) in Rahmen-Zoll liegt auf der Folie bei
    (shape.left + xi*914400, shape.top + yi*914400) EMU, weil die Rahmenbreite
    in Zoll exakt shape.width/914400 entspricht (1:1-Abbildung).
    """
    import math
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.dml.color import RGBColor

    root = _root(chart)
    fw = shape.width / 914400.0
    fh = shape.height / 914400.0
    left, top = shape.left, shape.top

    pa = root.find(".//" + _q("plotArea") + "/" + _q("layout")
                   + "/" + _q("manualLayout"))
    if pa is None:
        return 0

    def _g(t):
        e = pa.find(_q(t))
        return float(e.get("val")) if e is not None else None

    px, py, pw, ph = _g("x"), _g("y"), _g("w"), _g("h")
    if None in (px, py, pw, ph):
        return 0

    l, r = px * fw, (px + pw) * fw
    t_, b = py * fh, (py + ph) * fh
    cx, cy = (l + r) / 2, (t_ + b) / 2
    R_out = min(r - l, b - t_) / 2
    # Leader-Ansatz-Radius: start_tiefe schiebt den Ansatz vom Außenrand (0.0)
    # nach INNEN ins Ringband (0.5 = Mitte der Ringdicke, 1.0 = Innenkante).
    # So kommen die Linien sichtbar AUS dem farbigen Segment heraus.
    _hs_el = root.find(".//" + _q("holeSize"))
    _hs = float(_hs_el.get("val")) / 100.0 if _hs_el is not None else 0.68
    R_in = R_out * _hs
    R_start = R_out - start_tiefe * (R_out - R_in)

    vals = [float(v.text)
            for v in root.findall(".//" + _q("val") + "//" + _q("pt")
                                  + "/" + _q("v"))]
    if not vals:
        return 0
    tot = sum(vals) or 1.0
    mids, kum = [], 0.0
    for v in vals:
        f = v / tot
        mids.append((kum + f / 2) * 360)
        kum += f

    ser = root.find(".//" + _q("ser"))
    HALB_W, HALB_H = 0.33, 0.10

    # Idempotenz: alte eigene Leader dieses Charts entfernen
    praefix = "RingLeader_%s_" % shape.name
    dot_praefix = "RingLeaderDot_%s_" % shape.name
    for sp in list(slide.shapes):
        if sp.name and (sp.name.startswith(praefix)
                        or sp.name.startswith(dot_praefix)):
            sp._element.getparent().remove(sp._element)

    def _ex(xi):
        return int(left + xi * 914400)

    def _ey(yi):
        return int(top + yi * 914400)

    gezeichnet = 0
    for d in ser.findall(".//" + _q("dLbl")):
        ix = d.find(_q("idx"))
        ml = d.find(".//" + _q("manualLayout"))
        if ix is None or ml is None:
            continue
        ixv = int(ix.get("val"))
        if ixv >= len(mids):
            continue
        xe = ml.find(_q("x"))
        ye = ml.find(_q("y"))
        if xe is None or ye is None:
            continue
        mx = float(xe.get("val")) * fw + HALB_W       # Label-Box-Mitte
        my = float(ye.get("val")) * fh + HALB_H
        sm = math.radians(mids[ixv])
        sx = cx + R_start * math.sin(sm)              # Ansatz IM Ringband
        sy = cy - R_start * math.cos(sm)

        # GEKNICKTE Führung — radialer Teil + horizontaler Stub zur Zahl:
        #   Seite = relativ zur RING-MITTE (robust; der frühere Vergleich mit
        #   dem Radialstub kippte bei oberen Segmenten die Knick-Richtung).
        #   e_x = die dem Ring ZUGEWANDTE Seitenkante der Zahl-Box → der Stub
        #   setzt immer auf der richtigen Seite an.
        side = 1.0 if mx >= cx else -1.0
        e_x = mx - side * HALB_W
        e_y = my

        # Knick auf dem RADIALSTRAHL des Segments, auf Label-Höhe: so ist Teil 1
        # (S→Knick) echt radial und Teil 2 (Knick→Zahl) exakt horizontal.
        # Punkt auf dem Strahl mit y = my:  r = (cy - my) / cos(mid)
        # gerade=True: bewusst KEIN Knick → eine ruhige, gerade Linie vom Band
        # zur Zahl (gleichmäßiger, nicht wie harte technische Haken).
        cos_m = math.cos(sm)
        punkte = [(sx, sy), (e_x, e_y)]        # Default: gerade Linie
        if not gerade and abs(cos_m) > 0.30:
            r = (cy - my) / cos_m
            if R_out < r < R_out + 1.5:
                k_x = cx + r * math.sin(sm)
                # Stub muss vom Knick ZUR Zahl zeigen (dem Ring zugewandte
                # Seite). Nur dann knicken; sonst bleibt es die gerade Linie.
                if side * (e_x - k_x) > _LEADER_MIN_STUB:
                    punkte = [(sx, sy), (k_x, my), (e_x, e_y)]

        # Segmente als einzelne gerade Connectoren zeichnen (durchgehende
        # Polylinie, volle Kontrolle über den Knick).
        for teil, ((ax, ay), (bx, by)) in enumerate(
                zip(punkte[:-1], punkte[1:])):
            if abs(ax - bx) < 1e-4 and abs(ay - by) < 1e-4:
                continue
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, _ex(ax), _ey(ay), _ex(bx), _ey(by))
            conn.line.color.rgb = RGBColor(int(farbe[0:2], 16),
                                           int(farbe[2:4], 16),
                                           int(farbe[4:6], 16))
            conn.line.width = Emu(int(breite_emu))
            conn.name = "%s%d_%d" % (praefix, ixv, teil)

        # Punkt am LABEL-ENDE der Führungslinie (Wunsch 20.07.2026) — direkt vor
        # der Prozentzahl, NICHT am Ring. punkte[-1] ist die dem Ring zugewandte
        # Kante der Zahl-Box, also das äußere Ende der Linie. Nur wenn vom
        # Aufrufer gewünscht (z.B. Branchen-Ring der Thema-Familie).
        if punkt_zeichnen:
            pr = RGBColor(int(punkt_farbe[0:2], 16), int(punkt_farbe[2:4], 16),
                          int(punkt_farbe[4:6], 16))
            halb = punkt_durchmesser / 2.0
            ende_x, ende_y = punkte[-1]          # Linien-Ende beim Label
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                _ex(ende_x - halb), _ey(ende_y - halb),
                Emu(int(punkt_durchmesser * 914400)),
                Emu(int(punkt_durchmesser * 914400)))
            dot.fill.solid()
            dot.fill.fore_color.rgb = pr
            dot.line.fill.background()
            dot.shadow.inherit = False
            dot.name = "%s%d" % (dot_praefix, ixv)
        gezeichnet += 1
    return gezeichnet


def ring_label_schriftfarbe(chart, farbe=LABEL_SCHRIFTFARBE, fett=False):
    """Setzt die Schriftfarbe ALLER Prozent-Labels einheitlich (Default schwarz).

    Die Vorlagen setzen keine explizite Textfarbe (die Labels erben Schwarz aus
    dem Theme). Wir setzen sie trotzdem explizit — dann bleibt der Text schwarz,
    egal welches Theme die Vorlage mitbringt.

    fett=True setzt die Zahlen zusätzlich FETT (b="1"); die Farbe bleibt davon
    unberührt (schwarz). Bei fett=False wird das Fett-Attribut NICHT angefasst
    → Familien ohne Fett-Wunsch bleiben unverändert.

    Die Schriftfarbe hängt bewusst NICHT von der Segmentfarbe ab: Die Zuordnung
    Label→Segment erfolgt über die Position und die Führungslinie. Farbe gehört
    an den Ring, nicht an die Zahl.
    """
    from lxml import etree
    root = _root(chart)
    ser = root.find(".//" + _q("ser"))
    if ser is None:
        return 0
    dLbls = ser.find(_q("dLbls"))
    if dLbls is None:
        return 0

    n = 0
    for dLbl in dLbls.findall(_q("dLbl")):
        if dLbl.find(_q("idx")) is None:
            continue
        # CT_DLbl-Reihenfolge: idx, layout, tx, numFmt, spPr, txPr, dLblPos, …
        txPr = dLbl.find(_q("txPr"))
        if txPr is None:
            txPr = etree.Element(_q("txPr"))
            etree.SubElement(txPr, _A + "bodyPr")
            etree.SubElement(txPr, _A + "lstStyle")
            p = etree.SubElement(txPr, _A + "p")
            pPr = etree.SubElement(p, _A + "pPr")
            etree.SubElement(pPr, _A + "defRPr")
            etree.SubElement(p, _A + "endParaRPr")
            spPr = dLbl.find(_q("spPr"))
            (spPr.addnext(txPr) if spPr is not None
             else dLbl.find(_q("idx")).addnext(txPr))
        for rpr in txPr.iter(_A + "defRPr"):
            for tag in _FILL_TAGS:
                for el in rpr.findall(_A + tag):
                    rpr.remove(el)
            fill = etree.Element(_A + "solidFill")
            clr = etree.SubElement(fill, _A + "srgbClr")
            clr.set("val", farbe)
            rpr.insert(0, fill)
            if fett:
                rpr.set("b", "1")   # Prozentzahlen fett (nur wenn gewünscht)
        n += 1
    return n


def _enge_labelwinkel(chart, schwelle_deg=30.0):
    """True, wenn zwei benachbarte Segment-MITTELWINKEL naeher als schwelle_deg
    beieinander liegen.

    Ersetzt die alte Prozent-Schwelle ("Segment < 8 %"). Entscheidend ist nicht,
    wie klein ein Segment ist, sondern wie nah seine Label-Richtung an der des
    Nachbarn liegt. Ein 8,19 %-Segment fiel durch die alte Schwelle — und der
    ganze Entzerrungs-Pass wurde uebersprungen.
    """
    root = _root(chart)
    ser = root.find(".//" + _q("ser"))
    val = ser.find(_q("val")) if ser is not None else None
    if val is None:
        return False
    zahlen = [float(v.text) for v in val.iter(_q("v")) if v.text]
    s = sum(zahlen)
    if s <= 0 or len(zahlen) < 2:
        return False
    mitten, lauf = [], 0.0
    for z in zahlen:
        mitten.append((lauf + z / 2.0) / s * 360.0)
        lauf += z
    mitten.sort()
    for k in range(len(mitten)):
        luecke = (mitten[(k + 1) % len(mitten)] - mitten[k]) % 360.0
        if 0 < luecke < schwelle_deg:
            return True
    return False

def nachbearbeiten(prs, hole_size=79, label_gap_in=0.14,
                   min_gap_deg=24.0, min_gap_deg_klein=60.0,
                   tangential_in=0.14, tangential_klein=0.24,
                   rand_oben_klein=0.52, kopf_frei_klein=0.60,
                   leader_farbe=LEADER_FARBE,
                   label_schriftfarbe=LABEL_SCHRIFTFARBE):
    """EINE Funktion, die alle Charts einer fertigen Präsentation
    datenbasiert nachzieht — am Ende von generate_portfolioanalyse_pptx
    aufrufen, DIREKT VOR prs.save(...).

    Pro Chart-Typ:
      • Doughnut (Ringe): holeSize=hole_size (dünner Original-Look) +
        Außen-Labels radial aus dem Segmentwinkel.
      • Linie MIT Datums-Achse (Wertentwicklung): Achse auf die echte
        Datenspanne (kein Leerraum). Balken (catAx) bleiben unberührt.

    Rührt NICHTS an der Download-Logik an — arbeitet nur an Chart-XML.
    Gibt eine kleine Statistik zurück (für optionales Logging).
    """
    stat = {"ringe": 0, "linien": 0, "ringe_gefaerbt": 0, "labels_schwarz": 0,
            "punkte": 0, "stub_fix": 0}
    # Familie EINMAL bestimmen (für die Punkt-Regel: Punkte nur in Thema).
    ist_thema = _ist_thema_familie(prs)
    # Familienspezifische Ring-Optik EINMAL bestimmen (Wunsch 27.07.: CVV
    # kräftiger). Für Nicht-CVV liefert _ring_format die Defaults → deren
    # Ringe bleiben exakt wie bisher.
    _fam = _familie_aus_prs(prs)
    _fmt = _ring_format(_fam, hole_size, label_gap_in)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            chart = shape.chart
            typ = chart.chart_type.name if chart.chart_type else ""
            try:
                if "DOUGHNUT" in typ:
                    ring_holesize(chart, _fmt["hole"])

                    # NEU 10.07.2026 — Segmentfarben NAMENSBASIERT setzen.
                    # Gibt {} zurück, wenn es kein Assetklassen-Ring ist
                    # (Sektoren/Regionen bleiben damit unangetastet).
                    gesetzt = ring_segmentfarben(chart)
                    if gesetzt:
                        stat["ringe_gefaerbt"] += 1

                    # Rahmenmaße in Zoll (EMU/914400) für die geometrische
                    # Label-Platzierung.
                    _fw = shape.width / 914400.0
                    _fh = shape.height / 914400.0

                    # Adaptives Entzerren: NUR bei Assetklassen-Ringen mit
                    # mehreren kleinen Segmenten (dort drängeln sich die
                    # Labels oben). Sektoren-Ringe behalten min_gap_deg=24.
                    #
                    # Drei Stellschrauben gleichzeitig (10.07.2026):
                    #   min_gap_deg  spreizt die Winkel → DER wirksame Hebel.
                    #                r_use = max(r_use, R_out+0.05) verhindert,
                    #                dass ein Label radial nach unten in den Ring
                    #                rutscht — ein größerer rand_oben_in allein
                    #                bewegt oben stehende Labels daher kaum.
                    #                Nur der Winkel dreht sie vom Balken weg.
                    #   rand_oben_in Untergrenze der Label-MITTE (Oberkante =
                    #                Mitte - 0.10"). Muss in allen Pässen gelten.
                    #   tangential_in längerer, weniger steiler Führungsstrich
                    # (1) KOPF-FREIHALTUNG: BEDINGUNGSLOS fuer jeden Ring.
                    #     Die Sperre wird aus dem Chart GEMESSEN (der Balken ist
                    #     ein chartUserShape). Pass 6d verschiebt nur Labels,
                    #     die sie tatsaechlich verletzen — keine Prozentschwelle,
                    #     kein Sonderfall pro Strategie.
                    _kopf = kopf_sperre_aus_usershapes(chart, _fh)
                    if _kopf is None and gesetzt:
                        _kopf = kopf_frei_klein      # Fallback, falls kein Balken

                    # (2) STAERKERE ENTZERRUNG: strukturell + geometrisch statt
                    #     prozentual. Assetklassen-Ring UND benachbarte
                    #     Mittelwinkel unter 30 Grad.
                    _gap = min_gap_deg
                    _rand_oben = None
                    _tang = tangential_in
                    if gesetzt and _enge_labelwinkel(chart):
                        _gap = min_gap_deg_klein
                        _rand_oben = rand_oben_klein
                        _tang = tangential_klein

                    ring_labels_aussen_dynamisch(chart, _fw, _fh,
                                                 gap_in=_fmt["label_gap_in"],
                                                 min_gap_deg=_gap,
                                                 tangential_in=_tang,
                                                 rand_oben_in=_rand_oben,
                                                 kopf_frei_in=_kopf)

                    # Führungslinien: PowerPoints Auto-Leader ABSCHALTEN und
                    # stattdessen EIGENE Linien als Connector zeichnen — die
                    # erscheinen zuverlässig, unabhängig von PowerPoints
                    # undurchschaubarer Auto-Regel (Naht/Loch-Fälle). Labeltext
                    # bleibt schwarz; Zuordnung über Position + eigene Linie.
                    if leader_farbe:
                        ring_leaderlines_aus(chart)
                        # Führungslinien-Richtung reparieren: obere Labels, die
                        # sonst nur eine richtungslose gerade Linie bekämen (Zahl
                        # fast senkrecht überm Segment), minimal nach außen
                        # schieben → sauberer horizontaler Stub. Muss VOR dem
                        # Zeichnen laufen (verschiebt die Label-Position).
                        stat["stub_fix"] += ring_labels_stub_fix(
                            chart, _fw, _fh,
                            kopf_frei_in=(_kopf if _kopf is not None else 0.54))
                        # Punkt-Regel: auf den erlaubten Ringtypen, wenn Thema
                        # ODER die Familie punkte=True gesetzt hat (CVV, Wunsch
                        # 27.07.) — konfigurierbar im CONFIG-Block. Andere
                        # Familien (ESG/ETF/comdirect) bleiben ohne Punkte.
                        _typ = _ring_typ(chart, shape)
                        _punkt = (PUNKT_AN
                                  and _typ in PUNKT_RINGTYPEN
                                  and ((not PUNKT_NUR_THEMA or ist_thema)
                                       or _fmt["punkte"]))
                        ring_leader_zeichnen(slide, shape, chart,
                                             farbe=leader_farbe,
                                             breite_emu=_fmt["leader_breite_emu"],
                                             punkt_zeichnen=_punkt,
                                             punkt_durchmesser=_fmt["punkt_durchmesser"],
                                             start_tiefe=_fmt["leader_start_tiefe"],
                                             gerade=_fmt["leader_gerade"])
                        if _punkt:
                            stat["punkte"] += 1
                    if label_schriftfarbe:
                        stat["labels_schwarz"] += ring_label_schriftfarbe(
                            chart, farbe=label_schriftfarbe, fett=_fmt["label_fett"])
                    stat["ringe"] += 1
                elif "LINE" in typ and _hat_dateax(chart):
                    datumsachse_an_daten(chart)
                    stat["linien"] += 1
            except Exception:
                # Ein einzelnes problematisches Chart darf den Export nie
                # abbrechen — schlimmstenfalls bleibt dieses Chart wie gehabt.
                pass
    return stat

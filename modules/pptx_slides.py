"""
modules/pptx_slides.py — Slide-Befüllungs-Logik für die FFPB-Broschüre.

Domain-spezifische Funktionen für die vier Folien der Broschüre:
- Slide 7:  Anlagevorschlag (Tabelle + Allokations-Ring)
- Slide 8:  Wertentwicklung/Kurzübersicht (NEU Juli 2026 — alte cVV-Folie:
            Kennzahlen-Tabelle + Perf-p.a.-Balken + Wertentwicklungs-Linie)
- Slide 9:  Performance (Kennzahlen + 2 Charts, mit Benchmark-Vergleich)
- Slide 10: Aktuelle Portfoliozusammenstellung (2 Ring-Charts)

Dieses Modul kennt die FFPB-Vorlage, die Shape-Namen, die Asset-Klassen-
Klassifizierung und die Tabellen-Layouts. Es nutzt aber NUR die generischen
Module pptx_helpers und pptx_charts — kein direkter Streamlit-Zugriff.

Architektur:
    pptx_helpers (Shape/Text/Table/Slide-Manipulation)
    pptx_charts  (Chart-XML mit Bug-Workaround)
        ↑
    pptx_slides  (DIESE Datei — Domain-Logik)
        ↑
    pptx_export  (Orchestrierung der Broschüre)

Juni 2026 — Rückbau auf native Ring-Charts:
    Die zwischenzeitliche matplotlib-PNG-Lösung (modules/png_charts.py,
    replace_donut_chart) wurde wieder entfernt. Die Ringe auf Slide 7 und 9
    werden wieder über native PowerPoint-Donuts befüllt (replace_chart_data),
    d.h. das Template-Styling (Banner, Legende, Quelle, Datenlabels) bleibt
    unverändert und es werden nur die Chart-Daten ausgetauscht.
    => png_charts.py wird von diesem Modul NICHT mehr importiert/benötigt.

Juni 2026 — Kapazitäts-Fix Anlagevorschlag-Tabelle (Slide 7):
    Vorher: Bei mehr als 34 Datenzeilen (Gruppen-Header + Positionen +
    Liquidität) wurden überschüssige Zeilen in fill_table_with_positions
    STILL abgeschnitten (`if i >= max_data_rows: break`) — Positionen
    verschwanden kommentarlos aus dem Compliance-Dokument.
    Jetzt: ensure_table_capacity() klont bei Bedarf zusätzliche <a:tr>-Zeilen
    in die Tabellen-XML, fit_shape_to_table() staucht anschließend ALLE
    Zeilen proportional in den verfügbaren Platz (mit Untergrenze
    MIN_ROW_H_EMU; wird selbst die nicht mehr eingehalten, schrumpft
    zusätzlich die Schrift bis MIN_FONT_PT). Kalibriert an einem echten
    34-Zeilen-Export (Zeilenhöhe 0.1424", Schrift 6pt fix in der Vorlage).
    Realer Extremfall (35 Titel, 39 Zeilen) braucht dabei KEINE
    Schriftverkleinerung — nur minimale Zeilenstauchung (~92%).
    Nur im pathologischen Fall (deutlich >40 Titel, jenseits der Grenze wo
    selbst MIN_FONT_PT nicht mehr reicht) wird eine Warnung zurückgegeben,
    NIE mehr still Daten verworfen.

Juli 2026 — Wertentwicklungs-Folie (alte cVV-Folie) als neue Slide 8:
    Die Vorlage enthält jetzt 26 Slides. Die aus dem alten VBA-Tool
    übernommene Folie "Anlagestrategie {Name} | Wertentwicklung" wurde per
    ZIP-Slide-Copy in die Vorlage integriert (Template-Position 11, direkt
    nach der Performance-Folie). fill_wertentwicklung_slide() befüllt:
    Titel, 4 Kennzahlen (kumulierte WE, Rendite p.a., WE seit 01.01., 
    Duration), Balken-Chart (Perf p.a. vs. Benchmark, volle Kalenderjahre)
    und Linien-Chart (gesamte Historie, Index 1.0-basiert), plus die
    dynamische ***-Benchmark-Fußnote.
"""

import pandas as pd
from typing import Optional
from copy import deepcopy

from pptx.util import Pt, Emu
from pptx.oxml.ns import qn

# Generische PPTX-Helpers (Shape-Lookup, Text, Tabellen)
try:
    from modules.pptx_helpers import (
        find_shape_by_name, replace_text_in_shape,
        set_cell_text, set_cell_text_preserve_format,
        safe_float,
    )
except ImportError:
    from pptx_helpers import (
        find_shape_by_name, replace_text_in_shape,
        set_cell_text, set_cell_text_preserve_format,
        safe_float,
    )

# Chart-Manipulation (XML-basiert, mit Bug-Workaround)
try:
    from modules.pptx_charts import (
        replace_chart_data, replace_chart_data_safe,
        set_value_axis_min,
        set_date_axis_base_unit, set_series_line_width,
    )
except ImportError:
    from pptx_charts import (
        replace_chart_data, replace_chart_data_safe,
        set_value_axis_min,
        set_date_axis_base_unit, set_series_line_width,
    )

# Format-Helpers
try:
    from modules.formats import fmt_pct, fmt_ratio, fmt_date_de, PCT_FORMAT_CODE
except ImportError:
    from formats import fmt_pct, fmt_ratio, fmt_date_de, PCT_FORMAT_CODE


# ═══════════════════════════════════════════════════════════════════════════
# KONSTANTEN
# ═══════════════════════════════════════════════════════════════════════════

# ─── Strategienamen-Bereinigung ─────────────────────────────────────────────
STRATEGY_PREFIXES = ["cVV", "Muster", "Stiftung"]
"""Diese Präfixe werden in clean_strategy_name() vom Strategienamen entfernt."""

STRATEGIEENTWURF_TITLE = "Strategieentwurf im Rahmen einer Vermögensverwaltung"
"""Compliance-Anforderung Juni 2026: Slide 7 trägt diesen festen Titel
(statt 'Anlagevorschlag – Konservativ' o.ä.)."""


# ─── Shape-Namen in der Vorlage ─────────────────────────────────────────────
SHAPE_CHART_ALLOCATION = "C_Kennzahlen"    # Ring-Diagramm (Slides 7, 8)
SHAPE_TABLE = "T_Kennzahlen"               # Positionen-Tabelle (Slides 7, 8)
SHAPE_CHART_LEFT = "C_Kennzahlen1"         # Linkes Ring-Diagramm (Slides 9, 10)
SHAPE_CHART_RIGHT = "C_Kennzahlen2"        # Rechtes Ring-Diagramm (Slide 9)
SHAPE_TITLE = "Titel"
SHAPE_TITLE_ALT = "Titel 2"

# ─── Anlagekriterien-Kasten der Struktur-Folie (NEU 10.08.2026) ─────────────
# Tabelle 5x3 links oben: Spalte 0 Bezeichnung, Spalte 1 LEERE Abstandsspalte,
# Spalte 2 Wert. Zeile 0 ist die Kopfzeile ("Anlagekriterien" | | Strategie).
KRIT_KOPFZELLE = "Anlagekriterien"
KRIT_COL_BEZ = 0
KRIT_COL_WERT = 2
KRIT_ROW_KOPF = 0
"""Die Abstandsspalte 1 wird NIE angefasst — sie trägt die Optik des Kastens."""

# Shape-Namen der Wertentwicklungs-Folie (alte cVV-Folie, NEU Juli 2026).
# "Tabelle"/"Diagramm links"/"Diagramm rechts" heißen auf der Performance-
# Folie genauso — die Lookups sind aber immer per-Slide, daher kein Konflikt.
SHAPE_WE_TABLE = "Tabelle"
SHAPE_WE_CHART_BAR = "Diagramm links"      # Säulen: Perf p.a. im Benchmarkvergleich
SHAPE_WE_CHART_LINE = "Diagramm rechts"    # Linie: Wertentwicklung (Index)
SHAPE_WE_FUSSNOTE = "Fußnote"
SHAPE_WE_QUELLE = "Quelle"
SHAPE_WE_LEGENDE = "Legende Diagramm links"
"""Statische Legenden-Textbox NEBEN dem Säulen-Chart (kein Chart-Element).
Enthält je Eintrag Farbquadrat + Abstand + Beschriftung als eigene Runs —
siehe remove_legend_entry."""

WE_TITLE_FORMAT = "Anlagestrategie {name} | Wertentwicklung"
"""Titel-Muster der Wertentwicklungs-Folie (wie in der alten cVV-Broschüre)."""

WE_SERIES_PORTFOLIO = "Musterdepot"
WE_SERIES_BENCHMARK = "Benchmark"
"""Series-Namen im Balken-Chart der Wertentwicklungs-Folie.

10.08.2026 (Philip): ZURÜCK auf 'Musterdepot' — das ist der Begriff, der in
ALLEN Vorlagen original steht (verifiziert in der Slide-XML: ETF 17/19,
cVV 8/10/12/14/16, ESG 17/19/21/23, comdirect 7/9/11, Thema 12, FFPB 11).
Damit entfällt auch die Umschreibung der statischen Legenden-Textbox: die
Vorlage bleibt unangetastet, der Serienname folgt ihr.

Vorgeschichte: Am 02.07.2026 (Punkt 3) wurde 'Musterdepot' → 'Referenz-
portfolio' vereinheitlicht, weil die performance-Folie diesen Begriff
original führt (FFPB slide10, auch im Folientitel). Das Argument trägt aber
nur für die FFPB-Standardvorlage — die Rolle 'performance' kommt in KEINER
Familien-Konfiguration vor (siehe vorlagen_config.VORLAGEN_FAMILIEN), nur in
pptx_export.DEFAULT_TEMPLATE_CONFIG. In den Familien-Broschüren gab es also
gar keine Folie, zu der die Angleichung hätte passen können. Die
performance-Folie behält 'Referenzportfolio' bewusst (dort in der Vorlage
statisch, inkl. Titel)."""

# ─── Fußnoten-/Disclaimer-Umschreibung der Wertentwicklungs-Folie ───────────
# Juli 2026: Die YTD-Kennzahl folgt jetzt der Tool-Konvention (nach Kosten,
# taggenauer Honorarabzug) statt der alten VBA-Regel ("vor Kosten, ab 30.06.
# abzüglich halbjährigen Honorarsatz"). Die statischen Vorlagen-Texte, die
# noch die alte Regel beschreiben, werden beim Befüllen ersetzt. Die neue
# Formulierung übernimmt bewusst den Wortlaut des bereits freigegebenen
# Tool-Disclaimers ("...in eine äquivalente tägliche Belastung umgerechnet
# und ... taggenau ... abgezogen; eine halbjährliche Berücksichtigung
# erfolgt nicht").
# WICHTIG: Der Disclaimer ist in der Vorlage HART umbrochen (die Sätze
# verteilen sich mit festen Zeilenumbrüchen über mehrere Absätze). Die
# Ersetzung arbeitet daher absatzweise über eindeutige Präfixe; die neuen
# Zeilen sind auf ähnliche Länge kalibriert, damit das Layout hält.
WE_LINE_WIDTH_PT = 1.5
"""Linienstärke des F8-Wertentwicklungs-Charts (03.07.2026): angeglichen an
F9 (Vorlage F8: 0,75pt — für 211 Monatspunkte ausgelegt, bei Tagesdaten
zu unruhig)."""

WE_FOOTNOTE_STAR3_PREFIX = "***"
"""Präfix der Benchmark-Zeile in der Fußnote der Wertentwicklungs-Folie.
Matcht NUR die ***-Zeile: '** ' (mit Leerzeichen) greift dort nicht, weil an
Position 3 ein Stern steht. Bei Strategien OHNE Benchmark wird die Zeile
nicht umgeschrieben, sondern ganz entfernt (11.08.2026)."""

WE_FOOTNOTE_STAR2_PREFIX = "** "
WE_FOOTNOTE_STAR2_NEW = "** nach Kosten (taggenauer Honorarabzug)"

WE_FOOTNOTE_STAR1_PREFIX = "* "
WE_FOOTNOTE_STAR1_NEW = ("* nach Kosten; vollständige Performancehistorie kann "
                         "auf Anfrage eingesehen werden")
"""02.07.2026 (Punkt 7): Kennzahlen 1+2 laufen jetzt bis zum LETZTEN
Datenpunkt (volle Historie, identische Basis wie F9/Tool-UI) statt bis
31.12. des Vorjahres → der Zusatz 'bis zum 31.12. des Vorjahres' in der
*-Fußnote entfällt. Präfix '* ' matcht NUR die *-Zeile ('**…' hat an
Position 2 einen Stern, kein Leerzeichen)."""

WE_DISCLAIMER_REPLACEMENTS = [
    # (Absatz-Präfix in der Vorlage, neuer Absatz-Text)
    ("Der unterjährige Performance Ausweis",
     "Sämtliche Performance Angaben wurden nach Kosten berechnet; der jährliche "
     "Honorarsatz wird in eine äquivalente tägliche Belastung umgerechnet und "),
    ("Kosten berechnet.",
     "taggenau abgezogen (keine halbjährliche Berücksichtigung). Sowohl das VV "
     "Honorar als auch fremde Spesen und evtl. Produktkosten wurden "
     "berücksichtigt. Die Inflation kann negative Auswirkun-"),
]


# ─── Asset-Gruppen (Reihenfolge in Tabelle + Ring) ──────────────────────────
GROUP_AKTIEN = "AKTIEN"
GROUP_RENTEN = "RENTEN"
GROUP_EDELMETALLE = "EDELMETALLE"
GROUP_LIQUIDITAET = "LIQUIDITÄT"
GROUP_SONSTIGE = "SONSTIGE"

GROUP_ORDER = [GROUP_AKTIEN, GROUP_RENTEN, GROUP_EDELMETALLE, GROUP_LIQUIDITAET, GROUP_SONSTIGE]
"""Standard-Reihenfolge der Asset-Gruppen für Tabellen und Allokations-Ring."""


# ─── Tabellen-Spalten-Indizes (Anlagevorschlag-Tabelle, Slide 7) ────────────
# Die Tabelle hat 11 Spalten: 6 Daten-Spalten + 5 Spacer dazwischen
COL_WERTPAPIER = 0
COL_KUPON = 2
COL_FAELLIGKEIT = 4
COL_WKN = 6
COL_ANTEIL = 8
COL_RATING = 10

COL_SPACERS = [1, 3, 5, 7, 9]
"""Spalten-Indizes der Spacer-Spalten (immer leer)."""

# ─── Spalten-Layout als Map (NEU 20.07.2026, für die ETF-Vorlage) ───────────
# Die Standard-/Themen-/CVV-/ESG-Tabellen haben 11 Spalten (obige Konstanten).
# Die ETF-Vorlage hat nur 7 Spalten: WERTPAPIER, WKN, Anteil, Marktrisikowert
# — OHNE KUPON/FÄLLIGKEIT (ETFs sind keine Renten). Damit dieselbe Fill-Logik
# beide Layouts bedienen kann, wird das Layout als Map übergeben; fehlt sie,
# gilt DEFAULT_SPALTEN_MAP → exakt das bisherige 11-Spalten-Verhalten.
# Schlüssel None (kupon/faelligkeit bei ETF) → die Fill-Logik überspringt sie.
DEFAULT_SPALTEN_MAP = {
    "wertpapier": COL_WERTPAPIER,   # 0
    "kupon":       COL_KUPON,        # 2
    "faelligkeit": COL_FAELLIGKEIT,  # 4
    "wkn":         COL_WKN,          # 6
    "anteil":      COL_ANTEIL,       # 8
    "rating":      COL_RATING,       # 10  (Marktrisikowert)
    "spacers":     COL_SPACERS,      # [1, 3, 5, 7, 9]
}
ETF_SPALTEN_MAP = {
    "wertpapier": 0, "kupon": None, "faelligkeit": None,
    "wkn": 2, "anteil": 4, "rating": 6, "spacers": [1, 3, 5],
}


def _belegte_spalten(sm: dict):
    """Nicht-None Datenspalten des Layouts (zum Leeren/zur Leer-Prüfung).
    Für DEFAULT_SPALTEN_MAP genau [0, 2, 4, 6, 8, 10] wie bisher."""
    return [sm[k] for k in ("wertpapier", "kupon", "faelligkeit",
                            "wkn", "anteil", "rating") if sm[k] is not None]


# ─── Kennzahlen-Tabelle der Wertentwicklungs-Folie (7×3) ────────────────────
# Row 0: Header "KENNZAHLEN" | Row 1: Spacer | Rows 2-5: Kennzahlen | Row 6: Spacer
# Spalte 0 = Label, Spalte 1 = Spacer, Spalte 2 = Wert (hellblaue Box)
WE_COL_LABEL = 0
WE_COL_VALUE = 2
WE_ROW_KUMULIERT = 2
WE_ROW_PA = 3
WE_ROW_YTD = 4
WE_ROW_DURATION = 5


# ─── Positionen-Verteilung auf Slides ───────────────────────────────────────
SLIDE_7_DATA_ROWS = 34
"""Slide 7: 36 Zeilen - 1 Header - 1 Summen-Zeile = 34 Daten-Zeilen."""

SLIDE_8_DATA_ROWS = 12
"""Historisch: alte Anlagevorschlag-Teil-2-Folie. Wird nicht mehr genutzt
(Folie wird beim Export entfernt), Konstante bleibt für Import-Kompatibilität."""


# ─── Ring-Chart Konsolidierung (Slide 9: Regionen + Branchen) ───────────────
SMALL_SEGMENT_THRESHOLD = 0.03
"""Kategorien unter 3% werden zu 'Sonstige' zusammengefasst."""

MAX_SEGMENTS_IN_CHART = 7
"""Maximal so viele Kategorien im Ring (alle weiteren → 'Sonstige').
Liquidität wird ggf. NACH dieser Konsolidierung angehängt."""


# ═══════════════════════════════════════════════════════════════════════════
# DOMAIN-HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def finde_anlagekriterien_tabelle(slide):
    """Findet den Anlagekriterien-Kasten INHALTSBASIERT.

    Bewusst nicht über den Shape-Namen: Der Kasten heißt in der Vorlage
    schlicht "Tabelle" — und so heißt auf der Wertentwicklungs-Folie die
    KENNZAHLEN-Tabelle. Wer über den Namen sucht, erwischt irgendwann die
    falsche und überschreibt Kundenzahlen mit Anlageregionen.

    Erkennungsmerkmal ist deshalb die Kopfzelle (0,0) = "Anlagekriterien".
    Siehe Transferwissen #45.

    Returns: die Tabelle oder None.
    """
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        tabelle = shape.table
        if len(tabelle.rows) < 2 or len(tabelle.columns) <= KRIT_COL_WERT:
            continue
        kopf = tabelle.cell(KRIT_ROW_KOPF, 0).text.strip().lower()
        if KRIT_KOPFZELLE.lower() in kopf:
            return tabelle
    return None


def fill_anlagekriterien_slide(prs, slide_idx: int, kriterien,
                               anzeigename: Optional[str] = None) -> bool:
    """Schreibt die Anlagekriterien in den Kasten der Struktur-Folie.

    NEU 10.08.2026: Bis hierher stand der Kasten NUR statisch in der Vorlage —
    je Familie unterschiedlich geschrieben, mit Tippfehlern („FPFB Strategie
    30", „AUsgewogen") und uneinheitlicher Prozent-Schreibweise. Jetzt kommt
    er aus ``Mapping_Anlagekriterien.xlsx``, derselben Datei, die den Banner
    im Tool speist. Eine Änderung dort wirkt auf beides.

    Das ist die bewusste Ausnahme von der Regel „statischer Vorlagentext wird
    in der VORLAGE geändert" (CLAUDE.md): Der Text muss an ZWEI Stellen
    erscheinen, also braucht er EINE Quelle. Anders als beim Musterdepot-Fall
    widerspricht der Code der Vorlage hier nicht heimlich — er ersetzt sie
    als Quelle, nachvollziehbar über eine gepflegte Datei.

    Args:
        kriterien: [(Bezeichnung, Wert), …] in Zeilenreihenfolge (aus
            ``anlagekriterien.fuer``). Leer → die Folie bleibt unberührt.
        anzeigename: Strategiename für die Kopfzeile. None → Vorlage bleibt.

    Returns: True, wenn geschrieben wurde.
    """
    if not kriterien:
        return False
    if slide_idx < 0 or slide_idx >= len(prs.slides):
        return False
    tabelle = finde_anlagekriterien_tabelle(prs.slides[slide_idx])
    if tabelle is None:
        return False

    # Kopfzeile: Strategiename rechts neben "Anlagekriterien"
    if anzeigename:
        set_cell_text_preserve_format(
            tabelle.cell(KRIT_ROW_KOPF, KRIT_COL_WERT), anzeigename)

    # Datenzeilen ab 1. Mehr Kriterien als Zeilen → der Rest wird ignoriert,
    # statt die Tabelle zu sprengen (die Vorlage gibt die Zeilenzahl vor).
    for zeile, (bez, wert) in enumerate(kriterien, start=KRIT_ROW_KOPF + 1):
        if zeile >= len(tabelle.rows):
            break
        set_cell_text_preserve_format(tabelle.cell(zeile, KRIT_COL_BEZ), bez)
        set_cell_text_preserve_format(tabelle.cell(zeile, KRIT_COL_WERT), wert)
    return True


def clean_strategy_name(name: str) -> str:
    """Bereinigt einen Strategienamen für die Anzeige in der Broschüre.

    - Entfernt die in STRATEGY_PREFIXES definierten Präfixe (z.B. 'cVV',
      'Muster', 'Stiftung') sowohl am Anfang als auch am Ende.
    - Ersetzt Underscores durch Leerzeichen (Datenquellen-Konvention: 
      `ETF_Wachstum` → `ETF Wachstum`).
    - Erster Buchstabe wird großgeschrieben.

    Examples:
        >>> clean_strategy_name("cVV Konservativ")
        'Konservativ'
        >>> clean_strategy_name("Stiftung Konservativ")
        'Konservativ'
        >>> clean_strategy_name("Muster Konservativ cVV")
        'Konservativ'
        >>> clean_strategy_name("ETF_Wachstum")
        'ETF Wachstum'
    """
    if not name:
        return ""
    # Underscores zu Leerzeichen — Datenquellen-Konvention
    cleaned = str(name).strip().replace("_", " ")
    # Mehrfach iterieren, falls mehrere Präfixe vorhanden sind
    changed = True
    while changed:
        changed = False
        for prefix in STRATEGY_PREFIXES:
            if cleaned.lower().startswith(prefix.lower() + " "):
                cleaned = cleaned[len(prefix) + 1:].strip()
                changed = True
                break
            if cleaned.lower().endswith(" " + prefix.lower()):
                cleaned = cleaned[:-len(prefix) - 1].strip()
                changed = True
                break
    # Ersten Buchstaben großschreiben
    if cleaned:
        cleaned = cleaned[0].upper() + cleaned[1:]
    return cleaned


def set_title_with_autoscale(title_shape, text: str):
    """Setzt einen Folien-Titel mit automatischer Schriftgrößen-Anpassung.

    Die Titel-Box ist nur ~0.39" hoch (1 Zeile) und ~10.67" breit.
    Bei langem Strategienamen würde der Text in 2 Zeilen umbrechen.

    Strategie (kombiniert):
    1. Manuelle, aggressive Schwellen (empirisch kalibriert in Juni 2026)
    2. Auto-Fit als zusätzliche Sicherheit (PowerPoint skaliert ggf. nach)

    Schwellen (für Standard-Bold-Schrift, 10.67" Box-Breite):
    - ≤ 66 Zeichen → Layout-Default (~32 pt)
    - 67-72 Zeichen → 26 pt
    - 73-80 Zeichen → 22 pt
    - 81-88 Zeichen → 20 pt
    - 89-96 Zeichen → 18 pt
    - 97-108 Zeichen → 16 pt
    - > 108 Zeichen → 14 pt
    """
    replace_text_in_shape(title_shape, text)

    char_count = len(text)
    if char_count <= 66:
        font_size_pt = None  # Layout-Default beibehalten
    elif char_count <= 72:
        font_size_pt = 26
    elif char_count <= 80:
        font_size_pt = 22
    elif char_count <= 88:
        font_size_pt = 20
    elif char_count <= 96:
        font_size_pt = 18
    elif char_count <= 108:
        font_size_pt = 16
    else:
        font_size_pt = 14

    tf = title_shape.text_frame

    if font_size_pt is not None:
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(font_size_pt)

    # Auto-Fit aktivieren als Sicherheits-Netz
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = True
    except Exception:
        pass  # Nicht verfügbar in alten python-pptx-Versionen


def replace_paragraph_text_by_prefix(text_frame, prefix: str, new_text: str) -> bool:
    """Ersetzt den Text GENAU EINES Absatzes in einem Text-Frame — des ersten,
    dessen Text (getrimmt) mit `prefix` beginnt. Alle anderen Absätze bleiben
    unangetastet (NEU Juli 2026, für die dynamische ***-Benchmark-Fußnote der
    Wertentwicklungs-Folie).

    Formaterhaltung: Der neue Text wird in den ERSTEN Run des Absatzes
    geschrieben (dessen Formatierung — Schriftgröße, Farbe — bleibt erhalten),
    alle weiteren Runs des Absatzes werden geleert. Bewusst NICHT
    replace_text_in_shape verwenden — das würde die übrigen Absätze
    (Disclaimer, Fußnoten * und **) mit plattmachen.

    Returns:
        True wenn ein passender Absatz gefunden und ersetzt wurde, sonst False.
    """
    for para in text_frame.paragraphs:
        if para.text.strip().startswith(prefix):
            runs = para.runs
            if not runs:
                continue  # Absatz ohne Runs — nicht beschreibbar, weitersuchen
            runs[0].text = new_text
            for r in runs[1:]:
                r.text = ""
            return True
    return False


def replace_substring_in_runs(text_frame, old: str, new: str) -> bool:
    """Ersetzt einen Teilstring INNERHALB einzelner Runs eines Text-Frames
    (NEU 02.07.2026, für die Legenden-Textbox der Wertentwicklungs-Folie).

    Rührt nur Runs an, die `old` enthalten — alle übrigen Runs (z.B. die
    Wingdings-Farbquadrate der Legende) bleiben samt Formatierung unberührt.
    Voraussetzung: `old` liegt komplett in EINEM Run (in der Vorlage der Fall:
    'Musterdepot ' ist ein eigener Run).

    HINWEIS (10.08.2026): Seit der Rücknahme der Legenden-Umschreibung hat
    diese Funktion KEINEN Aufrufer mehr. Sie bleibt als generischer Helfer
    stehen, weil sie die einzige Stelle im Modul ist, die formaterhaltend in
    einzelnen Runs ersetzt — wer sie nicht braucht, kann sie entfernen.

    Returns:
        True wenn mindestens ein Run ersetzt wurde.
    """
    replaced = False
    for para in text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                replaced = True
    return replaced


def remove_paragraph_by_prefix(text_frame, prefix: str) -> bool:
    """Entfernt GANZ EINEN Absatz — den ersten, dessen Text (getrimmt) mit
    `prefix` beginnt (NEU 11.08.2026, für die ***-Benchmark-Fußnote bei
    Strategien ohne Vergleichsmaßstab).

    Unterschied zu replace_paragraph_text_by_prefix: dort bleibt die Zeile
    stehen und bekommt neuen Text. Hier verschwindet sie samt Absatzmarke —
    gewollt, wenn es zur Fußnote schlicht nichts zu sagen gibt (eine leere
    ***-Zeile hinterließe eine Lücke und einen Verweis ins Nichts).

    Der letzte verbliebene Absatz wird NIE entfernt: ein Text-Frame ohne
    Absatz ist in der Datei ungültig.

    Returns:
        True wenn ein Absatz entfernt wurde, sonst False.
    """
    absaetze = list(text_frame.paragraphs)
    if len(absaetze) <= 1:
        return False
    for para in absaetze:
        if para.text.strip().startswith(prefix):
            para._p.getparent().remove(para._p)
            return True
    return False


def _ist_schmuck_run(run) -> bool:
    """True, wenn ein Run nur Leerraum oder ein Symbol-Zeichen enthält.

    Die Legenden-Boxen der Vorlagen bestehen je Eintrag aus drei Runs:
    Farbquadrat + Abstand + Beschriftung. Das Quadrat ist KEIN Rechteck-Shape,
    sondern ein Zeichen aus dem Private-Use-Bereich (U+F0A2, Wingdings-
    Herkunft) — sichtbar nur an seiner Schriftfarbe. Beides zählt hier als
    Schmuck: es gehört zum Eintrag, trägt aber keine Aussage.
    """
    text = (run.text or "").strip()
    if not text:
        return True
    return all(0xE000 <= ord(zeichen) <= 0xF8FF for zeichen in text)


def remove_legend_entry(text_frame, label: str) -> bool:
    """Entfernt einen Eintrag aus einer statischen Legenden-Textbox — den Run
    mit der Beschriftung `label` samt zugehörigem Farbquadrat und Abstand
    (NEU 11.08.2026, für Strategien ohne Benchmark).

    Vorgehen: den Run finden, der `label` enthält, und von dort RÜCKWÄRTS
    alle Schmuck-Runs (Leerraum, Farbquadrat) mitnehmen, bis echter Text
    kommt — der gehört zum vorherigen Eintrag und bleibt stehen.

    Beispiel Vorlage_Thema, Box "Legende Diagramm links":
        ['■', ' ', 'Musterdepot ', '     ', '■', ' ', 'Benchmark***']
        remove_legend_entry(tf, "Benchmark") →
        ['■', ' ', 'Musterdepot ']

    Bewusst rückwärts statt über feste Indizes: die Vorlagen unterscheiden
    sich in der Zahl der Abstands-Runs, die Reihenfolge Quadrat-Abstand-Text
    ist dagegen überall gleich. Formatierung der bleibenden Runs wird nicht
    angefasst.

    Returns:
        True wenn ein Eintrag entfernt wurde, sonst False.
    """
    for para in text_frame.paragraphs:
        runs = list(para.runs)
        treffer = next((i for i, r in enumerate(runs) if label in (r.text or "")),
                       None)
        if treffer is None:
            continue
        start = treffer
        while start > 0 and _ist_schmuck_run(runs[start - 1]):
            start -= 1
        for run in runs[start:treffer + 1]:
            run._r.getparent().remove(run._r)
        return True
    return False


# XML-Namespace für DrawingML (für set_shape_text_static)
_NSA = "http://schemas.openxmlformats.org/drawingml/2006/main"


def set_shape_text_static(shape, text: str):
    """Setzt den Text eines Shapes STATISCH und entfernt dabei dynamische
    Felder (NEU 02.07.2026, Punkt 6 — 'Quelle … Stand'-Boxen).

    Hintergrund: Die 'Quelle'-Textboxen der Wertentwicklungs- und der
    Performance-Folie enthalten ein PowerPoint-DATUMSFELD (<a:fld
    type="datetime1">) — das rendert immer das ÖFFNUNGS-Datum der Datei,
    nicht den Datenstand. Öffnet der Kunde die Broschüre eine Woche später,
    steht dort ein anderes Datum. Fix: Feld-Elemente aus dem XML entfernen
    und den kompletten Text statisch setzen (Format des ersten Runs bleibt
    erhalten; existiert kein Run, übernimmt der erste Absatz das Format des
    Feld-Laufs nicht — in der Vorlage ist immer ein Text-Run vorhanden).

    No-op wenn das Shape kein Text-Frame hat.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 1) Alle <a:fld>-Elemente entfernen (Datums-/Foliennummern-Felder)
    txBody = tf._txBody
    for fld in txBody.findall(f".//{{{_NSA}}}fld"):
        fld.getparent().remove(fld)
    # 2) Text statisch setzen: erster Run behält Format, Rest wird geleert,
    #    überzählige Absätze entfernt (gleiches Muster wie replace_text_in_shape)
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    para = tf.paragraphs[0]
    runs = list(para.runs)
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        para.text = text


def safe_marktrisikowert(value) -> str:
    """Konvertiert die CSV-Spalte 'Marktrisikowert' zu einem Display-String.

    Float-Werte werden als Integer dargestellt (3.0 → '3'), damit in der
    Tabelle keine Nachkommastellen erscheinen. Fallback '-' bei None/NaN.
    """
    if value is None:
        return "-"
    try:
        if pd.isna(value):
            return "-"
    except (TypeError, ValueError):
        pass
    # Versuch: als ganze Zahl darstellen (3.0 → '3')
    try:
        return str(int(float(value)))
    except (ValueError, TypeError):
        s = str(value).strip()
        return s if s else "-"


def classify_gattung(gattung) -> str:
    """Ordnet eine Gattung einer der 5 Hauptgruppen zu.

    Heuristik:
    - "aktie" / "equity" → AKTIEN
    - "rente" / "anleihe" / "bond" → RENTEN
    - "edelmetall" / "gold" / "silber" → EDELMETALLE
    - "liquid" / "cash" → LIQUIDITÄT
    - sonst → SONSTIGE
    """
    if gattung is None:
        return GROUP_SONSTIGE
    try:
        if pd.isna(gattung):
            return GROUP_SONSTIGE
    except (TypeError, ValueError):
        pass
    g = str(gattung).lower()
    if "aktie" in g or "equity" in g:
        return GROUP_AKTIEN
    if "rente" in g or "anleihe" in g or "bond" in g:
        return GROUP_RENTEN
    if "edelmetall" in g or "gold" in g or "silber" in g:
        return GROUP_EDELMETALLE
    if "liquid" in g or "cash" in g:
        return GROUP_LIQUIDITAET
    return GROUP_SONSTIGE


def group_portfolio_positions(df: pd.DataFrame) -> dict:
    """Gruppiert Portfoliopositionen nach GROUP_ORDER.

    Innerhalb jeder Gruppe sind Positionen alphabetisch nach Wertpapier-Name
    sortiert (seit Juni 2026 — vorher Sortierung nach Gewicht).

    Positionen werden ausgefiltert wenn:
    - Kein Wertpapier-Name vorhanden ist
    - Gewicht = 0 oder NaN ist
    - Wertpapier-Name "nan", "NaT", "None" oder leer ist (Müll aus CSV)

    Wenn die Summe aller Position-Gewichte < 1.0 ist, wird die Differenz
    implizit als Liquidität ergänzt.

    Returns:
        {
            "AKTIEN": [{"wertpapier": ..., "wkn": ..., "gewicht": 0.02, ...}, ...],
            "RENTEN": [...],
            ...
        }
        Leere Gruppen werden weggelassen.
    """
    groups = {g: [] for g in GROUP_ORDER}

    # Junk-Strings die wir als "leer" behandeln
    JUNK_STRINGS = {"", "nan", "NaN", "NaT", "None", "null"}

    for _, row in df.iterrows():
        wertpapier = str(row.get("Wertpapier", "")).strip()
        gewicht = safe_float(row.get("Gewicht", 0.0), 0.0)

        # Müll-Zeilen rausfiltern
        if wertpapier in JUNK_STRINGS:
            continue
        if gewicht <= 0.0001:
            continue

        gruppe = classify_gattung(row.get("Gattung"))

        # WKN auch auf Müll checken
        wkn = str(row.get("WKN", "")).strip()
        if wkn in JUNK_STRINGS:
            wkn = ""

        pos = {
            "wertpapier": wertpapier,
            "wkn": wkn,
            "gewicht": gewicht,
            "kupon": row.get("Kupon"),
            "faelligkeit": row.get("Fälligkeit_parsed") if "Fälligkeit_parsed" in row.index else None,
            "rating": safe_marktrisikowert(row.get("Marktrisikowert")),
            "waehrung": str(row.get("Währung", "")).strip() if "Währung" in row.index else "",
        }
        groups[gruppe].append(pos)

    # Innerhalb jeder Gruppe alphabetisch sortieren
    for g in groups:
        groups[g] = sorted(groups[g], key=lambda p: str(p["wertpapier"]).lower())

    # Liquidität aus Differenz berechnen (falls nicht explizit in Daten)
    if "Gewicht" in df.columns:
        total_weight = safe_float(df["Gewicht"].sum(skipna=True), 0.0)
    else:
        total_weight = 0.0
    liq_from_positions = sum(safe_float(p["gewicht"], 0.0) for p in groups[GROUP_LIQUIDITAET])
    implicit_liq = max(0.0, 1.0 - total_weight)
    if implicit_liq > 0.0001 and liq_from_positions < 0.0001:
        groups[GROUP_LIQUIDITAET].append({
            "wertpapier": "",
            "wkn": "",
            "gewicht": implicit_liq,
            "kupon": None,
            "faelligkeit": None,
            "rating": "",
            "waehrung": "",
        })

    # Leere Gruppen entfernen
    return {g: ps for g, ps in groups.items() if ps}


def distribute_positions_to_slides(groups: dict) -> list:
    """Verteilt gruppierte Positionen auf die Tabellen-Slides.

    Seit Juni 2026 (Performance-Folie als neue Slide 8):
    - Alle Positionen kommen auf Slide 7
    - Bei mehr als SLIDE_7_DATA_ROWS Positionen erweitert
      fill_table_with_positions die Tabelle (ensure_table_capacity)

    Reihenfolge der Zeilen:
    - Asset-Gruppen nach Gewicht absteigend (AKTIEN, RENTEN, EDELMETALLE, ...)
    - LIQUIDITÄT IMMER am Ende als eigene Zeile

    Returns: Liste mit 2 Einträgen:
        [
            {"rows": [...alle Positionen...], "is_last_slide": True},
            {"rows": [], "is_last_slide": False},
        ]
    """
    non_liq = [(n, ps) for n, ps in groups.items() if n != GROUP_LIQUIDITAET]
    non_liq.sort(
        key=lambda kv: sum(safe_float(p["gewicht"], 0.0) for p in kv[1]),
        reverse=True,
    )
    liq_positions = groups.get(GROUP_LIQUIDITAET, [])
    has_liq = bool(liq_positions) and sum(
        safe_float(p["gewicht"], 0.0) for p in liq_positions
    ) > 0.0001

    if not non_liq and not has_liq:
        return [
            {"rows": [], "is_last_slide": True},
            {"rows": [], "is_last_slide": False},
        ]

    all_rows = []
    for group_name, positions in non_liq:
        all_rows.append({"type": "group_header", "data": {"name": group_name}})
        for pos in positions:
            all_rows.append({"type": "position", "data": pos})

    if has_liq:
        total_liq = sum(safe_float(p["gewicht"], 0.0) for p in liq_positions)
        all_rows.append({
            "type": "liquidity",
            "data": {"name": GROUP_LIQUIDITAET, "liq_value": total_liq},
        })

    return [
        {"rows": all_rows, "is_last_slide": True},
        {"rows": [], "is_last_slide": False},
    ]


# ─── Tabellen-Kapazität: Zeilen-Klonen + Stauchung (Juni 2026) ──────────────
# Kalibriert an einem echten Vollkapazitäts-Export (32 Titel, 1 Gruppe,
# genau 34 Datenzeilen = Kapazitätsgrenze der Original-Vorlage).
MAX_TABLE_BOTTOM_INCH = 6.60
SHAPE_PADDING_EMU = 50000  # ~0.05" Puffer für Rahmen

ORIGINAL_DATA_ROW_H_EMU = int(0.1424 * 914400)   # reale Vorlagen-Zeilenhöhe
ORIGINAL_FONT_PT = 6.0                            # reale Vorlagen-Schriftgröße

MIN_ROW_H_INCH = 0.115
"""Untergrenze der Zeilenhöhe, bevor zusätzlich die Schrift schrumpft.
6pt-Text braucht ca. 0.10" Zeilenhöhe (Faktor ~1.2 Line-Height);
0.115" lässt ~0.015" Puffer, bei Sichtprüfung noch komfortabel lesbar."""
MIN_ROW_H_EMU = int(MIN_ROW_H_INCH * 914400)

MIN_FONT_PT = 5.5
"""Absolute Schrift-Untergrenze. Wird selbst hiermit nicht genug Platz frei,
gibt fit_shape_to_table eine Warnung zurück statt weiter zu schrumpfen oder
(schlimmer) Positionen abzuschneiden."""


def remove_empty_table_rows(table, spalten_map=None):
    """Entfernt leere Daten-Zeilen aus der Anlagevorschlag-Tabelle.

    Eine Zeile gilt als 'leer' wenn alle relevanten Daten-Spalten leer sind
    (WERTPAPIER, KUPON, FÄLLIGKEIT, WKN, ANTEIL, RATING).
    Header (Zeile 0) bleibt immer erhalten.

    spalten_map: Layout-Map (Default = 11-Spalten-Layout). ETF (7 Spalten,
    ohne Kupon/Fälligkeit) übergibt sein eigenes Map — dann werden nur die
    tatsächlich belegten Spalten geprüft.

    WICHTIG: Anschließend muss fit_shape_to_table aufgerufen werden, damit
    die Shape-Höhe an die jetzt geringere Zeilenanzahl angepasst wird
    (sonst stretcht LibreOffice die verbleibenden Zeilen).
    """
    sm = spalten_map or DEFAULT_SPALTEN_MAP
    n_rows = len(table.rows)
    if n_rows <= 1:
        return

    indices_to_remove = []
    for i in range(1, n_rows):  # Header (Zeile 0) immer behalten
        row = table.rows[i]
        is_empty = True
        for col_idx in _belegte_spalten(sm):
            text = row.cells[col_idx].text_frame.text.strip()
            # NBSP wird auch als leer betrachtet
            if text and text != "\u00a0":
                is_empty = False
                break
        if is_empty:
            indices_to_remove.append(i)

    if not indices_to_remove:
        return

    # Aus dem XML entfernen — rückwärts, damit Indizes vorderer Zeilen stabil bleiben
    tbl_elem = table._tbl
    tr_elements = tbl_elem.findall(qn('a:tr'))

    for idx in sorted(indices_to_remove, reverse=True):
        tr_to_remove = tr_elements[idx]
        tbl_elem.remove(tr_to_remove)


def _clone_last_data_row(table):
    """Klont die letzte Datenzeile (direkt vor der Summenzeile) und fügt die
    Kopie an derselben Stelle ein. Struktur/Zellformatierung bleiben erhalten
    (Inhalt wird danach von fill_table_with_positions überschrieben — dabei
    IMMER mit explizitem is_bold, damit geklonte Zeilen keine geerbte
    Fett-Formatierung vom Quell-Row behalten).

    Returns:
        Das neu eingefügte <a:tr>-Element.
    """
    tbl_elem = table._tbl
    tr_elements = tbl_elem.findall(qn('a:tr'))
    template_row = tr_elements[-2]  # letzte Datenzeile (vor Summenzeile)
    new_row = deepcopy(template_row)
    summary_row = tr_elements[-1]
    summary_row.addprevious(new_row)
    return new_row


def ensure_table_capacity(table, n_needed_data_rows: int) -> int:
    """Stellt sicher, dass die Tabelle genug Datenzeilen für n_needed_data_rows
    hat. Klont bei Bedarf zusätzliche Zeilen (Struktur/Format der letzten
    Datenzeile) direkt vor die Summenzeile — ERSETZT das frühere stille
    Abschneiden überzähliger Positionen.

    Args:
        table: pptx.table.Table (Header + Datenzeilen + Summenzeile)
        n_needed_data_rows: Anzahl der benötigten Datenzeilen (ohne Header/Summe)

    Returns:
        Anzahl der neu hinzugefügten Zeilen (0 wenn Kapazität schon reichte).
    """
    n_rows_initial = len(table.rows)
    current_capacity = n_rows_initial - 2  # minus Header minus Summe
    n_missing = n_needed_data_rows - current_capacity
    if n_missing <= 0:
        return 0
    for _ in range(n_missing):
        _clone_last_data_row(table)
    return n_missing


def _set_all_font_sizes(table, font_pt: float):
    """Setzt die Schriftgröße aller Runs in allen Zellen der Tabelle."""
    for row in table.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_pt)


# Kalibriert an echter Vorlage (deck_pro.pptx, T_Kennzahlen-Spaltenbreiten)
_ORIGINAL_KUPON_W_EMU = 571500
_ORIGINAL_FAELLIGKEIT_W_EMU = 571500
_MIN_BOND_COL_W_EMU = 91440  # ~0.1" — schmale, aber sichere Nicht-Null-Breite


def maybe_narrow_bond_columns(table_shape, has_renten: bool,
                              spalten_map=None) -> None:
    """Macht die KUPON- und FÄLLIGKEIT-Spalten schmal, wenn kein RENTEN-Anteil
    im Portfolio vorhanden ist, und gibt den freiwerdenden Platz an die
    WERTPAPIER-Spalte weiter (Juni 2026, auf Wunsch: bei reinen Aktien-/
    Nicht-Renten-Strategien bleiben diese Spalten sonst komplett leer und
    wirken wie unnötige Lücken im Tabellenkopf).

    Rührt NUR Spaltenbreiten an (horizontal) — unabhängig von der
    Zeilen-Kapazitätslogik (fit_shape_to_table, vertikal). Reihenfolge der
    Aufrufe spielt daher keine Rolle.

    Args:
        table_shape: Die Tabellen-Shape (mit .table)
        has_renten: True wenn das Portfolio RENTEN-Positionen enthält.
            Bei True: Original-Spaltenbreiten bleiben unverändert.
        spalten_map: Layout-Map. Hat das Layout keine Kupon-/Fälligkeit-Spalten
            (z.B. ETF, 7 Spalten), passiert gar nichts.
    """
    sm = spalten_map or DEFAULT_SPALTEN_MAP
    if sm["kupon"] is None or sm["faelligkeit"] is None:
        return  # Layout ohne Bond-Spalten (z.B. ETF) → nichts zu verschmälern
    if has_renten:
        return  # Original-Layout beibehalten — Kupon/Fälligkeit werden gebraucht

    table = table_shape.table
    freed = (_ORIGINAL_KUPON_W_EMU - _MIN_BOND_COL_W_EMU) + \
            (_ORIGINAL_FAELLIGKEIT_W_EMU - _MIN_BOND_COL_W_EMU)

    table.columns[sm["kupon"]].width = Emu(_MIN_BOND_COL_W_EMU)
    table.columns[sm["faelligkeit"]].width = Emu(_MIN_BOND_COL_W_EMU)
    table.columns[sm["wertpapier"]].width = Emu(
        table.columns[sm["wertpapier"]].width + freed
    )


def fit_shape_to_table(table_shape, max_row_scale: float = 3.0,
                       max_bottom_inch: Optional[float] = None,
                       original_row_h_inch: Optional[float] = None) -> Optional[str]:
    """Passt die Höhe der Tabellen-Shape UND die Zeilenhöhen an die tatsächlich
    genutzte Zeilenanzahl an. Symmetrisch:

    - WENIGE Zeilen (< 70% Auslastung): Zeilenhöhen proportional hochskalieren,
      max. `max_row_scale`× Original (wie bisher — unverändertes Verhalten).
    - VIELE Zeilen (Summe > verfügbarer Platz): NEU — Zeilenhöhen proportional
      runterskalieren, bis MIN_ROW_H_EMU. Reicht das nicht, wird zusätzlich
      die Schrift bis MIN_FONT_PT verkleinert. Reicht selbst das nicht
      (pathologischer Fall, weit jenseits normaler Portfoliogrößen), wird
      NICHT mehr still abgeschnitten — stattdessen ein Warnhinweis
      zurückgegeben, den der Aufrufer (z.B. Streamlit-UI) anzeigen kann.

    Args:
        table_shape: Die Tabellen-Shape (mit .table-Property)
        max_row_scale: Maximaler Hochskalierungsfaktor pro Zeile (Default 3.0)

    Returns:
        Warnhinweis (str) falls selbst bei Minimalgröße nicht alles passt,
        sonst None.
    """
    table = table_shape.table

    # max_bottom_inch=None → modulweiter Default (Standard-Vorlage, 6.60").
    # Andere Vorlagen der Familie (z.B. CVV: Abschlusslinie bei 6.38") reichen
    # ihre eigene Unterkante durch.
    MAX_TABLE_BOTTOM_INCH_LOCAL = (MAX_TABLE_BOTTOM_INCH if max_bottom_inch is None
                                   else float(max_bottom_inch))
    # Vorlagen-Zeilenhöhe: ORIGINAL_DATA_ROW_H_EMU ist auf die STANDARD-Vorlage
    # kalibriert (0.1424"). Andere Vorlagen der Familie haben andere Zeilen
    # (CVV: 0.192") — sonst greift die Stauchungs-Untergrenze zu früh und die
    # Tabelle läuft über die Abschlusslinie.
    ORIG_ROW_H_EMU = (ORIGINAL_DATA_ROW_H_EMU if original_row_h_inch is None
                      else int(float(original_row_h_inch) * 914400))
    shape_top_inch = table_shape.top / 914400
    max_available_h_emu = int((MAX_TABLE_BOTTOM_INCH_LOCAL - shape_top_inch) * 914400)

    total_row_h = sum(row.height for row in table.rows)
    warning = None

    if total_row_h < max_available_h_emu * 0.7 and total_row_h > 0:
        # ── Wenige Zeilen: hochskalieren (unverändertes Verhalten) ──
        target_h = min(
            max_available_h_emu - SHAPE_PADDING_EMU,
            int(total_row_h * max_row_scale)
        )
        scale = target_h / total_row_h
        for row in table.rows:
            row.height = int(row.height * scale)
        total_row_h = sum(row.height for row in table.rows)

    elif total_row_h > max_available_h_emu - SHAPE_PADDING_EMU:
        # ── NEU: viele Zeilen -> proportional stauchen, mit Untergrenze ──
        target_h = max_available_h_emu - SHAPE_PADDING_EMU
        scale = target_h / total_row_h
        implied_data_row_h = ORIG_ROW_H_EMU * scale

        if implied_data_row_h < MIN_ROW_H_EMU:
            # An Zeilenhöhen-Untergrenze klemmen statt weiter zu stauchen
            floor_scale = MIN_ROW_H_EMU / ORIG_ROW_H_EMU
            for row in table.rows:
                row.height = int(row.height * floor_scale)
            total_row_h = sum(row.height for row in table.rows)

            if total_row_h > max_available_h_emu - SHAPE_PADDING_EMU:
                # Selbst an der Zeilenhöhen-Untergrenze reicht der Platz nicht:
                # zusätzlich Schrift verkleinern (bis MIN_FONT_PT).
                overflow_ratio = total_row_h / (max_available_h_emu - SHAPE_PADDING_EMU)
                font_pt = max(MIN_FONT_PT, ORIGINAL_FONT_PT / overflow_ratio)
                _set_all_font_sizes(table, font_pt)
                if font_pt <= MIN_FONT_PT and overflow_ratio > 1.02:
                    n_data_rows = len(table.rows) - 2
                    warning = (
                        f"Anlagevorschlag-Tabelle hat {n_data_rows} Zeilen — selbst "
                        f"bei minimaler Zeilenhöhe/Schrift ({MIN_FONT_PT}pt) reicht "
                        f"der Platz nicht ganz. Geringer optischer Überlauf über den "
                        f"Footer möglich. Bitte Folie 7 manuell prüfen."
                    )
        else:
            for row in table.rows:
                row.height = int(row.height * scale)
            total_row_h = sum(row.height for row in table.rows)

    # Shape-Höhe auf die (ggf. skalierte) Summe der Zeilenhöhen setzen
    table_shape.height = total_row_h + SHAPE_PADDING_EMU
    return warning


def adjust_table_shape_height(prs, table_shape, n_data_rows: int, needs_summary: bool):
    """Passt die Höhe der Tabellen-Shape an die tatsächlich benötigte Zeilenanzahl an.

    Kann die Shape auch vergrößern (nach unten), aber nur bis max. 6.60" Bottom
    (vor Footer bei 6.76").

    Args:
        prs: Presentation
        table_shape: Die Tabellen-Shape
        n_data_rows: Anzahl Daten-Zeilen die wir befüllen (inkl. Gruppen-Header)
        needs_summary: True wenn Summen-Zeile benötigt wird
    """
    ORIGINAL_HEADER_H = 0.236
    ORIGINAL_DATA_ROW_H = 0.142
    ORIGINAL_SUMMARY_H = 0.142
    MAX_TABLE_BOTTOM = 6.60  # inches

    n_buffer_rows = 2 if needs_summary else 0

    needed_h = ORIGINAL_HEADER_H + (n_data_rows * ORIGINAL_DATA_ROW_H) + (n_buffer_rows * ORIGINAL_DATA_ROW_H)
    if needs_summary:
        needed_h += ORIGINAL_SUMMARY_H

    shape_top_inch = table_shape.top / 914400
    shape_current_h_inch = table_shape.height / 914400
    max_available_h = MAX_TABLE_BOTTOM - shape_top_inch

    new_h_inch = min(needed_h, max_available_h)

    # Nur ändern wenn Änderung signifikant (>0.05" Differenz)
    if abs(new_h_inch - shape_current_h_inch) > 0.05:
        table_shape.height = int(new_h_inch * 914400)


def consolidate_small_segments(agg_series: pd.Series,
                                threshold: float = SMALL_SEGMENT_THRESHOLD,
                                max_segments: int = MAX_SEGMENTS_IN_CHART) -> pd.Series:
    """Fasst kleine Kategorien zu 'Sonstige' zusammen.

    Regel:
    - Alle Kategorien unter threshold werden zu 'Sonstige' gruppiert
    - Wenn nach Konsolidierung noch mehr als max_segments Kategorien da sind,
      werden die kleinsten zusätzlich in Sonstige verschoben bis max_segments
      erreicht ist

    Args:
        agg_series: Pandas Series (Index = Kategorie-Name, Werte = Gewicht)
        threshold: Schwellwert für 'kleine' Kategorie
        max_segments: Maximale Anzahl Segmente im Chart

    Returns:
        Konsolidierte Series, absteigend sortiert.
    """
    agg = agg_series.sort_values(ascending=False)

    big = agg[agg >= threshold]
    small = agg[agg < threshold]

    # Maximale Anzahl Segmente beachten
    if len(big) > max_segments - 1:  # -1 weil Platz für 'Sonstige' nötig
        keep = big.head(max_segments - 1)
        move_to_small = big.tail(len(big) - (max_segments - 1))
        big = keep
        small = pd.concat([small, move_to_small])

    # Sonstige zusammenfassen
    if len(small) > 0:
        sonstige_sum = small.sum()
        if sonstige_sum > 0.0001:
            existing = float(big["Sonstige"]) if "Sonstige" in big.index else 0.0
            big["Sonstige"] = existing + sonstige_sum
            big = big.sort_values(ascending=False)

    return big


def build_ring_series(df: pd.DataFrame, dim_col: str) -> pd.Series:
    """Baut die Werte-Serie für einen Ring auf Slide 9 (Regionen oder Branchen).

    Regionen- und Branchen-Ringe zeigen ausschliesslich die INVESTIERTEN
    Anteile (Aktien/Renten/…), skaliert auf 100%. Liquidität hat weder eine
    Region noch eine Branche und wird deshalb HERAUSGERECHNET — analog zur
    VBA-Makrolösung (Stand 10.07.2026, Option B). Der frühere Ansatz, die
    Liquidität als eigenes Segment anzuhängen, führte je nach Daten dazu, dass
    der Ring auf ~98% statt 100% summierte (Liquiditätszeile trug teils eine
    Region → wurde als klassifiziert gewertet, fiel aber unter den 3%-Threshold
    und verschwand in 'Sonstige' oder ganz).

    Ablauf:
    1. Liquidität per Gattung erkennen und komplett entfernen
       (nicht per leerer Region raten — das war die alte, fragile Annahme).
    2. Nach `dim_col` aggregieren; Positionen ohne Eintrag ignorieren.
    3. Auf 100% normieren (Summe der verbleibenden, investierten Gewichte).
    4. Kleine Kategorien zu 'Sonstige' konsolidieren.

    Die Normierung passiert VOR der Konsolidierung, damit die 'Sonstige'-
    Schwelle sich auf die bereits skalierten Anteile bezieht.
    """
    if dim_col not in df.columns or "Gewicht" not in df.columns:
        return pd.Series(dtype=float)

    work = df.copy()

    # ── 1) Liquidität herausrechnen (Option B) ────────────────────────────
    # Primär über die Gattung (robust, unabhängig von der Region-Spalte).
    # Fallback: fehlt die Gattung-Spalte, greifen wir auf leere dim_col-Werte
    # zurück (altes Verhalten der Nicht-Klassifizierung).
    if "Gattung" in work.columns:
        ist_liquide = work["Gattung"].apply(
            lambda g: classify_gattung(g) == GROUP_LIQUIDITAET)
        work = work[~ist_liquide]

    # ── 2) Nach Dimension aggregieren ─────────────────────────────────────
    col = work[dim_col].astype(str).replace(["nan", "NaT", "None"], "")
    has_value = col.str.strip() != ""
    classified = work[has_value]
    if classified.empty:
        return pd.Series(dtype=float)

    agg = classified.groupby(col[has_value])["Gewicht"].sum()
    agg = agg[agg > 0.0001]
    if agg.empty:
        return pd.Series(dtype=float)

    # ── 3) Auf 100% normieren (nur investierte Anteile) ───────────────────
    gesamt = float(agg.sum())
    if gesamt > 0.0001:
        agg = agg / gesamt

    # ── 4) Kleine Kategorien konsolidieren ────────────────────────────────
    agg = consolidate_small_segments(agg)

    return agg


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE-BEFÜLLUNG (Hauptfunktionen pro Folie)
# ═══════════════════════════════════════════════════════════════════════════

def fill_table_with_positions(table, slide_data: dict, total_weight: float = 1.0,
                              shape_height: int = 0, spalten_map=None):
    """Befüllt die Anlagevorschlag-Tabelle (Slide 7) mit Positionen.

    Die Tabellen-Struktur der Vorlage bleibt UNVERÄNDERT (keine Zeilen entfernt,
    keine Höhen geändert). Nicht benötigte Zeilen bleiben leer sichtbar.

    Args:
        table: Die Tabelle (shape.table)
        slide_data: {"rows": [...], "is_last_slide": bool}
        total_weight: Summe aller Gewichte (für Summen-Zeile, default 100%)
        shape_height: Höhe der Tabellen-Shape in EMU (aus Kompat-Gründen in der
                      Signatur belassen, wird nicht mehr verwendet)
    """
    n_rows_initial = len(table.rows)
    rows = slide_data["rows"]
    is_last = slide_data["is_last_slide"]

    # NEU (Juni 2026): Tabelle bei Bedarf erweitern statt still abzuschneiden.
    # Vorher wurde hier über max_data_rows = n_rows_initial - 2 hart begrenzt
    # und überzählige Positionen mit `if i >= max_data_rows: break` verworfen.
    ensure_table_capacity(table, len(rows))
    n_rows_initial = len(table.rows)  # ggf. jetzt größer

    # Summen-Zeile ist immer die letzte Zeile in der Vorlage
    summary_row_idx = n_rows_initial - 1
    max_data_rows = n_rows_initial - 2

    # Erst alle Datenzeilen leeren (nur belegte Datenspalten - Spacer bleiben)
    sm = spalten_map or DEFAULT_SPALTEN_MAP
    for row_idx in range(1, n_rows_initial):
        row = table.rows[row_idx]
        for col_idx in _belegte_spalten(sm):
            set_cell_text(row.cells[col_idx], "")

    # Zeilen befüllen
    for i, row_def in enumerate(rows):
        if i >= max_data_rows:
            break  # Kein Platz mehr

        target_row_idx = i + 1  # +1 weil Zeile 0 der Tabellen-Header ist
        row = table.rows[target_row_idx]

        if row_def["type"] in ("group_header", "liquidity"):
            # Gruppen-Header: Name in Spalte 0, alle anderen leer, fett
            name = row_def["data"]["name"]
            set_cell_text(row.cells[sm["wertpapier"]], name, is_bold=True)
            # Bei RENTEN: "KUPON" und "FÄLLIGKEIT" als Sub-Header (nur wenn das
            # Layout diese Spalten hat — ETF hat sie nicht)
            if name == GROUP_RENTEN and sm["kupon"] is not None:
                set_cell_text(row.cells[sm["kupon"]], "KUPON", is_bold=True)
                set_cell_text(row.cells[sm["faelligkeit"]], "FÄLLIGKEIT", is_bold=True)
            # Bei LIQUIDITÄT: Wert direkt in der Header-Zeile
            if name == GROUP_LIQUIDITAET and "liq_value" in row_def["data"]:
                set_cell_text(row.cells[sm["anteil"]], fmt_pct(row_def["data"]["liq_value"]), is_bold=True)

        elif row_def["type"] == "position":
            data = row_def["data"]
            # Alle Felder einer Position: explizit NICHT BOLD
            set_cell_text(row.cells[sm["wertpapier"]], data["wertpapier"], is_bold=False)
            set_cell_text(row.cells[sm["wkn"]], data["wkn"], is_bold=False)
            set_cell_text(row.cells[sm["anteil"]], fmt_pct(data["gewicht"]), is_bold=False)
            set_cell_text(row.cells[sm["rating"]], data.get("rating", "-"), is_bold=False)
            # Kupon/Fälligkeit nur, wenn das Layout diese Spalten hat (ETF: nein)
            if sm["kupon"] is not None:
                if data.get("kupon") is not None and not pd.isna(data["kupon"]) and data["kupon"] != 0:
                    set_cell_text(row.cells[sm["kupon"]], fmt_pct(data["kupon"]), is_bold=False)
                else:
                    set_cell_text(row.cells[sm["kupon"]], "", is_bold=False)
            if sm["faelligkeit"] is not None:
                if data.get("faelligkeit") is not None and not pd.isna(data["faelligkeit"]):
                    set_cell_text(row.cells[sm["faelligkeit"]], fmt_date_de(data["faelligkeit"]), is_bold=False)
                else:
                    set_cell_text(row.cells[sm["faelligkeit"]], "", is_bold=False)

    # Summen-Zeile: nur auf letzter Slide
    summary_row = table.rows[summary_row_idx]
    for col_idx in _belegte_spalten(sm):
        set_cell_text(summary_row.cells[col_idx], "")

    if is_last:
        set_cell_text(summary_row.cells[sm["anteil"]], fmt_pct(total_weight))

    # WICHTIG: Tabellen-Struktur der Vorlage bleibt UNVERÄNDERT.
    # Frühere Versuche, leere Zeilen zu entfernen, haben LibreOffice
    # zum Vergrößern der Zeilen veranlasst → Überlauf am Slide-Rand.


def fill_anlagevorschlag_slides(prs, slide_7_idx: int,
                                 df: pd.DataFrame, strategy_name: str,
                                 eval_date=None,
                                 titel_text: Optional[str] = None,
                                 max_bottom_inch: Optional[float] = None,
                                 original_row_h_inch: Optional[float] = None,
                                 spalten_map: Optional[dict] = None) -> Optional[str]:
    """Befüllt Slide 7 (Anlagevorschlag/Strategieentwurf) mit Portfolio-Daten.

    Seit Juni 2026 (Performance-Folie als Slide 8): Es gibt nur noch EINE
    Anlagevorschlag-Slide. Alle Positionen kommen auf Slide 7. Die Tabelle
    wird bei Bedarf um zusätzliche Zeilen erweitert (ensure_table_capacity)
    und anschließend proportional in den verfügbaren Platz eingepasst
    (fit_shape_to_table) — es werden NIE mehr Positionen still abgeschnitten.

    Args:
        prs: Presentation
        slide_7_idx: 0-indexed Index der Anlagevorschlag-Slide
        df: DataFrame mit Positionen (Wertpapier, WKN, Gewicht, Gattung, Kupon,
            Fälligkeit_parsed, Marktrisikowert)
        strategy_name: Name der Strategie für den Titel (schon bereinigt)
        eval_date: Auswertungsdatum (für Source-Annotation im Ring-Chart).
            Optional — das Quelle-Datum wird zentral über
            pptx_helpers.update_quelle_datum gesetzt (steht statisch im Template).
        titel_text: NEU (09.07.2026) — expliziter Folientitel. None (Default)
            = bisheriges Verhalten ("Strategieentwurf … - <Name>").
            "" = Titel gar nicht anfassen (CVV: die Vorlage trägt bereits
            "Anlagestrategie Konservativ" o.ä. und soll ihn behalten).
        max_bottom_inch: NEU (09.07.2026) — Unterkante der Tabelle in Zoll.
            None (Default) = modulweites MAX_TABLE_BOTTOM_INCH (6.60,
            Standard-Vorlage). Die CVV-Vorlage hat ihre Abschlusslinie schon
            bei 6.38" und braucht deshalb einen kleineren Wert.

    Returns:
        Warnhinweis (str) falls die Tabelle selbst bei Minimalgröße nicht
        vollständig passt (sehr seltener Extremfall, siehe fit_shape_to_table),
        sonst None. Aufrufer (z.B. Streamlit-UI) sollte das dem Nutzer anzeigen.
    """
    # 1. Daten vorbereiten
    groups = group_portfolio_positions(df)
    slide_distribution = distribute_positions_to_slides(groups)

    # 2. Allokations-Daten für Ring-Chart (nach Gruppen)
    alloc_labels = []
    alloc_values = []
    for g in GROUP_ORDER:
        if g in groups:
            total = sum(safe_float(p["gewicht"], 0.0) for p in groups[g])
            if total > 0.0001:
                alloc_labels.append(g)
                alloc_values.append(float(total))

    # Gesamt-Gewicht (für Summen-Zeile)
    total_weight = sum(alloc_values)

    # 3. Slide 7 befüllen
    slide_7 = prs.slides[slide_7_idx]
    # Titel: Strategieentwurf-Hinweis (Email-Anforderung Juni 2026, Compliance).
    # titel_text="" → Titel der Vorlage unangetastet lassen (CVV).
    _titel = (f"{STRATEGIEENTWURF_TITLE} - {strategy_name}"
              if titel_text is None else titel_text)
    if _titel:
        title = find_shape_by_name(slide_7, SHAPE_TITLE_ALT) or find_shape_by_name(slide_7, SHAPE_TITLE)
        if title:
            set_title_with_autoscale(title, _titel)

    # Ring-Chart: NATIVER PowerPoint-Donut (Daten ersetzen, Template-Styling
    # bleibt: Banner, Legende, Quelle, Datenlabels). Werte sind Anteile (0..1);
    # das Zahlenformat des Charts stellt sie als Prozent dar.
    chart = find_shape_by_name(slide_7, SHAPE_CHART_ALLOCATION)
    if chart and getattr(chart, "has_chart", False):
        if sum(alloc_values) > 0:
            replace_chart_data(
                chart,
                categories=list(alloc_labels),
                values=list(alloc_values),
                series_name="Anteil",
            )

    # Tabelle befüllen
    table_shape = find_shape_by_name(slide_7, SHAPE_TABLE)
    capacity_warning = None
    if table_shape:
        # NEU (Juni 2026): Kupon/Fälligkeit-Spalten schmal machen, wenn keine
        # RENTEN-Positionen vorhanden sind — sonst bleiben sie bei reinen
        # Aktien-Strategien komplett leer und wirken wie unnötige Lücken.
        maybe_narrow_bond_columns(table_shape, has_renten=(GROUP_RENTEN in groups),
                                  spalten_map=spalten_map)

        fill_table_with_positions(table_shape.table, slide_distribution[0], total_weight,
                                  shape_height=table_shape.height, spalten_map=spalten_map)
        # Leere Zeilen entfernen (nur relevant falls Kapazität > benötigt)
        remove_empty_table_rows(table_shape.table, spalten_map=spalten_map)
        # Trennstriche an die TATSÄCHLICHEN Kategoriegrenzen setzen (NEU
        # 07.08.2026). Muss NACH dem Entfernen leerer Zeilen laufen, sonst
        # stimmen die Zeilenindizes nicht mehr. Reine Kosmetik — ein Fehler
        # hier darf den Export nicht abbrechen.
        try:
            tabelle_kategorie_trennlinien(table_shape.table)
        except Exception as exc:
            EINZELTITEL_WARNUNGEN.append(
                f"Trennstriche der Positionstabelle konnten nicht gesetzt "
                f"werden: {type(exc).__name__}: {exc}")
        # Shape-Höhe an tatsächliche Zeilenanzahl anpassen (staucht/streckt je
        # nach Bedarf; ensure_table_capacity in fill_table_with_positions hat
        # vorher bereits ggf. zusätzliche Zeilen geklont)
        capacity_warning = fit_shape_to_table(
            table_shape, max_bottom_inch=max_bottom_inch,
            original_row_h_inch=original_row_h_inch)

    return capacity_warning


def fill_kennzahlen_table(table, kz: dict):
    """Befüllt die KENNZAHLEN-Tabelle auf der Performance-Folie.

    Tabellen-Layout (7 rows × 5 cols, mit Spacer-Spalten):
      Row 0: Header   (KENNZAHLEN | _ | REFERENZ | _ | BENCHMARK)
      Row 1: leer/Spacer
      Row 2: Performance p.a.
      Row 3: Volatilität
      Row 4: Sharpe Ratio
      Row 5: Max Drawdown
      Row 6: leer/Spacer

    Wert-Spalten: 2 (REFERENZ), 4 (BENCHMARK)
    """
    metric_rows = [
        ("performance_pa_ref",  "performance_pa_bench",   2, True),   # row 2, Prozent
        ("volatilitaet_ref",    "volatilitaet_bench",     3, True),   # row 3, Prozent
        ("sharpe_ref",          "sharpe_bench",           4, False),  # row 4, Dezimal
        ("max_drawdown_ref",    "max_drawdown_bench",     5, True),   # row 5, Prozent
    ]
    for ref_key, bench_key, row_idx, is_pct in metric_rows:
        if row_idx >= len(table.rows):
            continue
        row = table.rows[row_idx]
        ref_val = kz.get(ref_key)
        bench_val = kz.get(bench_key)
        if is_pct:
            ref_str = fmt_pct(ref_val)
            bench_str = fmt_pct(bench_val)
        else:
            ref_str = fmt_ratio(ref_val)
            bench_str = fmt_ratio(bench_val)
        # Spalte 2 = REFERENZ, Spalte 4 = BENCHMARK
        set_cell_text_preserve_format(row.cells[2], ref_str)
        set_cell_text_preserve_format(row.cells[4], bench_str)


# ─────────────────────────────────────────────────────────────────────────
# Generische Tabellen-Layout-Helfer (NEU 09.07.2026)
#
# Bewusst GENERISCH gehalten (Daten-Spalten werden als Parameter übergeben),
# damit sie für JEDE Vorlage der PowerPoint-Familie wiederverwendbar sind —
# unabhängig von der Spaltenanzahl. Die alten Helfer (remove_empty_table_rows,
# fit_shape_to_table) sind fest auf das 11-spaltige Anlagevorschlag-Layout
# verdrahtet und würden an der 7-spaltigen Einzeltitel-Tabelle mit
# IndexError sterben — deshalb diese Parallel-Implementierung.
#
# WICHTIG: Diese Funktionen werden AUSSCHLIESSLICH von
# fill_einzeltitel_themen_slide aufgerufen. Die Standard-Anlagevorschlag-
# Folie bleibt unangetastet (eigener, kalibrierter Pfad).
# ─────────────────────────────────────────────────────────────────────────

# Untere Grenze des Tabellenbereichs der Einzeltitel-Folie.
# Die dunkelblaue Abschlusslinie des Layouts ("Linie rechts") liegt bei
# 6.38"; 6.20" lässt ~0.18" Sicherheitsabstand (Renderer legen Zeilen minimal
# höher aus als berechnet — gemessen ~0.1" Drift über 37 Zeilen). NICHT mit dem globalen
# MAX_TABLE_BOTTOM_INCH (6.60, Anlagevorschlag-Folie) verwechseln.
EINZELTITEL_MAX_BOTTOM_INCH = 6.20

EINZELTITEL_MIN_ROW_H_INCH = 0.125
"""Untergrenze Datenzeilenhöhe (6pt-Text braucht ca. 0.10")."""

EINZELTITEL_MAX_ROW_H_INCH = 0.19
"""Obergrenze: bei wenigen Titeln sollen die Zeilen nicht aufgeblasen
werden — der freie Bereich bleibt bewusst leer (statt leerer Linien)."""

EINZELTITEL_MIN_FONT_PT = 6.0
EINZELTITEL_MAX_FONT_PT = 7.0

EINZELTITEL_GAP_H_INCH = 0.10
"""Höhe der linienlosen Abstandszeile zwischen Tabelle und Gesamt-Zeile."""

EINZELTITEL_SUMMARY_H_INCH = 0.17
"""Höhe der Gesamt-Zeile (etwas mehr Luft, damit '100,00 %' vertikal
mittig steht und nicht an der Linie klebt)."""

_ROW_H_PER_PT_INCH = 0.0237
"""Kalibriert an der Vorlage: 0.142"-Zeile trägt 6pt-Text."""

_LINE_TAGS = ("lnL", "lnR", "lnT", "lnB")


def _zelle_rahmen_entfernen(cell):
    """Entfernt alle Rahmenlinien einer Tabellenzelle (setzt sie auf noFill).

    Nötig für die Abstandszeile: eine geklonte Datenzeile bringt sonst ihre
    0.25pt-Linien mit und es stünden drei Linien übereinander.

    WICHTIG — Schema-Reihenfolge: In <a:tcPr> müssen lnL, lnR, lnT, lnB VOR
    den Füll-Elementen (noFill/solidFill/…) stehen. Hängt man sie hinten an,
    ignorieren die Renderer sie und zeichnen stattdessen die Standard-Rahmen
    des Tabellen-Styles (sichtbarer leerer Kasten). Deshalb: vorne einfügen.
    """
    from lxml import etree
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in _LINE_TAGS:
        for el in tcPr.findall(qn(f"a:{tag}")):
            tcPr.remove(el)
    # rückwärts an Position 0 einfügen → Endreihenfolge lnL, lnR, lnT, lnB
    for tag in reversed(_LINE_TAGS):
        ln = etree.Element(qn(f"a:{tag}"))
        etree.SubElement(ln, qn("a:noFill"))
        tcPr.insert(0, ln)


def _zelle_leeren_kompakt(cell, font_pt: float = 4.0):
    """Leert eine Zelle UND begrenzt die Höhe ihres leeren Absatzes.

    text_frame.clear() lässt einen leeren Absatz zurück, dessen Schriftgröße
    nicht gesetzt ist → er erbt die Default-Größe (~18pt) und bläht die Zeile
    auf, egal welche Höhe im XML steht. Deshalb defRPr UND endParaRPr klein
    setzen (Renderer nutzen unterschiedliche der beiden).
    """
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.font.size = Pt(font_pt)                       # defRPr
    end = p._p.get_or_add_endParaRPr()              # endParaRPr
    end.sz = int(font_pt * 100)


def _zeile_schrift_setzen(row, font_pt: float, bold: Optional[bool] = None):
    """Setzt Schriftgröße (und optional Fettung) ALLER Runs einer Zeile.

    Behebt u.a. den Fall, dass ein Run in der Vorlage KEIN sz-Attribut hat
    und dadurch die Default-Schriftgröße (~18pt) erbt — genau das ließ die
    'Gesamt'-Zelle über die Zeile hinauswachsen und die Rahmenlinie
    durchkreuzen.
    """
    for cell in row.cells:
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(font_pt)
                if bold is not None:
                    run.font.bold = bold


def tabelle_leere_zeilen_entfernen(table, daten_spalten,
                                   kopf_zeilen: int = 1,
                                   fuss_zeilen: int = 1) -> int:
    """Entfernt leere Datenzeilen zwischen Kopf und Fuß (Summenzeile).

    GENERISCH: `daten_spalten` ist die Liste der Spalten-Indizes, die über
    'leer' entscheiden — dadurch für jede Vorlage der Familie nutzbar.
    Indizes außerhalb der Tabelle werden ignoriert (statt IndexError).

    Returns: Anzahl entfernter Zeilen.
    """
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    if n_rows <= kopf_zeilen + fuss_zeilen:
        return 0

    spalten = [c for c in daten_spalten if c < n_cols]
    zu_entfernen = []
    for i in range(kopf_zeilen, n_rows - fuss_zeilen):
        row = table.rows[i]
        leer = True
        for c in spalten:
            text = row.cells[c].text_frame.text.strip()
            if text and text != "\u00a0":       # NBSP zählt als leer
                leer = False
                break
        if leer:
            zu_entfernen.append(i)

    if not zu_entfernen:
        return 0

    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    for idx in sorted(zu_entfernen, reverse=True):
        tbl.remove(trs[idx])
    return len(zu_entfernen)


def _zelle_rahmen_uebernehmen(ziel_cell, quell_cell, ziel_tag: str, quell_tag: str):
    """Kopiert eine Rahmenlinie (z.B. lnT) einer Quellzelle als andere Linie
    (z.B. lnB) in die Zielzelle — inklusive Stärke, Farbe und Strichart.

    Bewusst KOPIEREN statt Attribute nachzubauen: so bleibt exakt das
    Corporate-Design der Vorlage erhalten (0.75pt, Farbe 14355C, solid).
    """
    from lxml import etree
    q_tcPr = quell_cell._tc.get_or_add_tcPr()
    quelle = q_tcPr.find(qn(f"a:{quell_tag}"))
    if quelle is None:
        return False
    z_tcPr = ziel_cell._tc.get_or_add_tcPr()
    for el in z_tcPr.findall(qn(f"a:{ziel_tag}")):
        z_tcPr.remove(el)
    neu = deepcopy(quelle)
    neu.tag = qn(f"a:{ziel_tag}")
    # Schema-Reihenfolge in tcPr: lnL, lnR, lnT, lnB, ...
    reihenfolge = ["lnL", "lnR", "lnT", "lnB"]
    ziel_pos = reihenfolge.index(ziel_tag)
    eingefuegt = False
    for kind in list(z_tcPr):
        name = etree.QName(kind).localname
        if name in reihenfolge and reihenfolge.index(name) > ziel_pos:
            kind.addprevious(neu)
            eingefuegt = True
            break
    if not eingefuegt:
        # hinter die letzte vorhandene Linie, sonst an den Anfang
        letzte = None
        for kind in list(z_tcPr):
            if etree.QName(kind).localname in reihenfolge:
                letzte = kind
        if letzte is not None:
            letzte.addnext(neu)
        else:
            z_tcPr.insert(0, neu)
    return True


def _linien_breite_pt(cell, tag: str):
    """Stärke einer Rahmenlinie in Punkt, oder None wenn keine/abgeschaltet."""
    tcPr = cell._tc.find(qn("a:tcPr"))
    if tcPr is None:
        return None
    el = tcPr.find(qn("a:" + tag))
    if el is None or el.find(qn("a:noFill")) is not None:
        return None
    w = el.get("w")
    return int(w) / 12700.0 if w else None


def _ist_gruppenzeile(cell) -> bool:
    """Kategorie-Kopfzeile? = erste Spalte fett UND nicht leer.

    fill_table_with_positions schreibt Gruppennamen (RENTEN, AKTIEN,
    EDELMETALLE, LIQUIDITÄT) mit is_bold=True, Positionen mit is_bold=False —
    die Fettung ist damit das verlässliche Merkmal direkt am Ergebnis.
    """
    if not " ".join(cell.text_frame.text.split()):
        return False
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            if run.font.bold:
                return True
    return False


def tabelle_kategorie_trennlinien(table, kopf_zeilen: int = 1,
                                  letzte_zeilen: int = 1) -> int:
    """Setzt den DICKEN Trennstrich genau unter jede Kategorie-Überschrift.

    PROBLEM (gefunden 07.08.2026 an der CVV-Broschüre "Defensiv"):
    fill_table_with_positions schreibt zwar Text und Fettung in die Zeilen,
    fasst die Rahmenlinien aber nie an. Die dicken Trennstriche blieben
    deshalb dort stehen, wo die VORLAGE ihre Gruppen hatte — unabhängig
    davon, wo die Kategorien nach dem Befüllen tatsächlich beginnen. Bei
    "Defensiv" lief der dicke Strich mitten durch die Rentenliste (zwischen
    Fraport und Fresenius), während der Übergang Würth → AKTIEN gar keinen
    bekam. Über alle Familien waren es 80 falsch platzierte Striche.

    REGEL (festgelegt mit Philip, 07.08.2026): dicker Strich NUR unter der
    Kategorie-Überschrift; zwischen den Positionen und vor der nächsten
    Überschrift steht die dünne Linie.

        Würth Finance IHS 3 %     ← dünn darunter
        AKTIEN                    ← DICK darunter
        Future of Defence ETF     ← dünn darunter

    MECHANIK: Die beiden Linienarten werden aus der Tabelle selbst geerntet
    (dickste bzw. dünnste vorhandene Unterkante) und per
    _zelle_rahmen_uebernehmen kopiert — nicht nachgebaut. So bleiben Stärke,
    Farbe und Strichart der Vorlage exakt erhalten, spaltenweise.
    Wie in tabelle_abstandszeile_einfuegen dokumentiert, wird an jeder
    Zeilengrenze BEIDSEITIG gesetzt (lnB oben, lnT unten), weil Renderer
    angrenzende Zellrahmen zusammenführen.

    Args:
        table: die Tabelle (shape.table)
        kopf_zeilen: Anzahl Kopfzeilen, die unberührt bleiben
        letzte_zeilen: Anzahl Schlusszeilen (Summe/Abstand), die unberührt
            bleiben — die regelt tabelle_abschlusslinie_sichern

    Returns:
        Anzahl der gesetzten Zeilengrenzen (0 = nichts getan).
    """
    zeilen = list(table.rows)
    erste = kopf_zeilen
    letzte = len(zeilen) - 1 - letzte_zeilen      # letzte Datenzeile
    if letzte - erste < 1:
        return 0

    n_spalten = len(list(zeilen[erste].cells))

    # ── Linienarten spaltenweise aus der Tabelle ernten ──────────────────
    # Pro Spalte die Zelle mit der dicksten und der dünnsten Unterkante.
    dick_quelle, duenn_quelle = {}, {}
    for spalte in range(n_spalten):
        breiteste = schmalste = None
        b_max = b_min = None
        for r in range(erste, letzte + 1):
            zellen = list(zeilen[r].cells)
            if spalte >= len(zellen):
                continue
            cell = zellen[spalte]
            w = _linien_breite_pt(cell, "lnB")
            if w is None:
                continue
            if b_max is None or w > b_max:
                b_max, breiteste = w, cell
            if b_min is None or w < b_min:
                b_min, schmalste = w, cell
        # Nur brauchbar, wenn es WIRKLICH zwei verschiedene Stärken gibt
        if breiteste is not None and schmalste is not None and b_max > b_min:
            dick_quelle[spalte] = breiteste
            duenn_quelle[spalte] = schmalste

    if not dick_quelle:
        return 0        # Tabelle kennt keine zwei Linienarten → nichts tun

    # ── Grenzen setzen ──────────────────────────────────────────────────
    gesetzt = 0
    for r in range(erste, letzte):
        oben = list(zeilen[r].cells)
        unten = list(zeilen[r + 1].cells)
        # Dick, wenn DIESE Zeile eine Kategorie-Überschrift ist
        quellen = dick_quelle if _ist_gruppenzeile(oben[0]) else duenn_quelle
        for spalte in range(n_spalten):
            quelle = quellen.get(spalte)
            if quelle is None:
                continue
            if spalte < len(oben):
                _zelle_rahmen_uebernehmen(oben[spalte], quelle, "lnB", "lnB")
            if spalte < len(unten):
                _zelle_rahmen_uebernehmen(unten[spalte], quelle, "lnT", "lnB")
        gesetzt += 1
    return gesetzt


def tabelle_abschlusslinie_sichern(table, kopf_zeilen: int = 1) -> bool:
    """Vereinheitlicht die Unterkante der LETZTEN Datenzeile auf die normale,
    dünne Trennlinie einer mittleren Datenzeile.

    Hintergrund: In der Vorlage trägt die physisch letzte Datenzeile eine dicke
    0.75pt-Linie (weil dort früher direkt die Summe folgte). Werden leere
    Zeilen entfernt, rutscht mal eine dicke, mal eine dünne Zeile ans Ende —
    je nach Strategie. Wir setzen sie einheitlich auf DÜNN; die klare Trennung
    zur Summe übernimmt die dicke Linie über der Gesamt-Zeile.
    """
    n = len(table.rows)
    if n < kopf_zeilen + 3:
        return False
    letzte_daten = table.rows[n - 2]
    # Muster: eine mittlere Datenzeile (die erste trägt oft eine dicke Linie
    # unter der Gruppen-Überschrift)
    muster_idx = max(kopf_zeilen + 1, (kopf_zeilen + n - 2) // 2)
    muster = table.rows[muster_idx]
    ok = False
    for i, cell in enumerate(letzte_daten.cells):
        if i < len(muster.cells):
            ok = _zelle_rahmen_uebernehmen(cell, muster.cells[i], "lnB", "lnB") or ok
    return ok


def tabelle_abstandszeile_einfuegen(table):
    """Fügt direkt VOR der letzten Zeile (Summe) eine leere Abstandszeile ein.

    Rahmen-Logik (WICHTIG): Renderer führen angrenzende Zellrahmen zusammen —
    setzt man die Abstandszeile rundum auf noFill, verschwinden auch die
    Linien der Nachbarzeilen. Deshalb wird an jeder Zeilengrenze BEIDSEITIG
    dasselbe gesetzt:
        oben  (letzte Datenzeile ↔ Abstandszeile): dünne Trennlinie
        unten (Abstandszeile ↔ Gesamt-Zeile):      dicke Summenlinie
    Seitlich (lnL/lnR) bleibt die Zeile rahmenlos, sonst entsteht ein
    sichtbarer leerer Kasten.

    Returns: die neue Zeile (pptx-Row) oder None.
    """
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    if len(trs) < 3:
        return None
    vorlage = trs[-2]                    # letzte Datenzeile
    neu = deepcopy(vorlage)
    trs[-1].addprevious(neu)             # vor die Summenzeile

    idx = len(table.rows) - 2
    row = table.rows[idx]
    letzte_daten = table.rows[idx - 1]
    summe = table.rows[idx + 1]

    for i, cell in enumerate(row.cells):
        _zelle_leeren_kompakt(cell)          # leert + hält den Absatz klein
        _zelle_rahmen_entfernen(cell)        # erst mal alle vier weg
        # Oberkante = Unterkante der letzten Datenzeile (dünn)
        if i < len(letzte_daten.cells):
            _zelle_rahmen_uebernehmen(cell, letzte_daten.cells[i], "lnT", "lnB")
        # Unterkante = Oberkante der Gesamt-Zeile (dicke Summenlinie)
        if i < len(summe.cells):
            _zelle_rahmen_uebernehmen(cell, summe.cells[i], "lnB", "lnT")
    return row


def tabelle_dynamisch_skalieren(table_shape, max_bottom_inch: float,
                                min_row_h: float = EINZELTITEL_MIN_ROW_H_INCH,
                                max_row_h: float = EINZELTITEL_MAX_ROW_H_INCH,
                                min_font_pt: float = EINZELTITEL_MIN_FONT_PT,
                                max_font_pt: float = EINZELTITEL_MAX_FONT_PT,
                                gap_h: float = EINZELTITEL_GAP_H_INCH,
                                summary_h: float = EINZELTITEL_SUMMARY_H_INCH,
                                kopf_zeilen: int = 1) -> dict:
    """Skaliert Zeilenhöhen + Schriftgröße dynamisch auf den verfügbaren Platz.

    Erwartete Struktur: [Kopf] + n Datenzeilen + [Abstandszeile] + [Summe].
    Die Datenzeilen teilen sich den Rest zwischen Kopf und `max_bottom_inch`,
    begrenzt durch min_row_h/max_row_h. Die Schriftgröße folgt der Zeilenhöhe
    (kalibriert: 0.142" trägt 6pt), gedeckelt auf [min_font_pt, max_font_pt].

    - VIELE Zeilen → eng (bis min_row_h / min_font_pt)
    - WENIGE Zeilen → ruhiger (bis max_row_h), Rest bleibt bewusst leer

    Returns: Dict mit den gewählten Werten (Diagnose) + ggf. "warnung".
    """
    table = table_shape.table
    n_rows = len(table.rows)
    # letzte zwei Zeilen = Abstandszeile + Summe
    n_daten = max(0, n_rows - kopf_zeilen - 2)
    if n_daten == 0:
        return {"n_daten": 0}

    kopf_h = sum(table.rows[i].height for i in range(kopf_zeilen)) / 914400.0
    top_inch = table_shape.top / 914400.0
    verfuegbar = max_bottom_inch - top_inch - kopf_h - gap_h - summary_h

    roh = verfuegbar / n_daten
    row_h = min(max_row_h, max(min_row_h, roh))
    font_pt = min(max_font_pt, max(min_font_pt, round(row_h / _ROW_H_PER_PT_INCH * 2) / 2))

    warnung = None
    if roh < min_row_h:
        warnung = (f"Einzeltitel-Tabelle: {n_daten} Zeilen — bei minimaler "
                   f"Zeilenhöhe ({min_row_h}\") und {min_font_pt}pt reicht der "
                   f"Platz nicht ganz. Folie bitte prüfen.")

    # Zeilenhöhen setzen: Datenzeilen, Abstandszeile, Summenzeile
    for i in range(kopf_zeilen, kopf_zeilen + n_daten):
        table.rows[i].height = int(row_h * 914400)
    table.rows[n_rows - 2].height = int(gap_h * 914400)     # Abstandszeile
    table.rows[n_rows - 1].height = int(summary_h * 914400)  # Summe

    # Schrift der Datenzeilen (Fettung der Gruppen-Header bleibt erhalten)
    for i in range(kopf_zeilen, kopf_zeilen + n_daten):
        _zeile_schrift_setzen(table.rows[i], font_pt)

    # Shape-Höhe = Summe der Zeilenhöhen
    total = sum(r.height for r in table.rows)
    table_shape.height = total

    return {"n_daten": n_daten, "row_h": round(row_h, 4),
            "font_pt": font_pt, "bottom": round(top_inch + total / 914400.0, 3),
            "warnung": warnung}


def ensure_ring_legend_fits(chart_shape, categories,
                            per_entry_in: float = 0.22,
                            pad_in: float = 0.06,
                            char_w_in: float = 0.078,
                            swatch_in: float = 0.22,
                            wpad_in: float = 0.14) -> bool:
    """Vergrößert die Legenden-Box eines Ring-Charts, falls sie zu klein ist,
    um ALLE Kategorien vollständig (ohne Zeilenumbruch) anzuzeigen.

    Hintergrund (BUGFIX 20.07.2026): Die Themen-Einzeltitel-Vorlage hat eine
    Legenden-manualLayout, die nur für die 2 Platzhalter-Kategorien dimensioniert
    war (0,99" breit × 0,49" hoch). Nachdem der Ring datengetrieben befüllt wird:
      • HÖHE zu klein  → PowerPoint schnitt untere Einträge ab (LIQUIDITÄT fehlte)
      • BREITE zu klein → lange Labels wie 'EDELMETALLE' (~1,07") brachen auf zwei
        Zeilen um; der Umbruch fraß die Höhe → wieder ein Eintrag zu wenig.

    Fix: Box bei Bedarf vergrößern —
      • HÖHE nach OBEN (Unterkante bleibt am Rahmenboden → bleibt im Bild),
      • BREITE nach RECHTS (linke Kante bleibt bei x; horizontal ist links viel
        Platz bis zum Ring/zur Quelle-Notiz).
    Vergrößert NUR, verkleinert nie → bereits ausreichende Legenden (ESG/CVV)
    bleiben unangetastet. Generisch: skaliert mit Kategorienzahl UND längstem
    Label, keine hartcodierten Kategorien.

    per_entry_in: Höhe pro Eintrag (9pt-Font ≈ 0,21"; ESG/CVV: 4 Einträge in 0,84").
    char_w_in/swatch_in/wpad_in: Breitenschätzung Text + Farb-Swatch + Rand.
    """
    try:
        root = chart_shape.chart._chartSpace
    except Exception:
        return False
    C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
    leg = root.find(".//" + C + "legend")
    if leg is None:
        return False
    ml = leg.find(".//" + C + "manualLayout")
    if ml is None:
        return False                      # ohne manualLayout sizet PP selbst
    ye, he = ml.find(C + "y"), ml.find(C + "h")
    xe, we = ml.find(C + "x"), ml.find(C + "w")
    if ye is None or he is None or not categories:
        return False
    frame_w_in = chart_shape.width / 914400.0
    frame_h_in = chart_shape.height / 914400.0
    if frame_w_in <= 0 or frame_h_in <= 0:
        return False
    n = len(categories)
    geaendert = False

    # ── HÖHE: alle Einträge (je eine Zeile) müssen passen ──────────────────
    y = float(ye.get("val"))
    h = float(he.get("val"))
    noetig_h = (n * per_entry_in + pad_in) / frame_h_in
    if h < noetig_h:
        unterkante = y + h                 # bleibt (Legende bleibt im Rahmen)
        neu_y = max(0.0, unterkante - noetig_h)
        ye.set("val", "%.5f" % neu_y)
        he.set("val", "%.5f" % (unterkante - neu_y))
        geaendert = True

    # ── BREITE: längstes Label ohne Zeilenumbruch ──────────────────────────
    if xe is not None and we is not None:
        x = float(xe.get("val"))
        w = float(we.get("val"))
        max_chars = max(len(str(c).strip()) for c in categories)
        noetig_w = (swatch_in + max_chars * char_w_in + wpad_in) / frame_w_in
        if w < noetig_w:
            neu_w = min(noetig_w, 1.0 - x)  # nach rechts, im Rahmen bleiben
            we.set("val", "%.5f" % neu_w)
            geaendert = True

    return geaendert


def fill_einzeltitel_themen_slide(prs, slide_idx: int, df, strategy_name: str,
                                   eval_date=None):
    """Befüllt die Einzeltitel-Folie der THEMEN-Broschüren (NEU 06.07.2026,
    Pro / Pro Dividende / Offensiv).

    Layout dieser Vorlagen-Tabelle (7 Spalten, abweichend vom Standard-
    Anlagevorschlag!):
        Spalte 0: Wertpapier (bzw. Gruppen-Überschrift AKTIEN/RENTEN/…)
        Spalte 2: Währung
        Spalte 4: WKN
        Spalte 6: Anteil %
        Spalten 1/3/5: schmale Spacer (unverändert lassen)
    Zeile 0 = Kopf, letzte Zeile = Summe (100,00 %). Gruppen kommen in
    GROUP_ORDER (AKTIEN, RENTEN, EDELMETALLE, LIQUIDITÄT, SONSTIGE), Titel
    je Gruppe alphabetisch — dieselbe group_portfolio_positions-Logik wie
    der Standard, nur anderes Spaltenlayout und mit Währung.

    Anders als der Standard-Anlagevorschlag verteilt diese Funktion NICHT auf
    mehrere Folien — die Themen-Broschüre hat genau eine Einzeltitel-Folie.
    Passen nicht alle Positionen, wird die Tabelle erweitert (ensure_table_
    capacity) und eine Warnung über LAST_BUILD_ERRORS gemeldet (nie stilles
    Abschneiden).
    """
    slide = prs.slides[slide_idx]

    # Titel
    title = find_shape_by_name(slide, "Titel") or find_shape_by_name(slide, "Titel 2")
    if title and title.has_text_frame:
        replace_text_in_shape(title, "Einzeltitel")

    tabelle = (find_shape_by_name(slide, "T_Kennzahlen")
               or find_shape_by_name(slide, "Tabelle")
               or find_shape_by_name(slide, "Tabelle 2"))
    if tabelle is None or not getattr(tabelle, "has_table", False):
        return
    t = tabelle.table

    # Spaltenindizes dieses Layouts
    C_NAME, C_WAEHRUNG, C_WKN, C_ANTEIL = 0, 2, 4, 6

    groups = group_portfolio_positions(df)

    # ── Assetklassen-Ring (C_Kennzahlen) befüllen ──────────────────────────
    # BUGFIX 20.07.2026: Diese Funktion füllte bisher NUR die Tabelle. Der
    # Assetklassen-Ring der Themen-Broschüre wurde von KEINER Funktion mit
    # Daten versorgt (die Orchestrierung ruft für die Rolle 'einzeltitel_themen'
    # nur diese Funktion auf) und behielt daher die PLATZHALTER-Daten der
    # Vorlage (z.B. AKTIEN 98,09% / LIQUIDITÄT 1,91%). Folge: EDELMETALLE fehlte
    # und die Ring-Werte passten nicht zur Tabelle.
    #
    # Fix: den Ring mit DERSELBEN Aggregation wie der Standard-Anlagevorschlag
    # befüllen (GROUP_ORDER, inkl. EDELMETALLE) — dieselben `groups`, also
    # garantiert konsistent mit der Tabelle. Datengetrieben und generisch:
    # jede in den Positionen vorhandene Assetklasse erscheint als eigenes
    # Segment; keine Strategie ist hartcodiert.
    alloc_labels, alloc_values = [], []
    for g in GROUP_ORDER:
        if g in groups:
            gsum = sum(safe_float(p["gewicht"], 0.0) for p in groups[g])
            if gsum > 0.0001:
                alloc_labels.append(g)
                alloc_values.append(float(gsum))
    ring = find_shape_by_name(slide, SHAPE_CHART_ALLOCATION)
    if ring and getattr(ring, "has_chart", False) and sum(alloc_values) > 0:
        replace_chart_data(
            ring,
            categories=list(alloc_labels),
            values=list(alloc_values),
            series_name="Anteil",
        )
        # Legende an Kategorienzahl UND längstes Label anpassen: sonst schnitt
        # die kleine Vorlagen-Box LIQUIDITÄT ab bzw. brach 'EDELMETALLE' um.
        ensure_ring_legend_fits(ring, alloc_labels)

    # Flache Zeilenliste bauen: pro Gruppe eine Header-Zeile + Positionen
    zeilen = []  # (typ, dict)
    total = 0.0
    for g in GROUP_ORDER:
        if g not in groups:
            continue
        positionen = groups[g]
        gruppen_gewicht = sum(safe_float(p["gewicht"], 0.0) for p in positionen)
        total += gruppen_gewicht
        if g == GROUP_LIQUIDITAET:
            # Liquidität: nur Header-Zeile mit Wert (keine Einzelpositionen)
            zeilen.append(("header_liq", {"name": g, "wert": gruppen_gewicht}))
        else:
            zeilen.append(("header", {"name": g}))
            for p in positionen:
                zeilen.append(("position", p))

    # Kapazität sicherstellen (Header-Zeile 0 + Zeilen + Summenzeile)
    benoetigt = len(zeilen) + 2
    if benoetigt > len(t.rows):
        try:
            ensure_table_capacity(t, len(zeilen))
        except Exception:
            pass
    n_rows = len(t.rows)
    summen_idx = n_rows - 1
    max_daten = summen_idx - 1  # Zeilen 1 .. summen_idx-1

    if len(zeilen) > max_daten:
        _record_einzeltitel_warnung(strategy_name, len(zeilen), max_daten)

    # Datenzeilen leeren (nur Daten-Spalten, Spacer unangetastet)
    for r in range(1, n_rows):
        for c in (C_NAME, C_WAEHRUNG, C_WKN, C_ANTEIL):
            if c < len(t.columns):
                set_cell_text(t.rows[r].cells[c], "")

    # Zeilen schreiben
    for i, (typ, data) in enumerate(zeilen):
        if i >= max_daten:
            break
        row = t.rows[i + 1]  # +1 wegen Kopfzeile
        if typ == "header":
            set_cell_text(row.cells[C_NAME], data["name"], is_bold=True)
        elif typ == "header_liq":
            set_cell_text(row.cells[C_NAME], data["name"], is_bold=True)
            set_cell_text(row.cells[C_ANTEIL], fmt_pct(data["wert"]), is_bold=True)
        else:  # position
            set_cell_text(row.cells[C_NAME], data["wertpapier"], is_bold=False)
            set_cell_text(row.cells[C_WAEHRUNG], data.get("waehrung", ""), is_bold=False)
            set_cell_text(row.cells[C_WKN], data.get("wkn", ""), is_bold=False)
            set_cell_text(row.cells[C_ANTEIL], fmt_pct(data["gewicht"]), is_bold=False)

    # Summenzeile aufräumen und befüllen:
    # - "Gesamt"-Label links (C_NAME), 100%-Wert rechts (C_ANTEIL)
    # - die Vorlagen-Platzhalter "0" in den Spacer-Spalten (1/3/5) UND in
    #   Währung/WKN entfernen, damit die Zeile sauber aussieht.
    if summen_idx >= 1:
        summen_row = t.rows[summen_idx]
        # alle Zellen der Summenzeile bis auf Label + Wert leeren
        for c in range(len(t.columns)):
            if c not in (C_NAME, C_ANTEIL):
                set_cell_text(summen_row.cells[c], "")
        set_cell_text(summen_row.cells[C_NAME], "Gesamt", is_bold=True)
        if C_ANTEIL < len(t.columns):
            set_cell_text(summen_row.cells[C_ANTEIL],
                          fmt_pct(total if total > 0 else 1.0), is_bold=True)

    # ── Layout (NEU 09.07.2026) ────────────────────────────────────────────
    # Reihenfolge ist wichtig: erst leere Zeilen raus, dann Abstandszeile,
    # dann skalieren (die Skalierung rechnet mit der finalen Zeilenzahl).
    try:
        # 1. Leere Datenzeilen entfernen (7-Spalten-Layout: 0/2/4/6).
        #    Sonst bleiben bei kurzen Portfolios leere Zeilen MIT Linien stehen.
        tabelle_leere_zeilen_entfernen(t, daten_spalten=(C_NAME, C_WAEHRUNG,
                                                         C_WKN, C_ANTEIL))

        # 1b. Unterkante der (nun) letzten Datenzeile vereinheitlichen. In der
        #     Vorlage trägt die physisch letzte Zeile eine dicke Linie; nach dem
        #     Entfernen leerer Zeilen rutscht je nach Strategie mal eine dicke,
        #     mal eine dünne ans Ende. Wir setzen sie einheitlich auf DÜNN —
        #     die klare Trennung übernimmt die dicke Linie über "Gesamt".
        #     VOR der Abstandszeile (Indizes!).
        tabelle_abschlusslinie_sichern(t)

        # 1c. Dicke Trennstriche an die TATSÄCHLICHEN Kategoriegrenzen setzen
        #     (NEU 07.08.2026). Ohne diesen Schritt bleiben sie dort stehen,
        #     wo die Vorlage ihre Gruppen hatte. Nach 1b, damit die
        #     vereinheitlichte Unterkante der letzten Datenzeile erhalten
        #     bleibt (letzte_zeilen=1 lässt sie unberührt).
        tabelle_kategorie_trennlinien(t)

        # 2. Linienlose Abstandszeile vor die Gesamt-Zeile.
        tabelle_abstandszeile_einfuegen(t)

        # 3. Zeilenhöhen + Schrift dynamisch an den Platz anpassen; die
        #    Tabelle endet garantiert oberhalb der blauen Abschlusslinie.
        info = tabelle_dynamisch_skalieren(
            tabelle, max_bottom_inch=EINZELTITEL_MAX_BOTTOM_INCH)

        # 4. Gesamt-Zeile EXPLIZIT setzen. In der Vorlage hat die Zelle
        #    "Gesamt" kein sz-Attribut und erbt sonst ~18pt → der Text quillt
        #    über die Zeile und durchkreuzt die blaue Rahmenlinie.
        summen_font = min(EINZELTITEL_MAX_FONT_PT,
                          (info.get("font_pt") or EINZELTITEL_MIN_FONT_PT) + 0.5)
        _zeile_schrift_setzen(t.rows[len(t.rows) - 1], summen_font, bold=True)

        if info.get("warnung"):
            EINZELTITEL_WARNUNGEN.append(info["warnung"])
    except Exception as exc:      # Layout-Kosmetik darf den Export nie killen
        EINZELTITEL_WARNUNGEN.append(
            f"Einzeltitel {strategy_name}: Layout-Anpassung übersprungen ({exc})")


# Sammelt Kapazitäts-Warnungen der Einzeltitel-Themen-Folie (analog zum
# LAST_BUILD_ERRORS-Muster in pptx_export). Wird von dort ausgelesen.
EINZELTITEL_WARNUNGEN = []


def _record_einzeltitel_warnung(strategy_name, n_zeilen, kapazitaet):
    EINZELTITEL_WARNUNGEN.append(
        f"Einzeltitel {strategy_name}: {n_zeilen} Zeilen, aber nur "
        f"{kapazitaet} Tabellenzeilen — {n_zeilen - kapazitaet} abgeschnitten. "
        f"Tabelle in der Vorlage vergrößern.")



# ─────────────────────────────────────────────────────────────────────────
# Strategie-Übersichtstabelle (CVV Folie 17) — NEU 10.07.2026
#
# Eine Tabelle über ALLE Strategien: Zeiträume in den Zeilen, Strategien in
# den Spalten. Läuft EINMAL pro Broschüre, nicht je Strategie.
#
# Vorlagen-Layout (8x13):
#   Zeile 0      Kopf ("Zeitraum" | 5 Strategienamen) — STATISCH, bleibt
#   Zeile 1      leer (Abstand)
#   Zeilen 2..6  YTD / 1 / 3 / 5 / 10 Jahre
#   Zeile 7      leer (Abschluss)
#   Spalte 0     Zeitraum-Label (statisch)
#   Spalten 4/6/8/10/12  Strategiewerte (1/2/3/5/7/9/11 = schmale Spacer)
# ─────────────────────────────────────────────────────────────────────────

UEBERSICHT_PERIODEN = ["ytd", "1J", "3J", "5J", "10J"]
"""Reihenfolge der Datenzeilen — MUSS zu den Zeilenbeschriftungen der Vorlage
passen (YTD, 1-, 3-, 5-, 10-Jahres Performance). Keys wie in
pptx_export.compute_rollierend_data()."""

UEBERSICHT_ERSTE_DATENZEILE = 2
UEBERSICHT_WERTSPALTEN = [4, 6, 8, 10, 12]


def fill_uebersicht_slide(prs, slide_idx: int, rollierend_liste,
                          stand_date_str=None,
                          perioden=None, erste_zeile=None, spalten=None,
                          shape_name="Tabelle", fussnote_name="Fußnote"):
    """Befüllt die strategieübergreifende Wertentwicklungs-Tabelle.

    Args:
        rollierend_liste: Liste (in Spaltenreihenfolge!) von Dicts
            {"ytd": float|None, "1J": …, "3J": …, "5J": …, "10J": …}
            — genau die Rückgabe von compute_rollierend_data(). Ein Eintrag
            darf None sein (keine Zeitreihe) → ganze Spalte bekommt "-".
        stand_date_str: z.B. "07.07.2026". Wird als "Stand: …" an die Fußnote
            gehängt (idempotent — ein vorhandener Stand wird ersetzt).

    Fehlende Historie (z.B. Dynamic ohne 10 Jahre) → "-", nie 0,00 %.
    Kopfzeile, Zeitraum-Labels und Disclaimer bleiben unangetastet: sie stehen
    statisch in der Vorlage.

    Returns: Anzahl geschriebener Zellen.
    """
    perioden = perioden or UEBERSICHT_PERIODEN
    erste_zeile = UEBERSICHT_ERSTE_DATENZEILE if erste_zeile is None else erste_zeile
    spalten = spalten or UEBERSICHT_WERTSPALTEN

    slide = prs.slides[slide_idx]
    shape = find_shape_by_name(slide, shape_name)
    if shape is None or not shape.has_table:
        _record_einzeltitel_warnung(f"Übersicht (Folie {slide_idx + 1})", 0, 0)
        return 0
    t = shape.table
    n_rows, n_cols = len(t.rows), len(t.columns)

    geschrieben = 0
    for s_i, spalte in enumerate(spalten):
        if spalte >= n_cols:
            continue
        daten = rollierend_liste[s_i] if s_i < len(rollierend_liste) else None
        for p_i, periode in enumerate(perioden):
            zeile = erste_zeile + p_i
            if zeile >= n_rows:
                continue
            wert = (daten or {}).get(periode)
            # None / NaN → "-" (fehlende Historie), NIE 0,00 % anzeigen
            if wert is None or (isinstance(wert, float) and wert != wert):
                text = "-"
            else:
                text = fmt_pct(wert)
            set_cell_text_preserve_format(t.rows[zeile].cells[spalte], text)
            geschrieben += 1

    if stand_date_str:
        _fussnote_stand_setzen(slide, fussnote_name, stand_date_str)
    return geschrieben


def _fussnote_stand_setzen(slide, shape_name, stand_date_str):
    """Hängt 'Stand: TT.MM.JJJJ' an die Fußnote — idempotent.

    Ein bereits vorhandenes 'Stand: …' am Ende wird ersetzt, damit sich der
    Zusatz bei wiederholtem Befüllen derselben Präsentation nicht stapelt.
    """
    import re as _re
    shape = find_shape_by_name(slide, shape_name)
    if shape is None or not shape.has_text_frame:
        return False
    # Rückwärts den letzten Absatz suchen, der überhaupt Text trägt — der
    # letzte Absatz der Vorlage ist oft leer (0 Runs).
    run = None
    for par in reversed(list(shape.text_frame.paragraphs)):
        treffer = [r for r in par.runs if r.text and r.text.strip()]
        if treffer:
            run = treffer[-1]
            break
    if run is None:
        return False
    basis = _re.sub(r"\s*Stand:\s*\d{2}\.\d{2}\.\d{4}\s*$", "", run.text)
    run.text = f"{basis.rstrip()} Stand: {stand_date_str}"
    return True

def fill_rollierend_slide(prs, slide_idx: int, strategy_name: str,
                          rollierend_data: Optional[dict] = None,
                          stand_date_str: Optional[str] = None):
    """Befüllt die rollierende Wertentwicklungs-Tabelle der Themen-Broschüren
    (NEU 06.07.2026 — "Wertentwicklung der Strategie {Name}", Tabelle 8x7).

    Struktur der Vorlagen-Tabelle (Muster Pro):
      Zeile 0: Kopf ("Zeitraum" | ... | "Strategie {Name}")
      Zeile 2: YTD Performance    → Wert in Spalte 4
      Zeile 3: 1 - Jahres Perf.   → Spalte 4
      Zeile 4: 3 - Jahres Perf.   → Spalte 4
      Zeile 5: 5 - Jahres Perf.   → Spalte 4
      Zeile 6: 10 - Jahres Perf.  → Spalte 4
    Werte nach Kosten (dieselbe Logik wie die Streamlit-Rolltabelle);
    fehlende Historie (< n Jahre) → "–". None-Daten → Vorlagen-Platzhalter
    bleiben (Titel wird dennoch gesetzt).

    Args:
        rollierend_data: Dict aus compute_rollierend_data
            {"ytd","1J","3J","5J","10J"} (Dezimal) oder None.
        stand_date_str: optionales Datum für eine "Quelle"/"Stand"-Box.
    """
    slide = prs.slides[slide_idx]

    # Titel: "Wertentwicklung der Strategie {Name}" (Shape "Titel" ODER "Titel 2")
    title = find_shape_by_name(slide, "Titel") or find_shape_by_name(slide, "Titel 2")
    if title and title.has_text_frame:
        replace_text_in_shape(title, f"Wertentwicklung der Strategie {strategy_name}")

    # Kopf-Spaltentitel "Strategie {Name}" aktualisieren (Zeile 0, Spalte 4)
    tabelle = find_shape_by_name(slide, "Tabelle") or find_shape_by_name(slide, "Tabelle 2")
    if tabelle is None or not getattr(tabelle, "has_table", False):
        return  # kein Table-Shape → nichts zu tun
    t = tabelle.table
    WERT_SPALTE = 4
    if len(t.rows) > 0 and len(t.columns) > WERT_SPALTE:
        kopf = t.rows[0].cells[WERT_SPALTE]
        if kopf.text_frame.text.strip():
            set_cell_text_preserve_format(kopf, f"Strategie {strategy_name}")

    if rollierend_data is None:
        return  # Platzhalter-Modus: nur Titel/Kopf

    def _fmt(v):
        if v is None:
            return "–"
        return f"{v * 100:.2f}%".replace(".", ",")

    # Zeilen 2-6 = YTD/1/3/5/10 Jahre in Spalte 4
    zuordnung = [
        (2, "ytd"), (3, "1J"), (4, "3J"), (5, "5J"), (6, "10J"),
    ]
    for row_idx, key in zuordnung:
        if row_idx < len(t.rows) and len(t.columns) > WERT_SPALTE:
            zelle = t.rows[row_idx].cells[WERT_SPALTE]
            set_cell_text_preserve_format(zelle, _fmt(rollierend_data.get(key)))

    # Optionale Quelle/Stand-Box (falls in der Vorlage vorhanden)
    if stand_date_str:
        quelle = find_shape_by_name(slide, "Quelle")
        if quelle:
            set_shape_text_static(
                quelle, f"Quelle: Eigene Berechnung, Stand {stand_date_str}")


def fill_performance_slide(prs, slide_idx: int, strategy_name: str,
                            performance_data: Optional[dict] = None,
                            stand_date_str: Optional[str] = None):
    """Befüllt die Performance-Slide (Kennzahlen-Vergleich + 2 Charts mit BM).

    Änderungen 02.07.2026:
    - Punkt 2: Linien-Chart bekommt eine DATENBASIERTE Achsen-Untergrenze
      (Minimum über BEIDE Serien, 10%-Schritt darunter) statt des im Template
      fixierten 70%-Werts — vorher wurde die Portfoliolinie bei Strategien
      mit weit gelaufener Benchmark unnötig gestaucht bzw. bei jungen
      Strategien Platz verschenkt.
    - Punkt 6: Die 'Quelle'-Box enthielt ein Live-DATUMSFELD (zeigte das
      Öffnungs-Datum der Datei) → wird durch statischen Text mit dem
      Datenstand ersetzt.
    - Punkt 5 (nachbeauftragt 03.07.2026): Die Disclaimer-Fußnote der
      Vorlage enthält denselben veralteten Satz wie die cVV-Folie
      ("…erfolgt vor Kosten (ab 30.06. abzüglich halbjährigen
      Honorarsatz)…") — wortidentische Absätze, daher werden dieselben
      WE_DISCLAIMER_REPLACEMENTS angewendet. Läuft UNABHÄNGIG von
      performance_data (Textkorrektur gilt auch im Platzhalter-Modus),
      damit F8 und F9 nie widersprüchliche Kosten-Aussagen zeigen.

    Args:
        prs: Presentation
        slide_idx: 0-indexed Index der Performance-Slide
        strategy_name: Name der Strategie für den Titel
        performance_data: Dict mit Performance-Daten (siehe
            modules.analytics.compute_performance_data). Wenn None: nur Titel
            wird gesetzt, Charts/Tabelle bleiben mit Vorlagen-Platzhaltern.
        stand_date_str: Datum ("DD.MM.YYYY") für die statische Quelle-Zeile.
            None = Quelle-Box unangetastet lassen.
    """
    slide = prs.slides[slide_idx]

    # Titel anpassen: "{Strategy} | Wertentwicklung (mit Benchmark)"
    title = find_shape_by_name(slide, "Titel")
    if title and title.has_text_frame:
        new_title = f"{strategy_name} | Wertentwicklung (mit Benchmark)"
        replace_text_in_shape(title, new_title)

    # ── Fußnote (03.07.2026, Punkt 5): veralteten Kosten-Satz ersetzen ──
    # Die Vorlagen-Fußnote beschreibt noch die alte VBA-Honorarregel
    # ("vor Kosten, ab 30.06. …") — im Widerspruch zur F8 und zur
    # tatsächlichen Berechnung (nach Kosten, taggenau). Die Absätze sind
    # wortidentisch zur cVV-Folie → gleiche Ersetzungs-Konstanten.
    # BEWUSST vor dem Platzhalter-Return: Textkorrektur ist daten-unabhängig.
    fn = find_shape_by_name(slide, "Fußnote")
    if fn and fn.has_text_frame:
        for prefix, new_text in WE_DISCLAIMER_REPLACEMENTS:
            replace_paragraph_text_by_prefix(fn.text_frame, prefix, new_text)

    if performance_data is None:
        return  # Phase 1: nur Titel setzen

    # ── KENNZAHLEN-Tabelle befüllen ──
    kz = performance_data.get("kennzahlen", {})
    tab = find_shape_by_name(slide, "Tabelle")
    if tab and tab.has_table:
        fill_kennzahlen_table(tab.table, kz)

    # ── PERFORMANCE P.A. Chart (Säulen) ──
    # ACHTUNG, bewusste Abweichung zur Wertentwicklungs-Folie (10.08.2026):
    # Hier heißt die Serie weiter "Referenzportfolio", weil DIESE Folie den
    # Begriff in der Vorlage statisch führt (Vorlage_FFPB.pptx slide10 —
    # Legenden-Box UND Folientitel "Referenzportfolio| Wertentwicklung").
    # Die Wertentwicklungs-Folie nutzt WE_SERIES_PORTFOLIO ("Musterdepot"),
    # weil ihre Vorlagen das so vorgeben. Kein Versehen: jede Folie folgt
    # ihrer eigenen Vorlage. Wer das vereinheitlichen will, ändert die
    # Vorlagen, nicht diese Konstanten.
    pa = performance_data.get("performance_pa", {})
    # Standard TRUE — Begruendung siehe fill_wertentwicklung_slide.
    hat_benchmark = bool(performance_data.get("has_benchmark", True))
    chart_links = find_shape_by_name(slide, "Diagramm links")
    if chart_links and chart_links.has_chart and pa.get("jahre"):
        # Ohne Benchmark keine zweite Serie (11.08.2026) — analog zur
        # Wertentwicklungs-Folie, siehe dort. Diese Folie wird aktuell von
        # keiner Familien-Konfiguration verwendet (nur von
        # DEFAULT_TEMPLATE_CONFIG); die Behandlung steht hier trotzdem,
        # damit beide Folien nicht auseinanderlaufen.
        serien = [("Referenzportfolio", pa.get("referenz", []))]
        if hat_benchmark:
            serien.append(("Benchmark", pa.get("benchmark", [])))
        replace_chart_data_safe(
            chart_links,
            categories=[str(y) for y in pa["jahre"]],
            series_data=serien,
            data_label_format=PCT_FORMAT_CODE,
            # NEU (Juni 2026, Bug 4): Ohne diesen Parameter zeigte die
            # Y-Achse Rohwerte (0.05, 0.1, ...) statt Prozent (5%, 10%, ...),
            # obwohl die Daten-Labels über den Balken korrekt formatiert
            # waren. Bewiesen an echter Chart-XML: <c:valAx><c:numFmt
            # formatCode="General" sourceLinked="1"/> statt "0%"/sourceLinked=0.
            value_axis_format="0%",
        )

    # ── WERTENTWICKLUNG Chart (Linien) ──
    we = performance_data.get("wertentwicklung", {})
    chart_rechts = find_shape_by_name(slide, "Diagramm rechts")
    if chart_rechts and chart_rechts.has_chart and we.get("dates"):
        linien = [("Referenzportfolio", we.get("referenz", []))]
        if hat_benchmark:
            linien.append(("Benchmark", we.get("benchmark", [])))
        replace_chart_data_safe(
            chart_rechts,
            categories=we["dates"],
            series_data=linien,
            data_label_format=None,  # Linien-Chart hat keine Daten-Labels
        )
        # 02.07.2026 (Punkt 2): datenbasierte Untergrenze über BEIDE Serien
        # (Template-Fixwert 70% ersetzt; Details siehe _line_axis_min).
        all_vals = list(we.get("referenz", [])) + list(we.get("benchmark", []))
        set_value_axis_min(chart_rechts, _line_axis_min(all_vals))

    # ── Quelle (02.07.2026, Punkt 6): Datumsfeld → statischer Datenstand ──
    if stand_date_str:
        quelle = find_shape_by_name(slide, "Quelle")
        if quelle:
            set_shape_text_static(
                quelle, f"Quelle: Eigene Berechnung, Stand {stand_date_str}")


def _line_axis_min(values) -> float:
    """Berechnet die Achsen-Untergrenze für einen Index-Linien-Chart
    (NEU 02.07.2026, Punkt 2): 10%-Schritt UNTER dem Datenminimum,
    gedeckelt bei 1.0 (Achse beginnt nie über 100%, damit der Startpunkt
    der Kurve sichtbar bleibt) und nie unter 0.

    Beispiele: min(Daten)=0.94 → 0.9;  min=1.0 (nur steigend) → 0.9
    (eine Stufe Luft unter dem Start);  min=0.62 → 0.6.
    """
    try:
        data_min = min(float(v) for v in values)
    except (ValueError, TypeError):
        return 0.7  # Fallback: bisheriger Template-Wert
    import math
    floor10 = math.floor(data_min * 10.0) / 10.0
    if floor10 >= data_min - 1e-9:      # Datenminimum liegt exakt auf der Stufe
        floor10 -= 0.1                  # → eine Stufe Luft darunter
    return max(0.0, min(floor10, 0.9))


def fill_wertentwicklung_slide(prs, slide_idx: int, strategy_name: str,
                                we_data: Optional[dict] = None,
                                stand_date_str: Optional[str] = None):
    """Befüllt die Wertentwicklungs-/Kurzübersichts-Folie (NEU Juli 2026).

    Das ist die aus dem alten VBA-Tool übernommene Folie
    "Anlagestrategie {Name} | Wertentwicklung" (cVV-Broschüre), die beim
    Export als Folie 8 eingereiht wird. Bestandteile:

    - Titel: "Anlagestrategie {Name} | Wertentwicklung"
    - KENNZAHLEN-Tabelle (7×3): 4 Kennzahlen mit DYNAMISCHEN Labels
      (Auflagejahr + laufendes Jahr werden in die Label-Texte geschrieben):
        Row 2: "Wertentwicklung seit {Auflagejahr} kumuliert*"   → Wert %
        Row 3: "Rendite p.a. seit {Auflagejahr} nach Kosten"     → Wert %
        Row 4: "Wertentwicklung seit 01.01.{Jahr}**"             → Wert %
        Row 5: "Duration"                                        → Dezimal
    - "Diagramm links" (Säulen): Performance p.a. nach Kosten im
      Benchmarkvergleich, volle Kalenderjahre (max. 5); Achsenformat "0%"
      (02.07.2026, Punkt 1 — vereinheitlicht mit F9; Daten-Labels bleiben
      2-stellig "0.00%")
    - "Diagramm rechts" (Linie): Wertentwicklung als Index (Start = 1.0),
      gesamte Historie, EINE Serie; Achsen-Untergrenze DATENBASIERT
      (02.07.2026, Punkt 2 — 10%-Schritt unter Datenminimum, deterministisch
      in PowerPoint UND LibreOffice, statt Template-Fixwert/Auto)
    - Legende: bleibt unangetastet ("Musterdepot     Benchmark***" wie in der
      Vorlage). Der Angleich an F9 vom 02.07.2026 ist am 10.08.2026
      zurückgenommen — Begründung siehe WE_SERIES_PORTFOLIO.

    OHNE Benchmark (has_benchmark=False, z.B. "Muster SCHWEIZ Substanz")
    verschwindet der Vergleichsmaßstab VOLLSTÄNDIG von der Folie
    (11.08.2026): keine Benchmark-Serie im Säulen-Chart, kein
    "Benchmark***" in der Legenden-Box, keine ***-Zeile in der Fußnote.
    Vorher blieb an allen drei Stellen etwas stehen — im Chart Null-Balken,
    in der Fußnote sogar die Benchmark einer FREMDEN Strategie (der
    Vorlagentext von "Pro").
    - Quelle: Datumsfeld → statischer Text mit Datenstand (02.07.2026,
      Punkt 6 — vorher zeigte die Box das ÖFFNUNGS-Datum der Datei)
    - Fußnote: *-, **- und ***-Zeilen sowie Disclaimer-Satz dynamisch
      (siehe WE_FOOTNOTE_* / WE_DISCLAIMER_REPLACEMENTS)

    Args:
        prs: Presentation
        slide_idx: 0-indexed Index der Wertentwicklungs-Folie
        strategy_name: bereinigter Strategiename
        we_data: Dict aus pptx_export.compute_wertentwicklung_data(). Keys:
            auflage_jahr, laufendes_jahr, kum_nach_kosten, pa_nach_kosten,
            ytd, duration, benchmark_text, performance_pa, wertentwicklung.
            Wenn None: nur Titel wird gesetzt, Rest bleibt Vorlagen-Platzhalter.
        stand_date_str: Datum (String "DD.MM.YYYY") für die statische
            Quelle-Zeile. None = Quelle-Box unangetastet lassen.
    """
    slide = prs.slides[slide_idx]

    # ── Titel ──
    title = find_shape_by_name(slide, SHAPE_TITLE) or find_shape_by_name(slide, SHAPE_TITLE_ALT)
    if title:
        set_title_with_autoscale(title, WE_TITLE_FORMAT.format(name=strategy_name))

    if we_data is None:
        return  # Platzhalter-Modus (analog fill_performance_slide)

    # ── KENNZAHLEN-Tabelle: Labels UND Werte dynamisch ──
    aj = we_data.get("auflage_jahr")
    lj = we_data.get("laufendes_jahr")
    tab = find_shape_by_name(slide, SHAPE_WE_TABLE)
    if tab and getattr(tab, "has_table", False) and aj is not None:
        t = tab.table
        rows_spec = [
            (WE_ROW_KUMULIERT, f"Wertentwicklung seit {aj} kumuliert*",
             fmt_pct(we_data.get("kum_nach_kosten"))),
            (WE_ROW_PA, f"Rendite p.a. seit {aj} nach Kosten",
             fmt_pct(we_data.get("pa_nach_kosten"))),
            (WE_ROW_YTD, f"Wertentwicklung seit 01.01.{lj}**",
             fmt_pct(we_data.get("ytd"))),
            (WE_ROW_DURATION, "Duration",
             fmt_ratio(we_data.get("duration"))),
        ]
        for row_idx, label, value in rows_spec:
            if row_idx >= len(t.rows):
                continue
            row = t.rows[row_idx]
            set_cell_text_preserve_format(row.cells[WE_COL_LABEL], label)
            set_cell_text_preserve_format(row.cells[WE_COL_VALUE], value)

    # ── Säulen-Chart: Performance p.a. im Benchmarkvergleich ──
    # Datenbasis identisch zur Performance-Folie: nur VOLLE Kalenderjahre
    # (Fußnote *: "nach Kosten bis zum 31.12. des Vorjahres").
    pa = we_data.get("performance_pa", {})
    # Standard TRUE, nicht False (11.08.2026): Nur ein AUSDRÜCKLICHES
    # has_benchmark=False laesst Serie, Legendeneintrag und Fußnotenzeile
    # verschwinden. Fehlt der Schluessel — etwa weil ein Aufrufer ein
    # unvollstaendiges Dict baut —, bleibt alles stehen. Der umgekehrte
    # Standard waere gefaehrlich: er wuerde bei jedem solchen Aufruf still
    # Inhalte aus einer Kundenbroschuere entfernen, auch bei Strategien MIT
    # Vergleichsmassstab. compute_wertentwicklung_data setzt den Schluessel
    # immer, der Normalfall ist davon also unberuehrt.
    hat_benchmark = bool(we_data.get("has_benchmark", True))
    chart_bar = find_shape_by_name(slide, SHAPE_WE_CHART_BAR)
    if chart_bar and getattr(chart_bar, "has_chart", False) and pa.get("jahre"):
        # Ohne Benchmark wird die zweite Serie GAR NICHT übergeben
        # (11.08.2026). replace_chart_data_safe reicht das an python-pptx
        # weiter, das überzählige <c:ser> entfernt — die Vorlagen-Serie
        # verschwindet also samt Legendeneintrag des Charts, statt als Reihe
        # von Null-Balken stehen zu bleiben.
        serien = [(WE_SERIES_PORTFOLIO, pa.get("referenz", []))]
        if hat_benchmark:
            serien.append((WE_SERIES_BENCHMARK, pa.get("benchmark", [])))
        replace_chart_data_safe(
            chart_bar,
            categories=[str(y) for y in pa["jahre"]],
            series_data=serien,
            # Daten-Labels wie cVV-Original 2-stellig ("6,04%"); die ACHSE
            # dagegen "0%" ("5%, 10%, …") — 02.07.2026 (Punkt 1) angeglichen
            # an F9, vorher zeigte F8 klobige "5,00%/10,00%"-Achsenticks.
            data_label_format="0.00%",
            value_axis_format="0%",
        )

    # ── Linien-Chart: Wertentwicklung (Index, EINE Serie) ──
    we = we_data.get("wertentwicklung", {})
    chart_line = find_shape_by_name(slide, SHAPE_WE_CHART_LINE)
    if chart_line and getattr(chart_line, "has_chart", False) and we.get("dates"):
        replace_chart_data_safe(
            chart_line,
            categories=we["dates"],
            series_data=[
                ("Wertentwicklung", we.get("referenz", [])),
            ],
            data_label_format=None,   # Linie hat keine Daten-Labels
            value_axis_format="0%",   # Achse als Prozent (100%, 110%, ...)
        )
        # 02.07.2026 (Punkt 2): Achsen-Untergrenze DATENBASIERT setzen —
        # 10%-Schritt unter dem Datenminimum (max. 100%). Grund: Der
        # Template-Fixwert (70%) passt nicht für alle Strategien, und
        # PowerPoint-AUTO wählt bei weit über 100% laufenden Indizes gerne
        # 0% als Minimum und staucht die Kurve. Datenbasiert = identisches,
        # vorhersagbares Rendering in PowerPoint und LibreOffice.
        set_value_axis_min(chart_line, _line_axis_min(we.get("referenz", [])))
        # 03.07.2026: Das cVV-Chart war für MONATS-Daten gebaut —
        # baseTimeUnit="months" bündelt Tagespunkte monatsweise und die
        # Linie wirkt in POWERPOINT zerhackt (in LibreOffice unsichtbar!).
        # Auf Tages-Granularität stellen + Linienstärke an F9 angleichen
        # (0,75pt → 1,5pt; bei 6000+ Tagespunkten wirkt dünn unruhig).
        set_date_axis_base_unit(chart_line, "days")
        set_series_line_width(chart_line, WE_LINE_WIDTH_PT)

    # ── Fußnote: dynamische Zeilen ersetzen ──
    fn = find_shape_by_name(slide, SHAPE_WE_FUSSNOTE)
    if fn and fn.has_text_frame:
        # ***-Zeile: Benchmark-Zusammensetzung der Strategie.
        # WICHTIG: Reihenfolge — *** VOR ** ersetzen ist nicht nötig, denn
        # Präfix "** " (mit Leerzeichen) matcht die ***-Zeile NICHT
        # ("***…" hat an Position 3 einen Stern, kein Leerzeichen).
        #
        # OHNE Benchmark fliegt die Zeile RAUS (BUGFIX 11.08.2026): Bis dahin
        # wurde sie nur dann überschrieben, wenn ein Benchmark-Text vorlag —
        # bei "Muster SCHWEIZ Substanz"/"…Aktien" blieb deshalb der
        # Vorlagentext stehen und die Kundenbroschüre nannte die Benchmark
        # der Strategie "Pro" ("50% EuroStoxx 50; 50% MSCI World Euro").
        # Am echten Artefakt reproduziert.
        bm_text = we_data.get("benchmark_text")
        if not hat_benchmark:
            remove_paragraph_by_prefix(fn.text_frame, WE_FOOTNOTE_STAR3_PREFIX)
        elif bm_text:
            replace_paragraph_text_by_prefix(
                fn.text_frame, WE_FOOTNOTE_STAR3_PREFIX, f"*** {bm_text}")
        # *-Zeile (02.07.2026, Punkt 7): Kennzahlen 1+2 laufen jetzt bis zum
        # letzten Datenpunkt → "bis zum 31.12. des Vorjahres" entfällt.
        replace_paragraph_text_by_prefix(
            fn.text_frame, WE_FOOTNOTE_STAR1_PREFIX, WE_FOOTNOTE_STAR1_NEW)
        # **-Zeile + Disclaimer-Satz: alte VBA-Honorarregel → Tool-Konvention
        # (nach Kosten, taggenauer Abzug). Siehe WE_DISCLAIMER_REPLACEMENTS.
        replace_paragraph_text_by_prefix(
            fn.text_frame, WE_FOOTNOTE_STAR2_PREFIX, WE_FOOTNOTE_STAR2_NEW)
        for prefix, new_text in WE_DISCLAIMER_REPLACEMENTS:
            replace_paragraph_text_by_prefix(fn.text_frame, prefix, new_text)

    # ── Legende: nur der Benchmark-Eintrag, und nur wenn es keine gibt ─────
    # Die Textbox "Legende Diagramm links" steht in jeder Vorlage auf
    # "Musterdepot     Benchmark***" — genau so, wie sie gedruckt werden soll,
    # SOLANGE die Strategie eine Benchmark hat. Hat sie keine, nennt die
    # Legende einen Balken, den es im Chart nicht mehr gibt; dann wird der
    # Eintrag samt Farbquadrat entfernt (11.08.2026).
    #
    # 'Musterdepot' bleibt in JEDEM Fall unangetastet: Bis 10.08.2026 schrieb
    # hier eine Ersetzung 'Musterdepot ' → 'Referenzportfolio ' und kürzte zum
    # Ausgleich den 5-Leerzeichen-Run auf 3 (02.07.2026, Punkt 3). Das ist
    # zurückgenommen; Begründung siehe WE_SERIES_PORTFOLIO. Wer den BEGRIFF
    # ändern will, ändert die VORLAGE — nicht den Code.
    if not hat_benchmark:
        legende = find_shape_by_name(slide, SHAPE_WE_LEGENDE)
        if legende and legende.has_text_frame:
            remove_legend_entry(legende.text_frame, WE_SERIES_BENCHMARK)

    # ── Quelle (02.07.2026, Punkt 6): Datumsfeld → statischer Datenstand ──
    if stand_date_str:
        quelle = find_shape_by_name(slide, SHAPE_WE_QUELLE)
        if quelle:
            set_shape_text_static(
                quelle, f"Quelle: Eigene Berechnung, Stand {stand_date_str}")


def fill_zusammenstellung_slide(prs, slide_idx: int, df: pd.DataFrame,
                                 strategy_name: str, eval_date=None):
    """Befüllt Slide 9 mit 2 Ringen: Regionen (links) + Branchen/Segment (rechts).

    Kleine Kategorien (<3%) werden zu "Sonstige" zusammengefasst, maximal 8
    Segmente angezeigt. Nicht-klassifizierte Positionen (typischerweise
    Liquidität) erscheinen als eigenes Segment "Liquidität", damit der Ring
    auf 100% summiert.

    Die Ringe sind native PowerPoint-Donuts; es werden nur die Chart-Daten
    ersetzt, das Template-Styling (Banner "REGIONEN"/"Branchen", Legende,
    Quelle) bleibt erhalten.

    Args:
        prs: Presentation
        slide_idx: 0-indexed Slide-Position
        df: Portfolio-DataFrame
        strategy_name: bereinigter Strategiename
        eval_date: Auswertungsdatum. Optional (Quelle-Datum wird zentral gesetzt).
    """
    slide = prs.slides[slide_idx]

    # Titel
    title = find_shape_by_name(slide, SHAPE_TITLE) or find_shape_by_name(slide, SHAPE_TITLE_ALT)
    if title:
        replace_text_in_shape(title, f"Aktuelle Portfoliozusammenstellung – {strategy_name}")

    # Defensive Vorbereitung: Gewicht muss sauberer Float sein
    df_clean = df.copy()
    if "Gewicht" in df_clean.columns:
        df_clean["Gewicht"] = pd.to_numeric(df_clean["Gewicht"], errors="coerce").fillna(0.0).astype(float)

    # Regionen (links) — nativer PowerPoint-Donut (Daten ersetzen)
    region_agg = build_ring_series(df_clean, "Region")
    chart_left = find_shape_by_name(slide, SHAPE_CHART_LEFT)
    if not region_agg.empty and chart_left and getattr(chart_left, "has_chart", False):
        values = [float(v) for v in region_agg.values]
        if sum(values) > 0:
            replace_chart_data(
                chart_left,
                categories=region_agg.index.tolist(),
                values=values,
                series_name="Anteil",
            )

    # Segmente/Branchen (rechts) — nativer PowerPoint-Donut (Daten ersetzen)
    segment_agg = build_ring_series(df_clean, "Segment")
    chart_right = find_shape_by_name(slide, SHAPE_CHART_RIGHT)
    if not segment_agg.empty and chart_right and getattr(chart_right, "has_chart", False):
        values = [float(v) for v in segment_agg.values]
        if sum(values) > 0:
            replace_chart_data(
                chart_right,
                categories=segment_agg.index.tolist(),
                values=values,
                series_name="Anteil",
            )
